import gc
gc.collect()
import os
import glob
import numpy as np
import pandas as pd
import warnings
import time
import logging
import matplotlib.pyplot as plt
import seaborn as sns
import plotly.express as px
import openpyxl
import calendar
from scipy.stats import norm
from sklearn import metrics
from statsmodels.stats.outliers_influence import variance_inflation_factor
from statsmodels.formula.api import ols
from datetime import datetime,timedelta
from openpyxl import load_workbook
from openpyxl.chart import BarChart, Reference
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE,MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN,MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from io import BytesIO
import docx
from docx import Document
from docx.oxml import OxmlElement as word_oxml
from docx.shared import Inches,Pt, RGBColor as word_rgb
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml, OxmlElement as word_oxml
from docx.oxml.ns import qn,nsdecls
from docx.enum.table import WD_ALIGN_VERTICAL,WD_TABLE_ALIGNMENT,WD_CELL_VERTICAL_ALIGNMENT, WD_ROW_HEIGHT_RULE
from PIL import Image, ImageOps

warnings.filterwarnings('ignore')
warnings.simplefilter('ignore')
pd.set_option("display.max_columns",150)
pd.set_option("display.max_rows",200)
pd.set_option('display.float_format', lambda x: '%.9f' % x)
presentation=Presentation()
doc=Document()
heading_counters = [0, 0, 0, 0]
#...................................................................................................................................................................................#

logger=logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)

formatter=logging.Formatter('%(asctime)s---%(funcName)s---%(lineno)d---%(levelname)s---%(message)s')

timestamp=datetime.now().strftime("%Y%m%d_%H%M%S")
# path=fr's3://model-validation-monitoring-gbs/Validation/final_main_log{timestamp}.log'

file_handler=logging.FileHandler(filename="test.log",mode='w')
file_handler.setFormatter(formatter)

logger.addHandler(file_handler)

# ...................................................................................................................................................................................#
import google.ai.generativelanguage as glm
import pathlib
import google.generativeai as genai
genai.configure(api_key="AIzaSyCl24r6ILYtv3x0GW4RwPCHJ-92bsOOHfA")
def get_insights(path,text):
    model = genai.GenerativeModel("gemini-pro-vision")
    resepose = model.generate_content(glm.Content(
               parts = [glm.Part(text =text),
                       glm.Part(inline_data =glm.Blob(
                       mime_type = "image/png",
                       data = pathlib.Path(path).read_bytes())
                               ),
                       ],),
                                stream = True)
    resepose.resolve()
    return resepose.text

# ...............................................................................................................................................#
def save_table_as_image(table, image_path):
    data = []
    max_font_size = 0
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            paragraph = cell.paragraphs[0]
            run = paragraph.runs[0]
            font_size = run.font.size.pt
            max_font_size = max(max_font_size, font_size)
            row_data.append((paragraph.text, font_size))
        data.append(row_data)

    fig, ax = plt.subplots(figsize=(max_font_size * 0.05, max_font_size * 0.05))  # Adjust the image size based on font size
    ax.axis('off')  # Hide axis
    ax.table(cellText=data, loc='center')

    plt.savefig(image_path, dpi=300, bbox_inches='tight', pad_inches=0.05)  # Save as image with higher resolution


def read_column_name(portfolio_code):
    
    logger.debug(f"Reading the path where column names for the raw dataset of {portfolio_code} portfolio is stored starts.")
    
    path=r"\\zebra\GIRM\Group Portfolio Analytics and Risk Methodology\Named folders\GBS COE\Sharat\BSCR_L3 data\{}\columns.txt"
    path=path.format(portfolio_code)
    
    logger.debug(f"Reading the path where column names for the raw dataset of {portfolio_code} portfolio is stored ends.")
    logger.debug("#######################################################################################")
    
    return path

#.....................................................................................................................................................................#


def portfolio_variable_mapper(portfolio_code):
    
    logger.debug(f"Reading the excel mapper file containing SAS columns names of {portfolio_code} portfolio starts.")
    
    mapper_sheet=pd.read_excel(r"\\zebra\GIRM\Group Portfolio Analytics and Risk Methodology\Named folders\GBS COE\Sharat\BSCR_L3 data\Variable Mapper SAS to L3.xlsx",sheet_name=portfolio_code)
    
    logger.debug(f"Reading the excel mapper file containing SAS columns names of {portfolio_code} portfolio ends.")
    logger.debug("#######################################################################################")

    return mapper_sheet

#.....................................................................................................................................................................#


def read_l3_datas(portfolio_code):
    
    logger.debug(f"Reading the path where initial first draft of the l3 data was stored of {portfolio_code} portfolio starts.")    
    
    l3_data_path=r"\\zebra\GIRM\Group Portfolio Analytics and Risk Methodology\Named folders\GBS COE\Sharat\BSCR_L3 data\{}\l3_org_data"
    l3_data_path=l3_data_path.format(portfolio_code)
    extension="txt"
    os.chdir(l3_data_path)
    filename_list=glob.glob('*.{}'.format(extension))
    
    logger.debug(f"There are {len(filename_list)} files stored at the given path.")
    logger.debug(f"Reading the path where initial first draft of the l3 data was stored of {portfolio_code} portfolio ends.")    
    logger.debug("#######################################################################################")

    return filename_list,l3_data_path

#.....................................................................................................................................................................#



def temprory_data(filename_list,l3_data_path):
    temp_data=pd.DataFrame()
    
    logger.info("Concating the initial first draft of individual raw files into a single combined data file starts.")
    logger.debug(f"Dataframe temp_data has currently {temp_data.shape[0]} rows and {temp_data.shape[1]} columns.")
    logger.debug("#######################################################################################")
    
    for i in range(len(filename_list)):
        
        logger.info(f"Reading the file no {i+1} starts.")
        
        os.chdir(l3_data_path)
        
        if(i!=0):
        
            temp_data_1=pd.read_csv(filename_list[i],sep="~")
            
            
            logger.debug(f"The file no {i+1} has {temp_data_1.shape[0]} rows and {temp_data_1.shape[1]} columns.")
            
            temp_data=pd.concat([temp_data,temp_data_1.iloc[1:,:]])
        
            logger.debug(f"The file no {i+1} has been concatenated to the dataframe temp_data and now it has total {temp_data.shape[0]} rows and {temp_data.shape[1]} columns.")
            logger.debug("#######################################################################################")
            
        else:
            
            temp_data_1=pd.read_csv(filename_list[i],sep="~")
            
            
            logger.debug(f"The file no {i+1} has {temp_data_1.shape[0]} rows and {temp_data_1.shape[1]} columns.")
            
            temp_data=pd.concat([temp_data,temp_data_1],ignore_index=True)

            
            logger.debug(f"The file no {i+1} has been concatenated to the dataframe temp_data and now it has total {temp_data.shape[0]} rows and {temp_data.shape[1]} columns.")
            logger.debug("#######################################################################################")
      
    
            
            
    new_cols=[]
    for x in temp_data.columns.str.split("."):
        new_cols.append(x[1])
                
                
    temp_data.columns=new_cols
        
    logger.info("Concating the initial first draft of individual raw files into a single combined data file ends")
    logger.debug(f"Dataframe temp_data has finally {temp_data.shape[0]} rows and {temp_data.shape[1]} columns.")
    logger.debug("#######################################################################################")
    
    return temp_data

#.....................................................................................................................................................................#


def new_temprory_data(filename_list,l3_data_path):
    
    temp_data=pd.DataFrame()
    
    logger.info("Concating the second draft of individual raw files into a single combined data file starts.")
    logger.debug(f"Dataframe temp_data has currently {temp_data.shape[0]} rows and {temp_data.shape[1]} columns.")
    logger.debug("#######################################################################################")
    
    for i in range(len(filename_list)):
        
        logger.info(f"Reading the file no {i+1} starts.")
        
        os.chdir(l3_data_path)
        
        if(i!=0):
        
            temp_data_1=pd.read_csv(filename_list[i],sep="~",header=None)
            
            logger.debug(f"The file no {i+1} has {temp_data_1.shape[0]} rows and {temp_data_1.shape[1]} columns.")
            
            temp_data=pd.concat([temp_data,temp_data_1],ignore_index=True)
        
            logger.debug(f"The file no {i+1} has been concatenated to the dataframe temp_data and now it has total {temp_data.shape[0]} rows and {temp_data.shape[1]} columns.")
            logger.debug("#######################################################################################")
            
        else:
            
            temp_data_1=pd.read_csv(filename_list[i],sep="~")
            
            logger.debug(f"The file no {i+1} has {temp_data_1.shape[0]} rows and {temp_data_1.shape[1]} columns.")
            
            temp_data=pd.concat([temp_data,temp_data_1],ignore_index=True)
            
            new_cols=[]
            for x in temp_data.columns.str.split("_"):
                new_cols.append("_".join(x[4:]).split(".")[1])
                
                
            temp_data.columns=new_cols
        
            logger.debug(f"The file no {i+1} has been concatenated to the dataframe temp_data and now it has total {temp_data.shape[0]} rows and {temp_data.shape[1]} columns.")
            logger.debug("#######################################################################################")
        
    logger.info("Concating the second draft of individual raw files into a single combined data file ends")
    logger.debug(f"Dataframe temp_data has finally {temp_data.shape[0]} rows and {temp_data.shape[1]} columns.")
    logger.debug("#######################################################################################")
    
    return temp_data

#.....................................................................................................................................................................#



def change_to_Date_Time(df,list_):
    for i in list_:
        df[i]=pd.to_datetime(df[i])
    return df

#.....................................................................................................................................................................#
def remove_incomplete_per_data(df,start_dt,end_dt):
    
    logger.info("Removing the incomplete performance data to avoid duplicates starts.")
    logger.debug(f"Dataframe Combined data has currently {df.shape[0]} rows and {df.shape[1]} columns.")
    logger.debug("#######################################################################################")
    
    removed_data= ((df["SCORE_DATE"]>=(start_dt)) & (df["SCORE_DATE"]<=(end_dt)))
    
    logger.debug(f"Incomplete Performance data is to be removed from {start_dt} to {end_dt}.")
    
    df_F=df[~(removed_data)]
    
    logger.debug(f"After removing the incomplete performance data the combined data has {df_F.shape[0]} rows and {df_F.shape[1]} columns.")
    logger.info("Removing the incomplete performance data to avoid duplicates ends.")
    logger.debug("#######################################################################################")
    
    return df_F

#.....................................................................................................................................................................#


def new_dec_2021_data(portfolio_code,col_data):
    
    logger.info("Adding the new DEC'21 excel file which has the complete performance starts")
    logger.debug("#######################################################################################")    
    
    path=r"\\zebra\GIRM\Group Portfolio Analytics and Risk Methodology\Named folders\GBS COE\Sharat\BSCR_L3 data\{}\l3_org_data\l3_pb_uae_{}_bs_20230206.csv"
    path=path.format(portfolio_code,portfolio_code)
    Dec_21_Updated_df=pd.read_csv(path,sep="~",header=None,names=col_data.columns)
    
    logger.debug(f"The new DEC'21 file has {Dec_21_Updated_df.shape[0]} rows and {Dec_21_Updated_df.shape[1]} columns")
    logger.info("Adding the new DEC'21 excel file which has the complete performance ends.")
    logger.debug("#######################################################################################")    
    
    return Dec_21_Updated_df

#.....................................................................................................................................................................#


def read_new_l3_datas(portfolio_code):
    
    logger.debug(f"Reading the path where second draft of the l3 data was stored of {portfolio_code} portfolio starts.")    
    
    l3_data_path=r"\\zebra\GIRM\Group Portfolio Analytics and Risk Methodology\Named folders\GBS COE\Sharat\BSCR_L3 data\{}\l3_new_org_data"
    l3_data_path=l3_data_path.format(portfolio_code)
    extension="txt"
    os.chdir(l3_data_path)
    filename_list=glob.glob('*.{}'.format(extension))
    
    logger.debug(f"There are {len(filename_list)} files stored at the given path.")
    logger.debug(f"Reading the path where second draft of the l3 data was stored of {portfolio_code} portfolio ends.")    
    logger.debug("#######################################################################################")
    
    return filename_list,l3_data_path


#.....................................................................................................................................................................#


def New_Original_Bscore(df,portfolio_code):
    
    if(portfolio_code=='pl'):
        
        df[['SCORE_DATE','ACCOUNT_ID','BSCORE_ORIGINAL','SCORE_MODEL_CODE','RATIO_PRIN_BAL_ORIG_LOAN_WT', 'SLRY_AVGSLRY_1TO3_4TO9_WT', 'SLRY_MAXNUMMISNG_WT', 'SLRY_PCTTOTDUE_SLRY_1TO6_WT', 'LN2704LNTRM_WT', 'LN2711_TOB_ACCT_UNADJ_WT', 'MOSSNCDLQGT0_LAG_1TO24_WT', 'LN2809PMTC1TO6PCTBALC1TO6_WT', 'NATIONALITY_GROUP_WT', 'SLRY_AVGSLRY_1TO3_4TO12_WT', 'SLRY_PCT_MISNG_OCCR_LST3M_WT', 'SLRY_PCTTOTDUE_SLRY_1TO12_WT',  'LN2767PCTDLQGT1C1TO3_WT', 'LN2807PMTC1TO3PCTBALC1TO3_WT', 'LN2813PASTDUEC1PCTBALC1_RNG', 'LN2830MAXCMOSDLQGT1C1TO12_WT']].fillna(0,inplace=True)
        df['BSCORE_ORIGINAL_NEW']=df[['RATIO_PRIN_BAL_ORIG_LOAN_WT', 'SLRY_AVGSLRY_1TO3_4TO9_WT', 'SLRY_MAXNUMMISNG_WT', 'SLRY_PCTTOTDUE_SLRY_1TO6_WT', 'LN2704LNTRM_WT', 'LN2711_TOB_ACCT_UNADJ_WT', 'MOSSNCDLQGT0_LAG_1TO24_WT', 'LN2809PMTC1TO6PCTBALC1TO6_WT', 'NATIONALITY_GROUP_WT', 'SLRY_AVGSLRY_1TO3_4TO12_WT', 'SLRY_PCT_MISNG_OCCR_LST3M_WT', 'SLRY_PCTTOTDUE_SLRY_1TO12_WT',  'LN2767PCTDLQGT1C1TO3_WT', 'LN2807PMTC1TO3PCTBALC1TO3_WT', 'LN2813PASTDUEC1PCTBALC1_RNG', 'LN2830MAXCMOSDLQGT1C1TO12_WT']].sum(axis=1)
        
    if(portfolio_code=='al'):
        
        df[['DPST_AVGDPST_1TO6_7TO12_WT', 'SLRY_PCT_MISNG_OCCR_LST12M_WT', 'AUTO_TOTDUE_DPST_1TO12_WT', 'LN2748MAXDLQC2TO12_WT', 'LN2756PCTDLQGT0C1TO3_WT', 'LN2807PMTC1TO3PCTBALC1TO3_WT', 'LN2824MAXCMOSDLQGT0C_1TO6_WT', 'RATIO_PRIN_BAL_ORIG_LOAN_WT', 'USED_CAR_WT',
           'RATIO_SLRY_DPST_L6M_WT','SLRY_AVGSLRY_1TO3_4TO9_WT', 'SLRY_PCT_MISNG_OCCR_LST6M_WT', 'LN2809PMTC1TO6PCTBALC1TO6_WT','LN2704LNTRM_WT','AUTO_TOTDUE_DPST_1TO6_WT','NATIONALITY_GROUP_WT',
           'SLRY_AVGSLRY_1TO6_7TO12_WT', 'SLRY_PCTTOTDUE_SLRY_1TO3_WT', 'LN2718MOSSNCPMTGT0C1TO12_WT', 'LN2727DLQC1_WT', 'LN2734NUMDLQGT1C1TO6_WT']].fillna(0,inplace=True)
        df['BSCORE_ORIGINAL_NEW']=df[['DPST_AVGDPST_1TO6_7TO12_WT', 'SLRY_PCT_MISNG_OCCR_LST12M_WT', 'AUTO_TOTDUE_DPST_1TO12_WT', 'LN2748MAXDLQC2TO12_WT', 'LN2756PCTDLQGT0C1TO3_WT', 'LN2807PMTC1TO3PCTBALC1TO3_WT', 'LN2824MAXCMOSDLQGT0C_1TO6_WT', 'RATIO_PRIN_BAL_ORIG_LOAN_WT', 'USED_CAR_WT',
           'RATIO_SLRY_DPST_L6M_WT','SLRY_AVGSLRY_1TO3_4TO9_WT', 'SLRY_PCT_MISNG_OCCR_LST6M_WT', 'LN2809PMTC1TO6PCTBALC1TO6_WT','LN2704LNTRM_WT','AUTO_TOTDUE_DPST_1TO6_WT','NATIONALITY_GROUP_WT',
           'SLRY_AVGSLRY_1TO6_7TO12_WT', 'SLRY_PCTTOTDUE_SLRY_1TO3_WT', 'LN2718MOSSNCPMTGT0C1TO12_WT', 'LN2727DLQC1_WT', 'LN2734NUMDLQGT1C1TO6_WT']].sum(axis=1)
        
    if(portfolio_code=='nl'):
        
        df1=df[(df['SCORE_DATE']>='2022-07-31') & (df['SCORE_DATE']<='2022-12-31')]
        df2=df[~((df['SCORE_DATE']>='2022-07-31') & (df['SCORE_DATE']<='2022-12-31'))]
        
        july_dec_data=pd.read_excel(r"\\zebra\GIRM\Group Portfolio Analytics and Risk Methodology\Named folders\GBS COE\Sharat\BSCR_L3 data\nl\l3_new_org_data\NL_data_Jun22_dec22.xlsx")
        del df1['BSCORE_ORIGINAL']
        df1=pd.merge(df1,july_dec_data[['BSCORE_ORIGINAL','ACCOUNT_ID','SCORE_DATE']],how='left',left_on=['ACCOUNT_ID','SCORE_DATE'],right_on=['ACCOUNT_ID','SCORE_DATE'])
        
        df1['BSCORE_ORIGINAL_NEW']=df1['BSCORE_ORIGINAL']
        
        
        df2[['SLRY_PCT_MISNG_OCCR_LST12M_WT', 'LN2705RMNGLNTRM_WT', 'LN2730NUMDLQGT0C1_12_WT', 'DPST_AVGDPST_1TO3_4TO12_WT', 'RATIO_PRIN_BAL_ORIG_LOAN_WT','CURR_DELQ_WT', 'MAX_DELQ_L9M_WT', 'COUNT_DELQ_INC_L12M_WT', 'COUNT_PAY_L9M_WT', 'CUST_AGE_WT']].fillna(0,inplace=True)
        df2['BSCORE_ORIGINAL_NEW']=df2[['SLRY_PCT_MISNG_OCCR_LST12M_WT', 'LN2705RMNGLNTRM_WT', 'LN2730NUMDLQGT0C1_12_WT', 'DPST_AVGDPST_1TO3_4TO12_WT', 'RATIO_PRIN_BAL_ORIG_LOAN_WT','CURR_DELQ_WT', 'MAX_DELQ_L9M_WT', 'COUNT_DELQ_INC_L12M_WT', 'COUNT_PAY_L9M_WT', 'CUST_AGE_WT']].sum(axis=1)
        df2['BSCORE_ORIGINAL_NEW']=np.where(df2['SCORE_MODEL_CODE']=='PB_UAE_BSCR12',df2['BSCORE_ORIGINAL_NEW']+633,df2['BSCORE_ORIGINAL_NEW'])
        df=pd.concat([df1,df2],ignore_index=True)
        
    if(portfolio_code=='mort'):
        
        df1=df[(df['SCORE_DATE']>='2022-07-31') & (df['SCORE_DATE']<='2022-12-31')]
        df2=df[~((df['SCORE_DATE']>='2022-07-31') & (df['SCORE_DATE']<='2022-12-31'))]
#         del df1['BSCORE_ORIGINAL']
        df1['BSCORE_ORIGINAL']=np.where(df1['SCORE_VALUE']==492,429,df1['BSCORE_ORIGINAL'])
        df1['BSCORE_ORIGINAL']=np.where(df1['SCORE_VALUE']==510,462,df1['BSCORE_ORIGINAL'])
        df1['BSCORE_ORIGINAL']=np.where(df1['SCORE_VALUE']==535,509,df1['BSCORE_ORIGINAL'])
        df1['BSCORE_ORIGINAL']=np.where(df1['SCORE_VALUE']==556,549,df1['BSCORE_ORIGINAL'])
        df1['BSCORE_ORIGINAL']=np.where(df1['SCORE_VALUE']==563,562,df1['BSCORE_ORIGINAL'])
        df1['BSCORE_ORIGINAL']=np.where(df1['SCORE_VALUE']==577,587,df1['BSCORE_ORIGINAL'])
        df1['BSCORE_ORIGINAL']=np.where(df1['SCORE_VALUE']==578,589,df1['BSCORE_ORIGINAL'])
        df1['BSCORE_ORIGINAL']=np.where(df1['SCORE_VALUE']==587,606,df1['BSCORE_ORIGINAL'])
        df1['BSCORE_ORIGINAL']=np.where(df1['SCORE_VALUE']==588,609,df1['BSCORE_ORIGINAL'])
        df1['BSCORE_ORIGINAL']=np.where(df1['SCORE_VALUE']==590,612,df1['BSCORE_ORIGINAL'])
        # df['BSCORE_ORIGINAL']=np.where(df['SCORE_VALUE']==590,613,df['BSCORE_ORIGINAL'])
        df1['BSCORE_ORIGINAL']=np.where(df1['SCORE_VALUE']==607,643,df1['BSCORE_ORIGINAL'])
        df1['BSCORE_ORIGINAL']=np.where(df1['SCORE_VALUE']==608,646,df1['BSCORE_ORIGINAL'])
        df1['BSCORE_ORIGINAL']=np.where(df1['SCORE_VALUE']==627,681,df1['BSCORE_ORIGINAL']) 

        df=pd.concat([df1,df2],ignore_index=True)   
    
        
    return df


#.....................................................................................................................................................................#

def Islamic_pit_mob_6(df,portfolio_code):
    
    ''' This function outputs the dataframe after applying Islamic PiT mob 6 exclusions.
    
        input params:

        @ df - The required combined data after preprocessing and before calibration.

    '''

    logger.info(f"Defining conditions for Islamic PiT mob 6 flag for Islamic {portfolio_code} portfolio")
    logger.debug("#######################################################################################")    

    COND_PIT_MOB_LT_6_EXCLUSION=[
                
#MORTGAGES    
                #((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.INTERNATIONAL.isin([1]))),
                #(df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_02',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_05',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_04',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"]))),
#PL-STL
               # ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.INTERNATIONAL.isin([1]))),
                #((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_02',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_04',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_07',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_06',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_03',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])))]
    
    logger.info(f" Conditions defined successfuly for Islamic PiT mob 6 flag for Islamic {portfolio_code} portfolio")
    logger.debug("#######################################################################################")    
 

    logger.info(f" Applying condtions to create Islamic PiT mob 6 flag for Islamic {portfolio_code} portfolio")
    logger.debug("#######################################################################################") 
    
    logger.info(f" Islamic PiT mob 6 flag for Islamic for  {portfolio_code} portfolio has: {df['PIT_MOB_LT_6_EXCLUSION_REASON'].value_counts()}")
    logger.debug("#######################################################################################")
    
    
    df['ISLAMIC_PIT_MOB_LT_6_EXCLUSION_REASON']=np.select(COND_PIT_MOB_LT_6_EXCLUSION,[
#MORTGAGES
                          'A. OVERDRAFT','B. INTERNATIONAL','C. DPD > 90','D. POST_EXC_MOB <= 6','E. TOTAL OUTSTANDING = 1', 'F. POST_EXC_MOB > 6',
#PL-STL
                          'A. OVERDRAFT','B. INTERNATIONAL','C. ALREADY BAD','D. DPD > 90','E. POST_EXC_MOB <= 6', 'F. NO_SAL_LAST_12MON','G. TOTAL OUTSTANDING = 1', 'H. POST_EXC_MOB > 6'])
  
    logger.info(f" Islamic PiT mob 6 flag created successfuly for Islamic {portfolio_code} portfolio")
    logger.debug("#######################################################################################") 
    
    logger.info(f" PiT mob 6 flag for  {portfolio_code} portfolio has: {df['ISLAMIC_PIT_MOB_LT_6_EXCLUSION_REASON'].nunique()} unique exclusions i.e. {df['ISLAMIC_PIT_MOB_LT_6_EXCLUSION_REASON'].unique()}"  )
    logger.debug("#######################################################################################")
    
    return df

#.....................................................................................................................................................................#



#PORTFOILIO WISE PIT_EXCLUSION_REASON 
    
def islamic_pit_model_exclusion(df,portfolio_code):
    
    
    ''' This function outputs the dataframe after applying Islamic PiT exclusions.
    
        input params:

        @ df - The required combined data after preprocessing and before calibration.

    '''
    
    logger.info(f"Defining conditions for PiT model exclusion flag for Islamic  {portfolio_code} portfolio")
    logger.debug("#######################################################################################") 

    COND_PIT_EXCLUSION=[
                
#MORTGAGES    
                #((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.INTERNATIONAL.isin([1]))),
                #(df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_02',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_05',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_04',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"]))),

#PL-STL
                #((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.INTERNATIONAL.isin([1]))),
                #((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_02',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_04',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_07',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_06',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_03',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])))]
    

    logger.info(f"Conditions defined successfuly for Islamic PiT exclusion for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")   
    
  
    logger.info(f"Applying conditions to create PiT exclusion flag for Islamic  {portfolio_code} portfolio")
    logger.debug("#######################################################################################") 


    df['ISLAMIC_PIT_EXCLUSION_REASON']=np.select(COND_PIT_EXCLUSION,[
#MORTGAGES
                         'A. OVERDRAFT','B. INTERNATIONAL','C. DPD > 90','D. MOB <= 6','E. TOTAL OUTSTANDING = 1', 'F. POST_EXCLUSION',
#PL-STL
                         'A. OVERDRAFT','B. INTERNATIONAL','C. ALREADY BAD','D. DPD > 90','E. MOB <= 6', 'F. NO_SAL_LAST_12MON','G. TOTAL OUTSTANDING = 1', 'H. POST_EXCLUSION'])
 
    logger.info(f" Conditions applied successfuly for PiT exclusion flag for Islamic {portfolio_code} portfolio")
    logger.debug("#######################################################################################")
    
    
    logger.info(f" PiT exclusion flag for Islamic  {portfolio_code} portfolio has: {df['ISLAMIC_PIT_EXCLUSION_REASON'].nunique()} unique exclusions i.e. {df['ISLAMIC_PIT_EXCLUSION_REASON'].unique()}"  )
    logger.debug("#######################################################################################")
    
    return df
#.....................................................................................................................................................................#


def islamic_mis_exclusion(df,portfolio_code):
    
    ''' This function outputs the dataframe after applying Islamic MIS exclusions.
    
        input params:

        @ df - The required combined data after preprocessing and before calibration.

    '''
    
    
    logger.info(f"Defining conditions for Misexclusion flag for Islamic  {portfolio_code} portfolio")
    logger.debug("#######################################################################################") 

    COND_MIS_EXCLUSION=[
          
#MORTGAGES 
        
        #FOR MORTGAGES IN MIS RECON WE NEED TO REMOVE SOME PRODUCT IDs (4102, 21055, 21064, 21065, 3042) as SHARED BY DEV. TEAM
          #((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.ISLAMIC.isin([1]))) ,
          ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.OVERDRAFT.isin([1]))),
          ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.INTERNATIONAL.isin([1]))),
          #(df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),        
          ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"]))),

#PL-STL         
          #((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.ISLAMIC.isin([1]))) ,
          ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.OVERDRAFT.isin([1]))),
          ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.INTERNATIONAL.isin([1]))),
          #((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
          ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) &(df.PORTFOLIO_CODE.isin(["Personal Loans"])))]                   

    
    
    logger.info(f"Conditions defined successfuly for MIS exclusion flag for Islamic  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")
    
    logger.info(f"Applying conditions to create MIS exclusion flag for Islamic  {portfolio_code} portfolio")
    logger.debug("#######################################################################################") 
    
    df['ISLAMIC_MIS_RECON_EXCLUSION_REASON']= np.select(COND_MIS_EXCLUSION,[
                          
#MORTGAGES
                                    'A. OVERDRAFT','B. INTERNATIONAL','C. POST_EXCLUSION',
#PL-STL
                                   'A. OVERDRAFT','B. INTERNATIONAL','C. POST_EXCLUSION'])
    
    
    logger.info(f" Conditions applied successfuly for MIS exclusion flag for Islamic {portfolio_code} portfolio")
    logger.debug("#######################################################################################")
    
    
    logger.info(f" Mis exclusion flag for Islamic {portfolio_code} portfolio has: {df['ISLAMIC_MIS_RECON_EXCLUSION_REASON'].nunique()} unique exclusions i.e. {df['ISLAMIC_MIS_RECON_EXCLUSION_REASON'].unique()}"  )
    logger.debug("#######################################################################################")
    
    return df

#.....................................................................................................................................................................#

def islamic_bscore_model_exclusion(df,portfolio_code):
    
    
    
    ''' This function outputs the dataframe after applying Islamic Bscore exclusions.
    
        input params:

        @ df - The required combined data after preprocessing and before calibration.

    '''
    
    
    logger.info(f"Defining conditions for Bscore exclusion flag for Islamic {portfolio_code} portfolio")
    logger.debug("#######################################################################################") 
    
    #We are taking data from SEP'19 as the model went in production from this time.

    COND_EXCLUSION=[
                
#MORTGAGES    
                #((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.INTERNATIONAL.isin([1]))),
                #(df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_02',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_05',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_04',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Mortgage Loans"]))),

#PL-STL          
        
        
        

        
            
        
                #((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.INTERNATIONAL.isin([1]))),
                #((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_02',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_04',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_07',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_06',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_03',regex=False))),
                ((df['SCORE_DATE']>=pd.to_datetime("2021-09-30")) & (df['ISLAMIC'].isin([1])) & (df.PORTFOLIO_CODE.isin(["Personal Loans"])))]


    

    logger.info(f"Conditions defined successfuly for Bscore exclusion flag for Islamic {portfolio_code} portfolio")
    logger.debug("#######################################################################################")
    
    logger.info(f"Applying conditions to create Bscore exclusion flag for Islamic  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")     

    #ASSIGNING BSCORE_EXCLUSION_REASON TO ALL THE BSCORE DATASETS
    df['ISLAMIC_BSCORE_EXCLUSION_REASON']=np.select(COND_EXCLUSION,[
#MORTGAGES
                                'A. OVERDRAFT','B. INTERNATIONAL','C. DPD > 90','D. MOB <= 6','E. TOTAL OUTSTANDING = 1', 'F. POST_EXCLUSION',
#PL-STL
                                 'A. OVERDRAFT','B. INTERNATIONAL','C. ALREADY BAD','D. DPD > 90','E. MOB <= 6', 'F. NO_SAL_LAST_12MON','G. TOTAL OUTSTANDING = 1','H. POST_EXCLUSION'])
    
    
    

    logger.info(f" Conditions applied successfuly for Bscore exclusion flag for Islamic  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")
    
    
    logger.info(f" PiT mob 6 flag for  {portfolio_code} portfolio has: {df['ISLAMIC_BSCORE_EXCLUSION_REASON'].nunique()} unique exclusions i.e. {df['ISLAMIC_BSCORE_EXCLUSION_REASON'].unique()}"  )
    logger.debug("#######################################################################################")
    
    return df

#.....................................................................................................................................................................#


#PORTFOILIO WISE PIT_MOB_LT_6_EXCLUSION_REASON 
     
def pit_mob_6(df,portfolio_code):
    
    
    ''' This function outputs the dataframe after applying PiT mob 6 exclusions for non-islamic.
    
        input params:

        @ df - The required combined data after preprocessing and before calibration.

    '''
    
    
    logger.info(f"Defining conditions for PiT mob 6 flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")    
    
#AUTO LOANS
    COND_PIT_MOB_LT_6_EXCLUSION=[
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.EXCLUSION_CODES.str.contains('AL_02',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.EXCLUSION_CODES.str.contains('AL_04',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.EXCLUSION_CODES.str.contains('AL_03',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.EXCLUSION_CODES.str.contains('AL_05',regex=False))),                
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"]))),
#CREDIT CARDS 
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.INTERNATIONAL.isin([1]))) ,
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.EXCLUSION_CODES.str.contains('CC_01',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.EXCLUSION_CODES.str.contains('CC_08',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.EXCLUSION_CODES.str.contains('CC_02',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.EXCLUSION_CODES.str.contains('CC_03',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"]))),
#MORTGAGES    
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.INTERNATIONAL.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_02',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_05',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_04',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"]))),
#NL
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.EXCLUSION_CODES.str.contains('NL_05',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.EXCLUSION_CODES.str.contains('NL_03',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.EXCLUSION_CODES.str.contains('NL_12',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["National Loans"]))),
#PL-NSTL
                #((df.PORTFOLIO_CODE.isin(["PL-NSTL"])) & (df.ISLAMIC.isin([1]))) ,
                #((df.PORTFOLIO_CODE.isin(["PL-NSTL"])) & (df.OVERDRAFT.isin([1]))),
                #((df.PORTFOLIO_CODE.isin(["PL-NSTL"])) & (df.INTERNATIONAL.isin([1]))),
                #((df.PORTFOLIO_CODE.isin(["PL-NSTL"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                #((df.PORTFOLIO_CODE.isin(["PL-NSTL"])) & (df.EXCLUSION_CODES.str.contains('M_02',regex=False))),
                #((df.PORTFOLIO_CODE.isin(["PL-NSTL"])) & (df.EXCLUSION_CODES.str.contains('M_05',regex=False))),
                #((df.PORTFOLIO_CODE.isin(["PL-NSTL"])) & (df.EXCLUSION_CODES.str.contains('M_04',regex=False))),
                #((df.PORTFOLIO_CODE.isin(["PL-NSTL"]))),
#PL-STL
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.INTERNATIONAL.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_02',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_04',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_07',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_06',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_03',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])))]
    
    
    logger.info(f"Conditions defined successfuly for PiT mob 6 flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")
    
    logger.info(f"Applying conditions to create PiT mob 6 flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################") 
    
    df['PIT_MOB_LT_6_EXCLUSION_REASON']=np.select(COND_PIT_MOB_LT_6_EXCLUSION,[
#AUTO LOANS
                         'A. ISLAMIC','B. OVERDRAFT','C. NOT_IN_INCLUSION_LIST','D. ALREADY BAD','E. DPD > 90','F. TOTAL OUTSTANDING = 1','G. POST_EXC_MOB <= 6','H. POST_EXC_MOB > 6',                                                                       
#CREDIT CARDS
                         'A. INTERNATIONAL','B. NOT_IN_INCLUSION_LIST','C. CLOSED','D. WRITE-OFF','E. DPD > 90','F. POST_EXC_MOB <= 6','G. POST_EXC_MOB > 6',
#MORTGAGES
                         'A. ISLAMIC','B. OVERDRAFT','C. INTERNATIONAL','D. NOT_IN_INCLUSION_LIST','E. DPD > 90','F. POST_EXC_MOB <= 6','G. TOTAL OUTSTANDING = 1', 'H. POST_EXC_MOB > 6',
#NL
                         'A. ISLAMIC','B. OVERDRAFT','C. NOT_IN_INCLUSION_LIST','D. DPD > 90','E. POST_EXC_MOB <= 6','F. TOTAL OUTSTANDING = 1', 'G. POST_EXC_MOB > 6',
#PL-NSTL
                         #'A. ISLAMIC','B. OVERDRAFT','C. INTERNATIONAL','D. NOT_IN_INCLUSION_LIST','E. DPD > 90','F. POST_EXC_MOB <= 6','G. TOTAL OUTSTANDING = 1', 'H. POST_EXC_MOB > 6',
#PL-STL
                         'A. ISLAMIC','B. OVERDRAFT','C. INTERNATIONAL','D. NOT_IN_INCLUSION_LIST','E. ALREADY BAD','F. DPD > 90','G. POST_EXC_MOB <= 6', 'H. NO_SAL_LAST_12MON','I. TOTAL OUTSTANDING = 1', 'J. POST_EXC_MOB > 6'])

    
    logger.info(f" Conditions applied successfuly for PiT mob 6 flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")
    
    
    logger.info(f" PiT mob 6 flag for  {portfolio_code} portfolio has: {df['PIT_MOB_LT_6_EXCLUSION_REASON'].nunique()} unique exclusions i.e. {df['PIT_MOB_LT_6_EXCLUSION_REASON'].unique()}"  )
    logger.debug("#######################################################################################")
    
    
    return df

#.....................................................................................................................................................................#

#PORTFOILIO WISE PIT_EXCLUSION_REASON 
    
def pit_model_exclusion(df,portfolio_code):
    
    ''' This function outputs the dataframe after applying  PiT model exclusions for non-islamic.
    
        input params:

        @ df - The required combined data after preprocessing and before calibration.

    '''
    
    logger.info(f"Defining conditions for PiT exclusion flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")

#AUTO LOANS 
    COND_PIT_EXCLUSION=[
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.EXCLUSION_CODES.str.contains('AL_02',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.EXCLUSION_CODES.str.contains('AL_04',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.EXCLUSION_CODES.str.contains('AL_05',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.EXCLUSION_CODES.str.contains('AL_03',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"]))),
#CREDIT CARDS 
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.INTERNATIONAL.isin([1]))) ,
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.EXCLUSION_CODES.str.contains('CC_01',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.EXCLUSION_CODES.str.contains('CC_08',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.EXCLUSION_CODES.str.contains('CC_02',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.EXCLUSION_CODES.str.contains('CC_03',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.EXCLUSION_CODES.str.contains('CC_09',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"]))),
#MORTGAGES    
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.INTERNATIONAL.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_02',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_05',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_04',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"]))),
#NL
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.EXCLUSION_CODES.str.contains('NL_05',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.EXCLUSION_CODES.str.contains('NL_03',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.EXCLUSION_CODES.str.contains('NL_12',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["National Loans"]))),

#PL-STL
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.INTERNATIONAL.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_02',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_04',regex=False))|(df.EXCLUSION_CODES.str.contains('PL-NSTL_02',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_07',regex=False))|(df.EXCLUSION_CODES.str.contains('PL-NSTL_07',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_03',regex=False))|(df['T_Out_NSTL<=1'].isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])))]

    

    logger.info(f"Conditions defined successfuly for PiT exclusion flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")
    
    logger.info(f"Applying conditions to create PiT exclusion flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")     
    

    df['PIT_EXCLUSION_REASON']=np.select(COND_PIT_EXCLUSION,[
#AUTO LOANS
                         'A. ISLAMIC','B. OVERDRAFT','C. NOT_IN_INCLUSION_LIST','D. ALREADY BAD','E. DPD > 90','F. MOB <= 6','G. TOTAL OUTSTANDING = 1','H. POST_EXCLUSION',                                                                       
#CREDIT CARDS
                         'A. INTERNATIONAL','B. NOT_IN_INCLUSION_LIST','C. CLOSED','D. WRITE-OFF','E. DPD > 90','F. MOB <= 6','G. INACTIVE CARDS','H. POST_EXCLUSION',
#MORTGAGES
                         'A. ISLAMIC','B. OVERDRAFT','C. INTERNATIONAL','D. NOT_IN_INCLUSION_LIST','E. DPD > 90','F. MOB <= 6','G. TOTAL OUTSTANDING = 1', 'H. POST_EXCLUSION',
#NL
                         'A. ISLAMIC','B. OVERDRAFT','C. NOT_IN_INCLUSION_LIST','D. DPD > 90','E. MOB <= 6','F. TOTAL OUTSTANDING = 1', 'G. POST_EXCLUSION',
#PL-STL
                         'A. ISLAMIC','B. OVERDRAFT','C. INTERNATIONAL','D. NOT_IN_INCLUSION_LIST','E. ALREADY BAD','F. DPD > 90','G. MOB <= 6','H. TOTAL OUTSTANDING = 1','I. POST_EXCLUSION'])


    logger.info(f" Conditions applied successfuly for PiT exclusion flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")
    
    
    logger.info(f" PiT exclusion flag for  {portfolio_code} portfolio has: {df['PIT_EXCLUSION_REASON'].nunique()} unique exclusions i.e. {df['PIT_EXCLUSION_REASON'].unique()}"  )
    logger.debug("#######################################################################################")
    
    
    return df
#.....................................................................................................................................................................#


def mis_exclusion(df,portfolio_code):
    
    
    ''' This function outputs the dataframe after applying  MIS exclusions for non-islamic.
    
        input params:

        @ df - The required combined data after preprocessing and before calibration.

    '''    
    
    logger.info(f"Defining conditions for Mis exclusoin flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")

#AUTO LOANS
    COND_MIS_EXCLUSION=[
          ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.ISLAMIC.isin([1]))) ,
          ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.OVERDRAFT.isin([1]))),
          ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
          ((df.PORTFOLIO_CODE.isin(["Auto Loans"]))),
#CREDIT CARDS
          ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.INTERNATIONAL.isin([1]))) ,
          ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
          ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.EXCLUSION_CODES.str.contains('CC_01',regex=False))),
          ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.EXCLUSION_CODES.str.contains('CC_08',regex=False))),          
          ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.DUBAI_FIRST.isin([1]))),
          ((df.PORTFOLIO_CODE.isin(["Credit Cards"]))),
#MORTGAGES 
        
        #FOR MORTGAGES IN MIS RECON WE NEED TO REMOVE SOME PRODUCT IDs (4102, 21055, 21064, 21065, 3042) as SHARED BY DEV. TEAM
          ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.ISLAMIC.isin([1]))) ,
          ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.OVERDRAFT.isin([1]))),
          ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.INTERNATIONAL.isin([1]))),
          ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),        
          ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"]))),
#NL          
          ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.ISLAMIC.isin([1]))) ,
          ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.OVERDRAFT.isin([1]))),
          ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),         
          ((df.PORTFOLIO_CODE.isin(["National Loans"]))),

#PL-STL         
          ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.ISLAMIC.isin([1]))) ,
          ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.OVERDRAFT.isin([1]))),
          ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.INTERNATIONAL.isin([1]))),
          ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
          ((df.PORTFOLIO_CODE.isin(["Personal Loans"])))]                   

    
    logger.info(f"Conditions defined successfuly for Mis exclusion flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")
    
    logger.info(f"Applying conditions to create Mis exclusion flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")    

    df['MIS_RECON_EXCLUSION_REASON']= np.select(COND_MIS_EXCLUSION,[
#AUTO LOANS    
                                    'A. ISLAMIC','B. OVERDRAFT','C. NOT_IN_INCLUSION_LIST','D. POST_EXCLUSION',
#CREDIT CARDS
                                    'A. INTERNATIONAL','B. NOT_IN_INCLUSION_LIST','C. CLOSED','D. WRITE-OFF','E. DUBAI_FIRST','F. POST_EXCLUSION',
#MORTGAGES
                                    'A. ISLAMIC','B. OVERDRAFT','C. INTERNATIONAL','D. NOT_IN_INCLUSION_LIST','E. POST_EXCLUSION',
#NL
                                    'A. ISLAMIC','B. OVERDRAFT','C. NOT_IN_INCLUSION_LIST','D. POST_EXCLUSION',
#PL-NSTL
                                    #'A. ISLAMIC','B. OVERDRAFT','C. INTERNATIONAL','D. NOT_IN_INCLUSION_LIST','E. POST_EXCLUSION',
#PL-STL
                                    'A. ISLAMIC','B. OVERDRAFT','C. INTERNATIONAL','D. NOT_IN_INCLUSION_LIST','E. POST_EXCLUSION'])
    
    
    
    logger.info(f" Conditions applied successfuly for Mis exclusion flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")
    
    
    logger.info(f" Mis exclusion flag for  {portfolio_code} portfolio has: {df['MIS_RECON_EXCLUSION_REASON'].nunique()} unique exclusions i.e. {df['MIS_RECON_EXCLUSION_REASON'].unique()}"  )
    logger.debug("#######################################################################################")    

    return df
#.....................................................................................................................................................................#


def bscore_model_exclusion(df,portfolio_code):
    
    
    ''' This function outputs the dataframe after applying  Bscore exclusions for non-islamic.
    
        input params:

        @ df - The required combined data after preprocessing and before calibration.

    '''
    
    
    logger.info(f"Defining conditions for Bscore exclusion flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")

#AUTO LOANS 
    COND_EXCLUSION=[
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.EXCLUSION_CODES.str.contains('AL_02',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.EXCLUSION_CODES.str.contains('AL_04',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.EXCLUSION_CODES.str.contains('AL_05',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.EXCLUSION_CODES.str.contains('AL_03',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Auto Loans"]))),
#CREDIT CARDS 
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.INTERNATIONAL.isin([1]))) ,
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.EXCLUSION_CODES.str.contains('CC_01',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.EXCLUSION_CODES.str.contains('CC_08',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.EXCLUSION_CODES.str.contains('CC_02',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.EXCLUSION_CODES.str.contains('CC_03',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.EXCLUSION_CODES.str.contains('CC_09',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Credit Cards"]))),
#MORTGAGES    
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.INTERNATIONAL.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_02',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_05',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.EXCLUSION_CODES.str.contains('M_04',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"]))),
#NL
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.EXCLUSION_CODES.str.contains('NL_05',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.EXCLUSION_CODES.str.contains('NL_03',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.EXCLUSION_CODES.str.contains('NL_12',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["National Loans"]))),

#PL-STL
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.ISLAMIC.isin([1]))) ,
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.OVERDRAFT.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.INTERNATIONAL.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.NOT_IN_INCLUSION_LIST.isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_02',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_04',regex=False))|(df.EXCLUSION_CODES.str.contains('PL-NSTL_02',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_07',regex=False))|(df.EXCLUSION_CODES.str.contains('PL-NSTL_07',regex=False))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.EXCLUSION_CODES.str.contains('PL-STL_03',regex=False))|(df['T_Out_NSTL<=1'].isin([1]))),
                ((df.PORTFOLIO_CODE.isin(["Personal Loans"])))]

    
    logger.info(f"Conditions defined successfuly for Bscore exclusion flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")
    
    logger.info(f"Applying conditions to create Bscore exclusion for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")    

#ASSIGNING BSCORE_EXCLUSION_REASON TO ALL THE BSCORE DATASETS
    df['BSCORE_EXCLUSION_REASON']=np.select(COND_EXCLUSION,[
#AUTO LOANS
                                'A. ISLAMIC','B. OVERDRAFT','C. NOT_IN_INCLUSION_LIST','D. ALREADY BAD','E. DPD > 90','F. MOB <= 6','G. TOTAL OUTSTANDING = 1','H. POST_EXCLUSION',                                                                       
#CREDIT CARDS
                                'A. INTERNATIONAL','B. NOT_IN_INCLUSION_LIST','C. CLOSED','D. WRITE-OFF','E. DPD > 90','F. MOB <= 6','G. INACTIVE CARDS','H. POST_EXCLUSION',
#MORTGAGES
                                'A. ISLAMIC','B. OVERDRAFT','C. INTERNATIONAL','D. NOT_IN_INCLUSION_LIST','E. DPD > 90','F. MOB <= 6','G. TOTAL OUTSTANDING = 1', 'H. POST_EXCLUSION',
#NL
                                'A. ISLAMIC','B. OVERDRAFT','C. NOT_IN_INCLUSION_LIST','D. DPD > 90','E. MOB <= 6','F. TOTAL OUTSTANDING = 1', 'G. POST_EXCLUSION',
#PL-STL
                                 'A. ISLAMIC','B. OVERDRAFT','C. INTERNATIONAL','D. NOT_IN_INCLUSION_LIST','E. ALREADY BAD','F. DPD > 90','G. MOB <= 6','H. TOTAL OUTSTANDING = 1','I. POST_EXCLUSION'])

    logger.info(f" Conditions applied successfuly for Bscore flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")
    
    
    logger.info(f" Bscore exclusion flag for  {portfolio_code} portfolio has: {df['BSCORE_EXCLUSION_REASON'].nunique()} unique exclusions i.e. {df['BSCORE_EXCLUSION_REASON'].unique()}"  )
    logger.debug("#######################################################################################")
    
    return df

#.....................................................................................................................................................................#

def policy_exclusions(df,portfolio_code):   
    
    ''' This function outputs the dataframe after creating Exclusion flags.
    
        input params:

        @ df - The required combined data after preprocessing and before calibration.

    '''     
    
    logger.info(f"Defining conditions for Policy exclusion flags flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")
    
    
    logger.info(f"Defining Islamic for Policy exclusion flags flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")

#CONDITIONS FOR ISLAMIC FLAG
    CONDLIST_ISLAMIC = [((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.ACCOUNT_ID.str.slice(0, 5).isin(["ASLLD","ISBLD", "KIBLD"]))) , 
                    ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.ACCOUNT_ID.str.slice(0, 5).isin(["ASLLD","ISBLD", "KIBLD"]))),
                    ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.ACCOUNT_ID.str.slice(0, 5).isin(["ASLLD","ISBLD", "KIBLD"]))),
                    #((df.PORTFOLIO_CODE.isin(["PL-NSTL"])) & (df.ACCOUNT_ID.str.slice(0, 5).isin(["ASLLD","ISBLD", "KIBLD"]))),
                    ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.ACCOUNT_ID.str.slice(0, 5).isin(["ASLLD","ISBLD", "KIBLD"]))),
                    ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.ACCOUNT_ID.str.slice(0, 5).isin(["ASLLD","ISBLD", "KIBLD"])))]
#CREATING ISLAMIC FLAG
    df['ISLAMIC'] = np.select(CONDLIST_ISLAMIC, [1,1,1,1,1], 0)
    
    logger.info(f" Islamic distribution for  {portfolio_code} portfolio is {df['ISLAMIC'].value_counts()}")
    logger.debug("#######################################################################################")
    
    
    logger.info(f"Defining Overdraft for Policy exclusion flags flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")

#CONDITIONS FOR OVERDRAFT FLAG   
    CONDLIST_OVERDRAFT = [((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (df.ACCOUNT_ID.str.slice(3, 4).isin(["0","1","2","3","4","5","6","7","8","9"]))) , 
                      ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (df.ACCOUNT_ID.str.slice(3, 4).isin(["0","1","2","3","4","5","6","7","8","9"]))),
                      ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.ACCOUNT_ID.str.slice(3, 4).isin(["0","1","2","3","4","5","6","7","8","9"]))),
                      #((df.PORTFOLIO_CODE.isin(["PL-NSTL"])) & (df.ACCOUNT_ID.str.slice(3, 4).isin(["0","1","2","3","4","5","6","7","8","9"]))),
                      ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.ACCOUNT_ID.str.slice(3, 4).isin(["0","1","2","3","4","5","6","7","8","9"])))]
#CREATING OVERDRAFT FLAG
    df['OVERDRAFT'] = np.select(CONDLIST_OVERDRAFT, [1,1,1,1], 0)
    
    
    logger.info(f" Overdraft distribution for  {portfolio_code} portfolio is {df['OVERDRAFT'].value_counts()}")
    logger.debug("#######################################################################################")
    
    
    logger.info(f"Defining International for Policy exclusion flags flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")
    
#CONDITIONS FOR INTERNATIONAL FLAG     
    CONDLIST_INTERNATIONAl = [((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (df.ACCOUNT_ID.str.slice(0, 5).isin(["KSCMG","KIBMG", "KIBLD"]))),
                              #((df.PORTFOLIO_CODE.isin(["PL-NSTL"])) & (df.ACCOUNT_ID.str.slice(0, 5).isin(["KSCMG","KIBMG", "KIBLD"]))),
                              ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (df.ACCOUNT_ID.str.slice(0, 5).isin(["KSCMG","KIBMG", "KIBLD"]))),
                              ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.ACCOUNT_ID.str.slice(1, 7).isin(["000721", "000717","000719"])))]
#CREATING INTERNATIONAL FLAG
    df['INTERNATIONAL'] = np.select(CONDLIST_INTERNATIONAl, [1,1,1], 0)
    
    
    logger.info(f" International distribution for  {portfolio_code} portfolio is {df['INTERNATIONAL'].value_counts()}")
    logger.debug("#######################################################################################")
    
    


#CONDITIONS FOR DUBAI_FIRST FLAG 
    CONDLIST_DUBAI_FIRST = [((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (df.ACCOUNT_ID.str.slice(1, 8).isin(["0007137"])))]
#CREATING DUBAI_FIRST FLAG


    df['DUBAI_FIRST'] = np.select(CONDLIST_DUBAI_FIRST, [1], 0)
    
    logger.info(f"Defining conditions for Exclusion flag for  {portfolio_code} portfolio")
    logger.debug("#######################################################################################")

#CONDITIONS FOR NOT_IN_INCLUSION_LIST FLAG-(CATCH ALL LOGIC)
    CONDLIST_NOT_IN_INCLUSION_LIST = [((df.PORTFOLIO_CODE.isin(["Auto Loans"])) & (-df.ACCOUNT_ID.str.slice(0, 5).isin(["BNKMG","BNKLD","BNKAA"]))) , 
                                  ((df.PORTFOLIO_CODE.isin(["National Loans"])) & (-df.ACCOUNT_ID.str.slice(0, 5).isin(["BNKMG","BNKLD","BNKAA"]))),
                                  ((df.PORTFOLIO_CODE.isin(["Personal Loans"])) & (-df.ACCOUNT_ID.str.slice(0, 5).isin(["BNKMG","BNKLD","BNKAA"]))),
                                  #((df.PORTFOLIO_CODE.isin(["PL-NSTL"])) & (-df.ACCOUNT_ID.str.slice(0, 5).isin(["BNKMG","BNKLD"]))),
                                  ((df.PORTFOLIO_CODE.isin(["Mortgage Loans"])) & (-df.ACCOUNT_ID.str.slice(0, 5).isin(["BNKMG","BNKLD","BNKAA"]))),
                                  ((df.PORTFOLIO_CODE.isin(["Credit Cards"])) & (-df.ACCOUNT_ID.str.slice(1, 7).isin(["000713", "000714"])))]
#CREATING NOT_IN_INCLUSION_LIST-(CATCH ALL LOGIC)
    df['NOT_IN_INCLUSION_LIST'] = np.select(CONDLIST_NOT_IN_INCLUSION_LIST, [1,1,1,1,1], 0)
    
    return df

#.....................................................................................................................................................................#


def calibration_func(df,is_islamic,portfolio_code,model_segment_filters):
    
    ''' This function outputs the dataframe after applying recalibration equation based on given portfolio and scoremodel code.
    
        input params:

        @ df - The required combined data after preprocessing and applying exclusions.
        @ is_islamic - This takes bool variable i.e 1 for islamic portfolios & 0 for rest.
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ model_segment_filters - This takes list value for the corresponding scorecards for the given portfolio.
        
    '''
    
    logger.info(f"Applying Recalibration equation on the previous existing non calibrated scores function starts for {portfolio_code} portfolio")
    logger.debug("#######################################################################################")
    
    if(is_islamic!=1 ):
        
        df=df[df['SCORE_MODEL_CODE'].isin(model_segment_filters) & df['BSCORE_EXCLUSION_REASON'].str.contains('POST_EXCLUSION',regex=False)]
        df['SCORE_DATE'] = pd.to_datetime(df['SCORE_DATE'])
        
        if (portfolio_code=='mort'):
            df1=df[df['SCORE_DATE']<='2023-06-30']
            df2=df[df['SCORE_DATE']>'2023-06-30']
         
            logger.debug(f"The non islamic input post exclusion data has {df.shape[0]} rows and {df.shape[1]} columns.")
            logger.debug("We are dividing the whole data in 2 parts one df2 from where correct recalibrated scores are coming and second df1 upto where the old scores are used.")
            logger.debug("#######################################################################################")
        
        else:
            
            df1=df[df['SCORE_DATE']<='2023-03-31']
            df2=df[df['SCORE_DATE']>'2023-03-31']
            
            logger.debug(f"The non islamic input post exclusion data has {df.shape[0]} rows and {df.shape[1]} columns.")
            logger.debug("We are dividing the whole data in 2 parts one df2 from where correct recalibrated scores are coming and second df1 upto where the old scores are used.")
            logger.debug("#######################################################################################")

            
        INTERCEPT = 500
        SLOPE = 20/np.log(2)
        
        if(portfolio_code=='pl'):
            
            logger.info(f"Applying recalibration equation on non recalibrated scores for {portfolio_code} portfolio starts")
            
            df1['BSCORE_ORIGINAL']=np.where(df1['BSCORE_ORIGINAL'].isna(),df1['BSCORE_ORIGINAL_NEW'],df1['BSCORE_ORIGINAL'])
            df1['SCORE_VALUE']=np.where(df1['SCORE_MODEL_CODE'] == 'PB_UAE_BSCR04',round(df1['BSCORE_ORIGINAL']*0.7647 +140.37491),round(df1['BSCORE_ORIGINAL']*0.8273 +90.22461))
            df1['A'] = (df1['SCORE_VALUE'] - INTERCEPT)/SLOPE
            df1['PRED_BR'] = 1/(1+np.exp(df1['A']))
            df=pd.concat([df1,df2],ignore_index=True)
            
            logger.debug(f"The non islamic input post exclusion data after applying recalibration has {df.shape[0]} rows and {df.shape[1]} columns.")
            logger.info(f"Applying recalibration equation on non recalibrated scores for {portfolio_code} portfolio ends")
            logger.debug("#######################################################################################")
            
            
        elif(portfolio_code=='al'):
            
            logger.info(f"Applying recalibration equation on non recalibrated scores for {portfolio_code} portfolio starts")
            
            df1['BSCORE_ORIGINAL']=np.where(df1['BSCORE_ORIGINAL'].isna(),df1['BSCORE_ORIGINAL_NEW'],df1['BSCORE_ORIGINAL'])
            df1['SCORE_VALUE']=np.where(df1['SCORE_MODEL_CODE'] == 'PB_UAE_BSCR08',round(df1['BSCORE_ORIGINAL']*0.8868 +70.17865),(np.where(df1['SCORE_MODEL_CODE'] == 'PB_UAE_BSCR07',round(df1['BSCORE_ORIGINAL']*0.6494 +223.07629),round(df1['BSCORE_ORIGINAL']*0.6892 +182.63231)) ))
            df1['A'] = (df1['SCORE_VALUE'] - INTERCEPT)/SLOPE
            df1['PRED_BR'] = 1/(1+np.exp(df1['A']))
            df=pd.concat([df1,df2],ignore_index=True)
            
            logger.debug(f"The non islamic input post exclusion data after applying recalibration has {df.shape[0]} rows and {df.shape[1]} columns.")
            logger.info(f"Applying recalibration equation on non recalibrated scores for {portfolio_code} portfolio ends")
            logger.debug("#######################################################################################")            
            
            
        elif(portfolio_code=='nl'):
            
            logger.info(f"Applying recalibration equation on non recalibrated scores for {portfolio_code} portfolio starts")
            
            df1['BSCORE_ORIGINAL']=np.where(df1['BSCORE_ORIGINAL'].isna(),df1['BSCORE_ORIGINAL_NEW'],df1['BSCORE_ORIGINAL'])
            df1['SCORE_VALUE']=np.where(df1['SCORE_MODEL_CODE'] == 'PB_UAE_BSCR12',round(df1['BSCORE_ORIGINAL']*0.7857 +143.53188),round(df1['BSCORE_ORIGINAL']*1.609 -332.47386))
            df1['A'] = (df1['SCORE_VALUE'] - INTERCEPT)/SLOPE
            df1['PRED_BR'] = 1/(1+np.exp(df1['A']))
            df=pd.concat([df1,df2],ignore_index=True)
            
            logger.debug(f"The non islamic input post exclusion data after applying recalibration has {df.shape[0]} rows and {df.shape[1]} columns.")
            logger.info(f"Applying recalibration equation on non recalibrated scores for {portfolio_code} portfolio ends")
            logger.debug("#######################################################################################")


            
        elif(portfolio_code=='mort'):
            
            logger.info(f"Applying recalibration equation on non recalibrated scores for {portfolio_code} portfolio starts")
            
            df1['SCORE_VALUE']=round(df1['BSCORE_ORIGINAL']*0.5377 +265.5356)
            df1['A'] = (df1['SCORE_VALUE'] - INTERCEPT)/SLOPE
            df1['PRED_BR'] = 1/(1+np.exp(df1['A']))
            df=pd.concat([df1,df2],ignore_index=True)
            
            logger.debug(f"The non islamic input post exclusion data after applying recalibration has {df.shape[0]} rows and {df.shape[1]} columns.")
            logger.info(f"Applying recalibration equation on non recalibrated scores for {portfolio_code} portfolio ends")
            logger.debug("#######################################################################################")

            
            
    else:
        
        df=df[df['SCORE_MODEL_CODE'].isin(model_segment_filters) 
                                  & df['is_islamic_BSCORE_EXCLUSION_REASON'].str.contains('POST_EXCLUSION',regex=False) & 
                                  (~B_SCORE_DATA['is_islamic_BSCORE_EXCLUSION_REASON'].isin([0,"0"]))]

        df['SCORE_DATE'] = pd.to_datetime(df['SCORE_DATE'])

        if (portfolio_code=='mort'):
            df1=df[df['SCORE_DATE']<='2023-06-30']
            df2=df[df['SCORE_DATE']>'2023-06-30']
            
            logger.debug(f"The islamic input post exclusion data has {df.shape[0]} rows and {df.shape[1]} columns.")
            logger.debug("We are dividing the whole data in 2 parts one df2 from where correct recalibrated scores are coming and second df1 upto where the old scores are used.")
            logger.debug("#######################################################################################")

            
        else:
            
            df1=df[df['SCORE_DATE']<='2023-03-31']
            df2=df[df['SCORE_DATE']>'2023-03-31']
            
            logger.debug(f"The islamic input post exclusion data has {df.shape[0]} rows and {df.shape[1]} columns.")
            logger.debug("We are dividing the whole data in 2 parts one df2 from where correct recalibrated scores are coming and second df1 upto where the old scores are used.")
            logger.debug("#######################################################################################")

    
        if(portfolio_code=='pl'):
            
            logger.info(f"Applying recalibration equation on non recalibrated scores for islamic {portfolio_code} portfolio starts")
            
            df1['BSCORE_ORIGINAL']=np.where(df1['BSCORE_ORIGINAL'].isna(),df1['BSCORE_ORIGINAL_NEW'],df1['BSCORE_ORIGINAL'])
            df1['SCORE_VALUE']=np.where(df1['SCORE_MODEL_CODE'] == 'PB_UAE_BSCR04',round(df1['BSCORE_ORIGINAL']*0.7647 +140.37491),round(df1['BSCORE_ORIGINAL']*0.8273 +90.22461))

            df1['SCORE_VALUE']= np.where(df1['SCORE_MODEL_CODE'] == 'PB_UAE_BSCR04', df1['SCORE_VALUE'] , df1['SCORE_VALUE'] + 23 )
            df1['A'] = (df1['SCORE_VALUE'] - INTERCEPT)/SLOPE
            df1['PRED_BR'] = 1/(1+np.exp(df1['A']))
            df=pd.concat([df1,df2],ignore_index=True)
            
            logger.debug(f"The islamic input post exclusion data after applying recalibration has {df.shape[0]} rows and {df.shape[1]} columns.")
            logger.info(f"Applying recalibration equation on non recalibrated scores for islamic {portfolio_code} portfolio ends")
            logger.debug("#######################################################################################")



        elif(portfolio_code=='mort'):
            
            logger.info(f"Applying recalibration equation on non recalibrated scores for islamic {portfolio_code} portfolio starts")
            
            df1['SCORE_VALUE']=round(df1['BSCORE_ORIGINAL']*0.5377 +265.5356)
            df1['SCORE_VALUE']= df1['SCORE_VALUE']-16
            df1['A'] = (df1['SCORE_VALUE'] - INTERCEPT)/SLOPE
            df1['PRED_BR'] = 1/(1+np.exp(df1['A']))
            df=pd.concat([df1,df2],ignore_index=True)
            
            logger.debug(f"The islamic input post exclusion data after applying recalibration has {df.shape[0]} rows and {df.shape[1]} columns.")
            logger.info(f"Applying recalibration equation on non recalibrated scores for islamic {portfolio_code} portfolio ends")
            logger.debug("#######################################################################################")

    logger.info(f"Whole calibration_func function compiled and applying recalibration equation on the previous existing non calibrated scores function ends for {portfolio_code} portfolio.")
    logger.debug("#######################################################################################")
     
    return df

#.....................................................................................................................................................................#



def read_portfolio_data(portfolio_code,is_islamic):
#     ******************************************************************

    logger.info(f"B score data preprocessing steps for {portfolio_code} portfolio starts.")
    logger.debug("#######################################################################################")

    path=read_column_name(portfolio_code)
    path
#     ******************************************************************   
    logger.debug("Reading path for column name done")
    
    col_data=pd.read_csv(path)
    col_data
#     ******************************************************************    
    logger.debug(f"Column name data for {portfolio_code} portfolio is read having {col_data.shape[0]} rows.")
    
    mapper_sheet=portfolio_variable_mapper(portfolio_code)
    mapper_sheet
    
#     ******************************************************************   

    logger.debug(f"Reading SAS variable mapper excel file for {portfolio_code} portfolio is read having {mapper_sheet.shape[0]} rows.")
    
    dictionary={}
    for i in range(col_data.shape[1]):
        if pd.notna(mapper_sheet['Name_of_Variable in SAS dataset'][i]):
            dictionary[mapper_sheet['Variable_Name'][i]]=mapper_sheet['Name_of_Variable in SAS dataset'][i]
          
    col_data.rename(columns=(dictionary),inplace=True)
    
    logger.debug(f"Variables renamed according to SAS format. Total {col_data.shape[0]} variables present in {portfolio_code} portfolio.")
    logger.debug("#######################################################################################")

    
#     ****************************************************************** 


    filename_list=read_l3_datas(portfolio_code)[0]
    
#     ******************************************************************    
    logger.debug(f"Reading the initial draft of raw data files, total {len(filename_list)} files present.")


    l3_data_path=read_l3_datas(portfolio_code)[1]
    
#     ******************************************************************    
    
    combined_data=temprory_data(filename_list,l3_data_path,col_data)

#     ******************************************************************  
    logger.debug(f"All initial raw data files are read. It has {combined_data.shape[0]} rows and {combined_data.shape[1]} columns.")

    
    col_to_change=["SCORE_DATE"]
    Combined_data=change_to_Date_Time(combined_data,col_to_change)
    

#     ******************************************************************  
    
    Combined_data_F=remove_incomplete_per_data(Combined_data,"2021/12/31","2022/09/30")

#     ******************************************************************  
    logger.debug(f"Incomplete performance data filterd and removed. Now the data has {Combined_data_F.shape[0]} rows and {Combined_data_F.shape[1]} columns.")

    
    
    Dec_21_Updated_df=new_dec_2021_data(portfolio_code,col_data)

#     ********************************************************************
    logger.debug(f"DEC_21 complete performance data read having {Dec_21_Updated_df.shape[0]} rows and {Dec_21_Updated_df.shape[1]} columns.")

    filename_list_new=read_new_l3_datas(portfolio_code)[0]
    
    logger.debug(f"Reading the second draft of raw data files, total {len(filename_list_new)} files present.")

    
    l3_data_path_new = read_new_l3_datas(portfolio_code)[1]
    
#     ******************************************************************    
    
    combined_data_new = new_temprory_data(filename_list_new,l3_data_path_new,col_data)

#     ******************************************************************  

    logger.debug(f"All second draft raw data files are read. It has {combined_data_new.shape[0]} rows and {combined_data_new.shape[1]} columns.")


    combined_data_new=change_to_Date_Time(combined_data_new,col_to_change)
    

    Dec_21_Updated_df=change_to_Date_Time(Dec_21_Updated_df,col_to_change)
    
    
    

#     ******************************************************************  

    Combined_data=pd.concat([Combined_data_F,Dec_21_Updated_df,combined_data_new],ignore_index=True)
    
    logger.debug(f"All individual dataframes are concatenated.Final combined data has {Combined_data.shape[0]} rows and {Combined_data.shape[1]} columns.")
    logger.debug("#######################################################################################")


    Combined_data["month_score"]=Combined_data["SCORE_DATE"].dt.month

    Combined_data["YEAR_SCORE"]=Combined_data["SCORE_DATE"].dt.year

    Combined_data['QUARTER']="Q"+Combined_data["SCORE_DATE"].dt.quarter.astype('str') + Combined_data["SCORE_DATE"].dt.year.astype('str')
    Combined_data['QTR']=Combined_data["SCORE_DATE"].dt.year.astype('str')+"Q"+Combined_data["SCORE_DATE"].dt.quarter.astype('str') 

    

#     ****************************************************************** 




    
#     ****************************************************************** 

    if (is_islamic==0):
        
        logger.debug(f"Non islamic exclusions starts for {portfolio_code} portfolio.")
        
        Combined_data=policy_exclusions(Combined_data,portfolio_code)
        Combined_data=bscore_model_exclusion(Combined_data,portfolio_code)
        Combined_data=mis_exclusion(Combined_data,portfolio_code)
        Combined_data=pit_model_exclusion(Combined_data,portfolio_code)
        Combined_data=pit_mob_6(Combined_data,portfolio_code)
        
        logger.debug(f"Non islamic exclusions ends for {portfolio_code} portfolio.")
        logger.debug("#######################################################################################")
        
    else:
        
        logger.debug(f"Islamic exclusions starts for islamic {portfolio_code} portfolio.")

        Combined_data=is_islamic_bscore_model_exclusion(Combined_data,portfolio_code)
        Combined_data=is_islamic_mis_exclusion(Combined_data,portfolio_code)
        Combined_data=is_islamic_pit_model_exclusion(Combined_data,portfolio_code)
        Combined_data=is_islamic_pit_mob_6(Combined_data,portfolio_code)
        
        logger.debug(f"Islamic exclusions ends for islamic {portfolio_code} portfolio.")
        logger.debug("#######################################################################################")

        
#     ****************************************************************** 


    
    Model_Segment_Filters=portfolio_wise_score_model_code(portfolio_code)
    
    Combined_data=New_Original_Bscore(Combined_data,portfolio_code)

    Combined_data=calibration_func(Combined_data,is_islamic,portfolio_code,Model_Segment_Filters)
    
    path=save_and_read_intermediate_data(portfolio_code)
        
    Combined_data.to_csv(path,index=False)
    
    logger.info(f"Final combined intermediate data for {portfolio_code} portfolio is stored with {Combined_data.shape[0]} rows and {Combined_data.shape[1]} columns.")
    logger.info("Whole read_portfolio_data function read and compiled successfully and combined data returned.")
    logger.debug("#######################################################################################")
        
    return Combined_data
    
    
#.....................................................................................................................................................................#
    


def save_and_read_intermediate_data(portfolio_code,path_file):
    
    '''
        This function returns the path of the intermediate data for each portfolio.
        
        input params:
        
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        
        output params:
        
        It returns the path where the intermediate data is kept
    
    '''
    if portfolio_code=="pl":
        path=path_file
    elif portfolio_code=="al":
        path=path_file
    elif portfolio_code=="nl":
        path=path_file
    else:
        path=""
    #path=r"\\zebra\GIRM\Group Portfolio Analytics and Risk Methodology\Named folders\GBS COE\Sharat\BSCR_L3 data\pl\Intermediate data\INTERMEDIATE_DATA.csv"
    #path=path.format(portfolio_code)
    return path

#.....................................................................................................................................................................#

def find_validation_imm_omm_dates(current_qtr):
    year,quarter=map(int,current_qtr.split('Q'))
    last_month=quarter*3
    last_day=calendar.monthrange(year,last_month)[1]
    end_date=datetime(year,last_month,last_day)
    end_date=end_date.strftime("%Y-%m-%d")
    end_date_2=str(int(end_date[:4])-1)+end_date[4:]
    
    end_date_1=pd.to_datetime(end_date)
    current_date = end_date_1
    start_date = current_date - timedelta(days=12*30)
    start_date=start_date+pd.offsets.MonthEnd(0)    
    start_date=str(start_date)[:10]
    start_date_2=str(int(start_date[:4])-1)+start_date[4:]
    
    return start_date,end_date,start_date_2,end_date_2

#.....................................................................................................................................................................#

def score_model_wise_variable_lists_wt_rng(score_model_code,score_model_wise_variable_dict):
# def score_model_wise_variable_lists_wt_rng(score_model_code,path_file):
    
    '''
    This function outputs the list of score_model_code wise variable's value of their range and weights
    
    input params.
    @ score_model_code - This takes string variable and accounts for the scorecard whose variables is to be generated.i.e., "PB_UAE_BSCR04"
    
    output params.
    
    This will return 2 parameters, first the given score_model's variables name for their range value second for their weight value.
    
    '''
  
    
    RNG_VAR=score_model_code+"_RNG"
    WT_VAR=score_model_code+"_WT"
    
    return score_model_wise_variable_dict[RNG_VAR],score_model_wise_variable_dict[WT_VAR]
    
   
    # df = pd.read_csv(path_file)
    # RNG_VAR = list(set([col  for col in df.columns if col.endswith("_RNG")]))
    # WT_VAR =  list(set([col  for col in df.columns if col.endswith("_WT")]))
  
    # return RNG_VAR,WT_VAR
#.....................................................................................................................................................................#

def model_variable_full_form(score_model_code,model_variable_full_form):
    model_variable_full_form1 = model_variable_full_form
    model_variable_full_form={
    "PB_UAE_BSCR04":{
        "LN2704LNTRM_RNG":"Loans Term or Tenure", "LN2711_TOB_ACCT_UNADJ_RNG":"Time on Books", "LN2809PMTC1TO6PCTBALC1TO6_RNG":"Ratio of payment to total Due (EMI this month + arrear last month) in the last 6 months", "MOSSNCDLQGT0_LAG_1TO24_RNG":"Number of months since the account has been delinquent, starting with the cycle in which the account could have been current", "NATIONALITY_GROUP_RNG":"Nationality", "RATIO_PRIN_BAL_ORIG_LOAN_RNG":"Ratio of principal balance to original loan amount", "SLRY_AVGSLRY_1TO3_4TO9_RNG":"Ratio of average salary in 1-3 to 4-9 months", "SLRY_MAXNUMMISNG_RNG":"No of times salary is missing in 1-12 months", "SLRY_PCTTOTDUE_SLRY_1TO6_RNG":"Ratio of Total Due to Salary in last 6 months"        
    },
     
    "PB_UAE_BSCR05":{
        "LN2704LNTRM_RNG":"Loans Term or Tenure", "LN2767PCTDLQGT1C1TO3_RNG":"Pct. of months where delinquency is greater than 1 in last 3 months", "LN2807PMTC1TO3PCTBALC1TO3_RNG":"Ratio of payment to total Due (EMI this month + Arrears last month) in last 3 months", "LN2813PASTDUEC1PCTBALC1_RNG":"Past due as Pct. of Total Due", "LN2830MAXCMOSDLQGT1C1TO12_RNG":"Max Consecutive Months where delinquency is greater than 1 in last 12M", "NATIONALITY_GROUP_RNG":"Nationality", "RATIO_PRIN_BAL_ORIG_LOAN_RNG":"Ratio of principal balance to original loan amount", "SLRY_AVGSLRY_1TO3_4TO12_RNG":"Ratio of Avg. Salary 1-3 Months to 4 -12 Months", "SLRY_PCT_MISNG_OCCR_LST3M_RNG":"Pct. of Months’ Salary Missing in last 3 months", "SLRY_PCTTOTDUE_SLRY_1TO12_RNG":"Ratio of Total Due to Salary 1-12 Months"
    }, 
        
    "PB_UAE_BSCR12":{
        "EMPLOYER_RNG":"EMPLOYER_RNG" , "SECURED_UNSECURED_RNG":"SECURED_UNSECURED_RNG" , "TotConsInc_bal5_CASA_6m_RNG":"TotConsInc_bal5_CASA_6m_RNG" ,"TotConsInc_util_12m_RNG":"TotConsInc_util_12m_RNG" ,
          "max_dpd_6m_RNG":"max_dpd_6m_RNG","nbr_0_3m_RNG":"nbr_0_3m_RNG","nbr_SALARY_by_min_due_GT150_12m_RNG":"nbr_SALARY_by_min_due_GT150_12m_RNG",
          "nbr_times_pay_grt_min_due_l6_RNG":"nbr_times_pay_grt_min_due_l6_RNG"
    },
        
    "PB_UAE_BSCR13":{
        "COUNT_DELQ_INC_L12M_RNG":"Count of consecutive increase in delinquency in last 12 months" , "COUNT_PAY_L9M_RNG":"Count of number of Payments in last 9 months " , "CURR_DELQ_RNG":"Current Delinquency Range" , "CUST_AGE_RNG":"Customer Age", "MAX_DELQ_L9M_RNG":"Maximum Delinquency in last 9 months"
    }
    
    
    }
    
    return model_variable_full_form1
#.....................................................................................................................................................................#
    

def score_model_wise_implemented_ttc_pd(score_model_code,score_model_wise_imp_pd_dict):
    
    '''
    This function outputs implemented ttc-pd of the score_model.
    
    input params.
    @ score_model_code - This takes string variable and accounts for the scorecard whose variables is to be generated.i.e., "PB_UAE_BSCR04"
    
    output params.
    
    This will return implemented ttc-pd of the score_model in absolute number.
    
    '''
    
    # score_model_wise_imp_pd_dict= {
    
    # "PB_UAE_BSCR01":0.12,
    # "PB_UAE_BSCR02":0.23,
    # "PB_UAE_BSCR03":0.32,
    # "PB_UAE_BSCR04":0.0167,
    # "PB_UAE_BSCR05":0.3635,
    # "PB_UAE_BSCR07":0.0226,
    # "PB_UAE_BSCR08":0.0060,
    # "PB_UAE_BSCR09":0.2707,
    # "PB_UAE_BSCR10":0.0639,
    # "PB_UAE_BSCR12":0.0082,
    # "PB_UAE_BSCR13":0.3452  
    
    # }
    
    return score_model_wise_imp_pd_dict[score_model_code]

#.....................................................................................................................................................................#


def score_model_wise_benchmark_gini(score_model_code,score_model_wise_bm_gini_dict):
    
    '''
    This function outputs the benchmark_gini of the score_model.
    
    input params.
    @ score_model_code - This takes string variable and accounts for the scorecard whose variables is to be generated.i.e., "PB_UAE_BSCR04"
    
    output params.
    
    This will return the benchmark_gini of the score_model in absolute number.
    
    '''
    
    # score_model_wise_bm_gini_dict= {
    
    # "PB_UAE_BSCR01":0.12,
    # "PB_UAE_BSCR02":0.23,
    # "PB_UAE_BSCR03":0.32,
    # "PB_UAE_BSCR04":0.4950,
    # "PB_UAE_BSCR05":0.6380,
    # "PB_UAE_BSCR07":0.5425,
    # "PB_UAE_BSCR08":0.5770,
    # "PB_UAE_BSCR09":0.6178,
    # "PB_UAE_BSCR10":0.7504,
    # "PB_UAE_BSCR12":0.2862,
    # "PB_UAE_BSCR13":0.4163    
    
    # }
    
    return score_model_wise_bm_gini_dict[score_model_code]

# #.....................................................................................................................................................................#

def score_model_wise_ppt_graphs_title(score_model_code,score_model_wise_ppt_graphs_title):
    
    '''
    This function outputs the title of the score_model's monitoring charts.
    
    input params.
    @ score_model_code - This takes string variable and accounts for the scorecard whose variables is to be generated.i.e., "PB_UAE_BSCR04"
    
    output params.
    
    This will return the benchmark_gini of the score_model in absolute number.
    
    '''
    
    # score_model_wise_ppt_graphs_title= {
    
   
    # "PB_UAE_BSCR12":"National Loan B-Score Current Segment",
    # "PB_UAE_BSCR13":"National Loan B-Score Delinquent Segment"    
    
    # }
    
    return score_model_wise_ppt_graphs_title[score_model_code]

#.....................................................................................................................................................................#

def portfolio_wise_score_model_code(segment):
    
    '''
    This function outputs the list of portfolio wise score_model_codes.
    
    input params.
    @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
    
    output params.
    
    This will return list of score_model_codes according to given portfolio.
    
    '''
    
    # portfolio_wise_score_model_code = {
    #     "cc":["PB_UAE_BSCR01","PB_UAE_BSCR02","PB_UAE_BSCR03"],
    #     "pl":["PB_UAE_BSCR04","PB_UAE_BSCR05"],
    #     "al":["PB_UAE_BSCR07","PB_UAE_BSCR08","PB_UAE_BSCR09"],
    #     "mort":["PB_UAE_BSCR10"],
    #     "nl":["PB_UAE_BSCR12","PB_UAE_BSCR13"]
    # }
    
    # return portfolio_wise_score_model_code[portfolio_code]

    return segment

#.....................................................................................................................................................................#


def portfolio_wise_full_form(portfolio_code):
    
    '''
    This function outputs the list of portfolio wise score_model_codes.
    
    input params.
    @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
    
    output params.
    
    This will return list of score_model_codes according to given portfolio.
    
    '''
    
    portfolio_wise_full_form = {
        "cc":"Credit Cards",
        "pl":"Personal Loans",
        "al":"Auto Loans",
        "mort":"Mortgage Loans",
        "nl":"National Loans"
    }
    
    return portfolio_wise_full_form[portfolio_code]

#.....................................................................................................................................................................#

def call_all_variables_portfolio_wise(score_model_wise_variable_dict,segment1):
    
    '''
    This function outputs the list of total variables with range and weight values for the whole portfolio.
    
    input params.
    @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
    
    output params.
    
    This will return list of total variables with range and weight values for the whole given portfolio code.
    
    '''    
    
    temp_1=[]
    for x in portfolio_wise_score_model_code(segment=segment1): #(score_model_code,score_model_wise_variable_dict)
        temp=(score_model_wise_variable_lists_wt_rng(x,score_model_wise_variable_dict)[0])
        temp+=(score_model_wise_variable_lists_wt_rng(x,score_model_wise_variable_dict)[1])
        temp_1+=temp
    # temp_1.extend(score_model_wise_variable_lists_wt_rng(path_file)[0])
    # temp_1.extend(score_model_wise_variable_lists_wt_rng(path_file)[1])
        
    return temp_1


#.....................................................................................................................................................................#


def IMM_islamic_non_islamic_columns(portfolio_code,is_islamic,val_start_date,val_end_date,path_file,score_model_wise_variable_dict,columns1,segment,bm_year=2020):
    
    ''' This function outputs the dataframe required for calculating IMM.
    
        input params:
        
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ is_islamic - This takes bool variable i.e 1 for islamic portfolios & 0 for rest.
        @ val_start_date -This takes string variable in the format (YYYY-MM-DD) representing start date (month end date) for the current validation/monitoring exercise.
        @ val_end_date - This takes string variable in the format (YYYY-MM-DD) representing end date (month end date) for the current validation/monitoring exercise.
        @ bm_year - This takess integer values for the benchmark year, default value is 2020.
        
        
        output params:
        
        It returns a DataFrame for IMM.
    
    '''

    logger.debug("Input Monitoring data extraction starts")

    portfolio_code=portfolio_code.lower()
    is_islamic = int(is_islamic)
    bm_year = int(bm_year)
    
    var_portfolio = call_all_variables_portfolio_wise(score_model_wise_variable_dict,segment)
  
    columns = columns1
    columns = columns+var_portfolio
    # if portfolio_code=="pl":
        # columns.append("NSTL_FLAG")

    logger.debug("All required columns for input monitoring are read for the given portfolio")
        
    val_start_date = pd.to_datetime(val_start_date)
    val_end_date   = pd.to_datetime(val_end_date)
    
    columns += ['BSCORE_EXCLUSION_REASON']
    path = save_and_read_intermediate_data(portfolio_code,path_file) # score_model_wise_variable_dict
    df = pd.read_csv(path,usecols=columns) # 
    rows = df.shape[0]
    columns = df.shape[1]
    
    logger.debug(f"Bscore data for the input monitoring metrics is read with {rows} rows and {columns} columns.")

    df=df[df['BSCORE_EXCLUSION_REASON'].str.contains('POST_EXCLUSION',regex=False)]
    df['SCORE_DATE'] = pd.to_datetime(df['SCORE_DATE'])
    df['BENCHMARK']= np.select([df['SCORE_DATE'].dt.year==bm_year],[1],0)
    df['VALIDATION_SAMPLE']= np.select([(df['SCORE_DATE'] >=  val_start_date) & (df['SCORE_DATE'] <= val_end_date)],[1],0)   
    # df=df[df["SCORE_DATE"]<=val_end_date]

    logger.debug("Non islamic b score post exclusion data filterd and benchmark and validation samples are marked")
           
    rows = df.shape[0]
    columns = df.shape[1]
    
        
    logger.info(f"Whole IMM_islamic_non_islamic_columns function compiled & input monitoring dataframe with {rows} rows and {columns} columns stored and returned successfully")
    logger.debug("#######################################################################################")
      

    return df


#.....................................................................................................................................................................#

def OMM_islamic_non_islamic_columns(portfolio_code,is_islamic,val_start_date,val_end_date,path_file,score_model_wise_variable_dict,columns2,segment,bm_year=2020):
    
    ''' This function outputs the dataframe required for calculating OMM.
    
        input params:
        
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ is_islamic - This takes bool variable i.e 1 for islamic portfolios & 0 for rest.
        @ val_start_date -This takes string variable in the format (YYYY-MM-DD) representing start date (month end date) for the current validation/monitoring exercise.
        @ val_end_date - This takes string variable in the format (YYYY-MM-DD) representing end date (month end date) for the current validation/monitoring exercise.
        @ bm_year - This takess integer values for the benchmark year, default value is 2020.
        
        
        output params:
        
        It returns a DataFrame for OMM.
    
    '''
    
    logger.debug("Output Monitoring data extraction starts")

    
    portfolio_code=portfolio_code.lower()
    is_islamic = int(is_islamic)
    bm_year=int(bm_year)
    
    columns_gini = columns2   
    # print("****************columns2",columns2)

    var_portfolio = call_all_variables_portfolio_wise(score_model_wise_variable_dict,segment)
    columns_gini += var_portfolio
    

    logger.debug("All required columns for output monitoring are read for the given portfolio")

    
    val_start_date = pd.to_datetime(val_start_date)
    val_end_date   = pd.to_datetime(val_end_date)
    
   
    columns_gini += ['BSCORE_EXCLUSION_REASON']
    path = save_and_read_intermediate_data(portfolio_code,path_file)
    df = pd.read_csv(path,usecols=columns_gini)
    rows = df.shape[0]
    columns = df.shape[1]

    logger.debug(f"Bscore data for the output monitoring metrics is read with {rows} rows and {columns} columns.")

    df=df[df['BSCORE_EXCLUSION_REASON'].str.contains('POST_EXCLUSION',regex=False)]
    df['SCORE_DATE'] = pd.to_datetime(df['SCORE_DATE'])
    df['BENCHMARK']= np.select([df['SCORE_DATE'].dt.year==bm_year],[1],0)
    df['VALIDATION_SAMPLE']= np.select([(df['SCORE_DATE'] >=  val_start_date) & (df['SCORE_DATE'] <= val_end_date)],[1],0)  
    
    logger.debug("Non islamic b score post exclusion data filterd and benchmark and validation samples are marked")
            
   

        
    df['TARGET_12'].replace(99,0,inplace=True)
    
        
    rows = df.shape[0]
    columns = df.shape[1]
    
    
        
    logger.info(f"Whole IMM_islamic_non_islamic_columns function compiled & input monitoring dataframe with {rows} rows and {columns} columns stored and returned successfully")
    logger.debug("#######################################################################################")

        
    return df



#.....................................................................................................................................................................#


def calculate_psi_score_band(df,portfolio_code,score_model_code,current_qtr,segment1,deciles=10,len_rolling_window=4):
    
    
    ''' This function outputs the score-band level PSI table for last 4 quarters including current quarter based on the approach and portfolio selected.

    
        input params:
           
        @ df - The required data that is to be used for generating the PSI table.
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        @ current_qtr - This takes string variable for the current monitoring/validation quarter in format "YYYYQ". i.e. "2022Q4"
        @ deciles - This takes integer variable for making deciles for PSI. By deafult the value is 10.
        @ len_rolling_window - This takes integer variable for the length of the window we want in the rolling psi approach only. By default the value is 4.
        
        
        output params:
        
        It returns a DataFrame for PSI.
    
    '''

    logger.info(f"Score band level calculating Rolling PSI function starts for {portfolio_code} portfolio and {score_model_code} scoremodel.")
    

    if score_model_code not in portfolio_wise_score_model_code(segment=segment1):
        
        print("Please check the corresponding portfolio and score model code")
        
        logger.info("PSI Can't be calculated for given portfolio and scoremodel code")
        
        return 
        
    else:
        
    
        rows = df.shape[0]
        columns = df.shape[1]
            
        logger.debug(f"Bscore data for the input monitoring metrics has {rows} rows and {columns} columns.")
        logger.debug("#######################################################################################")
        
        logger.debug(f"Score band calculation with {deciles} deciles based on score-value starts")
    
        
        BSCORE_DATA_BM= df[df['BENCHMARK']==1].iloc[:,:]
        BSCORE_DATA_BM = BSCORE_DATA_BM.sort_values(['SCORE_MODEL_CODE','SCORE_VALUE'],ascending=False)
        BSCORE_DATA_RANK= BSCORE_DATA_BM[BSCORE_DATA_BM.SCORE_MODEL_CODE.isin([score_model_code])].loc[:,:]       
        BSCORE_DATA_RANK['rank'] = BSCORE_DATA_RANK['SCORE_VALUE'].rank()    
        BSCORE_DATA_RANK['DECILE'] = np.floor(BSCORE_DATA_RANK['rank']*(deciles)/(len(BSCORE_DATA_RANK['SCORE_VALUE'])+1))
        BSCORE_DATA_RANK= BSCORE_DATA_RANK.groupby(['SCORE_MODEL_CODE', 'DECILE']).agg(                   
                            BM_MAX_SCORE_VALUE = pd.NamedAgg(column='SCORE_VALUE', aggfunc='max'),
                            BM_MIN_SCORE_VALUE = pd.NamedAgg(column='SCORE_VALUE', aggfunc='min'),
                            COUNT = pd.NamedAgg(column='SCORE_VALUE', aggfunc='count')).reset_index(drop=False)
        
        logger.debug("Minimum and Maximum score value for each decile has been calculated")
        
        
        BSCORE_DATA_MODEL= df.loc[(df.SCORE_MODEL_CODE.isin([score_model_code]))]
        cutoff=BSCORE_DATA_RANK['BM_MIN_SCORE_VALUE'].values.astype('float64')
        cutoff[0]=-np.inf
        cutoff=np.append(cutoff,[np.inf])
        BSCORE_DATA_MODEL['SCORE_RANGE']=pd.cut(BSCORE_DATA_MODEL['SCORE_VALUE'],bins=cutoff,right=False)
        
        unq_bands=BSCORE_DATA_MODEL['SCORE_RANGE'].nunique()
        
        logger.debug(f"Score bands with {unq_bands} unique bands created for calculating PSI")
        logger.debug("#######################################################################################")
            
        VAL_MODEL=pd.DataFrame()
        
        #Calculating PSI with rolling window of 4 qtrs i.e. 1 year
        
        logger.debug(f"Calculation of rolling window for {len_rolling_window} qtrs with current quarter as {current_qtr} starts")
    
        def rolling_window(current_qtr):
            year = int(current_qtr[:4])
            qtr = int(current_qtr[-1])
            window = [current_qtr]
            for i in range(len_rolling_window-1):
                if qtr>1:
                    qtr=qtr-1
                else:
                    qtr=4
                    year=year-1
                window.append(str(year)+'Q'+str(qtr))
                window=sorted(window)
            return window
                
        temp_window = rolling_window(current_qtr) 
        listi=[]
        for qtrs in temp_window[::-1]:
            listi.append(rolling_window(qtrs))
        
        logger.debug("IMM data aggregation for calculating PSI for current validation/monitoring period starts")    
        
        for i in range(len(listi)):
            
            temp=BSCORE_DATA_MODEL[BSCORE_DATA_MODEL['QTR'].isin(listi[i])]
            temp['temp']=listi[i][3]
            
        
            QTR_MODEL= temp.groupby(['temp','SCORE_RANGE']).agg(
                        COUNT_SNAPSHOT = pd.NamedAgg(column='ACCOUNT_ID', aggfunc='count'),
                        ).reset_index(drop=False)
        
        
            VAL_MODEL=  pd.concat([VAL_MODEL,QTR_MODEL],ignore_index=True)
        
        
        DATA_MODEL_1=VAL_MODEL
        
        
        logger.debug(f"IMM data aggregation for calculating PSI for current validation/monitoring period ends for {len(listi)} unique quarters.")    
        logger.debug("#######################################################################################")
    
        logger.debug("Benchmark snapshot aggregation with current validation/monitoring data starts")
        
        T2=BSCORE_DATA_MODEL[BSCORE_DATA_MODEL['BENCHMARK']==1].loc[:,:]
        
        T2= T2.groupby(['SCORE_RANGE']).agg(
                        COUNT_BENCHMARK = pd.NamedAgg(column='ACCOUNT_ID', aggfunc='count'),
                        ).reset_index(drop=False)
    
        PSI_MODEL = pd.merge(DATA_MODEL_1[['temp','SCORE_RANGE', 'COUNT_SNAPSHOT']], T2[['SCORE_RANGE', 'COUNT_BENCHMARK']],
                     how='outer', on=['SCORE_RANGE'])
        
        PSI_MODEL['COUNT_BENCHMARK'].fillna(1,inplace=True)
        
        PSI_CALC = PSI_MODEL.groupby(['temp']).agg(
                        TOTAL_COUNT_SNAPSHOT = pd.NamedAgg(column='COUNT_SNAPSHOT', aggfunc='sum'),
                        TOTAL_COUNT_BENCHMARK = pd.NamedAgg(column='COUNT_BENCHMARK', aggfunc='sum')
                        ).reset_index(drop=False)
        
        PSI_1= pd.merge(PSI_MODEL[['temp', 'SCORE_RANGE', 'COUNT_SNAPSHOT', 'COUNT_BENCHMARK']], PSI_CALC[['temp','TOTAL_COUNT_SNAPSHOT', 'TOTAL_COUNT_BENCHMARK']],
                      left_on = "temp", right_on = "temp", how = "left", suffixes=('_',''))
        
        logger.debug("Benchmark snapshot aggregation with current validation/monitoring data ends")
        logger.debug("#######################################################################################")
          
        logger.debug("PSI calcultion starts")
    
        PSI_1= PSI_1.rename(columns = {'temp':'YYYYMM'})
        PSI_1['SCORE_CARD']= score_model_code
        PSI_1['VARIABLE_NAME']= "SCORE_BAND"
        PSI_1 = PSI_1.rename(columns = {'SCORE_RANGE':'VALUE'})
        PSI_1['CONCENTRATION_SNAPSHOT']= PSI_1['COUNT_SNAPSHOT']/PSI_1['TOTAL_COUNT_SNAPSHOT']
        PSI_1['CONCENTRATION_BENCHMARK']= PSI_1['COUNT_BENCHMARK']/PSI_1['TOTAL_COUNT_BENCHMARK']
        PSI_1['DIFFERENCE']= PSI_1['CONCENTRATION_SNAPSHOT']- PSI_1['CONCENTRATION_BENCHMARK']
        PSI_1['LOG_CONC']= np.log(PSI_1['CONCENTRATION_SNAPSHOT']/ PSI_1['CONCENTRATION_BENCHMARK'])
        PSI_1['INDIVIDUAL_PSI']= PSI_1['DIFFERENCE'] * PSI_1['LOG_CONC']
        
        PSI_2 = PSI_1.sort_values(['YYYYMM','VALUE'],ascending=False).groupby(['YYYYMM' ]).agg(
                        PSI = pd.NamedAgg(column='INDIVIDUAL_PSI', aggfunc='sum'),
                        HCI = pd.NamedAgg(column='CONCENTRATION_SNAPSHOT', aggfunc='max')
                        ).reset_index(drop=False)
        
        logger.debug("PSI calcultion ends")    
        
        PSI = pd.merge(PSI_1[['YYYYMM','SCORE_CARD', 'VALUE', 'COUNT_SNAPSHOT', 'COUNT_BENCHMARK', 'VARIABLE_NAME',           
               'CONCENTRATION_SNAPSHOT', 'CONCENTRATION_BENCHMARK', 'DIFFERENCE',
               'LOG_CONC', 'INDIVIDUAL_PSI']], PSI_2[['YYYYMM','PSI', 'HCI']],
                      left_on = "YYYYMM", right_on = "YYYYMM", how = "left", suffixes=('_',''))
        
        logger.debug("QoQ calculated PSI-HCI merged with respective quarters")
        logger.debug("#######################################################################################")

        
        col=["YYYYMM","VARIABLE_NAME",'VALUE', 'COUNT_SNAPSHOT', 'COUNT_BENCHMARK',           
               'CONCENTRATION_SNAPSHOT', 'CONCENTRATION_BENCHMARK', 'DIFFERENCE',
               'LOG_CONC', 'INDIVIDUAL_PSI','PSI', 'HCI']
        PSI=PSI[col]
        PSI = PSI.sort_values(['YYYYMM','VALUE'],ascending=True).reset_index(drop=True)
        
        unq_qtrs=PSI['YYYYMM'].nunique()
        
        logger.info(f"Whole psi_score_band function compiled and score band level Rolling PSI calculated for {portfolio_code} portfolio and {score_model_code} scoremodel for {unq_qtrs} qtrs and returned successfully.")    
        logger.debug("#######################################################################################")
        # display(PSI)
        return PSI


#.............................................................................................................................................................................................................................#    


def calculate_psi_score_point(df,portfolio_code,score_model_code,current_qtr,segment,len_rolling_window=4):
    
        
    ''' This function outputs the score-point level PSI table for last 4 quarters including current quarter based on the approach and portfolio selected.

    
        input params:
           
        @ df - The required data that is to be used for generating the PSI table.
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        @ current_qtr - This takes string variable for the current monitoring/validation quarter. i.e. "2022Q4"
        @ len_rolling_window - This takes integer variable for the length of the window we want in the rolling psi approach only. By default the value is 4.
        
        
        output params:
        
        It returns a DataFrame for PSI.
    
    '''

    logger.info(f"Score band level calculating Rolling PSI function starts for {portfolio_code} portfolio and {score_model_code} scoremodel.")
    

    if score_model_code not in portfolio_wise_score_model_code(segment=segment):
        
        print("Please check the corresponding portfolio and score model code")
        
        logger.info("PSI Can't be calculated for given portfolio and scoremodel code")
        
        return 

    else:
    
        rows = df.shape[0]
        columns = df.shape[1]
            
        logger.debug(f"Bscore data for the input monitoring metrics has {rows} rows and {columns} columns.")
    
    
        BSCORE_DATA_MODEL= df.loc[(df.SCORE_MODEL_CODE.isin([score_model_code]))]
        VAL_MODEL=pd.DataFrame()

        unq_scores=BSCORE_DATA_MODEL['SCORE_VALUE'].nunique()
        logger.debug(f"Score point with {unq_scores} unique scores created for calculating PSI")

        #Calculating PSI with rolling window of 4 qtrs i.e. 1 year
        logger.debug("#######################################################################################")

        
        logger.debug(f"Calculation of rolling window for {len_rolling_window} qtrs with current quarter as {current_qtr} starts")
        

        def rolling_window(current_qtr):
            year = int(current_qtr[:4])
            qtr = int(current_qtr[-1])
            window = [current_qtr]
            for i in range(len_rolling_window-1):
                if qtr>1:
                    qtr=qtr-1
                else:
                    qtr=4
                    year=year-1
                window.append(str(year)+'Q'+str(qtr))
                window=sorted(window)
            return window
                
        temp_window = rolling_window(current_qtr) 
        listi=[]
        for qtrs in temp_window[::-1]:
            listi.append(rolling_window(qtrs))

        logger.debug("IMM data aggregation for calculating PSI for current validation/monitoring period starts")    

        for i in range(len(listi)):
            
            temp=BSCORE_DATA_MODEL[BSCORE_DATA_MODEL['QTR'].isin(listi[i])]
            temp['temp']=listi[i][3]
            
        
            QTR_MODEL= temp.groupby(['temp','SCORE_VALUE']).agg(
                        COUNT_SNAPSHOT = pd.NamedAgg(column='ACCOUNT_ID', aggfunc='count'),
                        ).reset_index(drop=False)

            
            VAL_MODEL = pd.concat([VAL_MODEL,QTR_MODEL],ignore_index=True)
            
        DATA_MODEL_1=VAL_MODEL

        logger.debug(f"IMM data aggregation for calculating PSI for current validation/monitoring period ends for {len(listi)} unique quarters.")    
        logger.debug("#######################################################################################")
    
        logger.debug("Benchmark snapshot aggregation with current validation/monitoring data starts")        
        
        
        MORT_BM_DATA=pd.read_sas(r"\\zebra\GIRM\Group Portfolio Analytics and Risk Methodology\Validation\PB Models\2020\B Score Validation\Data & Codes\Loans\Mortgage\mort_tree_jan18_dec19.sas7bdat",encoding='')
        MORT_BM_DATA=MORT_BM_DATA[['AS_AT','new_contract_id','NODE_SCORE']]
        MORT_BM_DATA.columns= ['SCORE_DATE','ACCOUNT_ID','BSCORE_ORIGINAL']
        MORT_BM_DATA['SCORE_VALUE']=round(265.5356 + 0.5377 * MORT_BM_DATA['BSCORE_ORIGINAL'])
        
        MORT_BM_DATA=MORT_BM_DATA[MORT_BM_DATA.SCORE_DATE.dt.year==2018]
        T2=MORT_BM_DATA
        
        n_rows=MORT_BM_DATA.shape[0]
        n_cols=MORT_BM_DATA.shape[1]

        logger.debug(f"Benchmark data for 2018 read successfully with {n_rows} rows and {n_cols} columns.")        
        
        T2= T2.groupby(['SCORE_VALUE']).agg(
                        COUNT_BENCHMARK = pd.NamedAgg(column='ACCOUNT_ID', aggfunc='count'),
                        ).reset_index(drop=False)
    
        PSI_MODEL = pd.merge(DATA_MODEL_1[['temp','SCORE_VALUE', 'COUNT_SNAPSHOT']], T2[['SCORE_VALUE', 'COUNT_BENCHMARK']],
                     how='outer', on=['SCORE_VALUE'])
        
        PSI_MODEL['COUNT_BENCHMARK'].fillna(1,inplace=True)
        
        PSI_CALC = PSI_MODEL.groupby(['temp']).agg(
                        TOTAL_COUNT_SNAPSHOT = pd.NamedAgg(column='COUNT_SNAPSHOT', aggfunc='sum'),
                        TOTAL_COUNT_BENCHMARK = pd.NamedAgg(column='COUNT_BENCHMARK', aggfunc='sum')
                        ).reset_index(drop=False)
        
        PSI_1= pd.merge(PSI_MODEL[['temp', 'SCORE_VALUE', 'COUNT_SNAPSHOT', 'COUNT_BENCHMARK']], PSI_CALC[['temp','TOTAL_COUNT_SNAPSHOT', 'TOTAL_COUNT_BENCHMARK']],
                      left_on = "temp", right_on = "temp", how = "left", suffixes=('_',''))

        logger.debug("Benchmark snapshot aggregation with current validation/monitoring data ends")
        logger.debug("#######################################################################################")
          
        logger.debug("PSI calcultion starts")
        
        
        PSI_1= PSI_1.rename(columns = {'temp':'YYYYMM'})
        PSI_1['SCORE_CARD']= score_model_code
        PSI_1['VARIABLE_NAME']= "SCORE_BAND"
        PSI_1 = PSI_1.rename(columns = {'SCORE_VALUE':'VALUE'})
        PSI_1['CONCENTRATION_SNAPSHOT']= PSI_1['COUNT_SNAPSHOT']/PSI_1['TOTAL_COUNT_SNAPSHOT']
        PSI_1['CONCENTRATION_BENCHMARK']= PSI_1['COUNT_BENCHMARK']/PSI_1['TOTAL_COUNT_BENCHMARK']
        PSI_1['DIFFERENCE']= PSI_1['CONCENTRATION_SNAPSHOT']- PSI_1['CONCENTRATION_BENCHMARK']
        PSI_1['LOG_CONC']= np.log(PSI_1['CONCENTRATION_SNAPSHOT']/ PSI_1['CONCENTRATION_BENCHMARK'])
        PSI_1['INDIVIDUAL_PSI']= PSI_1['DIFFERENCE'] * PSI_1['LOG_CONC']
        
        PSI_2 = PSI_1.sort_values(['YYYYMM','VALUE'],ascending=False).groupby(['YYYYMM' ]).agg(
                        PSI = pd.NamedAgg(column='INDIVIDUAL_PSI', aggfunc='sum'),
                        HCI = pd.NamedAgg(column='CONCENTRATION_SNAPSHOT', aggfunc='max')
                        ).reset_index(drop=False)
        
        logger.debug("PSI calcultion ends")            
        
        PSI = pd.merge(PSI_1[['YYYYMM','SCORE_CARD', 'VALUE', 'COUNT_SNAPSHOT', 'COUNT_BENCHMARK', 'VARIABLE_NAME',           
               'CONCENTRATION_SNAPSHOT', 'CONCENTRATION_BENCHMARK', 'DIFFERENCE',
               'LOG_CONC', 'INDIVIDUAL_PSI']], PSI_2[['YYYYMM','PSI', 'HCI']],
                      left_on = "YYYYMM", right_on = "YYYYMM", how = "left", suffixes=('_',''))
        
        logger.debug("QoQ calculated PSI-HCI merged with respective quarters")
        logger.debug("#######################################################################################")

        col=["YYYYMM","VARIABLE_NAME",'VALUE', 'COUNT_SNAPSHOT', 'COUNT_BENCHMARK',           
               'CONCENTRATION_SNAPSHOT', 'CONCENTRATION_BENCHMARK', 'DIFFERENCE',
               'LOG_CONC', 'INDIVIDUAL_PSI','PSI', 'HCI']
        PSI=PSI[col]
        PSI = PSI.sort_values(['YYYYMM','VALUE'],ascending=True).reset_index(drop=True)
        
        unq_qtrs=PSI['YYYYMM'].nunique()
        
        logger.info(f"Whole psi_score_point function compiled and score point level Rolling PSI calculated for {portfolio_code} portfolio and {score_model_code} scoremodel for {unq_qtrs} qtrs and returned successfully.")    
        logger.debug("#######################################################################################")
        
        return PSI
    
    


#......................................................................................................................................................................#

   
def calculate_csi_score_band(df,portfolio_code,score_model_code,segment,path_file,score_model_wise_variable_dict):
    
    
        
    ''' This function outputs the score-band level CSI table. 
    
        input params:
           
        @ df - The required data that is to be used for generating the CSI table.
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"


        
        output params:
        
        It returns a DataFrame for CSI.
    
    '''  
    
    logger.info(f"Score band level CSI function starts for {portfolio_code} portfolio and {score_model_code} scoremodel.")
    

    if score_model_code not in portfolio_wise_score_model_code(segment=segment):
        
        print("Please check the corresponding portfolio and score model code")
        
        logger.info("CSI Can't be calculated for given portfolio and scoremodel code")
        
        return 
    
    else:

        variable_list_1=score_model_wise_variable_lists_wt_rng(score_model_code,score_model_wise_variable_dict)[0]
        variable_list_2=score_model_wise_variable_lists_wt_rng(score_model_code,score_model_wise_variable_dict)[1]
       
        logger.debug(f"Scorecard variables called. Total {len(variable_list_1)} variables present in this scorecard.")
        
        
        B_SCORE_DATA_BM= df[df['BENCHMARK']==1].iloc[:,:]
        B_SCORE_DATA_BM = B_SCORE_DATA_BM.sort_values(['SCORE_MODEL_CODE'],ascending=False)
        
        logger.debug("Benchmark period data filtered")
        
        logger.debug("#######################################################################################")

        
        BF = []
        for j in range(len(variable_list_1)):
            
            logger.debug(f"CSI calculation for {variable_list_1[j]} variable starts")
                    
            BSCORE_DATA_MODEL= df[df.SCORE_MODEL_CODE.isin([score_model_code]) ].loc[:,:]
    
            
            VAL_MODEL=BSCORE_DATA_MODEL[BSCORE_DATA_MODEL['VALIDATION_SAMPLE']==1].loc[:,:]
            VAL_MODEL['QTR']= "Validation_Sample"
            
            VAL_MODEL= VAL_MODEL.groupby(['QTR',variable_list_1[j], variable_list_2[j]]).agg(
                            COUNT_SNAPSHOT = pd.NamedAgg(column='ACCOUNT_ID', aggfunc='count'),
                            ).reset_index(drop=False)
            
            QTR_MODEL= BSCORE_DATA_MODEL.groupby(['QTR',variable_list_1[j], variable_list_2[j]]).agg(
                            COUNT_SNAPSHOT = pd.NamedAgg(column='ACCOUNT_ID', aggfunc='count'),
                            ).reset_index(drop=False)
            
            logger.debug(f"QTR data and validataion sample data aggregated for {variable_list_1[j]} variable and snapshot count calculated at variable level.")
        
            DATA_MODEL_1= pd.concat([VAL_MODEL,QTR_MODEL],ignore_index=True)
            
            T2=BSCORE_DATA_MODEL[BSCORE_DATA_MODEL['BENCHMARK']==1].loc[:,:]
            
            T2= T2.groupby([variable_list_1[j]]).agg(
                            COUNT_BENCHMARK = pd.NamedAgg(column='ACCOUNT_ID', aggfunc='count'),
                            ).reset_index(drop=False)
            
            logger.debug(f"Benchmark count calculated for {variable_list_1[j]} variable at variable level")
            logger.debug("#######################################################################################")
        
            CSI_MODEL = pd.merge(DATA_MODEL_1[['QTR',variable_list_1[j], variable_list_2[j], 'COUNT_SNAPSHOT']], T2[[variable_list_1[j], 'COUNT_BENCHMARK']],
                         how='outer', on=[variable_list_1[j]])
    
            CSI_MODEL['COUNT_BENCHMARK'].fillna(1,inplace=True)
            
            CSI_CALC = CSI_MODEL.groupby(['QTR']).agg(
                            TOTAL_COUNT_SNAPSHOT = pd.NamedAgg(column='COUNT_SNAPSHOT', aggfunc='sum'),
                            TOTAL_COUNT_BENCHMARK = pd.NamedAgg(column='COUNT_BENCHMARK', aggfunc='sum')
                            ).reset_index(drop=False)
            
            logger.debug(f"Total Benchmark and Snapshot count calculated for {variable_list_1[j]} variable at QTR level")
            
            CSI_1= pd.merge(CSI_MODEL[['QTR', variable_list_1[j], variable_list_2[j] ,'COUNT_SNAPSHOT', 'COUNT_BENCHMARK']], CSI_CALC[['QTR','TOTAL_COUNT_SNAPSHOT', 'TOTAL_COUNT_BENCHMARK']],
                          left_on = "QTR", right_on = "QTR", how = "left", suffixes=('_',''))
            
            logger.debug("#######################################################################################")
       
            logger.debug(f"Variable CSI computation using formula starts for {variable_list_1[j]} variable")
            
            CSI_1= CSI_1.rename(columns = {'QTR':'YYYYMM'})
            CSI_1['SCORE_CARD']= score_model_code
            CSI_1['VARIABLE_NAME']= variable_list_1[j]
            CSI_1 = CSI_1.rename(columns = {variable_list_1[j]:'VALUE'})
            CSI_1 = CSI_1.rename(columns = {variable_list_2[j]:'BUCKET_SCORE'})
            CSI_1['CONCENTRATION_SNAPSHOT']= CSI_1['COUNT_SNAPSHOT']/CSI_1['TOTAL_COUNT_SNAPSHOT']
            CSI_1['CONCENTRATION_BENCHMARK']= CSI_1['COUNT_BENCHMARK']/CSI_1['TOTAL_COUNT_BENCHMARK']
            CSI_1['DIFFERENCE']= CSI_1['CONCENTRATION_SNAPSHOT']- CSI_1['CONCENTRATION_BENCHMARK']
            CSI_1['LOG_CONC']= np.log(CSI_1['CONCENTRATION_SNAPSHOT']/ CSI_1['CONCENTRATION_BENCHMARK'])
            CSI_1['INDIVIDUAL_CSI']= CSI_1['DIFFERENCE'] * CSI_1['LOG_CONC']
            
            CSI_2 = CSI_1.sort_values(['YYYYMM','VALUE'],ascending=False).groupby(['YYYYMM' ]).agg(
                            CSI = pd.NamedAgg(column='INDIVIDUAL_CSI', aggfunc='sum'),
                            HCI = pd.NamedAgg(column='CONCENTRATION_SNAPSHOT', aggfunc='max')
                            ).reset_index(drop=False)
            
            logger.debug(f"Variable CSI computation using formula ends for {variable_list_1[j]} variable")  
            
            logger.debug("#######################################################################################")
            
            
            CSI = pd.merge(CSI_1[['YYYYMM','SCORE_CARD', 'VALUE', 'BUCKET_SCORE', 'COUNT_SNAPSHOT', 'COUNT_BENCHMARK', 'VARIABLE_NAME',           
                   'CONCENTRATION_SNAPSHOT', 'CONCENTRATION_BENCHMARK', 'DIFFERENCE',
                   'LOG_CONC', 'INDIVIDUAL_CSI']], CSI_2[['YYYYMM','CSI', 'HCI']],
                          left_on = "YYYYMM", right_on = "YYYYMM", how = "left", suffixes=('_',''))
            
            CSI = CSI.sort_values(['YYYYMM','VALUE'],ascending=True)
            BF.append(CSI)
            
            logger.debug(f"QoQ calculated CSI for {variable_list_1[j]} variable is merged with its respective quarters")
            
            logger.debug("##########################################################################################################################################")
     
        
        bf1 = pd.concat(BF, axis=0, ignore_index=True)
        L= [ 'SCORE_CARD','YYYYMM','VARIABLE_NAME', 'VALUE','BUCKET_SCORE',  'COUNT_SNAPSHOT', 'COUNT_BENCHMARK',       
           'CONCENTRATION_SNAPSHOT', 'CONCENTRATION_BENCHMARK', 'DIFFERENCE',
           'LOG_CONC', 'INDIVIDUAL_CSI', 'CSI']
    
        df11 = bf1[L]    
        df12= df11.dropna()
        
        unq_qtrs=df12['YYYYMM'].nunique()
            
        logger.info(f"Whole csi_score_band function compiled and CSI calculated for {unq_qtrs} qtrs and returned successfully.")    
        logger.debug("##########################################################################################################################################")
    
        
        
        return df12  

#......................................................................................................................................................................#

def mort_adding_branch(df):
    
    Conditoins=[
        (df['LEAFNODE_ID'] == "leaf1" ),
        (df['LEAFNODE_ID'] == "leaf10"),
        (df['LEAFNODE_ID' ]== "leaf11"),
        (df['LEAFNODE_ID'] == "leaf12"),
        (df['LEAFNODE_ID'] == "leaf13"),
        (df['LEAFNODE_ID'] == "leaf14"),
        (df['LEAFNODE_ID'] == "leaf17"),
        (df['LEAFNODE_ID'] == "leaf18"),
        (df['LEAFNODE_ID'] == "leaf2" ),
        (df['LEAFNODE_ID'] == "leaf3" ),
        (df['LEAFNODE_ID'] == "leaf4" ),
        (df['LEAFNODE_ID'] == "leaf5" ),
        (df['LEAFNODE_ID'] == "leaf6" ),
        (df['LEAFNODE_ID'] == "leaf7" ),
        (df['LEAFNODE_ID'] == "leaf9" )
        
        ]
    
    df['NODE_VAL']=np.select(Conditoins,[
                
            "N1",
            "N6",
            "N7",
            "N8",
            "N9",
            "N10",
            "N14",
            "N15",
            "N4",
            "N2",
            "N3",
            "N11",
            "N12",
            "N13",
            "N5",
            
                
            ])
    
    df['BRANCH_0']="L1"
    
    branch_condition=[
    
        df['NODE_VAL'] == "N1" ,
        df['NODE_VAL'] == "N2" ,
        df['NODE_VAL'] == "N3" ,
        df['NODE_VAL'] == "N4" ,
        df['NODE_VAL'] == "N5" ,
        df['NODE_VAL'] == "N6" ,
        df['NODE_VAL'] == "N7" ,
        df['NODE_VAL'] == "N8" ,
        df['NODE_VAL'] == "N9" ,
        df['NODE_VAL'] == "N10",
        df['NODE_VAL'] == "N11",
        df['NODE_VAL'] == "N12",
        df['NODE_VAL'] == "N13",
        df['NODE_VAL'] == "N14",
        df['NODE_VAL'] == "N15"
        
        ]
    df['BRANCH_1']=np.select(branch_condition,[
    
            "L1.1",
            "L1.1",
            "L1.1",
            "L1.1",
            "L1.1",
            "L1.1",
            "L1.1",
            "L1.2",
            "L1.2",
            "L1.2",
            "L1.2",
            "L1.2",
            "L1.2",
            "L1.2",
            "L1.2"
            
            ])
    
    df['BRANCH_2']=np.select(branch_condition,[
    
        "L1.1.1",
        "L1.1.1",
        "L1.1.1",
        "L1.1.1",
        "L1.1.1",
        "L1.1.2",
        "L1.1.2",
        "L1.2.1",
        "L1.2.1",
        "L1.2.1",
        "L1.2.1",
        "L1.2.1",
        "L1.2.1",
        "L1.2.2",
        "L1.2.2"
        
        
        ])
    
    df['BRANCH_3']=np.select(branch_condition,[
    
        "L1.1.1.1",
        "L1.1.1.1",
        "L1.1.1.1",
        "N4",
        "N5",
        "N6",
        "N7",
        "L1.2.1.1",
        "L1.2.1.1",
        "L1.2.1.1",
        "L1.2.1.2",
        "L1.2.1.2",
        "L1.2.1.2",
        "N14",
        "N15"
        
        ])
    
    
    df['BRANCH_4']=np.select(branch_condition,[
    
        "N1", 
        "N2", 
        "N3", 
        "N4", 
        "N5", 
        "N6", 
        "N7", 
        "N8", 
        "N9", 
        "N10",
        "N11",
        "N12",
        "N13",
        "N14",
        "N15"
        
        
        ])
       
    return df



# ................................................................................................................................................................................................................................................................................................

def benchmark_data():
    MORT_BM_DATA=pd.read_csv(r"\\zebra\GIRM\Group Portfolio Analytics and Risk Methodology\Named folders\GBS COE\Deepak\REFERENCE\mort_SAS_DATA.tab",sep='\t')

    MORT_BM_DATA=MORT_BM_DATA[['AS_AT','new_contract_id','CAL_SCORE_VALUE','LEAFNODE_ID','final_product']]
    MORT_BM_DATA.columns= ['SCORE_DATE','ACCOUNT_ID','BSCORE_ORIGINAL','LEAFNODE_ID','final_product']
    
    final_bm_data=mort_adding_branch(MORT_BM_DATA)
    
    return final_bm_data
    



#

def Node_wise_count(data,Parent_Branch,Child_Branch,var3,period):
    name="{}_count"
    
    name1="{}_concentration"
    name1=name1.format(period)
    
    
    num=len(Parent_Branch)
    DF1=[]
    DF2=pd.DataFrame()
    for i in range(num):
        var1=Parent_Branch[i]
        var2=Child_Branch[i]
        agg_dict={name.format(period) : pd.NamedAgg(column='ACCOUNT_ID', aggfunc='count')}
        test=data.groupby([var3,var1 ,var2 ]).agg(**agg_dict
        
                      ).reset_index(drop=False).rename(columns={var1: "Parent",var2: "Child" })
        test=test[test['Parent']!=test['Child']]
#         test[name1]=test[name.format(period)].sum()
        
        DF1.append(test)
        
    DF2=pd.concat(DF1)
    DF2=DF2.sort_values([var3,'Parent'])
    return DF2


# ...................................................................................................................................................................................................


def MORT_CSI_FUNC(df,portfolio_code,score_model_code):
    
    bm_data=benchmark_data()
    
    val_data=mort_adding_branch(df)
    
    val_sample=val_data[val_data['VALIDATION_SAMPLE']==1]
    val_sample['QTR']='Validation'
    
    Parent_Branch=['BRANCH_0', 'BRANCH_1', 'BRANCH_2', 'BRANCH_3']

    Child_Branch=['BRANCH_1', 'BRANCH_2', 'BRANCH_3', 'BRANCH_4']
    
    val_with_node=Node_wise_count(val_data,Parent_Branch,Child_Branch,var3='QTR',period='Validation')
    
    val_with_node_Overall=Node_wise_count(val_sample,Parent_Branch,Child_Branch,var3='QTR',period='Validation')
    
    val_with_node=pd.concat([val_with_node,val_with_node_Overall])
    
    bm_with_node=Node_wise_count(bm_data,Parent_Branch,Child_Branch,var3='final_product',period='Benchmark')
    
    final=pd.merge(val_with_node,bm_with_node,left_on=['Parent','Child'],right_on=['Parent','Child'],how='outer')
    
    order=['QTR','final_product','Parent','Child','Validation_count','Benchmark_count']
    
    final=final.reindex(columns=order).copy()
    
    final1=final.sort_values(['QTR','Parent'])
    
    Total_Node_wise_count=final1.groupby(['QTR','Parent']).agg(
                Val_Node_Count=pd.NamedAgg(column='Validation_count',aggfunc='sum'),
                Bench_Node_Count=pd.NamedAgg(column='Benchmark_count',aggfunc='sum'))
    
    
    final_pk=pd.merge(final,Total_Node_wise_count,left_on=['QTR','Parent'],right_on=['QTR','Parent'],how='outer')
    
    
    final_pk['concentration_benchmark']=final_pk['Benchmark_count']/final_pk['Bench_Node_Count']
    final_pk['concentration_snapshot']=final_pk['Validation_count']/final_pk['Val_Node_Count']
    
    
    final_pk['difference']=final_pk['concentration_snapshot']-final_pk['concentration_benchmark']
    
    final_pk['log_conc']=np.log(final_pk['concentration_snapshot']/final_pk['concentration_benchmark'])
    
    final_pk['individual_csi']= final_pk['difference']*final_pk['log_conc']
    
    
    Total_Node_wise_csi=final_pk.groupby(['QTR','Parent']).agg(
                CSI=pd.NamedAgg(column='individual_csi',aggfunc='sum')
                )
    
    final_pk=final_pk.sort_values(['QTR','Parent'])
    
    final_pk=pd.merge(final_pk,Total_Node_wise_csi,left_on=['QTR','Parent'],right_on=['QTR','Parent'],how='outer')
        
    return final_pk
    
# ...................................................................................................................................................................................................

def calculate_population_summary(df,portfolio_code,score_model_code,aggregator_var,current_qtr,segment,deciles=10,bm_year=2020):

    
    ''' This function outputs the Population Summary table for the selected scoremodel.

    
        input params:
           
        @ df - The required data that is to be used for generating the Population summary table.
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        @ aggregator_var - This takes string input and accounts for the column to be used as the variable for grouping our population summary data. i.e. 'SCORE_RANGE' or 'SCORE_VALUE'.
        @ deciles - This takes integer variable for making deciles for Population summary table. By default the value is 10.
        
        output params:
        
        It returns a DataFrame for Population Summary.
    
    '''

    
    logger.info(f"Calculating Population Summary function starts for {portfolio_code} portfolio and {score_model_code} scorecard.")

    
    if score_model_code not in portfolio_wise_score_model_code(segment):
    
        print("Please check the corresponding portfolio and score model code")
    
        logger.info("Population Summary Can't be performed for given portfolio and scoremodel code")
    
        return    
    
    else:
    #     if portfolio_code!='mort':
            
        logger.debug(f"Creation of Score bands from benchmark data starts for {portfolio_code} portfolio and {score_model_code} scorecard starts")
    
        BSCORE_DATA_BM= df[df['BENCHMARK']==1].iloc[:,:]
        BSCORE_DATA_BM = BSCORE_DATA_BM.sort_values(['SCORE_MODEL_CODE','SCORE_VALUE'],ascending=False)
        BSCORE_DATA_RANK= BSCORE_DATA_BM[BSCORE_DATA_BM.SCORE_MODEL_CODE.isin([score_model_code])].loc[:,:]     

        logger.debug("Benchmark data filtered for making deciles and calculating population summary")

        BSCORE_DATA_RANK['rank'] = BSCORE_DATA_RANK['SCORE_VALUE'].rank()    
        BSCORE_DATA_RANK['DECILE'] = np.floor(BSCORE_DATA_RANK['rank']*(deciles)/(len(BSCORE_DATA_RANK['SCORE_VALUE'])+1))
        BSCORE_DATA_RANK= BSCORE_DATA_RANK.groupby(['SCORE_MODEL_CODE', 'DECILE']).agg(                   
                            BM_MAX_SCORE_VALUE = pd.NamedAgg(column='SCORE_VALUE', aggfunc='max'),
                            BM_MIN_SCORE_VALUE = pd.NamedAgg(column='SCORE_VALUE', aggfunc='min'),
                            COUNT = pd.NamedAgg(column='SCORE_VALUE', aggfunc='count')).reset_index(drop=False)
        
        logger.debug("Minimum and Maximum score value for each decile has been calculated")
        logger.debug("#######################################################################################")
        
        
        BSCORE_DATA_MODEL= df.loc[(df.SCORE_MODEL_CODE.isin([score_model_code]))]
        cutoff=BSCORE_DATA_RANK['BM_MIN_SCORE_VALUE'].values.astype('float64')
        cutoff[0]=-np.inf
        cutoff=np.append(cutoff,[np.inf])
        BSCORE_DATA_MODEL['SCORE_RANGE']=pd.cut(BSCORE_DATA_MODEL['SCORE_VALUE'],bins=cutoff,right=False)
        
        unq_bands=BSCORE_DATA_MODEL['SCORE_RANGE'].nunique()
        
        logger.debug(f"Score bands with {unq_bands} unique bands created for calculating Population Summary")
        logger.debug("#######################################################################################")            
 
            
        # else:
            
        #     logger.debug(f"Creation of Score bands from benchmark data starts for {portfolio_code} portfolio and {score_model_code} scorecard starts")
            
        #     BSCORE_DATA_BM= df[df['BENCHMARK']==1].iloc[:,:]
        #     BSCORE_DATA_RANK= BSCORE_DATA_BM.groupby(['SCORE_VALUE']).agg(                   
        #                         BM_MAX_SCORE_VALUE = pd.NamedAgg(column='SCORE_VALUE', aggfunc='max'),
        #                         BM_MIN_SCORE_VALUE = pd.NamedAgg(column='SCORE_VALUE', aggfunc='min'),
        #                         COUNT = pd.NamedAgg(column='SCORE_VALUE', aggfunc='count')).reset_index(drop=False)
            
            
        #     BSCORE_DATA_MODEL= df.loc[(df.SCORE_MODEL_CODE.isin([score_model_code]))]
            
        #     unq_bands=BSCORE_DATA_MODEL['SCORE_VALUE'].nunique()
            # logger.debug(f"Score bands with {unq_bands} unique bands created for calculating Population Summary")
            # logger.debug("#######################################################################################")

        def rolling_window(current_qtr):
            year = int(current_qtr[:4])
            qtr = int(current_qtr[-1])
            window = [current_qtr]
            for i in range(4-1):
                if qtr>1:
                    qtr=qtr-1
                else:
                    qtr=4
                    year=year-1
                window.append(str(year)+'Q'+str(qtr))
                window=sorted(window)
            return window
        
        temp_window = rolling_window("2022Q3") 
        listi=[]
        for qtrs in temp_window[::-1]:
            listi.append(rolling_window(qtrs))
    
        df1=[]
        df2=pd.DataFrame()
        TTD_data=BSCORE_DATA_MODEL.sort_values(by='YEAR_SCORE',inplace=True)
        unique_year=sorted(df['YEAR_SCORE'].unique())
        len_unique=len(unique_year)
        
        logger.debug(f"Final data created for calculating the Population summary, which will have {len_unique} unique years")
        logger.debug("#######################################################################################")
        
        logger.debug("Calculating Population distribution for Population Summary table starts")
        
        bm_pop_data=BSCORE_DATA_MODEL[BSCORE_DATA_MODEL['YEAR_SCORE']==bm_year]
        monitoring_data=BSCORE_DATA_MODEL[BSCORE_DATA_MODEL['QTR'].isin(listi[0])]
        performance_data=BSCORE_DATA_MODEL[BSCORE_DATA_MODEL['QTR'].isin(listi[1])]
        
        bm_pop_data_con=bm_pop_data.groupby(aggregator_var)['ACCOUNT_ID'].count()/bm_pop_data.groupby(aggregator_var)['ACCOUNT_ID'].count().sum()
        bm_pop_data_con=bm_pop_data_con.reset_index(drop=True)
        
        monitoring_data_con=monitoring_data.groupby(aggregator_var)['ACCOUNT_ID'].count()/monitoring_data.groupby(aggregator_var)['ACCOUNT_ID'].count().sum()
        monitoring_data_con=monitoring_data_con.reset_index(drop=True)
        
        val_dict={"bm_period":bm_year,"mon_period":listi[0],"perf_period":listi[1]}
 
 
        l=[]
        l.append(bm_pop_data_con)
        df1=pd.DataFrame(list(bm_pop_data_con),columns=["%pop_benchmark_period"])
        df2=pd.concat([df2,df1],axis=1) 
        
        l=[]
        l.append(monitoring_data_con)
        df1=pd.DataFrame(list(monitoring_data_con),columns=["%_pop_monitoring_period"])
        df2=pd.concat([df2,df1],axis=1)  
        
        Perf_count=performance_data.groupby(aggregator_var)['ACCOUNT_ID'].count().sort_index().reset_index()['ACCOUNT_ID']
        Perf_bads=performance_data.groupby([aggregator_var])['TARGET_12'].sum().sort_index().reset_index()['TARGET_12']
        Perf_bad_rate=Perf_bads/Perf_count
        dicti={'%_bad_rate_performance_period':Perf_bad_rate}
        df3=pd.DataFrame(dicti) #
        df2=pd.concat([df2,df3],axis=1)        
        
        
        # for i in unique_year:
#             
            # logger.debug(f"Population distribution for year {i} starts")
#             
            # TTD_data=BSCORE_DATA_MODEL[BSCORE_DATA_MODEL['YEAR_SCORE'].isin([i])]
            # TTD_con=TTD_data.groupby(aggregator_var)['ACCOUNT_ID'].count()/TTD_data.groupby(aggregator_var)['ACCOUNT_ID'].count().sum()
            # TTD_con=TTD_con.reset_index(drop=True)
#             
            # l=[]
            # l.append(TTD_con)
            # df1=pd.DataFrame(list(TTD_con),columns=['pop_'+str(i)])
            # df2=pd.concat([df2,df1],axis=1)
            
            # logger.debug(f"{aggregator_var} aggregated counts calculated for year {i} and stored")
            # logger.debug("#######################################################################################")
# 
        # logger.debug("Population distribution for overall data starts")
#         
        # TTD_data=BSCORE_DATA_MODEL
        # TTD_con=TTD_data.groupby(aggregator_var)['ACCOUNT_ID'].count()/TTD_data.groupby(aggregator_var)['ACCOUNT_ID'].count().sum()
        # TTD_con=TTD_con.reset_index(drop=True)
        # l=[]
        # l.append(TTD_con)
        # df1=pd.DataFrame(list(TTD_con),columns=['pop_overall'])
        # df2=pd.concat([df2,df1],axis=1)
#         
        # logger.debug("Population distribution for overall data ends and stored with year on year Population distribution")
        # logger.debug("#######################################################################################")
#        
#         
        # logger.debug("Calculating Performance distribution for Population Summary table starts")
#         
        # for j in unique_year[:len_unique-1]:
#             
            # logger.debug(f"Performance distribution for year {j} starts")
#             
            # Performance_data=BSCORE_DATA_MODEL[BSCORE_DATA_MODEL['YEAR_SCORE'].isin([j])]
            # Perf_count=Performance_data.groupby(aggregator_var)['ACCOUNT_ID'].count().sort_index().reset_index()['ACCOUNT_ID']
            # Perf_bads=Performance_data.groupby([aggregator_var])['TARGET_12'].sum().sort_index().reset_index()['TARGET_12']
            # Perf_bad_rate=Perf_bads/Perf_count
            # dicti={'per_'+str(j):Perf_bad_rate}
            # df3=pd.DataFrame(dicti) #
            # df2=pd.concat([df2,df3],axis=1)
#             
            # logger.debug(f"{aggregator_var} aggregated counts calculated for year {j} and stored")
            # logger.debug("#######################################################################################")
# 
        # logger.debug("Performance distribution for overall data starts")
#         
        # Performance_data=BSCORE_DATA_MODEL[BSCORE_DATA_MODEL['YEAR_SCORE'].isin(unique_year[:len_unique-1])]
        # Perf_count=Performance_data.groupby(aggregator_var)['ACCOUNT_ID'].count().sort_index().reset_index()['ACCOUNT_ID']
        # Perf_bads=Performance_data.groupby([aggregator_var])['TARGET_12'].sum().sort_index().reset_index()['TARGET_12']
        # Perf_bad_rate=Perf_bads/Perf_count
        # dicti={'per_overall':Perf_bad_rate}
        # df3=pd.DataFrame(dicti) #
        # df2=pd.concat([df2,df3],axis=1)
#         
        # logger.debug("Performance distribution for overall data ends and stored with year on year Performance distribution")
        # logger.debug("#######################################################################################")
#         

    #if portfolio_code!='mort':
        x=pd.DataFrame(list(BSCORE_DATA_MODEL[aggregator_var].unique().sort_values()),columns=[aggregator_var])
        x=pd.concat([x,df2],axis=1)        
        logger.info(f"Whole calculate_population_summary function compiled and Population summary calculated for {score_model_code} scorecard and returned successfully")    

    # else:
    #     x=pd.DataFrame(list(sorted(BSCORE_DATA_MODEL[aggregator_var].unique())),columns=[aggregator_var])
    #     x=pd.concat([x,df2],axis=1)        
    #     logger.info(f"Whole calculate_population_summary function compiled and Population summary calculated for {score_model_code} scorecard and returned successfully")    
    


        logger.debug("#######################################################################################")

    return x, val_dict



#...............................................................................................................................................................................#


def calculate_scorepoint_gini(df,portfolio_code,score_model_code,gini_type,bm_gini,segment):
    
    ''' This function outputs the SCORE-POINT GINI table depending upon the choice of result we want.
    
        input params:
           
        @ df - The required data that is to be used for generating the GINI table.
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        @ gini_type - This is a string variable which takes only 2 input i.e "monthly" or "quarterly"  depending upon the type of result we want.
        @ bm_gini - This is the float variable for the benchmark gini of the particular scorecard.
        
        output params:
        
        It returns a DataFrame for SCORE-POINT GINI.
    
    '''

    logger.info(f"Calculating Score Point GINI function starts for {portfolio_code} portfolio and {score_model_code} scorecard")
    

    if score_model_code not in portfolio_wise_score_model_code(segment=segment):
        
        print("Please check the corresponding portfolio and score model code")
        
        logger.info("GINI Can't be calculated for given portfolio and scoremodel code")
        
        return 
    
    else:
        
    
        if(gini_type=="quarterly"):
            var1="QUARTER"
            var2="QTR"
        elif(gini_type=="monthly"):
            var1="DATE"
            var2="SCORE_DATE"
        else:
            print("Please check the gini type again") 
            return 
        
        logger.debug(f"The GINI calculation type selected is {gini_type}.")
        
        GINI_val=pd.DataFrame()
        GINI=pd.DataFrame()
        gini=[]
        gini_val=[]
        B_SCORE_DATA_val=df[(df["VALIDATION_SAMPLE"].isin([1])) & (df['SCORE_MODEL_CODE'].isin([score_model_code]))]
        
        logger.debug("Validation sample data filtered for GINI calculation")
        
        logger.debug(f"TARGET_12 var contains {B_SCORE_DATA_val['TARGET_12'].nunique()} unq values.")
        
        fpr_val, tpr_val, thresholds_val = metrics.roc_curve(B_SCORE_DATA_val["TARGET_12"], B_SCORE_DATA_val["PRED_BR"])
        auc_val=metrics.auc(fpr_val, tpr_val)
        gini_val.append(2*auc_val-1)
        GINI_val[var1]=['Validation_Sample']
        GINI_val['GINI']=gini_val
        
        logger.debug("Validation sample GINI calculated and stored")
        
        logger.debug("#######################################################################################")
        
        agg_var=sorted(df[(df["VALIDATION_SAMPLE"].isin([1]))][var2].unique())
        
        logger.debug(f"{var2} aggregated GINI calculation starts for {len(agg_var)} {var2}.")
        
        for j in agg_var:
            B_SCORE_DATA=df[df[var2].isin([j]) & (df['SCORE_MODEL_CODE'].isin([score_model_code]))]
            fpr, tpr, thresholds = metrics.roc_curve(B_SCORE_DATA["TARGET_12"], B_SCORE_DATA["PRED_BR"])
            auc=metrics.auc(fpr, tpr)
            gini.append(2*auc-1)

        logger.debug(f"{var2} aggregated GINI calculation ends for {len(agg_var)} {var2}.")

        logger.debug("#######################################################################################")
            
        GINI[var1]=agg_var
        GINI['GINI']=gini
        GINI=pd.concat([GINI,GINI_val],ignore_index=True)  
        
        logger.info(f"Whole calculate_scorepoint_gini function compiled and GINI calculated for {GINI.shape[0]} {var2} and returned successfully.")    

        logger.debug("#######################################################################################")

        return GINI  


#...............................................................................................................................................................................#


def calculate_decile_gini(df,portfolio_code,score_model_code,gini_type,bm_gini,deciles=10):

    ''' This function outputs the DECILE GINI table depending upon the choice of result we want.
    
        input params:

        @ df - The required data that is to be used for generating the GINI table.
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        @ gini_type - This is a string variable which takes only 2 input i.e "monthly" or "quarterly"  depending upon the type of result we want.
        @ bm_gini - This is the float variable for the benchmark gini of the particular scorecard.
        @ deciles - This takes integer variable for making deciles for GINI. By deafult the value is 10.

        output params:
        
        It returns a DataFrame for DECILE GINI.
    
    '''
    
    
    
    if(gini_type=="quarterly"):
        var1="QUARTER"
        var2="QTR"
    elif(gini_type=="monthly"):
        var1="DATE"
        var2="SCORE_DATE"
    else:
        print("Please chcek the gini type again")
    
    df=df[(df["VALIDATION_SAMPLE"].isin([1])) & (df['SCORE_MODEL_CODE'].isin([score_model_code]))]
    
    DF4=[]
    DF_F=[]
    scr_dt=pd.DataFrame()
    scr_dt[var1]= sorted(df[var2].unique())
    
    # if(portfolio_code=="mort"):
    #     aggregate_variable_2='SCORE_VALUE'
    # else:
    aggregate_variable_2='SCORE_RANGE'

        
    for i in range(len(scr_dt)):
        
        
        DF_n1=df[df[var2]==scr_dt[var1][i]]
        
            
        DF_n1[aggregate_variable_2]=pd.qcut(DF_n1['SCORE_VALUE'],q=deciles,labels=False,duplicates='drop')
        
        DF_F.append(DF_n1)
            
    DF1=pd.concat(DF_F)
    
    DF_val = df[(df["VALIDATION_SAMPLE"].isin([1])) & (df['SCORE_MODEL_CODE'].isin([score_model_code]))]
    
    DF_val[var2] = "Validation_Sample"
    
    DF_val[aggregate_variable_2]=pd.qcut(DF_val['SCORE_VALUE'],q=deciles,labels=False,duplicates='drop')
    
    DF1=pd.concat([DF1,DF_val])
    
    DF1 = DF1.groupby([var2, aggregate_variable_2]).agg(
                                    MIN = pd.NamedAgg(column=aggregate_variable_2, aggfunc='min'),   
                                    MAX = pd.NamedAgg(column=aggregate_variable_2, aggfunc='max'),
                                    TOTAL = pd.NamedAgg(column='TARGET_12', aggfunc='count'),
                                    BADS = pd.NamedAgg(column='TARGET_12', aggfunc='sum')
                                    ).reset_index(drop=False)
    
    DF1['GOODS']= DF1['TOTAL'] - DF1['BADS']
    DF1['BAD_RATE']= DF1['BADS']/DF1['TOTAL']
    
    
    DF2 = DF1.groupby([var2]).agg(
                      TOTAL_GOOD = pd.NamedAgg(column='GOODS', aggfunc='sum'),
                      TOTAL_BAD = pd.NamedAgg(column='BADS', aggfunc='sum')
                      ).reset_index(drop=False)
    
    
    
    DF3 = pd.merge(DF1[[var2,aggregate_variable_2, 'MIN','MAX','TOTAL', 'GOODS', 'BADS', 'BAD_RATE']],
                   DF2[[var2, 'TOTAL_GOOD', 'TOTAL_BAD']],
                   left_on = [var2],
                   right_on = [var2], how = "left", suffixes=('_',''))
    
    
    for i in ((DF3[var2].unique())):
        
        DF_n=DF3[DF3[var2]==i]
        
            
        DF_n['PERCENTAGE_GOODS']=DF_n['GOODS']/DF_n['TOTAL_GOOD']
        DF_n['PERCENTAGE_BADS']=DF_n['BADS']/DF_n['TOTAL_BAD']
        
        
        
        DF_n['CUMM_GOODS']=DF_n['PERCENTAGE_GOODS'].cumsum()
        DF_n['CUMM_BADS']=DF_n['PERCENTAGE_BADS'].cumsum()
        DF_n['CUMM_POP']=DF_n['TOTAL'].cumsum()
        
        DF_n['AUC']=0.5*(DF_n['CUMM_GOODS']-DF_n['CUMM_GOODS'].shift(1).fillna(0))*(DF_n['CUMM_BADS']+DF_n['CUMM_BADS'].shift(1).fillna(0))
        DF_n['DIFFERENTIATION']=DF_n['CUMM_BADS']-DF_n['CUMM_GOODS']
        
        DF_n['DIFFERENTIATION']=DF_n['CUMM_BADS']-DF_n['CUMM_GOODS']
        
        DF_n['GINI_STATISTIC']=-1 + 2*DF_n['AUC'].sum()
        DF_n['KS_STATISTIC']=DF_n['DIFFERENTIATION'].max()
        
        DF4.append(DF_n)
                
    DF4=pd.concat(DF4)
    if(gini_type=="quarterly"):
        
        cols_req=['QTR','SCORE_RANGE','TOTAL','GOODS','BADS','BAD_RATE','PERCENTAGE_GOODS','PERCENTAGE_BADS','CUMM_GOODS','CUMM_BADS','CUMM_POP','AUC','DIFFERENTIATION','GINI_STATISTIC','KS_STATISTIC']
    elif(gini_type=="monthly"):
        
        cols_req=['SCORE_DATE','SCORE_RANGE','TOTAL','GOODS','BADS','BAD_RATE','PERCENTAGE_GOODS','PERCENTAGE_BADS','CUMM_GOODS','CUMM_BADS','CUMM_POP','AUC','DIFFERENTIATION','GINI_STATISTIC','KS_STATISTIC']
    
    #cols_req=['QTR','SCORE_RANGE','TOTAL','GOODS','BADS','BAD_RATE','PERCENTAGE_GOODS','PERCENTAGE_BADS','CUMM_GOODS','CUMM_BADS','CUMM_POP','AUC','DIFFERENTIATION','GINI_STATISTIC','KS_STATISTIC']
    return DF4[cols_req]



#...............................................................................................................................................................................#


def calculate_mape(df,portfolio_code,score_model_code,val_start_date,val_end_date,mape_type,prod_ttc_pd,segment):
    
    
    ''' This function outputs the MAPE table depending upon the choice of result we want.
    
        input params:
           
        @ df - The required data that is to be used for generating the MAPE table.
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        @ val_start_date -This takes string variable in the format (YYYY-MM-DD) representing start date for the current validation/monitoring exercise.
        @ val_end_date - This takes string variable in the format (YYYY-MM-DD) representing end date for the current validation/monitoring exercise.
        @ mape_type - This is a string variable which takes only 2 input i.e "monthly" or "quarterly"  depending upon the type of result we want.
        @ prod_ttc_pd - This takes a float variable which accounts for the implemented in-house production value of TTC PD for the given score-card.
        
        output params:
        
        It returns a DataFrame for MAPE.
    
    '''
    
    
    
    
    if(prod_ttc_pd>1):
        
        print("Take percentage value of TTC PD")
        logger.warning("Give percentage value of TTC PD")
        return 
        
    else:
        
        logger.info(f"Calculating MAPE function starts for {portfolio_code} portfolio and {score_model_code} scorecard.")
    

        if score_model_code not in portfolio_wise_score_model_code(segment=segment):
        
            print("Please check the corresponding portfolio and score model code")
        
            logger.info("MAPE Can't be calculated for given portfolio and scoremodel code")
        
            return 
        
        else:
        

            MAPE_data=df[df['SCORE_MODEL_CODE']==score_model_code]
                
            val_start_date = pd.to_datetime(val_start_date)
            val_end_date = pd.to_datetime(val_end_date)
                
            MAPE_data_validation=MAPE_data[(MAPE_data['SCORE_DATE']>=val_start_date) & (MAPE_data['SCORE_DATE']<=val_end_date)]
            
            logger.debug(f"Validation data for MAPE calculation filtered, ranging from score_date {val_start_date} to {val_end_date}.")
 

            logger.debug("#######################################################################################")

            if mape_type=='monthly':
                
                logger.debug(f"{mape_type} Expected and Actual defaults calculation starts")

                MAPE_data_grp=MAPE_data_validation.groupby(["SCORE_DATE"]).agg(
                    ExpectedDefaults = pd.NamedAgg(column='PRED_BR', aggfunc='sum'),
                    ObservedDefaults = pd.NamedAgg(column='TARGET_12', aggfunc='sum'),
                    Total_Acc = pd.NamedAgg(column='ACCOUNT_ID', aggfunc='count'))
                
                logger.debug(f"{mape_type} Expected and Actual defaults calculation ends")
                
                
                logger.debug("#######################################################################################")

                            
            elif mape_type=='quarterly':
                
                logger.debug(f"{mape_type} Expected and Actual defaults calculation starts")

                MAPE_data_grp=MAPE_data_validation.groupby(["QTR"]).agg(
                    ExpectedDefaults = pd.NamedAgg(column='PRED_BR', aggfunc='sum'),
                    ObservedDefaults = pd.NamedAgg(column='TARGET_12', aggfunc='sum'),
                    Total_Acc = pd.NamedAgg(column='ACCOUNT_ID', aggfunc='count'))
                
                logger.debug(f"{mape_type} Expected and Actual defaults calculation ends")
        
                logger.debug("#######################################################################################")


            else:
                print("check the mape type name again")
                logger.warning("check the mape type name again")
                return
            
            logger.debug("EDR and ODR calculation starts") 
            
            MAPE_data_grp["ExpectedDefaults"]=round(MAPE_data_grp["ExpectedDefaults"])
            MAPE_data_grp["ExpectedDefaults"]=MAPE_data_grp["ExpectedDefaults"].astype(int)
            MAPE_data_grp["EDR"]= MAPE_data_grp["ExpectedDefaults"]/MAPE_data_grp["Total_Acc"]
            MAPE_data_grp["ODR"]=MAPE_data_grp["ObservedDefaults"]/MAPE_data_grp["Total_Acc"]
            MAPE_data_grp=MAPE_data_grp.reset_index()
            MAPE_data_grp=MAPE_data_grp[~MAPE_data_grp['ObservedDefaults'].isin([0])]
            MAPE_data_grp['prod_ttc_pd']=prod_ttc_pd
            
            MAPE_data_grp['Forecasted_LTPD_LTODR']=sum(MAPE_data_grp['ExpectedDefaults'])/sum(MAPE_data_grp['ObservedDefaults'])
            
            logger.debug("EDR and ODR calculation ends") 

            logger.debug("#######################################################################################")
            
            logger.info(f"Whole calculate_mape function compiled and MAPE calculated, ranging from score_date {val_start_date} to {val_end_date}.")    

            logger.debug("#######################################################################################")

            return MAPE_data_grp

#...............................................................................................................................................................................#

def calculate_binomial_test_quarterly(df,portfolio_code,score_model_code,segment,no_of_deciles=10,alpha=0.95,two_tailed=True):

    '''  This function outputs the Binomial table depending upon the choice of result we want.
    
        input params:

        @ df - The required data that is to be used for generating the Binomial test output.  
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        @ no_of_deciles - This takes integer value for no of deciles we want to break our quarterly data for calculating binomial test. By default the value is 10.
       	@ alpha - This takes float value for making the Confidence Interval. By default the value is .95.
        @ two_tailed - This takes bool value True or False depending upon the type of test we want to perform. By default the value is True.
        
        output params:
        
        It returns a DataFrame for Binomial Test.
    
    '''
    
    logger.info(f"Calculating Binomial Test function starts for {portfolio_code} portfolio and {score_model_code} scorecard.")
    
    if score_model_code not in portfolio_wise_score_model_code(segment):
    
        print("Please check the corresponding portfolio and score model code")
    
        logger.info("Binomial Test Can't be performed for given portfolio and scoremodel code")
    
        return     
    
    else:
    
        df= df[df.SCORE_MODEL_CODE.isin([score_model_code])].loc[:,:]
    
        VALIDATION=df[df['VALIDATION_SAMPLE']==1].loc[:,:]
        VALIDATION['QTR']= "Validation_Sample"
    
        df_new=pd.concat([df,VALIDATION])
        unique_qtr=df_new['QTR'].unique()
        len_unique=len(unique_qtr)
        df_result=pd.DataFrame()
        
        logger.debug("Validation Sample data for Binomial Test calculation filtered and appended to QoQ data")
        logger.debug("#######################################################################################")

        logger.debug(f"Binomial test calculation starts with parameters,  {no_of_deciles} deciles, {alpha} CI and two tailed is {two_tailed}")
        logger.debug("#######################################################################################")
        
        for i in unique_qtr:
            
            logger.debug(f"Binomial Test calculation starts for {i} quarter")
            
            df_1 = df_new[df_new['QTR'].isin([i])]
            df_1 = df_1.sort_values(['PRED_BR'],ascending=False)
            df_1['RANK'] = df_1['SCORE_VALUE'].rank()    
            df_1['DECILE'] = np.floor(df_1['RANK']*(no_of_deciles)/(len(df_1['SCORE_VALUE'])+1))
    
            df_1['alpha_1'] = alpha
            temp='z_value_'+str(alpha*100)
            
            if (two_tailed==True):
                df_1[temp] = norm.ppf(df_1['alpha_1'] + (1-df_1['alpha_1'])/2)
            else:
                df_1[temp] = norm.ppf(df_1['alpha_1'])
                
            test_df = df_1.sort_values(['DECILE'],ascending=True).groupby(['QTR','DECILE', temp]).agg(
                                   TOTAL = pd.NamedAgg(column='TARGET_12', aggfunc='count'),
                                   BADS = pd.NamedAgg(column='TARGET_12', aggfunc='sum'),
                                   PD = pd.NamedAgg(column='PRED_BR', aggfunc='mean')
                                   ).reset_index(drop=False)
    
            test_df['GOODS'] = test_df['TOTAL'] - test_df['BADS']
        
            logger.debug(f"Total counts, Goods, Bads calculated for {i} quarter.")
    
            test_df['alpha_1'] = alpha
            temp='z_value_'+str(alpha*100)
            
            if (two_tailed==True):
                test_df[temp] = norm.ppf(test_df['alpha_1'] + (1-test_df['alpha_1'])/2)
                
                test_df['LL_CI_DEF'] = (test_df['TOTAL']*test_df['PD']) - (test_df[temp]*np.sqrt(test_df['TOTAL']*test_df['PD']*(1-test_df['PD'])))
                test_df['Round_LL_CI_DEF']= np.ceil(test_df['LL_CI_DEF'])
        
                test_df['UL_CI_DEF'] = (test_df['TOTAL']*test_df['PD']) + (test_df[temp]*np.sqrt(test_df['TOTAL']*test_df['PD']*(1-test_df['PD'])))
                test_df['Round_UL_CI_DEF']= np.ceil(test_df['UL_CI_DEF'])
                
                logger.debug("Upper and Lower limit calculated using the formula")
            
            else:
                
                test_df[temp] = norm.ppf(test_df['alpha_1'])  
                test_df['UL_CI_DEF'] = (test_df['TOTAL']*test_df['PD']) + (test_df[temp]*np.sqrt(test_df['TOTAL']*test_df['PD']*(1-test_df['PD'])))
                test_df['Round_UL_CI_DEF']= np.ceil(test_df['UL_CI_DEF'])
            
                logger.debug("Upper limit calculated using the formula")
            
            df_result=pd.concat([df_result,test_df])
            #df_result=df_result[~df_result['BADS'].isin([0])]
            
            logger.debug(f"Final result data for {i} quarter appended for final resut generation")
            logger.debug("#######################################################################################")
            
        
        df_result['LL_Pass/Fail @95%']=np.where(df_result['BADS']>df_result['Round_LL_CI_DEF'],'PASS','FAIL')
        df_result['UL_Pass/Fail @95%']=np.where(df_result['BADS']<df_result['Round_UL_CI_DEF'],'PASS','FAIL')
        
        logger.info(f"Whole calculate_binomial_test_quarterly function compiled and Binomial Test calculated for {df_result['QTR'].nunique()} unique quarters and returned.")    

        logger.debug("#######################################################################################")
        
        
        return df_result


#...............................................................................................................................................................................#


def calculate_vif(df,portfolio_code,score_model_code,target_var,score_model_variable_wt,val_start_date,val_end_date,segment,path_file):

    '''  This function outputs the VIF table for the given score_card.
    
        input params:
        
        @ df - The required data that is to be used for generating the VIF output.  
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        @ target_var - This takes the string input for the name of the Target variable column in the dataset. 
        @ score_model_variable_wt - This takes the list of weight variables of the given scorecard to calculate vif. i.e. PB_UAE_BSCR04_WT
        @ val_start_date -This takes string variable in the format (YYYY-MM-DD) representing start date for the current validation/monitoring exercise.
        @ val_end_date - This takes string variable in the format (YYYY-MM-DD) representing end date for the current validation/monitoring exercise.


        
        output params:
        
        It returns a DataFrame for VIF.
    
    '''
    
    logger.info(f"Calculating VIF Test function starts for {portfolio_code} portfolio and {score_model_code} scorecard.")
    logger.debug("#######################################################################################")
    
    if score_model_code not in portfolio_wise_score_model_code(segment):
    
        print("Please check the corresponding portfolio and score model code")
    
        logger.info("VIF Test Can't be performed for given portfolio and scoremodel code")
    
        return     
    
    else:
 
    
        df['VIF_VALIDATION_SAMPLE']= np.select([(df['SCORE_DATE'] >=  val_start_date) & (df['SCORE_DATE'] <= val_end_date)],[1],0)
        df_1=df[df['VIF_VALIDATION_SAMPLE']==1]
        
        logger.debug(f"Validation sample filtered from {val_start_date} to {val_end_date} and regression started to calculate VIF.")
        
        model_formula = f"{target_var} ~ {' + '.join(score_model_variable_wt)}"
        
        logger.debug(f"Total {len(score_model_variable_wt)} variables present in {score_model_code} scorecard, named {score_model_variable_wt}")
        
        model = ols(model_formula, data=df_1).fit()
        
        logger.debug("Regression fitted successfully")
        
        vif = pd.DataFrame()
        vif["Variable"] = model.params.index
        vif["VIF"] = [variance_inflation_factor(model.model.exog, i) for i in range(model.model.exog.shape[1])]
        
        logger.info(f"Whole calculate_vif function compiled and VIF test calculated for {score_model_code} scorecard having {len(score_model_variable_wt)} variables.")    
    
        logger.debug("#######################################################################################")
        
        
        
        return vif

#...............................................................................................................................................................................#



def calculate_iv(df,portfolio_code,score_model_code,val_start_date,val_end_date,segment,path_file,score_model_wise_variable_dict):

    '''  
        This function outputs the IV table for the given score_card.
        
        input params:
            
        @ df - The required data that is to be used for generating the VIF output.  
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"   
        @ val_start_date -This takes string variable in the format (YYYY-MM-DD) representing start date for the current validation/monitoring exercise.
        @ val_end_date - This takes string variable in the format (YYYY-MM-DD) representing end date for the current validation/monitoring exercise.
    

        
        output params:
    
        It returns a DataFrame for IV.
    
    '''
    
    logger.info(f"Calculating IV function starts for {portfolio_code} portfolio and {score_model_code} scorecard.")
    logger.debug("#######################################################################################")
    
    if score_model_code not in portfolio_wise_score_model_code(segment):
    
        print("Please check the corresponding portfolio and score model code")
    
        logger.info("IV Test Can't be performed for given portfolio and scoremodel code")
    
        return    
    
    else:
        
        variable_list_1=score_model_wise_variable_lists_wt_rng(score_model_code,score_model_wise_variable_dict)[0]
        variable_list_2=score_model_wise_variable_lists_wt_rng(score_model_code,score_model_wise_variable_dict)[1]
        
        logger.debug(f"Weight variables and Range variables for {score_model_code} scorecard is fetched.")
        logger.debug(f"Range variables are {variable_list_1}")
        logger.debug(f"Weight variables are {variable_list_2}")
        logger.debug("###############################################################################################")
        
        if portfolio_code=="pl":
            df=df[df["NSTL_FLAG"]==0]
        df['IV_VALIDATION_SAMPLE']= np.select([(df['SCORE_DATE'] >=  val_start_date) & (df['SCORE_DATE'] <= val_end_date)],[1],0)
        
        logger.debug(f"IV Validation sample filtered from {val_start_date} to {val_end_date}")
        logger.debug("###############################################################################################")
        
        BF = []
        
        BSCORE_DATA_MODEL= df[df.SCORE_MODEL_CODE.isin([score_model_code]) ].loc[:,:]
            
        VAL_MODEL=BSCORE_DATA_MODEL[BSCORE_DATA_MODEL['IV_VALIDATION_SAMPLE']==1].loc[:,:]
        VAL_MODEL['QTR']= "IV_Validation_Sample"
        
        logger.debug("IV validation sample data is filtered ")
        
        for j in range(len(variable_list_1)):
                    
            logger.debug(f"IV calculation for {variable_list_2[j]} variable starts")
           
            VAL_MODEL_T= VAL_MODEL.groupby(['QTR',variable_list_2[j]]).agg(
                            Total_count = pd.NamedAgg(column='ACCOUNT_ID', aggfunc='count'),
                            Bad_count = pd.NamedAgg(column='TARGET_12', aggfunc='sum'),
                            ).reset_index(drop=False)
            VAL_MODEL_T['Good_count']=VAL_MODEL_T['Total_count']-VAL_MODEL_T['Bad_count']
            
            VAL_MODEL_T['Bad_count'].replace(0,1,inplace=True)
            VAL_MODEL_T['Good_count'].replace(0,1,inplace=True)
            
            
            logger.debug("Variable weight aggregated total counts and bads are calculated.")
    
            QTR_MODEL= VAL_MODEL_T.groupby(['QTR']).agg(
                            TOTAL_COUNT_ALL = pd.NamedAgg(column='Total_count', aggfunc='sum'),
                            TOTAL_GOOD_COUNT = pd.NamedAgg(column='Good_count', aggfunc='sum'),
                            TOTAL_BAD_COUNT = pd.NamedAgg(column='Bad_count', aggfunc='sum'),
                            ).reset_index(drop=False)
        
            logger.debug("Overall variable aggregated total counts and bads are calculated.")
    
            VIF_MODEL = pd.merge(VAL_MODEL_T[['QTR',variable_list_2[j], 'Total_count','Bad_count','Good_count']], QTR_MODEL[['QTR','TOTAL_COUNT_ALL','TOTAL_GOOD_COUNT','TOTAL_BAD_COUNT']],
                         how='outer', on='QTR')
    
            VIF_MODEL_1=VIF_MODEL
            VIF_MODEL_1['SCORE_CARD']= score_model_code
            VIF_MODEL_1['VARIABLE_NAME']= variable_list_2[j]
            #VIF_MODEL_1 = VIF_MODEL_1.rename(columns = {variable_list_1[j]:'VALUE'})
            VIF_MODEL_1 = VIF_MODEL_1.rename(columns = {variable_list_2[j]:'BUCKET_SCORE'})
            VIF_MODEL_1['PERC_BAD']= VIF_MODEL_1['Bad_count']/VIF_MODEL_1['TOTAL_BAD_COUNT']
            VIF_MODEL_1['PERC_GOOD']= VIF_MODEL_1['Good_count']/VIF_MODEL_1['TOTAL_GOOD_COUNT']
            VIF_MODEL_1['PERC_POP']= VIF_MODEL_1['Total_count']/VIF_MODEL_1['TOTAL_COUNT_ALL']
            VIF_MODEL_1['DIFFERENCE']= VIF_MODEL_1['PERC_GOOD']- VIF_MODEL_1['PERC_BAD']
            VIF_MODEL_1['LOG_CONC']= np.log(VIF_MODEL_1['PERC_GOOD']/ VIF_MODEL_1['PERC_BAD'])
            VIF_MODEL_1['Individual_IV']= VIF_MODEL_1['DIFFERENCE'] * VIF_MODEL_1['LOG_CONC']

            
            logger.debug(f"Individual IV value for {VIF_MODEL_1['BUCKET_SCORE'].nunique()} bins of {variable_list_1[j]} calculated.")
            logger.debug("###############################################################################################")
            
    
            BF.append(VIF_MODEL_1)
        

        bf1 = pd.concat(BF, axis=0, ignore_index=True)
        
            
        
        L= ['SCORE_CARD','QTR','VARIABLE_NAME','BUCKET_SCORE','PERC_BAD','PERC_GOOD',       
           'PERC_POP','Total_count','Bad_count','Good_count','DIFFERENCE',
           'LOG_CONC','Individual_IV']
        
        df11 = bf1[L] 
        
        t=df11.groupby(['QTR','VARIABLE_NAME'])['Individual_IV'].sum().reset_index()
        t.rename(columns={'Individual_IV':'IV'},inplace=True)
    
        IV_Table=pd.merge(df11,t,on=['QTR','VARIABLE_NAME'],how='left')

        logger.info(f"Whole calculate_iv function compiled and IV calculated for {len(variable_list_1)} variables of {score_model_code} scorecard.")    

        logger.debug("#######################################################################################")
        
 
    
        return IV_Table

#...............................................................................................................................................................................#




def calculate_rank_order(df,portfolio_code,score_model_code,segment,deciles=10):
    
    
    ''' This function outputs the Rank order table for the selected scoremodel.

    
        input params:
           
        @ df - The required data that is to be used for generating the Rank Order table.
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        @ deciles - This takes integer variable for making deciles for Rank Order table. By default the value is 10.
        
        output params:
        
        It returns a DataFrame for Rank order.
    
    '''

    logger.info(f"Calculating Rank Order function starts for {portfolio_code} portfolio and {score_model_code} scorecard.")

    
    if score_model_code not in portfolio_wise_score_model_code(segment):
    
        print("Please check the corresponding portfolio and score model code")
    
        logger.info("Rank order Can't be performed for given portfolio and scoremodel code")
    
        return    
    
    else:
    
        BSCORE_DATA_BM = df[df['BENCHMARK']==1].iloc[:,:]
        BSCORE_DATA_BM = BSCORE_DATA_BM.sort_values(['SCORE_MODEL_CODE','SCORE_VALUE'],ascending=False)
        B_SCORE_DATA_RANK= BSCORE_DATA_BM[BSCORE_DATA_BM.SCORE_MODEL_CODE.isin([score_model_code])].loc[:,:] 
        
        logger.debug("Benchmark data filtered for making deciles and calculating rank order")
    
        B_SCORE_DATA_RANK['rank'] = B_SCORE_DATA_RANK['SCORE_VALUE'].rank()    
        B_SCORE_DATA_RANK['DECILE'] = np.floor(B_SCORE_DATA_RANK['rank']*(deciles)/(len(B_SCORE_DATA_RANK['SCORE_VALUE'])+1))
        B_SCORE_DATA_RANK= B_SCORE_DATA_RANK.groupby(['SCORE_MODEL_CODE', 'DECILE']).agg(                   
                            BM_MAX_SCORE_VALUE = pd.NamedAgg(column='SCORE_VALUE', aggfunc='max'),
                            BM_MIN_SCORE_VALUE = pd.NamedAgg(column='SCORE_VALUE', aggfunc='min'),
                            COUNT = pd.NamedAgg(column='SCORE_VALUE', aggfunc='count')).reset_index(drop=False)
        
        logger.debug("Minimum and Maximum score value for each decile has been calculated")
        logger.debug("#######################################################################################")

        
        DATA_RANGE= df.loc[(df.SCORE_MODEL_CODE.isin([score_model_code])) & (df.VALIDATION_SAMPLE.isin([1]))] 
        
#         cutoff=B_SCORE_DATA_RANK['BM_MIN_SCORE_VALUE'].values.astype('float64')
#         cutoff[0]=-np.inf
#         cutoff=np.append(cutoff,[np.inf])
#         DATA_RANGE['SCORE_RANGE']=pd.cut(DATA_RANGE['SCORE_VALUE'],bins=cutoff,right=False)      
        
        
        if portfolio_code != 'mort':
             
            logger.debug(f"Creation of Score bands for {portfolio_code} portfolio and {score_model_code} scorecard starts")
            
            cutoff=B_SCORE_DATA_RANK['BM_MIN_SCORE_VALUE'].values.astype('float64')
            cutoff[0]=-np.inf
            cutoff=np.append(cutoff,[np.inf])
            DATA_RANGE['SCORE_RANGE']=pd.cut(DATA_RANGE['SCORE_VALUE'],bins=cutoff,right=False)
    
            BENCHMARK=df.loc[(df.SCORE_MODEL_CODE.isin([score_model_code])) & (df.BENCHMARK.isin([1]))]
            BENCHMARK['SCORE_RANGE'] = pd.cut(BENCHMARK['SCORE_VALUE'],bins=cutoff,right=False)
            BENCHMARK['QTR']= "Benchmark"
            
            unq_bands=BENCHMARK['SCORE_RANGE'].nunique()
            
            logger.debug(f"Score bands with {unq_bands} unique bands created for calculating Rank order")
            logger.debug("#######################################################################################")
            
                        
        # else:
            
        # logger.debug(f"Creation of Score bands for {portfolio_code} portfolio and {score_model_code} scorecard starts")

        # DATA_RANGE['SCORE_RANGE'] = DATA_RANGE['SCORE_VALUE']

        # BENCHMARK=df.loc[(df.SCORE_MODEL_CODE.isin([score_model_code])) & (df.BENCHMARK.isin([1]))] 
        # BENCHMARK['SCORE_RANGE'] =  BENCHMARK['SCORE_VALUE']
        # BENCHMARK['QTR']= "Benchmark"
        
        # unq_bands=BENCHMARK['SCORE_RANGE'].nunique()
        
        # logger.debug(f"Score bands with {unq_bands} unique bands created for calculating Rank order")
        # logger.debug("#######################################################################################")

        
        BENCHMARK = BENCHMARK.groupby(['QTR','SCORE_RANGE']).agg(
                        BADS = pd.NamedAgg(column='TARGET_12', aggfunc='sum'),
                        TOTAL = pd.NamedAgg(column='SCORE_RANGE', aggfunc='count')
                        ).reset_index(drop=False)
        
        logger.debug("Bechmark data filtered and total counts and bads count calculated")
        
        VALIDATION=DATA_RANGE[DATA_RANGE['VALIDATION_SAMPLE']==1].loc[:,:]
        VALIDATION['QTR']= "Validation_Sample"
        
        VALIDATION = VALIDATION.groupby(['QTR','SCORE_RANGE']).agg(
                        BADS = pd.NamedAgg(column='TARGET_12', aggfunc='sum'),
                        TOTAL = pd.NamedAgg(column='SCORE_RANGE', aggfunc='count')
                        ).reset_index(drop=False)
        
        logger.debug("Validation Sample data filterd and total counts and bads count calculated")
        
        QTR_DATA = DATA_RANGE.groupby(['QTR','SCORE_RANGE']).agg(
                        BADS = pd.NamedAgg(column='TARGET_12', aggfunc='sum'),
                        TOTAL = pd.NamedAgg(column='SCORE_RANGE', aggfunc='count')
                        ).reset_index(drop=False)
        
        logger.debug(f"Quarterly data for {DATA_RANGE['QTR'].nunique()} quarters filtered and total counts and bads count calculated")
        logger.debug("#######################################################################################")
        
        
        RANK_ORDER_DATA = pd.concat([BENCHMARK, VALIDATION, QTR_DATA], axis = 0,sort = False)
        
        logger.debug("Benchmark, validation and quarterly data concatenated for rank order analysis")
        logger.debug("#######################################################################################")
        logger.debug("Rank order calculations starts")

        
        RANK_ORDER_DATA['BAD_RATE'] = RANK_ORDER_DATA['BADS']/RANK_ORDER_DATA['TOTAL']
        df8 = RANK_ORDER_DATA.groupby(['QTR']).agg(
                        TOTAL_SUM = pd.NamedAgg(column='TOTAL', aggfunc='sum')
                        ).reset_index(drop=False)
        
        RANK_ORDER_DATA_1 = pd.merge(RANK_ORDER_DATA[['QTR', 'SCORE_RANGE', 'TOTAL', 'BADS', 'BAD_RATE']], df8[['TOTAL_SUM', 'QTR']],
                      left_on = "QTR", right_on = "QTR", how = "left", suffixes=('_',''))
        RANK_ORDER_DATA_1['CONCENTRATION_POPULATION'] = RANK_ORDER_DATA_1['TOTAL']/RANK_ORDER_DATA_1['TOTAL_SUM']
        RANK_ORDER_DATA_1['VARIABLE_NAME']="score_band"
        
        logger.debug("Rank order calculations ends")
        logger.debug("#######################################################################################")        
        
        col=["QTR","VARIABLE_NAME","SCORE_RANGE","TOTAL","BADS","CONCENTRATION_POPULATION","BAD_RATE"]
        RANK_ORDER_DATA_1=RANK_ORDER_DATA_1[col]
        
        c1=["Benchmark"]
        x=QTR_DATA['QTR'].unique()
        x=sorted(list(x))  
        cols=["Validation_Sample"]
    
        RANK_ORDER_DATA_1 = RANK_ORDER_DATA_1.sort_values(by=['QTR','SCORE_RANGE'],ascending=[False,True])
        l1=RANK_ORDER_DATA_1[RANK_ORDER_DATA_1['QTR'].isin(c1)]
        l2=RANK_ORDER_DATA_1[RANK_ORDER_DATA_1['QTR'].isin(x)]
        l2=l2.sort_values(by=['QTR','SCORE_RANGE'])
        l3=RANK_ORDER_DATA_1[RANK_ORDER_DATA_1['QTR'].isin(cols)]
        temp_df=pd.concat([l1,l2,l3],ignore_index=True)
        
        logger.info(f"Whole calculate_rank_order function compiled and Rank order calculated for {score_model_code} scorecard having {temp_df['QTR'].nunique()} Quarters including benchmark and validation sample.")    
    
        logger.debug("#######################################################################################")


        return temp_df      
  
        
#...............................................................................................................................................................................#


def call_metric_data(portfolio_code,is_islamic,val_start_date_imm,val_end_date_imm,val_start_date_omm,val_end_date_omm,path_file,score_model_wise_variable_dict,columns1, columns2,segment,bm_year=2020):

    ''' 
        This function outputs the dataframe required for input and output monitoring.
    
        input params:
        
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ is_islamic - This takes bool variable i.e 1 for islamic portfolios & 0 for rest.
        @ val_start_date_imm -This takes string variable in the format (YYYY-MM-DD) representing start date for Input Monitoring KPIs of the current validation/monitoring exercise.
        @ val_end_date_imm - This takes string variable in the format (YYYY-MM-DD) representing end date for Input Monitoring KPIs of the current validation/monitoring exercise.
        @ val_start_date_omm -This takes string variable in the format (YYYY-MM-DD) representing start date for Output Monitoring KPIs of the current validation/monitoring exercise.
        @ val_end_date_omm - This takes string variable in the format (YYYY-MM-DD) representing end date for Output Monitoring KPIs of the current validation/monitoring exercise.
        @ bm_year - This takess integer values for the benchmark year, default value is 2020.
        
        
        output params:
        
        It returns a DataFrame for IMM and OMM.
    
    '''
    
    logger.info(f"Fetching of Bscore data for {portfolio_code} portfolio for input and output monitoring metrics starts")
    logger.debug("#######################################################################################")
    
    path=save_and_read_intermediate_data(portfolio_code,path_file)
    # print("*****************",columns1)
    
    B_score_imm=IMM_islamic_non_islamic_columns(portfolio_code=portfolio_code,is_islamic=is_islamic,val_start_date=val_start_date_imm,val_end_date=val_end_date_imm,path_file=path,score_model_wise_variable_dict=score_model_wise_variable_dict,columns1=columns1,segment=segment,bm_year=bm_year)
   
    # display(B_score_imm)


    logger.info(f"Input monitoring metrics data for {portfolio_code} portfolio is read. It has {B_score_imm.shape[0]} rows {B_score_imm.shape[1]} columns. Validation Sample starts from {val_start_date_imm} to {val_end_date_imm}.")

    B_score_omm=OMM_islamic_non_islamic_columns(portfolio_code=portfolio_code,is_islamic=is_islamic,val_start_date=val_start_date_omm,val_end_date=val_end_date_omm,path_file=path_file,score_model_wise_variable_dict=score_model_wise_variable_dict,columns2=columns2,segment=segment,bm_year=bm_year)   

    # display(B_score_omm)

    logger.info(f"Output monitoring metrics data for {portfolio_code} portfolio is read. It has {B_score_omm.shape[0]} rows {B_score_omm.shape[1]} columns. Validation Sample starts from {val_start_date_omm} to {val_end_date_omm}.")
    logger.info(f"Fetching of Bscore data for {portfolio_code} portfolio for input and output monitoring metrics ends")
    
    logger.debug("#######################################################################################")
    
    
    return B_score_imm,B_score_omm


#...................................................................................................................................................................................#


def call_psi(df,portfolio_code,score_model_code,current_qtr,segment1,deciles=10,len_rolling_window=4):

    ''' 
        This function outputs the PSI table for last 4 quarters including current quarter based on the approach and portfolio selected.
    
        input params:

        @ df -  The dataframe created for input monitoring. 
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        @ current_qtr - This takes string variable for the current monitoring/validation quarter. i.e. "2022Q4"
        @ deciles - This takes integer variable for making deciles for PSI. By default the value is 10.
        @ len_rolling_window - This takes integer variable for the length of the window we want in the rolling psi approach only. By default the value is 4.
        
        
        output params:
        
        It returns a DataFrame for PSI.
    
    '''
    logger.info(f"Calcuation of PSI-HCI for {portfolio_code} portfolio and {score_model_code} scoremodel starts")
        
  
  
    PSI = calculate_psi_score_band(df=df,portfolio_code=portfolio_code,score_model_code=score_model_code,current_qtr=current_qtr,segment1 =segment1, deciles=deciles,len_rolling_window=len_rolling_window)
    
    logger.info(f"Calcuation of PSI-HCI for {portfolio_code} portfolio and {score_model_code} scoremodel ends.")
    logger.info(f"PSI data has {PSI.shape[0]} rows and {PSI.shape[1]} columns.")
    logger.debug("#######################################################################################")
        
    logger.info(f"Whole call_psi function compiled and PSI calculated for {portfolio_code} portfolio and {score_model_code} scoremodel and returned successfully")    

    logger.debug("#######################################################################################")
    

    return PSI


#...................................................................................................................................................................................#


def call_csi(df,portfolio_code,score_model_code,current_qtr,segment,path_file,score_model_wise_variable_dict):

    '''  
        This function outputs the score-band level CSI table. 
    
        input params:

        @ df -  The dataframe created for input monitoring.         
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        @ current_qtr - This takes string variable for the current monitoring/validation quarter. i.e. "2022Q4"
       
        
        output params:
        
        It returns a DataFrame for CSI.
    
    '''
    logger.info(f"Calcuation of CSI for {portfolio_code} portfolio and {score_model_code} scoremodel starts")
    

    
    
   
        
    CSI = calculate_csi_score_band(df=df,portfolio_code=portfolio_code,score_model_code=score_model_code,segment=segment,path_file=path_file,score_model_wise_variable_dict=score_model_wise_variable_dict) 
    
    logger.info(f"Calcuation of CSI for {portfolio_code} portfolio and {score_model_code} scoremodel ends.")
    logger.info(f"CSI has {CSI.shape[0]} rows and {CSI.shape[1]} columns.")
    logger.debug("#######################################################################################")
    
    logger.info(f"Whole call_csi function compiled and CSI calculated for {portfolio_code} portfolio and {score_model_code} scoremodel and returned successfully")    

    logger.debug("#######################################################################################")
    
 
    return CSI

#........................................................................................................................................
def call_population_summary(df,portfolio_code,score_model_code,aggregator_var,current_qtr,segment,deciles=10,bm_year=2020):

    '''  
        This function outputs the Population Summary table for the selected scoremodel.
    
        input params:

        @ df -  The dataframe created for output monitoring.                         
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        @ aggregator_variable - This takes string input and accounts for the column to be used as the variable for grouping our population summary data. i.e. 'SCORE_RANGE' or 'SCORE_VALUE'.
        @ deciles - This takes integer variable for making deciles for PSI. By default the value is 10.
        
        output params:
        
        It returns a DataFrame for Population Summary.
    
    '''

    logger.info(f"Calcuation of Population Summary for {portfolio_code} portfolio and {score_model_code} scoremodel on {aggregator_var} level starts.")
    
    pop_summary_table,val_dict=calculate_population_summary(df=df,portfolio_code=portfolio_code,score_model_code=score_model_code,aggregator_var=aggregator_var,current_qtr=current_qtr,segment =segment,deciles=deciles,bm_year=bm_year)
    
    logger.info(f"Whole call_population_summary function compiled and Population Summary table calculated for {portfolio_code} portfolio and {score_model_code} scoremodel on {aggregator_var} level ends and returned successfully.")    

    logger.debug("#######################################################################################")

    
    return pop_summary_table,val_dict


#...................................................................................................................................................................................#

def call_gini(df,portfolio_code,score_model_code,gini_method,gini_type,bm_gini,segment,deciles=10):

    '''  
        This function outputs the GINI table depending upon the choice of result we want.
    
        input params:

        @ df -  The dataframe created for output monitoring.                 
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        @ gini_method - This is a string variable which takes only 2 input i.e "decile" or "scorepoint"  depending upon the type of result we want.
        @ gini_type - This is a string variable which takes only 2 input i.e "monthly" or "quarterly"  depending upon the type of result we want.
        @ bm_gini - This is float variable for the benchmark gini.
        @ deciles - This takes integer variable for making deciles for GINI. By deafult the value is 10.
       
        
        output params:
        
        It returns a DataFrame for GINI.
    
    '''
    
    logger.info(f"Calcuation of GINI for {portfolio_code} portfolio and {score_model_code} scoremodel with {gini_type} type GINI starts")
    

    if gini_method=="decile":
        
        GINI=calculate_decile_gini(df=df,portfolio_code=portfolio_code,score_model_code=score_model_code,gini_type=gini_type,bm_gini=bm_gini,deciles=deciles)
        logger.info(f"Calcuation of GINI for {portfolio_code} portfolio and {score_model_code} scoremodel with {gini_type} type GINI ends.")
        logger.info(f"GINI has {GINI.shape[0]} rows and {GINI.shape[1]} columns.")
        logger.debug("#######################################################################################")


    elif gini_method=="scorepoint":
        
        GINI=calculate_scorepoint_gini(df=df,portfolio_code=portfolio_code,score_model_code=score_model_code,gini_type=gini_type,bm_gini=bm_gini,segment=segment)
        logger.info(f"Calcuation of GINI for {portfolio_code} portfolio and {score_model_code} scoremodel with {gini_type} type GINI ends.")
        logger.info(f"GINI has {GINI.shape[0]} rows and {GINI.shape[1]} columns.")
        logger.debug("#######################################################################################")
    
    else:
        
        print("Please enter the correct word")
        return 
    
    logger.info(f"Whole call_gini function compiled and GINI calculated for {portfolio_code} portfolio and {score_model_code} scoremodel and returned successfully")    

    logger.debug("#######################################################################################")

    return GINI

#...................................................................................................................................................................................#


def call_mape(df,portfolio_code,val_start_date,val_end_date,score_model_code,mape_type,prod_ttc_pd,segment):

    '''  
        This function outputs the MAPE table depending upon the choice of result we want.
    
        input params:

        @ df -  The dataframe created for output monitoring.                         
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ val_start_date -This takes string variable in the format (YYYY-MM-DD) representing start date for the current validation/monitoring exercise.
        @ val_end_date - This takes string variable in the format (YYYY-MM-DD) representing end date for the current validation/monitoring exercise.
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        @ mape_type - This is a string variable which takes only 2 input i.e "monthly" or "quarterly"  depending upon the type of result we want.
        @ prod_ttc_pd - This takes a float variable which accounts for the implemented in-house production value of TTC PD for the given score-card.
       
        
        output params:
        
        It returns a DataFrame for MAPE.
    
    '''

    logger.info(f"Calcuation of MAPE for {portfolio_code} portfolio and {score_model_code} scoremodel with {mape_type} type MAPE starts")

    
    MAPE=calculate_mape(df=df,portfolio_code=portfolio_code,score_model_code=score_model_code,val_start_date=val_start_date,val_end_date=val_end_date,mape_type=mape_type,prod_ttc_pd=prod_ttc_pd,segment=segment)
    
    logger.info(f"Whole call_mape function compiled and MAPE calculated for {portfolio_code} portfolio and {score_model_code} scoremodel and returned successfully for time period {val_start_date} to {val_end_date}.")    

    logger.debug("#######################################################################################")

    
    return MAPE
    


#...................................................................................................................................................................................#


def call_binomial_test_quarterly(df,portfolio_code,score_model_code,segment,no_of_deciles=10,alpha=0.95,two_tailed=True):

    '''  
        This function outputs the Binomial table depending upon the choice of result we want.
    
        input params:

        @ df -  The dataframe created for output monitoring.                                 
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        @ no_of_deciles - This takes integer value for no of deciles we want to break our quarterly data for calculating binomial test. By default the value is 10.
       	@ alpha - This takes float value for making the Confidence Interval. By default the value is .95.
        @ two_tailed - This takes bool value True or False depending upon the type of test we want to perform. By default the value is True.
        
        output params:
        
        It returns a DataFrame for Binomial Test Results.
    
    '''

    logger.info(f"Calcuation of Binomial Test for {portfolio_code} portfolio and {score_model_code} scoremodel with two tailed being {two_tailed} and CI of {alpha*100} starts")

    Quarterly_Binomial=calculate_binomial_test_quarterly(df=df,portfolio_code=portfolio_code,score_model_code=score_model_code,segment=segment,no_of_deciles=10,alpha=0.95,two_tailed=True)
    
    logger.info(f"Whole call_binomial_test_quarterly function compiled and Binomial table calculated for {portfolio_code} portfolio and {score_model_code} scoremodel and returned successfully.")    

    logger.debug("#######################################################################################")

    
    return Quarterly_Binomial


#...................................................................................................................................................................................#


        
def call_vif(df,portfolio_code,val_start_date,val_end_date,score_model_code,target_var,segment,path_file,score_model_wise_variable_dict):

    '''  
        This function outputs the VIF table for the given score_card.
    
        input params:

        @ df -  The dataframe created for output monitoring.         
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ val_start_date -This takes string variable in the format (YYYY-MM-DD) representing start date for the current validation/monitoring exercise.
        @ val_end_date - This takes string variable in the format (YYYY-MM-DD) representing end date for the current validation/monitoring exercise.
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        @ target_var - This takes the string input for the name of the Target variable column in the dataset.
        
        output params:
        
        It returns a DataFrame for VIF.
    
    '''
    # if (portfolio_code=='mort'):
        
    #     logger.info(f"VIF can not be calculated for {portfolio_code} portfolio as it is a decision tree model.")   
    #     vif_table=pd.DataFrame()

    # else:    
    required_variable_wt = score_model_wise_variable_lists_wt_rng(score_model_code,score_model_wise_variable_dict)[1]
    
    logger.info(f"Calcuation of VIF Test for {portfolio_code} portfolio and {score_model_code} scoremodel having {len(required_variable_wt)} variables as {required_variable_wt} starts.")
    
    
    vif_table=calculate_vif(df=df,portfolio_code=portfolio_code,score_model_code=score_model_code,target_var=target_var,score_model_variable_wt=required_variable_wt ,val_start_date=val_start_date,val_end_date=val_end_date,segment=segment,path_file=path_file)
    
    logger.info(f"Whole call_vif function compiled and VIF table calculated for {portfolio_code} portfolio and {score_model_code} scoremodel having {len(required_variable_wt)} variables as {required_variable_wt} ends and returned successfully.")    

    logger.debug("#######################################################################################")

    
    return vif_table


#...................................................................................................................................................................................#


def call_iv(df,portfolio_code,val_start_date,val_end_date,score_model_code,segment,path_file,score_model_wise_variable_dict):

    '''  
        This function outputs the IV table for the given score_card.
    
        input params:

        @ df -  The dataframe created for output monitoring.                 
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ val_start_date -This takes string variable in the format (YYYY-MM-DD) representing start date for the current validation/monitoring exercise.
        @ val_end_date - This takes string variable in the format (YYYY-MM-DD) representing end date for the current validation/monitoring exercise.
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        
        output params:
        
        It returns a DataFrame for IV.
    
    '''
    logger.info(f"Calcuation of IV values for {portfolio_code} portfolio and {score_model_code} scoremodel starts.")
    
    # if (portfolio_code=='mort'):
        
    #     iv_table=calculate_mort_bad_rate(df=df,portfolio_code=portfolio_code,score_model_code=score_model_code,val_start_date=val_start_date,val_end_date=val_end_date,segment=segment)
        
    #     logger.info(f"Whole call_iv function compiled and IV table calculated for {portfolio_code} portfolio and {score_model_code} scoremodel ends and returned successfully.")    
        
    #     logger.debug("#######################################################################################")

    # else:
        
    iv_table=calculate_iv(df=df,portfolio_code=portfolio_code,score_model_code=score_model_code,val_start_date=val_start_date,val_end_date=val_end_date,segment=segment,path_file=path_file,score_model_wise_variable_dict=score_model_wise_variable_dict)
    
    logger.info(f"Whole call_iv function compiled and IV table calculated for {portfolio_code} portfolio and {score_model_code} scoremodel ends and returned successfully.")    
    
    logger.debug("#######################################################################################")
    

    return iv_table





#...................................................................................................................................................................................#



def call_rank_order(df,portfolio_code,score_model_code,segment,deciles=10):

    '''  
        This function outputs the Rank Order table for the selected scoremodel.
    
        input params:

        @ df -  The dataframe created for output monitoring.                                 
        @ portfolio_code - This takes string variable dpeneding upon the name of the portfolio i.e. "al","nl","pl","mort","cc".
        @ score_model_code - This takes string variable and accounts for the scorecard whose result is to be generated.i.e., "PB_UAE_BSCR04"
        @ deciles - This takes integer variable for making deciles for PSI. By default the value is 10.

        output params:
        
        It returns a DataFrame for Rank Order.
    
    '''
    logger.info(f"Calcuation of Rank Order for {portfolio_code} portfolio and {score_model_code} scoremodel starts.")
    
    rank_order_table=calculate_rank_order(df=df,portfolio_code=portfolio_code,score_model_code=score_model_code,segment=segment,deciles=deciles)

    logger.info(f"Whole call_rank_order function compiled and Rank Order table calculated for {portfolio_code} portfolio and {score_model_code} scoremodel ends and returned successfully.")    

    logger.debug("#######################################################################################")

    
    return rank_order_table


#...................................................................................................................................................................................#

def make_waterfall_exclusion(df,waterfall_type, val_start_date_omm,val_end_date_omm):
    
    def change_name_exclusion(exclusion):
        if "POST" in exclusion:
            exclusion_f=exclusion
        else:
            exclusion_f="(-) " + exclusion
            
        return exclusion_f 
    
    if waterfall_type=="imm":
        overall=df.groupby('YEAR_SCORE')['ACCOUNT_ID'].count().reset_index()
        overall.rename(columns={"ACCOUNT_ID":"Total Observations"},inplace=True)
        overall['Total Observations']=overall['Total Observations'].apply(lambda x: f'{x:,}')
                
        x=pd.pivot_table(data=df,index=["BSCORE_EXCLUSION_REASON"],columns=["YEAR_SCORE"],values=["ACCOUNT_ID"],aggfunc='count').reset_index(drop=False)
        Y1=x['ACCOUNT_ID']
        Y2=x['BSCORE_EXCLUSION_REASON']
        df2=pd.concat([Y2,Y1],axis=1)
        df2['BSCORE_EXCLUSION_REASON']=df2['BSCORE_EXCLUSION_REASON'].str[2:]
        df2=df2.fillna(0)
        for x in overall['YEAR_SCORE']:
            df2[x]=df2[x].astype('int').apply(lambda x: f'{x:,}')
        df2["BSCORE_EXCLUSION_REASON"]=df2["BSCORE_EXCLUSION_REASON"].apply(lambda x: change_name_exclusion(x))
            
        filter_df=df[df['BSCORE_EXCLUSION_REASON'].str.contains("POST_EXCLUSION")]  
        x_segment=pd.pivot_table(data=filter_df,index=["SCORE_MODEL_CODE"],columns=["YEAR_SCORE"],values=["ACCOUNT_ID"],aggfunc='count').reset_index(drop=False)
        Y1_seg=x_segment['ACCOUNT_ID']
        Y2_seg=x_segment['SCORE_MODEL_CODE']
        df2_seg=pd.concat([Y2_seg,Y1_seg],axis=1)
        for x in overall['YEAR_SCORE']:
            df2_seg[x]=df2_seg[x].astype('int').apply(lambda x: f'{x:,}')
        
        
    else:
        df['SCORE_DATE']=pd.to_datetime(df['SCORE_DATE'])
        df=df[(df['SCORE_DATE']>= val_start_date_omm) & (df['SCORE_DATE']<=val_end_date_omm)]

        overall=df.shape[0]
        overall=f'{overall:,}'

        df2=pd.pivot_table(data=df,index=["BSCORE_EXCLUSION_REASON"],values=["ACCOUNT_ID"],aggfunc='count').reset_index(drop=False)
        df2['BSCORE_EXCLUSION_REASON']=df2['BSCORE_EXCLUSION_REASON'].str[2:]
        df2['ACCOUNT_ID']=df2['ACCOUNT_ID'].astype('int').apply(lambda x: f'{x:,}')
        df2["BSCORE_EXCLUSION_REASON"]=df2["BSCORE_EXCLUSION_REASON"].apply(lambda x: change_name_exclusion(x))
        filter_df=df[df['BSCORE_EXCLUSION_REASON'].str.contains("POST_EXCLUSION")]
        df2_seg=pd.pivot_table(data=filter_df,index=["SCORE_MODEL_CODE"],values=["ACCOUNT_ID"],aggfunc='count').reset_index(drop=False)
        df2_seg['ACCOUNT_ID']=df2_seg['ACCOUNT_ID'].astype('int').apply(lambda x: f'{x:,}')   
    
    return overall,df2,df2_seg
#...................................................................................................................................................................................#
    

def df_psi_hci_chart(df,y):
    
        
    ''' 
        This function outputs bar chart for PSI & HCI.
    
        input params:
           
        @ df - The corresponding PSI-HCI data generated from the code.
        @ x - This takes string variable denoting the X axis of the graph.
        @ y - This takes string variable denoting the Y axis of the graph.
        @ title - This takes string variable denoting the Title of the graph.
        @ score_model_code - This takes string variable and accounts for the scorecard whose graph is to be saved.i.e., "PB_UAE_BSCR04"    
        output params:
        
        It returns a Bar Graph for the PSI or HCI.
    
    '''
    
    
    pd.set_option('display.float_format', lambda x: '%.5f' % x)
    df_1=df[['YYYYMM','PSI','HCI']].drop_duplicates().reset_index(drop=True)
    
    initial_color="#5B9BD5"
    count_blue=df_1['YYYYMM'].nunique()-1
    color_list=[initial_color]*count_blue
    
    if(y=='PSI'):
        
        val_value=df_1["PSI"].to_list()[-1]
        if(val_value<=0.1):
            final_color="#7F7F7F"
            document_color="#00B050"
        elif(val_value<=0.25):
            final_color="#7F7F7F"
            document_color="#FFC000"
        else:
            final_color="#7F7F7F"  
            document_color="#FF0000"
            
    if(y=='HCI'):
        
        val_value=df_1["HCI"].to_list()[-1]
        if(val_value<=0.2):
            final_color="#7F7F7F"
            document_color="#00B050"
        elif(val_value<=0.3):
            final_color="#7F7F7F"
            document_color="#FFC000"
        else:
            final_color="#7F7F7F"          
            document_color="#FF0000"
            
    color_list.append(final_color) 
    
    return df_1, color_list, document_color

#...................................................................................................................................................................................#


def plot_psi_hci_chart(df,x,y,title,score_model_code,path_folder):
    
        
    ''' 
        This function outputs bar chart for PSI & HCI.
    
        input params:
           
        @ df - The corresponding PSI-HCI data generated from the code.
        @ x - This takes string variable denoting the X axis of the graph.
        @ y - This takes string variable denoting the Y axis of the graph.
        @ title - This takes string variable denoting the Title of the graph.
        @ score_model_code - This takes string variable and accounts for the scorecard whose graph is to be saved.i.e., "PB_UAE_BSCR04"    
        output params:
        
        It returns a Bar Graph for the PSI or HCI.
    
    '''
    
    logger.info(f"Plotting of {y} graph begins for {score_model_code} scorecard.")
    
    pd.set_option('display.float_format', lambda x: '%.5f' % x)

    sns.set(style="white", palette="dark",font_scale=0.9,rc={"figure.figsize":(15, 6)})
    
    df_1,color_list,document_color=df_psi_hci_chart(df=df,y=y)    
    
    # bar_plot = sns.barplot(data=df_1,x=x, y=y, palette=color_list,width=0.4,dodge=0.15,linewidth=1)
    
    
    if y=='PSI':
        
        val_value=df_1[y].to_list()[-1]
        if(val_value<=0.1):
            line_plot = sns.lineplot(x=df_1[x], y=[0.1]*len(df_1),color="#002060",linestyle='dashed',linewidth=0.75)
            plt.text(1.5,0.15,'RAG Cut off 10%',horizontalalignment='center',verticalalignment='center',fontsize=10,fontweight=1000,fontstretch='semi-condensed',wrap=True,color='#002060')

        elif(val_value<=0.25):
            line_plot = sns.lineplot(x=df_1[x], y=[0.1]*len(df_1),color="#002060",linestyle='dashed',linewidth=0.75)
            line_plot = sns.lineplot(x=df_1[x], y=[0.25]*len(df_1),color="#FFC000",linestyle='dashed',linewidth=0.75)
            plt.text(1.5,.24,'RAG Cut off 10%-25%',horizontalalignment='center',verticalalignment='center',fontsize=7,fontweight=1000,fontstretch='condensed',wrap=True,color='#002060')
            
        else:
            line_plot = sns.lineplot(x=df_1[x], y=[0.25]*len(df_1),color="#002060",linestyle='dashed',linewidth=0.75) 
            plt.text(1.5,.3,'RAG Cut off 25%',horizontalalignment='center',verticalalignment='center',fontsize=7,fontweight=1000,fontstretch='condensed',wrap=True,color='#002060')
        
    else:
        val_value=df_1[y].to_list()[-1]
        if(val_value<=0.2):
            line_plot = sns.lineplot(x=df_1[x], y=[0.2]*len(df_1),color="#002060",linestyle='dashed',linewidth=0.75)
            plt.text(1.5,0.25,'RAG Cut off 20%',horizontalalignment='center',verticalalignment='center',fontsize=15,fontweight='bold',fontstretch='ultra-condensed',wrap=True,color='#002060')

        elif(val_value<=0.3):
            line_plot = sns.lineplot(x=df_1[x], y=[0.2]*len(df_1),color="#002060",linestyle='dashed',linewidth=0.75)
            line_plot = sns.lineplot(x=df_1[x], y=[0.3]*len(df_1),color="#FFC000",linestyle='dashed',linewidth=0.75)
            plt.text(1.5,.29,'RAG Cut off 20%-30%',horizontalalignment='center',verticalalignment='center',fontsize=15,fontweight='bold',fontstretch='ultra-condensed',wrap=True,color='#002060')
            
        else:
            line_plot = sns.lineplot(x=df_1[x], y=[0.3]*len(df_1),color="#002060",linestyle='dashed',linewidth=0.75) 
            plt.text(1.5,.35,'RAG Cut off 30%',horizontalalignment='center',verticalalignment='center',fontsize=15,fontweight='bold',fontstretch='ultra-condensed',wrap=True,color='#002060')
    
    logger.debug(f"Bar graph plotted for {len(df_1[x])} quarters.")

    # for p in bar_plot.patches:
    #     value=p.get_height()
    #     bar_plot.annotate(f'{value:.2%}', (p.get_x() + p.get_width() / 2., p.get_height()),
    #                 ha='center', va='center', xytext=(0, 4), textcoords='offset points', color='#002060',fontsize=10)
    

    # bar_plot.set(xlabel=None,ylabel=None)
    # bar_plot.yaxis.set_visible(False)
    # bar_plot.set_xticklabels(bar_plot.get_xticklabels(),fontsize=12,color='#404040')
    # bar_plot.set_ylim(top=max(df_1[y])+.15)
    # plt.title(title, {'fontsize':10,'color':"#002060",'fontweight':1000})
    # for spine in bar_plot.spines.values():
    #     spine.set_edgecolor("#BFBFBF")
    #     spine.set_linewidth(0.75)
    

    # plt.show()
    if y=='PSI':
        bar_plot = sns.barplot(data=df_1,x=x, y=y, palette=color_list,width=0.6,dodge=0.05,linewidth=1)
        for p in bar_plot.patches:
            value=p.get_height()
            bar_plot.annotate(f'{value:.2%}', (p.get_x() + p.get_width() / 2., p.get_height()),
                    ha='center', va='center', xytext=(0, 4), textcoords='offset points', color='#002060',fontsize=10)
        bar_plot.set(xlabel=None,ylabel=None)
        bar_plot.yaxis.set_visible(False)
        bar_plot.set_xticklabels(bar_plot.get_xticklabels(),fontsize=10,color='#404040')
        bar_plot.set_ylim(top=max(df_1[y])+.15)
        plt.title(title, {'fontsize':10,'color':"#002060",'fontweight':1000})
        for spine in bar_plot.spines.values():
            spine.set_edgecolor("#BFBFBF")
            spine.set_linewidth(0.75)
        plt.savefig(f'{path_folder}/nl/PSI_{score_model_code}.png',bbox_inches='tight')
    elif y=='HCI':
        bar_plot = sns.barplot(data=df_1,x=x, y=y, palette=color_list,width=0.6,dodge=0.05,linewidth=1)
        for p in bar_plot.patches:
            value=p.get_height()
            bar_plot.annotate(f'{value:.2%}', (p.get_x() + p.get_width() / 2., p.get_height()),
                    ha='center', va='center', xytext=(0, 4), textcoords='offset points', color='#002060',fontsize=15)
    

        bar_plot.set(xlabel=None,ylabel=None)
        bar_plot.yaxis.set_visible(False)
        bar_plot.set_xticklabels(bar_plot.get_xticklabels(),fontsize=20,color='#404040')
        bar_plot.set_ylim(top=max(df_1[y])+.15)
        plt.title(title, {'fontsize':20,'color':"#002060",'fontweight':1000})
        for spine in bar_plot.spines.values():
            spine.set_edgecolor("#BFBFBF")
            spine.set_linewidth(1)#.75
        plt.savefig(f'{path_folder}/nl/HCI_{score_model_code}.png',bbox_inches='tight')
    
    logger.info(f"Plotting of {y} graph ends for {score_model_code} scorecard.")
    logger.debug("#######################################################################################")
    
    return

#...................................................................................................................................................................................#

def df_csi_chart(df,current_qtr):
    
    def rolling_window(current_qtr):
        year = int(current_qtr[:4])
        qtr = int(current_qtr[-1])
        window = [current_qtr]
        for i in range(4-1):
            if qtr>1:
                qtr=qtr-1
            else:
                qtr=4
                year=year-1
            window.append(str(year)+'Q'+str(qtr))
            window=sorted(window)
        return window
    
    sample=rolling_window(current_qtr=current_qtr)
    sample+=["Validation_Sample"]
    
    df_1=df[df["YYYYMM"].isin(sample)]
    df_1.sort_values(by=["VARIABLE_NAME","YYYYMM"],inplace=True)
    del df_1['SCORE_CARD']
    df_1=df_1.reset_index(drop=True)
    df_1=df_1[["YYYYMM","VARIABLE_NAME","CSI"]].drop_duplicates()

    
    return df_1
    
#...................................................................................................................................................................................#

def plot_csi_chart(df,x,y,title,current_qtr,score_model_code,var_name,path_folder):
            

    logger.info(f"Plotting of {y} graph begins for {score_model_code} scorecard.")

    sns.set(style="white", palette="dark",font_scale=0.9,rc={"figure.figsize":(6, 3.5)})
        
    df=df_csi_chart(df=df,current_qtr=current_qtr)
    
    df_1=df[df["VARIABLE_NAME"]==var_name]
    uniq_qtrs=df_1['YYYYMM'].nunique()-1
    
    color_1="#5B9BD5"
    color_list=[color_1]*uniq_qtrs
    final_color="#7F7F7F"
    color_list.append(final_color)

    bar_plot = sns.barplot(data=df_1,x=x, y=y, palette=color_list,width=0.55,gap=0.2)
    
    # y_axis=max(df_1[y])+1
    
    logger.debug(f"Bar graph plotted for {len(df[x])} quarters.")

    for p in bar_plot.patches:
        value=p.get_height()
        bar_plot.annotate(f'{value:.2%}', (p.get_x() + p.get_width() / 2., p.get_height()),
                    ha='center', va='center', xytext=(0, 4), textcoords='offset points', color='#002060',fontsize=7)
    
    bar_plot.set(xlabel=None,ylabel=None)
    bar_plot.yaxis.set_visible(False)
    bar_plot.set_xticklabels(bar_plot.get_xticklabels(),fontsize=9,color='#404040')
    
    if (max(df_1[y])<0.1):
        bar_plot.set_ylim(top=max(df_1[y])+.05)        
    elif (max(df_1[y])<2):
        bar_plot.set_ylim(top=max(df_1[y])+.15)
    else:
        bar_plot.set_ylim(top=max(df_1[y])+1.2)
    plt.title(title, {'fontsize':8,'color':"#002060",'fontweight':1000})
    for spine in bar_plot.spines.values():
        spine.set_edgecolor("#BFBFBF")
        spine.set_linewidth(0.75)


    # plt.show()
    plt.savefig(f'{path_folder}/nl/CSI_{score_model_code}_{var_name}.png',bbox_inches='tight')
    logger.info(f"Plotting of {y} graph ends for {score_model_code} scorecard.")
    logger.debug("#######################################################################################")
    
    return 

#...................................................................................................................................................................................#

def df_gini_chart(df,score_model_code,score_model_wise_bm_gini_dict,z=None):
            
    ''' 
        This function outputs bar chart for GINI.
    
        input params:
           
        @ df - The corresponding GINI data generated from the code.
        @ x - This takes string variable denoting the X axis of the graph.
        @ y - This takes string variable denoting the Y axis of the graph.
        @ title - This takes string variable denoting the Title of the graph.
        @ score_model_code - This takes string variable and accounts for the scorecard whose benchmark gini is to be fetched and graph is to be saved.i.e., "PB_UAE_BSCR04"
        output params:
        
        It returns a Bar Graph for the GINI.
    
    '''
    
    bm_gini=score_model_wise_benchmark_gini(score_model_code,score_model_wise_bm_gini_dict)
    bm_gini_row={'QUARTER':'Benchmark_Sample','GINI':bm_gini}
    df_1=df.copy(deep=True)
    df_1.loc[df.shape[0]]=bm_gini_row
 
    initial_color="#5B9BD5"
    #count_blue = 0 # Assign default value
 
    if z == 'DATE':
        count_blue=df['DATE'].nunique()-1
    elif z == 'SCORE_DATE':
        count_blue=df['SCORE_DATE'].nunique()-1
    
    elif z=="QUARTER":
        count_blue=df['QUARTER'].nunique()-1
    
    color_list=[initial_color]*count_blue
    
    final_val=df['GINI'].to_list()[-1]
    
    per_change_gini=abs((final_val-bm_gini)/bm_gini)
    
    if(final_val-bm_gini>0):
        final_color="#7F7F7F"
        document_color="#00B050"
        per_change_gini=0-per_change_gini
    
    else:
    
        if(per_change_gini<=0.2):
            final_color="#7F7F7F"
            document_color="#00B050"
        elif(per_change_gini<=0.3):
            final_color="#7F7F7F"
            document_color="#FFC000"
        else:
            final_color="#7F7F7F"
            document_color="#FF0000" 
        
    color_list.append(final_color)
    
    return df_1, color_list, per_change_gini, document_color



#...................................................................................................................................................................................#


def plot_gini_chart(df,x,y,title,score_model_code,path_folder,score_model_wise_bm_gini_dict,z=None):
            
    ''' 
        This function outputs bar chart for GINI.
    
        input params:
           
        @ df - The corresponding GINI data generated from the code.
        @ x - This takes string variable denoting the X axis of the graph.
        @ y - This takes string variable denoting the Y axis of the graph.
        @ title - This takes string variable denoting the Title of the graph.
        @ score_model_code - This takes string variable and accounts for the scorecard whose benchmark gini is to be fetched and graph is to be saved.i.e., "PB_UAE_BSCR04"
        output params:
        
        It returns a Bar Graph for the GINI.
    
    '''
    logger.info(f"Plotting of {y} graph begins for {score_model_code} scorecard.")

    sns.set(style="white", palette="dark",font_scale=0.9,rc={"figure.figsize":(18,7.5)})
        
    df,color_list,per_change_gini,document_color=df_gini_chart(df=df,score_model_code=score_model_code,score_model_wise_bm_gini_dict=score_model_wise_bm_gini_dict,z=z)
    df_1=df.iloc[:-1,:]
    bar_plot = sns.barplot(data=df_1,x=x, y=y, palette=color_list,width=0.25,gap=0.2)
    
    bm_gini=score_model_wise_benchmark_gini(score_model_code,score_model_wise_bm_gini_dict)
    y_axis=max(df_1[y])+.055
    if z!= None:
         line_plot = sns.lineplot(x=df_1[x].astype(str), y=[bm_gini]*len(df_1),color="#002060",linestyle='dashed',linewidth=0.75)
    else:
         line_plot = sns.lineplot(x=df_1[x], y=[bm_gini]*len(df_1),color="#002060",linestyle='dashed',linewidth=0.75)

    #line_plot = sns.lineplot(x=df_1[x], y=[bm_gini]*len(df_1),color="#002060",linestyle='dashed',linewidth=0.75)
    text=f"Benchmark GINI {round(bm_gini*100,2)}%"
    plt.text(2.5,y_axis,text,horizontalalignment='center',verticalalignment='center',fontsize=16,fontweight='bold',fontstretch='ultra-condensed',wrap=True,color='#002060')
    
    logger.debug(f"Bar graph plotted for {len(df[x])} quarters.")

    for p in bar_plot.patches:
        value=p.get_height()
        bar_plot.annotate(f'{value:.2%}', (p.get_x() + p.get_width() / 2., p.get_height()),
                    ha='center', va='center', xytext=(0, 4), textcoords='offset points', color='#002060',fontsize=10)
    
    bar_plot.set(xlabel=None,ylabel=None)
    bar_plot.yaxis.set_visible(False)
    bar_plot.set_xticklabels(bar_plot.get_xticklabels(),fontsize=12,color='#404040')
    bar_plot.set_ylim(top=max(df_1[y])+.02)
    plt.title(title, {'fontsize':16,'color':"#002060",'fontweight':1000})
    for spine in bar_plot.spines.values():
        spine.set_edgecolor("#BFBFBF")
        spine.set_linewidth(0.75)


    # plt.show()
    plt.savefig(f'{path_folder}/nl/GINI_{score_model_code}.png',bbox_inches='tight')
    
    logger.info(f"Plotting of {y} graph ends for {score_model_code} scorecard.")
    logger.debug("#######################################################################################")
    
    return 

#...................................................................................................................................................................................#

def df_rank_order(df):
            
    ''' 
        This function outputs line graph for Rank-Order.
    
        input params:
           
        @ df - The corresponding Rank-Order data generated from the code.
        @ x - This takes string variable denoting the X axis of the graph.
        @ y - This takes string variable denoting the Y axis of the graph.
        @ title - This takes string variable denoting the Title of the graph.
        @ score_model_code - This takes string variable and accounts for the scorecard whose graph is to be saved.i.e., "PB_UAE_BSCR04"        
        output params:
        
        It returns a Line Graph for the Rank-Order.
    
    '''
    
    df_1=df.copy(deep=True)
    df_1=df_1[df_1['QTR']!='Benchmark']
    df_1['SCORE_RANGE']=df_1['SCORE_RANGE'].astype(str)
    df_1['BAD_RATE']=(df_1['BAD_RATE']*100)
 
    df_1_1=df_1[df_1['QTR']!='Validation_Sample']
    df_1_2=df_1[df_1['QTR']=='Validation_Sample']
    
    return df_1_1,df_1_2
    
#...................................................................................................................................................................................#

def plot_rank_order(df,x,y,title,score_model_code,path_folder):
            
    ''' 
        This function outputs line graph for Rank-Order.
    
        input params:
           
        @ df - The corresponding Rank-Order data generated from the code.
        @ x - This takes string variable denoting the X axis of the graph.
        @ y - This takes string variable denoting the Y axis of the graph.
        @ title - This takes string variable denoting the Title of the graph.
        @ score_model_code - This takes string variable and accounts for the scorecard whose graph is to be saved.i.e., "PB_UAE_BSCR04"        
        output params:
        
        It returns a Line Graph for the Rank-Order.

    '''

    logger.info(f"Plotting of Rank Order graph begins for {score_model_code} scorecard.")
    
    sns.set(style="white", palette="dark",font_scale=0.9,rc={"figure.figsize":(27, 17)})
    
    fig,ax1=plt.subplots()
    
    df_1_1,df_1_2=df_rank_order(df=df)
    df_1_2.rename(columns={'CONCENTRATION_POPULATION':'Population(%)'},inplace=True)
    df_1_2['Population(%)']=round(df_1_2['Population(%)']*100,2)
    
    bar_plot= sns.barplot(data=df_1_2,x=x, y='Population(%)',ax=ax1, color='#002060',width=.4,gap=0.15,linewidth=0.75,label="Population(%)")
    ax2=ax1.twinx()
    line_plot = sns.lineplot(data=df_1_2,x=x, y=y,ax=ax2,color="#C00000",linewidth=3,label=y)
    
    # for index, row in df_1_2.iterrows():
    #     plt.annotate(f'{row[y]:.2f}%',(row[x],row[y]),textcoords="offset points",xytext=(2,4), ha='center',fontsize=8,color='#002060')
        
    # line_plot = sns.lineplot(data=df_1_1,x=x, y=y,hue='QTR',palette="flare",linestyle='dashed',linewidth=0.8)
    
    logger.debug(f"Line graph plotted for {len(df_1_2[x])} score bands.")

    plt.title(title, {'fontsize':25,'color':"#002060",'fontweight':1000})
    ax1.set_xlabel(x,{'fontsize':25,'color':'#002060','fontweight':1000},labelpad=10)
    ax1.set_ylabel('Population(%)',{'fontsize':25,'color':'#002060','fontweight':1000},labelpad=10)
    ax2.set_ylabel(y+"(%)",{'fontsize':25,'color':'#002060','fontweight':1000},labelpad=10)   
     
    ax1.legend(loc='upper right',fontsize=16)
    ax2.legend(loc='upper right',fontsize=16,bbox_to_anchor=(1,0.925))
    ax2.tick_params(axis='y', labelsize=25) 
    ax2.set_ylabel(y +"(%)", {'fontsize':25,'color':'#002060','fontweight':1000}, labelpad=10) 
    bar_plot.set_yticklabels(bar_plot.get_yticklabels(),fontsize=25)
    bar_plot.set_xticklabels(bar_plot.get_xticklabels(),fontsize=25,rotation=90)
    #bar_plot.tick_params(axis='y',labelsize=25)

    for spine in line_plot.spines.values():
        spine.set_edgecolor("#BFBFBF")
        spine.set_linewidth(0.75)
    # plt.show()
    plt.savefig(f'{path_folder}/nl/RANK_ORDER_{score_model_code}.png',bbox_inches='tight')
    
    logger.info(f"Plotting of Rank Order graph ends for {score_model_code} scoremodel.")
    logger.debug("#######################################################################################")
    
    return 

#...................................................................................................................................................................................#

def df_mape_chart(df,x,y,score_model_code,score_model_wise_imp_pd_dict):
                
    ''' 
        This function outputs line graph for MAPE.
    
        input params:
           
        @ df - The corresponding MAPE data generated from the code.
        @ x - This takes string variable denoting the X axis of the graph.
        @ y - This takes string variable denoting the Y axis of the graph.
        @ title - This takes string variable denoting the Title of the graph.
        @ score_model_code - This takes string variable and accounts for the scorecard whose production ttc pd is to be fetched and graph to be saved.i.e., "PB_UAE_BSCR04"
        output params:
        
        It returns a Line Graph for the MAPE.
    
    '''
    
    df_1=df.copy(deep=True)
    df_1['LT_ODR']=df_1['ObservedDefaults'].sum()/df_1['Total_Acc'].sum()
    df_1['LT_EDR']=df_1['ExpectedDefaults'].sum()/df_1['Total_Acc'].sum()
    
    prod_ttc_pd=score_model_wise_implemented_ttc_pd(score_model_code,score_model_wise_imp_pd_dict)
    df_1['PROD_TTC_PD']=prod_ttc_pd
    
    df_1=pd.melt(frame=df_1,value_vars=['EDR','ODR','LT_EDR','LT_ODR','PROD_TTC_PD'],id_vars=[x],value_name=y)
    df_1[y]=(df_1[y]*100)
  
       
    df_1_1=df_1[df_1['variable'].isin(['EDR','ODR'])]
    df_1_2=df_1[df_1['variable'].isin(['LT_EDR','LT_ODR','PROD_TTC_PD'])]
    
    return df_1_1, df_1_2


#...................................................................................................................................................................................#


def plot_mape(df,x,y,title,score_model_code,path_folder,score_model_wise_imp_pd_dict):
                
    ''' 
        This function outputs line graph for MAPE.
    
        input params:
           
        @ df - The corresponding MAPE data generated from the code.
        @ x - This takes string variable denoting the X axis of the graph.
        @ y - This takes string variable denoting the Y axis of the graph.
        @ title - This takes string variable denoting the Title of the graph.
        @ score_model_code - This takes string variable and accounts for the scorecard whose production ttc pd is to be fetched and graph to be saved.i.e., "PB_UAE_BSCR04"
        output params:
        
        It returns a Line Graph for the MAPE.
    
    '''
    
    logger.info(f"Plotting of MAPE graph begins for {score_model_code} scorecard")
    sns.set(style="white", palette="dark",font_scale=0.9,rc={"figure.figsize":(12.45,4.2)})
    
    df[x]=df[x].astype(str)
    xlabels=df[x].to_list()
    
    

    df_1_1,df_1_2=df_mape_chart(df=df,x=x,y=y,score_model_code=score_model_code,score_model_wise_imp_pd_dict=score_model_wise_imp_pd_dict)
    print(df_1_1)
    line_plot = sns.lineplot(data=df_1_1,x=x, y=y,palette=["#000000","#002060"],hue='variable',linewidth=2.5)
    line_plot = sns.lineplot(data=df_1_2,x=x, y=y,hue='variable',palette=["#000000","#002060","#FF0000"],linestyle='dashed',linewidth=1.5)
    
    logger.debug(f"Line graph plotted for {len(df[x])} score dates.")
    
    plt.title("Long Term Default Rate Series", {'fontsize':12,'color':"#002060",'fontweight':1000})
    line_plot.set_xticklabels(xlabels,fontsize=9,rotation=90,ha='right',color='#404040')
    line_plot.set_yticklabels(line_plot.get_yticklabels(),fontsize=9,color='#404040')
    # line_plot.set(xlabel=None,ylabel=None)
    
    logger.debug(f"No of xticks is {len(xlabels)}")
    logger.debug(f"No of yticks is {len(line_plot.get_yticklabels())}")
    
    line_plot.set_xlabel(x,{'fontsize':10,'color':"#002060"},labelpad=10)
    line_plot.set_ylabel(y,{'fontsize':10,'color':"#002060"},labelpad=10)
    line_plot.legend(loc='upper right',fontsize=8)
    for spine in line_plot.spines.values():
        spine.set_edgecolor("#BFBFBF")
        spine.set_linewidth(0.75)

    # plt.show()
    plt.savefig(f'{path_folder}/nl/MAPE_{score_model_code}.png',bbox_inches='tight')
    
    logger.info(f"Plotting of MAPE graph ends for {score_model_code} scorecard.")
    logger.debug("#######################################################################################")
    
    return 

#...................................................................................................................................................................................#



def plotting_imm_monitoring_charts(df,y_psi,y_hci,title_name,score_model_code):
    
    ''' 
        This function pastes charts in the PPT.
    
        input params:
           
        @ title_name - This takes string variable denoting the Title of the slide.
        @ score_model_code - This takes string variable and accounts for the scorecard whose charts are to be pasted.i.e., "PB_UAE_BSCR04"
        @ ppt_file_path_name - This takes string variable for the path and file name of the ppt.Ex- "ppt/final_ppt_1.pptx"
        
        output params:
        
        It returns a saved presentation at given file path.
    
    '''
    logger.info(f"Pasting of monitoring charts in ppt begins for {score_model_code} scorecard")
        
    imm_title=f"{title_name} – Input Model Monitoring: Model Level"
    
    slide_width=Inches(13.33)
    slide_height=Inches(7.5)
    
    slide_layout = presentation.slide_layouts[5] # Title only
    slide = presentation.slides.add_slide(slide_layout)
    presentation.slide_width, presentation.slide_height=slide_width,slide_height
    
    logger.debug("New slide created")
    
    left_inch = Inches(0.30)
    top_inch = Inches(0.45)
    width=Inches(12.7)
    height=Inches(0.45)
    title = slide.shapes.title
    title.left=left_inch
    title.top=top_inch
    title.width=width
    title.height=height
    title.text = imm_title

    title_text_frame=title.text_frame
    title_text_frame.paragraphs[0].font.size=Pt(22)
    title_text_frame.paragraphs[0].font.bold=True
    title_text_frame.paragraphs[0].font.color.rgb=RGBColor(47,85,151)
    title_text_frame.paragraphs[0].alignment=PP_ALIGN.LEFT

    line_left=left_inch
    line_left_top=top_inch+height
    line_right=left_inch+width
    line_right_top=top_inch+height
    
    line = slide.shapes.add_connector(
    MSO_CONNECTOR.STRAIGHT, line_left,line_left_top, line_right, line_right_top)
    
    line.line.color.rgb=RGBColor(47,85,151)
    line.line.width=Pt(2.25)
    
    logger.debug("Title added")
    
    # Add first graph
    img_stream_1 = f"{path_folder}/nl/PSI_{score_model_code}.png"
    left_inch = Inches(0.50)
    top_inch = Inches(1)
    pic = slide.shapes.add_picture(img_stream_1, left=left_inch, top=top_inch, width=Inches(6), height=Inches(3.7))
    outline=pic.line
    outline.color.rgb=RGBColor(217,217,217)
    outline.width=Pt(.75)
    # Add second graph
    
    logger.debug("PSI graph pasted")
    
    img_stream_2 = F"{path_folder}/nl/HCI_{score_model_code}.png"
    left_inch = Inches(6.95)
    top_inch = Inches(1)
    pic = slide.shapes.add_picture(img_stream_2, left=left_inch, top=top_inch, width=Inches(6), height=Inches(3.7))
    outline=pic.line
    outline.color.rgb=RGBColor(217,217,217)
    outline.width=Pt(.75)
    
    logger.debug("HCI graph pasted")  
    logger.debug("###############################################")
    
    # Text Box No 1.
    #Adding Outline
    left_inch = Inches(0.50)
    top_inch = Inches(4.95)
    text_box_1 = slide.shapes.add_textbox(left=left_inch, top=top_inch, width=Inches(6), height=Inches(1.75))
    
    text_frame_1=text_box_1.text_frame    
    outline=text_box_1.line
    outline.color.rgb=RGBColor(217,217,217)
    outline.width=Pt(.75)
    
    #Adding Text Frame, and title of text box
    paragraph=text_frame_1.paragraphs[0]
    title="Comments"
    paragraph.text=title
    paragraph.font.size=Pt(11)
    paragraph.font.underline=True
    paragraph.alignment=PP_ALIGN.CENTER
    paragraph.font.color.rgb=RGBColor(0,0,0)
    # text_frame_1.margin_left=Pt(10)

    
    df_1_psi,color_list_psi,document_color=df_psi_hci_chart(df=df,y=y_psi) 
    df_1_hci,color_list_hci,document_color=df_psi_hci_chart(df=df,y=y_hci)
    val_value_psi=df_1_psi["PSI"].to_list()[-1]
    val_value_hci=df_1_hci["HCI"].to_list()[-1]
    ini_qtr=df_1_psi["YYYYMM"].to_list()[0]
    fin_qtr=df_1_psi["YYYYMM"].to_list()[-1]
    
    if (color_list_psi[-1]=="#FFC000"):
        status_psi="MEDIUM"
    elif (color_list_psi[-1]=="#FF0000"):
        status_psi="HIGH"
    else:
        status_psi="LOW"
        
    if (color_list_hci[-1]=="#FFC000"):
        status_hci="MEDIUM"
    elif (color_list_hci[-1]=="#FF0000"):
        status_hci="HIGH"
    else:
        status_hci="LOW"
    
    first_point=f" For the entire observation time period ({ini_qtr} to {fin_qtr}), PSI is {round(val_value_psi*100,2)}%."
    second_point=f" Based on the evidence, PSI for the model has been assigned a RAG status of {status_psi}."
    
    bullet_points = [first_point,
    second_point]
    
    paragraph.add_run().text="\n"
    
    text_frame=text_box_1.text_frame
    text_frame.word_wrap = True

    
    for points in bullet_points:
                
        paragraph=text_frame.add_paragraph()
        run = paragraph.add_run()
        run.text=u'\u2022'
        run.font.size=Pt(9)
        run.font.color.rgb=RGBColor(0,0,0)
        run.font.name='Calibri'

        
        if(points==second_point):
            words=points.split()
            for j, word in enumerate(words):
                run=paragraph.add_run()
                if j==0:
                    run.text=" " + word + " "
                else:
                    run.text=word + " "

                if j== len(words)-1:
                    if(words[j]=="LOW."):
                        run.font.color.rgb = RGBColor(0,176,80)
                    elif(words[j]=="MEDIUM."):
                        run.font.color.rgb = RGBColor(255,192,0)
                    else:
                        run.font.color.rgb = RGBColor(255,0,0)
                        
                    run.font.size=Pt(11)
                    run.font.underline=True
                    run.font.bold=True
                else:                    
                    run.font.color.rgb = RGBColor(0,0,0)
                    run.font.size=Pt(11)
                    
        
        else:
        
            comment_run=paragraph.add_run()
            comment_run.text=points           
            
            comment_run.font.size=Pt(11)
            comment_run.alignment=PP_ALIGN.LEFT
            comment_run.font.color.rgb=RGBColor(0,0,0) 
            comment_run.level=Pt(2)
            
            text_frame.margin_left=Pt(10)

    
    left_inch = Inches(6.95)
    top_inch = Inches(4.95)
    text_box_2 = slide.shapes.add_textbox(left=left_inch, top=top_inch, width=Inches(6), height=Inches(1.75))
    
    text_frame_1=text_box_2.text_frame
    outline=text_box_2.line
    outline.color.rgb=RGBColor(217,217,217)
    outline.width=Pt(.75)
    
    #Adding Text Frame, and title of text box
    paragraph=text_frame_1.paragraphs[0]
    title="Comments"
    paragraph.text=title
    paragraph.font.size=Pt(11)
    paragraph.font.underline=True
    paragraph.alignment=PP_ALIGN.CENTER
    paragraph.font.color.rgb=RGBColor(0,0,0)
    # text_frame_1.margin_left=Pt(10)
    
    first_point=f" For the entire observation time period ({ini_qtr} to {fin_qtr}), HCI is {round(val_value_hci*100,2)}%."
    second_point=f" Based on the evidence, HCI for the model has been assigned a RAG status of {status_hci}."
 
    bullet_points = [first_point,
    second_point]
    
    paragraph.add_run().text="\n"
    
    text_frame=text_box_2.text_frame
    text_frame.word_wrap = True
    
    for points in bullet_points:
                
        paragraph=text_frame.add_paragraph()
        
        run = paragraph.add_run()
        run.text=u'\u2022'
        run.font.size=Pt(9)
        run.font.color.rgb=RGBColor(0,0,0)
        run.font.name='Calibri'
        
        if(points==second_point):
            words=points.split()
            for j, word in enumerate(words):
                run=paragraph.add_run()
                if j==0:
                    run.text=" " + word + " "
                else:
                    run.text=word + " "

                if j== len(words)-1:
                    if(words[j]=="LOW."):
                        run.font.color.rgb = RGBColor(0,176,80)
                    elif(words[j]=="MEDIUM."):
                        run.font.color.rgb = RGBColor(255,192,0)
                    else:
                        run.font.color.rgb = RGBColor(255,0,0)
                        
                    run.font.size=Pt(11)
                    run.font.underline=True
                    run.font.bold=True
                else:                    
                    run.font.color.rgb = RGBColor(0,0,0)
                    run.font.size=Pt(11)
                    
        
        else:
        
            comment_run=paragraph.add_run()
            comment_run.text=points           
            
            comment_run.font.size=Pt(11)
            comment_run.alignment=PP_ALIGN.LEFT
            comment_run.font.color.rgb=RGBColor(0,0,0) 
            comment_run.level=Pt(2)
            
            text_frame.margin_left=Pt(10)

    return

#...................................................................................................................................................................................#

def plotting_omm_monitoring_charts_1(df_gini,df_rankorder,title_name,score_model_code):
    
    
    omm_title_1=f"{title_name} – Output Model Monitoring: Model Level (1/2)"
    
    slide_width=Inches(13.33)
    slide_height=Inches(7.5)

    slide_layout = presentation.slide_layouts[5] # Title only
    slide = presentation.slides.add_slide(slide_layout)
    presentation.slide_width, presentation.slide_height=slide_width,slide_height
    
    logger.debug("New slide created")
    
    left_inch = Inches(0.30)
    top_inch = Inches(0.45)
    width=Inches(12.7)
    height=Inches(0.45)
    title = slide.shapes.title
    title.left=left_inch
    title.top=top_inch
    title.width=width
    title.height=height
    title.text = omm_title_1
    title_text_frame=title.text_frame
    title_text_frame.paragraphs[0].font.size=Pt(22)
    title_text_frame.paragraphs[0].font.bold=True
    title_text_frame.paragraphs[0].font.color.rgb=RGBColor(47,85,151)
    title_text_frame.paragraphs[0].alignment=PP_ALIGN.LEFT
    
    line_left=left_inch
    line_left_top=top_inch+height
    line_right=left_inch+width
    line_right_top=top_inch+height
    
    line = slide.shapes.add_connector(
    MSO_CONNECTOR.STRAIGHT, line_left,line_left_top, line_right, line_right_top)
    
    line.line.color.rgb=RGBColor(47,85,151)
    line.line.width=Pt(2.25)
    
    logger.debug("Title added")

    # Add first graph
    img_stream_1 = f"{path_folder}/nl/GINI_{score_model_code}.png"
    left_inch = Inches(0.50)
    top_inch = Inches(1)
    pic = slide.shapes.add_picture(img_stream_1, left=left_inch, top=top_inch, width=Inches(6), height=Inches(3.7))
    outline=pic.line
    outline.color.rgb=RGBColor(217,217,217)
    outline.width=Pt(0.75)
    
    logger.debug("GINI graph pasted")

    # Add second graph
    img_stream_2 = f"{path_folder}/nl/RANK_ORDER_{score_model_code}.png"
    left_inch = Inches(6.95)
    top_inch = Inches(1)
    pic = slide.shapes.add_picture(img_stream_2, left=left_inch, top=top_inch, width=Inches(6), height=Inches(3.7))
    outline=pic.line
    outline.color.rgb=RGBColor(217,217,217)
    outline.width=Pt(0.75)
    
    logger.debug("Rank Order graph pasted")
    logger.debug("###############################################")

    left_inch = Inches(0.50)
    top_inch = Inches(4.95)
    text_box_1 = slide.shapes.add_textbox(left=left_inch, top=top_inch, width=Inches(6), height=Inches(1.75))

    text_frame_1=text_box_1.text_frame
    outline=text_box_1.line
    outline.color.rgb=RGBColor(217,217,217)
    outline.width=Pt(0.75)

     
    #Adding Text Frame, and title of text box
    paragraph=text_frame_1.paragraphs[0]
    title="Comments"
    paragraph.text=title
    paragraph.font.size=Pt(11)
    paragraph.font.underline=True
    paragraph.alignment=PP_ALIGN.CENTER
    paragraph.font.color.rgb=RGBColor(0,0,0)
    # text_frame_1.margin_left=Pt(10)
    
    df,color_list,per_change_gini,document_color=df_gini_chart(df=df_gini,score_model_code=score_model_code)
    
    if (color_list[-1]=="#FFC000"):
        status="MEDIUM"
    elif (color_list[-1]=="#FF0000"):
        status="HIGH"
    else:
        status="LOW"    
    
    ini_qtr=df["QUARTER"].to_list()[0]
    fin_qtr=df["QUARTER"].to_list()[-3]
    
    first_point=f" % Drop in GINI for current observation period ({ini_qtr} to {fin_qtr}) from the benchmark GINI ({round(df['GINI'].to_list()[-1]*100,2)}%) is {round(per_change_gini*100,2)}%."
    second_point=f" Based on the evidence, GINI for the model has been assigned a RAG status of {status}."
  
    bullet_points = [first_point,
    second_point]
    
    paragraph.add_run().text="\n"
    
    text_frame=text_box_1.text_frame
    text_frame.word_wrap = True
    for points in bullet_points:
                
        paragraph=text_frame.add_paragraph()
        
        run = paragraph.add_run()
        run.text=u'\u2022'
        run.font.size=Pt(9)
        run.font.color.rgb=RGBColor(0,0,0)
        run.font.name='Calibri'
        
        if(points==second_point):
            words=points.split()
            for j, word in enumerate(words):
                run=paragraph.add_run()
                if j==0:
                    run.text= " " + word + " "
                else:
                    run.text=word + " "

                if j== len(words)-1:
                    if(words[j]=="LOW."):
                        run.font.color.rgb = RGBColor(0,176,80)
                    elif(words[j]=="MEDIUM."):
                        run.font.color.rgb = RGBColor(255,192,0)
                    else:
                        run.font.color.rgb = RGBColor(255,0,0)
                        
                    run.font.size=Pt(11)
                    run.font.underline=True
                    run.font.bold=True
                else:                    
                    run.font.color.rgb = RGBColor(0,0,0)
                    run.font.size=Pt(11)
                    
        
        else:
        
            comment_run=paragraph.add_run()
            comment_run.text=points           
            
            comment_run.font.size=Pt(11)
            comment_run.alignment=PP_ALIGN.LEFT
            comment_run.font.color.rgb=RGBColor(0,0,0) 
            comment_run.level=Pt(2)
            
            text_frame.margin_left=Pt(10)        

    left_inch = Inches(6.95)
    top_inch = Inches(4.95)
    text_box_2 = slide.shapes.add_textbox(left=left_inch, top=top_inch, width=Inches(6), height=Inches(1.75))

    text_frame_2=text_box_2.text_frame
    outline=text_box_2.line
    outline.color.rgb=RGBColor(217,217,217)
    outline.width=Pt(0.75)

    #Adding Text Frame, and title of text box
    paragraph=text_frame_2.paragraphs[0]
    title="Comments"
    paragraph.text=title
    paragraph.font.size=Pt(9)
    paragraph.font.underline=True
    paragraph.alignment=PP_ALIGN.CENTER
    paragraph.font.color.rgb=RGBColor(0,0,0)
    # text_frame_1.margin_left=Pt(10)
    
    val_df=df_rank_order(df=df_rankorder)[1]
    val_df["%change"]=val_df['BAD_RATE'].pct_change()
    no_breaks=len(val_df[val_df['%change']>0])
    
    if (no_breaks<=0):
        first_point=f" Strong rank order demonstrated by the scorecard throughout high, medium and low risk bands in all the quarters under monitoring ({ini_qtr} to {fin_qtr})."
        
    elif(no_breaks<=2):
        first_point=f" Mostly stable and satisfactory risk Rank ordering demonstrated by the scorecard in all the quarters under monitoring ({ini_qtr} to {fin_qtr})."
    
    else:
        first_point=f" Poor rank order demonstrated by the scorecard. Multiple trend breaks have been observed in all the quarters under monitoring ({ini_qtr} to {fin_qtr})."
    
    bullet_points = [first_point,
    ]
    
    paragraph.add_run().text="\n"
    
    text_frame=text_box_2.text_frame
    text_frame.word_wrap = True
    for points in bullet_points:
        
        paragraph=text_frame.add_paragraph()
        
        run = paragraph.add_run()
        run.text=u'\u2022'
        run.font.size=Pt(11)
        run.font.color.rgb=RGBColor(0,0,0)
        run.font.name='Calibri'
        
        comment_run=paragraph.add_run()
        comment_run.text=points
        comment_run.font.size=Pt(11)
        comment_run.alignment=PP_ALIGN.LEFT
        comment_run.font.color.rgb=RGBColor(0,0,0) 
        comment_run.level=Pt(2)
        
        text_frame.margin_left=Pt(10)    
        
    return
        
#...................................................................................................................................................................................#

    
def plotting_omm_monitoring_charts_2(df,y,title_name,score_model_code,ppt_file_path_name):
    
    omm_title_2=f"{title_name} – Output Model Monitoring: Model Level (2/2)"
    slide_width=Inches(13.33)
    slide_height=Inches(7.5)
    
    slide_layout = presentation.slide_layouts[5] # Title only
    slide = presentation.slides.add_slide(slide_layout)
    presentation.slide_width, presentation.slide_height=slide_width,slide_height
    
    logger.debug("New slide created")
    
    left_inch = Inches(0.30)
    top_inch = Inches(0.45)
    width=Inches(12.7)
    height=Inches(0.45)
    title = slide.shapes.title
    title.left=left_inch
    title.top=top_inch
    title.width=width
    title.height=height
    title.text = omm_title_2
    title_text_frame=title.text_frame
    title_text_frame.paragraphs[0].font.size=Pt(22)
    title_text_frame.paragraphs[0].font.bold=True
    title_text_frame.paragraphs[0].font.color.rgb=RGBColor(47,85,151)
    title_text_frame.paragraphs[0].alignment=PP_ALIGN.LEFT
    
    line_left=left_inch
    line_left_top=top_inch+height
    line_right=left_inch+width
    line_right_top=top_inch+height
    
    line = slide.shapes.add_connector(
    MSO_CONNECTOR.STRAIGHT, line_left,line_left_top, line_right, line_right_top)
    
    line.line.color.rgb=RGBColor(47,85,151)
    line.line.width=Pt(2.25)
    
    logger.debug("Title added")
    
    img_stream_1 = f"{path_folder}/nl/MAPE_{score_model_code}.png"
    left_inch = Inches(0.50)
    top_inch = Inches(1)
    pic = slide.shapes.add_picture(img_stream_1, left=left_inch, top=top_inch, width=Inches(12.45), height=Inches(4.2))
    outline=pic.line
    outline.color.rgb=RGBColor(217,217,217)
    outline.width=Pt(0.75)
    
    logger.debug("MAPE graph pasted")
    logger.debug("###############################################")
    
    left_inch = Inches(0.50)
    top_inch = Inches(5.40)
    text_box_1 = slide.shapes.add_textbox(left=left_inch, top=top_inch, width=Inches(12.45), height=Inches(1.5))
    
    text_frame_1=text_box_1.text_frame
    outline=text_box_1.line
    outline.color.rgb=RGBColor(217,217,217)
    outline.width=Pt(0.75)
    
    #Adding Text Frame, and title of text box
    paragraph=text_frame_1.paragraphs[0]
    title="Comments"
    paragraph.text=title
    paragraph.font.size=Pt(11)
    paragraph.font.underline=True
    paragraph.alignment=PP_ALIGN.CENTER
    paragraph.font.color.rgb=RGBColor(0,0,0)
    # text_frame_1.margin_left=Pt(10)
    
    df_1_1,df_1_2=df_mape_chart(df=df,y=y,score_model_code=score_model_code)

    lt_odr=df_1_2[df_1_2['variable']=='LT_ODR'][y].to_list()[0]
    lt_edr=df_1_2[df_1_2['variable']=='LT_EDR'][y].to_list()[0]

    ini_date=sorted(df_1_2['SCORE_DATE'].unique())[0]
    fin_date=sorted(df_1_2['SCORE_DATE'].unique())[-1]
    mape=abs((lt_edr-lt_odr)/lt_odr)
    
    if (mape<=0.2):
        status="LOW"
    elif (mape<=0.3):
        status="MEDIUM"
    else:
        status="HIGH"
    
    first_point=f" This segment has an Observed Long term default Rate of {round(lt_odr,2)}% while the Predicted Long term default rate for the same period ({ini_date} to {fin_date}) was {round(lt_edr,2)}%."
    second_point=f" The MAPE against the TTC PD is {round(mape*100,2)}%."
    third_point=f" Based on the evidence, MAPE for the model has been assigned a RAG status of {status}."
    
    bullet_points = [first_point,
    second_point,
    third_point]
    
    paragraph.add_run().text="\n"
    
    text_frame=text_box_1.text_frame
    text_frame.word_wrap = True
    for points in bullet_points:
                
        paragraph=text_frame.add_paragraph()
        
        run = paragraph.add_run()
        run.text=u'\u2022'
        run.font.size=Pt(9)
        run.font.color.rgb=RGBColor(0,0,0)
        run.font.name='Calibri'
        
        if(points==third_point):
            words=points.split()
            for j, word in enumerate(words):
                run=paragraph.add_run()
                if j==0:
                    run.text=" " + word + " "
                else:
                    run.text=word + " "

                if j== len(words)-1:
                    if(words[j]=="LOW."):
                        run.font.color.rgb = RGBColor(0,176,80)
                    elif(words[j]=="MEDIUM."):
                        run.font.color.rgb = RGBColor(255,192,0)
                    else:
                        run.font.color.rgb = RGBColor(255,0,0)
                        
                    run.font.size=Pt(11)
                    run.font.underline=True
                    run.font.bold=True
                else:                    
                    run.font.color.rgb = RGBColor(0,0,0)
                    run.font.size=Pt(11)
                    
        
        else:
        
            comment_run=paragraph.add_run()
            comment_run.text=points           
            
            comment_run.font.size=Pt(11)
            comment_run.alignment=PP_ALIGN.LEFT
            comment_run.font.color.rgb=RGBColor(0,0,0) 
            comment_run.level=Pt(2)
            
            text_frame.margin_left=Pt(10)
            
    presentation.save(ppt_file_path_name)
    
    return
#...................................................................................................................................................................................#


def plotting_population_summary_table(df,val_dict,title_name,score_model_code):
    
    ''' 
        This function pastes charts in the PPT.
    
        input params:
           
        @ title_name - This takes string variable denoting the Title of the slide.
        @ score_model_code - This takes string variable and accounts for the scorecard whose charts are to be pasted.i.e., "PB_UAE_BSCR04"
        @ ppt_file_path_name - This takes string variable for the path and file name of the ppt.Ex- "ppt/final_ppt_1.pptx"
        
        output params:
        
        It returns a saved presentation at given file path.
    
    '''
    imm_title=f"{title_name} – Population Summary"
    
    slide_width=Inches(13.33)
    slide_height=Inches(7.5)
    
    slide_layout = presentation.slide_layouts[5] # Title only
    slide = presentation.slides.add_slide(slide_layout)
    presentation.slide_width, presentation.slide_height=slide_width,slide_height
    
    logger.debug("New slide created")
    
    left_inch = Inches(0.30)
    top_inch = Inches(0.45)
    width=Inches(12.7)
    height=Inches(0.45)
    title = slide.shapes.title
    title.left=left_inch
    title.top=top_inch
    title.width=width
    title.height=height
    title.text = imm_title

    title_text_frame=title.text_frame
    title_text_frame.paragraphs[0].font.size=Pt(22)
    title_text_frame.paragraphs[0].font.bold=True
    title_text_frame.paragraphs[0].font.color.rgb=RGBColor(47,85,151)
    title_text_frame.paragraphs[0].alignment=PP_ALIGN.LEFT
    
    logger.debug("Title added")
    
    line_left=left_inch
    line_left_top=top_inch+height
    line_right=left_inch+width
    line_right_top=top_inch+height
    
    line = slide.shapes.add_connector(
    MSO_CONNECTOR.STRAIGHT, line_left,line_left_top, line_right, line_right_top)
    
    line.line.color.rgb=RGBColor(47,85,151)
    line.line.width=Pt(2.25)
    
    df_1=df.copy(deep=True)
    df_1['%change']=df_1['%_bad_rate_performance_period'].pct_change()
    no_breaks=len(df_1[df_1['%change']>0])
    
    for cols in df_1.columns[1:-1]:
        df_1[cols] = df[cols].map('{:.2%}'.format)
        
    df_1=df_1.iloc[:,:-1]
    
    val_dict
    ini_qtr=val_dict["perf_period"][0]
    fin_qtr=val_dict["perf_period"][-1]
    
    if (no_breaks<=0):
        first_point=f" Strong rank order demonstrated by the scorecard throughout high, medium and low risk bands in all the quarters under monitoring ({ini_qtr} to {fin_qtr})."
        
    elif(no_breaks<=2):
        first_point=f" Mostly stable and satisfactory risk Rank ordering demonstrated by the scorecard in all the quarters under monitoring ({ini_qtr} to {fin_qtr})."
    
    else:
        first_point=f" Poor rank order demonstrated by the scorecard. Multiple trend breaks have been observed in all the quarters under monitoring ({ini_qtr} to {fin_qtr})."
    
    # Add first table
    left_inch = Inches(0.50)
    top_inch = Inches(1)
    table = slide.shapes.add_table(rows=df_1.shape[0]+1,cols=df_1.shape[1], left=left_inch, top=top_inch, width=Inches(9.6), height=Inches(5.225)).table
    
    first_row=table.rows[0]
    first_row.height=Inches(0.75)
    
    first_column=table.columns[0]
    first_column.width=Inches(1.6)
    
    def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element
    
    def _set_cell_border(cell, border_color="#000000", border_width='3175'):
        """ Hack function to enable the setting of border width and border color
            - left border
            - right border
            - top border
            - bottom border
        """
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
    
        # Left Cell Border
        for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
            
            # Every time before a node is inserted, the nodes with the same tag should be removed.
            tag = lines.split(":")[-1]
            for e in tcPr.getchildren():
                if tag in str(e.tag):
                    tcPr.remove(e)
            # end
            
            ln = SubElement(tcPr, lines , w=border_width, cap='flat', cmpd='sng', algn='ctr')
            solidFill = SubElement(ln, 'a:solidFill')
            srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
            prstDash = SubElement(ln, 'a:prstDash', val='solid')
            round_ = SubElement(ln, 'a:round')
            headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
            tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')
    
        return cell
    
    #loop through all cells. apply before setting any cell colors/fonts/..
    for cell in table.iter_cells():
         _set_cell_border(cell)

    for col, col_name in enumerate(df_1.columns):
        cell=table.cell(0,col)
        cell.text=col_name
        cell.text_frame.paragraphs[0].font.size=Pt(11)
        cell.text_frame.paragraphs[0].font.name='Calibri'
        cell.text_frame.paragraphs[0].font.bold=True
        cell.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb=RGBColor(47,85,151)
        cell.text_frame.paragraphs[0].font.color.rgb=RGBColor(255,255,255)

        
    for row in range(df_1.shape[0]):
        for col in range(df_1.shape[1]):
            cell=table.cell(row+1,col)
            cell.text=str(df_1.iloc[row,col])
            cell.text_frame.paragraphs[0].font.size=Pt(11)
            cell.text_frame.paragraphs[0].font.name='Calibri'
            cell.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER
            cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb=RGBColor(255,255,255)
            # cell.border.line.fill.solid()
            # cell.border.line.color.rgb=RGBColor(217,217,217)            
    
    left_inch = Inches(9.5)
    top_inch = Inches(1)
    text_box_2 = slide.shapes.add_textbox(left=left_inch, top=top_inch, width=Inches(3.5), height=Inches(5.5))
    
    text_frame_1=text_box_2.text_frame
    outline=text_box_2.line
    outline.color.rgb=RGBColor(217,217,217)
    outline.width=Pt(.75)
    
    #Adding Text Frame, and title of text box
    paragraph=text_frame_1.paragraphs[0]
    title="Comments"
    paragraph.text=title
    paragraph.font.size=Pt(11)
    paragraph.font.underline=True
    paragraph.alignment=PP_ALIGN.CENTER
    paragraph.font.color.rgb=RGBColor(0,0,0)
    # text_frame_1.margin_left=Pt(10)
    
    # first_point=f".  For the entire observation time period ({ini_qtr} to {fin_qtr}), HCI is {round(val_value_hci*100,2)}%."
    # second_point=f".  Based on the evidence, HCI for the model has been assigned a RAG status of {status_hci}."
 
    bullet_points = [first_point,
    ]
    
    paragraph.add_run().text="\n"
    
    text_frame=text_box_2.text_frame
    text_frame.word_wrap = True
    for points in bullet_points:
        
        paragraph=text_frame.add_paragraph()
        
        run = paragraph.add_run()
        run.text=u'\u2022'
        run.font.size=Pt(9)
        run.font.color.rgb=RGBColor(0,0,0)
        run.font.name='Calibri'
        
        comment_run=paragraph.add_run()
        comment_run.text=points
        # paragraph.add_run().text="\n"
        comment_run.font.size=Pt(11)
        comment_run.alignment=PP_ALIGN.LEFT
        comment_run.font.color.rgb=RGBColor(0,0,0) 
        comment_run.level=Pt(2)
        
        text_frame.margin_left=Pt(10)
        
    return
#...................................................................................................................................................................................#

# Some commonly used function for report writing.

    def paragraph_add(text,bold,font_pt,font_color,font_name):      
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        run=paragraph.add_run()
        run.text= text
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return
        
    def add_paragraph_space(space_after,space_before):
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)   
            
        return
        
    def add_auto_numbered_heading(level, text,bold,font_pt,font_color,font_name):
        if level==1:
            heading_counters[1]=0
        if level==2:
            heading_counters[2]=0
        if level==3:
            heading_counters[3]=0
        heading_counters[level - 1] += 1 # Increment the counter for the current level
        number = '.'.join(str(heading_counters[i]) for i in range(level)) # Generate the numbering
        heading = doc.add_paragraph()
        run = heading.add_run()
        run.text=f"{number} {text}"
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        heading.style = f'Heading {level}'
        heading.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        if level==1:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(18)
        else:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(6)            
        return 

    def create_table(n_rows,n_cols):
        table=doc.add_table(rows=n_rows, cols=n_cols)
        table.style='Table Grid' 
        table.alignment=WD_TABLE_ALIGNMENT.CENTER
        table.autofit=True 
        
        for i in range(len(table.rows)):
            for j in range(len(table.columns)):
        
                table.cell(i,j).paragraphs[0].alignment  = WD_ALIGN_VERTICAL.CENTER
                table.cell(i,j).vertical_alignment=WD_CELL_VERTICAL_ALIGNMENT.CENTER
                
        return table

    def fill_table_color(table,row,col,color="FFFFFF"):
        shading_elm_1 = parse_xml(r'<w:shd {} w:fill="{}"/>'.format(nsdecls('w'),color))
        table.rows[row].cells[col]._tc.get_or_add_tcPr().append(shading_elm_1)
        return

    
    def fill_table(table,row,col,bold,text,r,g,b,alignment):
        
        cell_1=table.cell(row,col).paragraphs[0]
        run=cell_1.add_run()
        run.text=text
        run.bold=bold
        run.font.size=Pt(9)
        run.font.name='Arial'
        run.font.color.rgb=word_rgb(r,g,b)
        cell_1.alignment=alignment
        return
    
    def add_table_caption(caption):
        paragraph = doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)
        run=paragraph.add_run()
        run.text=f"Table {len(doc.tables)}: {caption}"
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.style = 'Caption'
        run.font.name='Arial'
        run.bold=True
        run.font.size=Pt(8)
        run.font.color.rgb=word_rgb(0,0,0)  
        
        return
    
    def set_table_col_width(table,width_arr):
        for i,x in enumerate(width_arr):
            for col in table.columns[i].cells:
                col.width=Inches(x)
        return       

#...................................................................................................................................................................................#

def document_part_1(exclusion_df,psi_hci_list,gini_list,mape_list,portfolio_code,val_start_date_imm,val_start_date_omm,val_end_date_imm,val_end_date_omm,segment,score_model_wise_bm_gini_dict,score_model_wise_imp_pd_dict,bank_name,z=None,x1=None,formatted_date=None):

    def paragraph_add(space_after,space_before,text,bold,font_pt,font_color,font_name):    
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        run=paragraph.add_run()
        run.text= text
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return
    
    def mixed_paragraph_add(space_after,space_before):
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        
        return paragraph
    
    def mixed_run_add(paragraph,text,bold,underline,italic,font_pt,font_color,font_name):
        run=paragraph.add_run(text)
        run.bold=bold
        run.underline=underline
        run.italic=italic
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return
    
    def b_u_paragraph_add(space_after,space_before,text,bold,underline,italic,font_pt,font_color,font_name):
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        run=paragraph.add_run()
        run.text= text
        run.bold=bold
        run.underline=underline
        run.italic=italic
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
    def add_auto_numbered_heading(level, text,bold,font_pt,font_color,font_name):
        if level==1:
            heading_counters[1]=0
        if level==2:
            heading_counters[2]=0
        if level==3:
            heading_counters[3]=0
        heading_counters[level - 1] += 1 # Increment the counter for the current level
        number = '.'.join(str(heading_counters[i]) for i in range(level)) # Generate the numbering
        heading = doc.add_paragraph()
        run = heading.add_run()
        run.text=f"{number}  {text}"
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        heading.style = f'Heading {level}'
        heading.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        if level==1:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(18)
        else:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(6)            
        return 

    def create_table(n_rows,n_cols):
        table=doc.add_table(rows=n_rows, cols=n_cols)
        table.style='Table Grid' 
        table.alignment=WD_TABLE_ALIGNMENT.CENTER
        table.autofit=True 
        
        for i in range(len(table.rows)):
            for j in range(len(table.columns)):
        
                table.cell(i,j).paragraphs[0].alignment  = WD_ALIGN_VERTICAL.CENTER
                table.cell(i,j).vertical_alignment=WD_CELL_VERTICAL_ALIGNMENT.CENTER
                
        return table
    
    
    def fill_table_color(table,row,col,color="FFFFFF"):
        shading_elm_1 = parse_xml(r'<w:shd {} w:fill="{}"/>'.format(nsdecls('w'),color))
        table.rows[row].cells[col]._tc.get_or_add_tcPr().append(shading_elm_1)
        return

    
    def fill_table(table,row,col,bold,text,r,g,b,alignment):
        
        cell_1=table.cell(row,col).paragraphs[0]
        run=cell_1.add_run()
        run.text=text
        run.bold=bold
        run.font.size=Pt(9)
        run.font.name='Arial'
        run.font.color.rgb=word_rgb(r,g,b)
        cell_1.alignment=alignment
        return
    
    def add_table_caption(caption):
        paragraph = doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)
        run=paragraph.add_run()
        run.text=f"Table {len(doc.tables)}: {caption}"
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.style = 'Caption'
        run.font.name='Arial'
        run.bold=True
        run.font.size=Pt(8)
        run.font.color.rgb=word_rgb(0,0,0)  
        
        return
    
    def set_table_col_width(table,width_arr):
        for i,x in enumerate(width_arr):
            for col in table.columns[i].cells:
                col.width=Inches(x)
        return      
    
    val_end_date_input=val_end_date_imm
    val_start_date_input=val_start_date_imm
    val_start_date_output=val_start_date_omm
    val_end_date_output=val_end_date_omm
    
    val_start_date_imm=pd.to_datetime(val_start_date_imm)
    val_start_date_omm=pd.to_datetime(val_start_date_omm)    
    val_end_date_imm=pd.to_datetime(val_end_date_imm)
    val_end_date_omm=pd.to_datetime(val_end_date_omm)

    sections=doc.sections
    for section in sections:
        section.top_margin=Inches(1)
        section.bottom_margin=Inches(1)
        section.left_margin=Inches(1)
        section.right_margin=Inches(1)

    add_auto_numbered_heading(level=1,text="Executive Summary",bold=True,font_pt=14,font_color=(0,0,0),font_name='Arial')
    

    # add_paragraph_space(0,8)
    text_="This section provides an overall summary of the model background, scope of validation and quantitative validation results."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    add_auto_numbered_heading(level=2,text="Model Background",bold=True,font_pt=12,font_color=(0,0,0),font_name='Arial')   
    
    # add_paragraph_space(0,8)
    full_form=portfolio_wise_full_form(portfolio_code=portfolio_code)
    text_1="This report contains" 
    text_2=" ‘Annual Model Validation’"
    text_3=" results of the" 
    text_4=f" B-score/TTC PD model for the {bank_name} {full_form} portfolio."
    
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=True,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_3,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_4,bold=True,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
        
#     paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    text_="The TTC PD estimates derived from the B-score are used as input in the PiT PD estimation. The PiT PD estimates are eventually used in the computation of Expected Credit Loss (ECL). As per the IFRS9 guidelines, Banks should calculate Expected Credit Loss (ECL) for each facility using a combination of Probability of Default (PD), Loss Given default (LGD) and Exposure at Default (EAD)."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    text_=f"Following is a model chronology of {full_form} B-Score model."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    year=val_end_date_imm.year
    # print("year>>>>>>>>>>>>>>>>>>",year)
    rows_t=(year%100)-12 # 21-13=8 #
    table=create_table(rows_t+1,3)
    add_table_caption("Model Lifecycle chronology")
    
    for row in table.rows:
        row.height = Inches(0.3)
    
    row_0=["Model History Log",	"Time Period",	"Performed by"]
    for i,x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)
        fill_table_color(table,0,i,"0070C0")
      
    # if portfolio_code!="nl":
    #     col_0=["B-Score Model Development",
    #             "Model Approval",
    #             "System Implementation Approval *For scoring and data collection only",
    #            "New Model Validation (First validation)",
    #            "Business Go-live",
    #         #    "2021 Annual Validation *Approved by MRMC in Mar-22",
    #         #    "2022 Annual Validation *Approved by MRMC in Apr-23",
    #            "Model Recalibration",
    #            "2023 Annual Validation "
                
    # else:
    col_0=["B-Score Model Development",
            "Model Approval",
              "Business Go-live",
            "Model Recalibration",
            "Annual Validation ",
              " ",
            " ",
             " "
            ]        
    
    for i,x in enumerate(col_0): # 6
        #print(i,x)
        fill_table(table,i+1,0,False,x,0,0,0,0) #fill_table(table,row,col,bold,text,r,g,b,alignment):
        fill_table_color(table,i+1,0,"FFFFFF")
        
    # if portfolio_code=="pl":
        
    col_1=[" ",
            " ",
            " ",
            " ",
            " ",
            " ", 
            " ",
            " "
               ]
    # if portfolio_code=="al":
    #     col_1=["2018",
    #             "Oct-2019",
    #             "Jun-2020",
    #            "Nov-2020",
    #            "Sep-2021",
    #            "Dec-2021",
    #            "Dec-2022",
    #            "Mar-2023",
    #            "Dec-2023"
    #             ]
    # if portfolio_code=="nl":
    #     col_1=["2019",
    #             "Sep-2020",
    #             "Oct-2020",
    #            "Feb-2021",
    #            "Sep-2021",
    #            "Dec-2022",
    #            "Mar-2023",
    #            "Dec-2023"
    #             ]        
    
    for i,x in enumerate(col_1):
        fill_table(table,i+1,1,False,x,0,0,0,1)
        fill_table_color(table,i+1,1,"FFFFFF")  
      
    # if portfolio_code!="nl":
    col_2=[" ",
           " ",
            " ",
            " ",
            " ", 
            " ",
            " ",
            " "
            ]
    # else:
    #     col_2=["FICO Consulting",
    #             "MRMC",
    #            "Auronova Consulting",
    #            "MRMC", 
    #            "Enterprise Risk Solutions",
    #            "FAB Model Validation",
    #            "FAB Model Development",
    #            "FAB Model Validation"
    #             ]        
        
    for i,x in enumerate(col_2):
        fill_table(table,i+1,2,False,x,0,0,0,0)
        fill_table_color(table,i+1,2,"FFFFFF") 
        
    table_width=[2.93,1.31,2.25]                
    set_table_col_width(table,table_width)
    
    full_form=portfolio_wise_full_form(portfolio_code=portfolio_code)
    scorecard_available=portfolio_wise_score_model_code(segment)
    no_segments=len(scorecard_available)
    text_=f"There are {no_segments} segments in {full_form} B-Score model. These are as follows."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    imm_end_date=val_end_date_imm.month_name()[:3]+"'"+str(val_end_date_imm.year)
    omm_start_date=val_start_date_omm.month_name()[:3]+"'"+str(val_start_date_omm.year)
    omm_end_date=val_end_date_omm.month_name()[:3]+"'"+str(val_end_date_omm.year)

    table=create_table(2+no_segments,4) 
    add_table_caption(f"Model and Portfolio Details (As of {imm_end_date})")  
    table.cell(1+no_segments,0).merge(table.cell(1+no_segments,1))    
     
    for row in table.rows:
        row.height = Inches(0.3)
    
    row_0=["Model Segments",	"Segment Description",	"#Accounts", "Exposure(In AED)"]
    for i,x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)
        fill_table_color(table,0,i,"0070C0")
        
        
    if portfolio_code=="pl":
        row_1=["Segment 1: STL - Current","Account in ‘CURRENT i.e., DPD = 0’ & Salary Account is linked to their loan account"]
        row_2=["Segment 2: STL - Delinquent","Account in ‘DELINQUENT i.e., DPD>0’ & Salary Account is linked to their loan account"]
        n_rows=[row_1,row_2]
        
    elif portfolio_code=="nl":
        if no_segments == 1:
            row_1=["Segment 1: Current","Account in ‘CURRENT’ i.e., DPD = 0 in their payment as of snapshot month"]
            n_rows=[row_1]
        elif no_segments == 2:
            row_1=["Segment 1: Current","Account in ‘CURRENT’ i.e., DPD = 0 in their payment as of snapshot month"]
            row_2=["Segment 2: Delinquent","Account in ‘DELINQUENT’ i.e., DPD>0 in their payment as of snapshot month"]
            n_rows=[row_1,row_2]
        
    elif portfolio_code=="al":
        row_1=[]
        row_2=[]
    else:
        row_1=[]
        row_2=[]        

        
    if ((portfolio_code=="pl") | (portfolio_code=="mort")):
        temp_df=exclusion_df[(exclusion_df['SCORE_DATE']==val_end_date_input) & (~exclusion_df["BSCORE_EXCLUSION_REASON"].str.contains("ISLAMIC"))]
    else:
        temp_df=exclusion_df[(exclusion_df['SCORE_DATE']==val_end_date_input)]        
    temp_df_grp=temp_df.groupby("SCORE_MODEL_CODE").agg(count=pd.NamedAgg(column='ACCOUNT_ID',aggfunc='count'),
                                           exposure=pd.NamedAgg(column='LOAN_BALANCE_00',aggfunc='sum')).reset_index(drop=True)
    
    overall_rag_temp_df=temp_df_grp.copy(deep=True)
    
    temp_df_grp["count"]=temp_df_grp["count"].apply(lambda x: f"{x:,}")
    temp_df_grp["exposure"]=round(temp_df_grp["exposure"]/1000000,0)
    temp_df_grp["exposure"]=temp_df_grp["exposure"].apply(lambda x: f"{x:,}")
    
    
    for i in range(temp_df_grp.shape[0]):
        new_list=temp_df_grp.iloc[i,:].to_list()
        n_rows[i]=n_rows[i]+new_list
        
    for j in range(len(n_rows)):
        for i,x in enumerate(n_rows[j]):
            if (i==len(n_rows[j])-1):
                x_1=f"AED {n_rows[j][i][:-2]} Mn"
                fill_table(table,j+1,i,False,x_1,0,0,0,1)
                fill_table_color(table,j+1,i,"FFFFFF")
            elif (i==0):
                fill_table(table,j+1,i,True,x,0,0,0,0)
                fill_table_color(table,j+1,i,"FFFFFF") 
            else:
                fill_table(table,j+1,i,False,x,0,0,0,1)
                fill_table_color(table,j+1,i,"FFFFFF")                 
            
    row_n=["Overall Population",""]
    n_total=temp_df.shape[0]
    expo_total=round(temp_df["LOAN_BALANCE_00"].sum()/1000000,0)
    n_total="{:,}".format(n_total)
    expo_total="{:,}".format(expo_total)[:-2]
    row_n.append(n_total)
    row_n.append(expo_total)    
    for i,x in enumerate(row_n):
        if (i==len(row_n)-1):
            x_1=f"AED {row_n[i]} Mn"
            fill_table(table,1+no_segments,i,True,x_1,0,0,0,1)
            fill_table_color(table,1+no_segments,i,"FFFFFF") 
        else:
            fill_table(table,1+no_segments,i,True,x,0,0,0,1)
            fill_table_color(table,1+no_segments,i,"FFFFFF")             
        
    table_width=[1.87,2.78,.79,1.05]                
    set_table_col_width(table,table_width)    
    
    
    text_=f"The portfolio details in Table 2 above are as of {imm_end_date}. This is inclusive of all policy and model exclusions (pre-exclusion accounts). {full_form} portfolio exposure is approximately AED {expo_total} Mn. There are over {n_total[:-4]}K active accounts on the book. ECL for the same period is AED (to be updated) This is approx.(to be updated) of the portfolio exposure."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')     
    
    add_auto_numbered_heading(level=2,text="Validation Approach",bold=True,font_pt=12,font_color=(0,0,0),font_name='Arial')
    text_="This section explains the qualitative and quantitative validation measures. The entire validation process was divided into two parts as shown below:"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    table=create_table(3,2)
    add_table_caption("Qualitative validation parameters for B-Score")

    for row in table.rows:
        row.height = Inches(0.3)
    
    row_0=["Qualitative Validation","Key checks"]
    for i,x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)
        fill_table_color(table,0,i,"0070C0")
        
    row_1=["Review of Validation Data Quality",
           ["Are there significant concerns on the data quality?",
           "Is the modelling data reconciling closely with business reports?"
           ]
    ]
    
    for i in range(len(row_1)):
        if(i==0):
            fill_table(table,1,i,False,row_1[i],0,0,0,0)
        else:

            for j in range(len(row_1[1])):
                if(j!=0):
                    paragraph=table.cell(1,i).add_paragraph()
                    paragraph.style='List Bullet'
                    paragraph.style.font.size=Pt(8)
                    paragraph.style.bold=True
                    run=paragraph.add_run(f"{row_1[1][j]}")
                    paragraph.alignment = 0
                    run.font.size=Pt(9)
                    run.font.name='Arial'
                else:
                    paragraph=table.cell(1,i).paragraphs[0]
                    paragraph.style='List Bullet'
                    paragraph.style.font.size=Pt(8)
                    paragraph.style.bold=True
                    run=paragraph.add_run(f"{row_1[1][j]}")
                    paragraph.alignment = 0
                    run.font.size=Pt(9)
                    run.font.name='Arial'                    
                
                             
                
    row_2=["Review of Data and Modelling Exclusions", "Are there any explained dips / rise in the volume of Modelling exclusions?"]
    
    for i,x in enumerate(row_2):   
        if i==0:
            fill_table(table,2,i,False,x,0,0,0,0)  
        else:
            paragraph=table.cell(2,i).paragraphs[0]
            run=paragraph.add_run(f"{x}")
            paragraph.alignment = 0
            paragraph.style='List Bullet'
            paragraph.style.font.size=Pt(8)
            paragraph.style.bold=True
            run.font.size=Pt(9)
            run.font.name='Arial'     
        
    table_width=[2.28,4.21]                
    set_table_col_width(table,table_width)    
    
    
    table=create_table(8,3)
    add_table_caption("Quantitative validation parameters for B-Score") 
    for row in table.rows:
        row.height = Inches(0.3)
    
    row_0=["Quantitative Validation","KPIs","Comment"]
    for i,x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)
        fill_table_color(table,0,i,"0070C0")
    
    col_0=["Model Discriminatory Power",
           "Model Stability",
           "Model Accuracy",
           "Model Concentration",
           "Variable Stability",
           "Variable Correlation",
           "Variable Discriminatory Power"
           ]
    for i,x in enumerate(col_0):
        fill_table(table,i+1,0,False,x,0,0,0,0)
        fill_table_color(table,i+1,0,"FFFFFF")    

    col_1=["GINI",
           "PSI",
           "MAPE",
           "HCI",
           "CSI",
           "VIF",
           "IV & Bad rate Rank Order"
           ]
    for i,x in enumerate(col_1):
        fill_table(table,i+1,1,False,x,0,0,0,0)
        fill_table_color(table,i+1,1,"FFFFFF") 
    
    col_2=[" ",
           " ",
           " ",
           " ",
           " ",
           " ",
           " "
           ]
    for i,x in enumerate(col_2):
        fill_table(table,i+1,2,False,x,0,0,0,0)
        fill_table_color(table,i+1,2,"FFFFFF") 

    table_width=[1.85,.96,3.68]                
    set_table_col_width(table,table_width) 
    
    table=create_table(2,3)
    add_table_caption("Implementation Testing for B-Score Annual Validation") 

    for row in table.rows:
        row.height = Inches(0.3)
    
    row_0=["Implementation Check","KPIs","Comment"]
    for i,x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)
        fill_table_color(table,0,i,"0070C0")

   
    row_1=["Reconciliation of B-Score (manually calculated) With B-Score",
           "Sample checks",
           f"{portfolio_code.upper()} B-score model implementation testing was performed for {formatted_date} to {imm_end_date}."
]
    for i,x in enumerate(row_1):
        fill_table(table,1,i,False,x,0,0,0,1)
        fill_table_color(table,1,i,"FFFFFF")
    
    table_width=[1.85,.96,3.68]                
    set_table_col_width(table,table_width) 
    
    table=create_table(4,2)
    add_table_caption("Validation Data Period")
    
    for row in table.rows:
        row.height = Inches(0.3)
    
    row_0=["Data Statistics","Validation Period"]
    for i,x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)
        fill_table_color(table,0,i,"0070C0")
        
    row_1=["Observation Snapshot Months in validation data",
           [f"{omm_start_date} - {omm_end_date}"]
    ]
    
    for i in range(len(row_1)):
        if(i==0):
            fill_table(table,1,i,False,row_1[i],0,0,0,0)
        else:

            for j in range(len(row_1[1])):
                if(j!=0):
                    paragraph=table.cell(1,i).add_paragraph()
                    paragraph.style='List Bullet'
                    paragraph.style.font.size=Pt(8)
                    paragraph.style.bold=True
                    run=paragraph.add_run(f"{row_1[1][j]}")
                    paragraph.alignment = 0
                    run.font.size=Pt(9)
                    run.font.name='Arial'
                else:
                    paragraph=table.cell(1,i).paragraphs[0]
                    paragraph.style='List Bullet'
                    paragraph.style.font.size=Pt(8)
                    paragraph.style.bold=True
                    run=paragraph.add_run(f"{row_1[1][j]}")
                    paragraph.alignment = 0
                    run.font.size=Pt(9)
                    run.font.name='Arial'

    row_2=["Observation window", ["All model variables look at up to historical 12 months of customer payment behaviour prior to observation snapshot."]]
    
    for i,x in enumerate(row_2):   
        if i==0:
            fill_table(table,2,i,False,x,0,0,0,0)  
        else:
            for j in range(len(row_2[1])):
                if(j!=0):
                    paragraph=table.cell(2,i).add_paragraph()
                    paragraph.style='List Bullet'
                    paragraph.style.font.size=Pt(8)
                    paragraph.style.bold=True
                    run=paragraph.add_run(f"{row_2[1][j]}")
                    paragraph.alignment = 0
                    run.font.size=Pt(9)
                    run.font.name='Arial'
                else:
                    paragraph=table.cell(2,i).paragraphs[0]
                    paragraph.style='List Bullet'
                    paragraph.style.font.size=Pt(8)
                    paragraph.style.bold=True
                    run=paragraph.add_run(f"{row_2[1][j]}")
                    paragraph.alignment = 0
                    run.font.size=Pt(9)
                    run.font.name='Arial'     
            
    from datetime import datetime, timedelta
    # def get_month_and_year(date,months_after):
    #     current_date = date
    #     target_date = current_date + timedelta(days=months_after*30)
    #     month_name = target_date.strftime('%B')
    #     year = target_date.year
    #     return month_name, year
    def get_month_and_year(date,months_after):
        date1 = pd.to_datetime(date).date()
        timestamp = f"{date1}"
        date_object = datetime.strptime(timestamp, '%Y-%m-%d')
        formatted_date = date_object.strftime("%b")
        
        if formatted_date == "Jan":
            current_date = date
            target_date = current_date + timedelta(days=months_after*28)
            month_name = target_date.strftime('%B')
            year = target_date.year
            return month_name, year
        else:
            current_date = date
            target_date = current_date + timedelta(days=months_after*30)
            month_name = target_date.strftime('%B')
            year = target_date.year
            return month_name, year
    months = 1 
    next_month_start, next_year_start = get_month_and_year(val_start_date_omm,months)
    start_1=next_month_start[:3]+"'"+str(next_year_start)
    
    months = 12
    twelve_month_start, twelve_month_year_start = get_month_and_year(val_start_date_omm,months)
    end_1=twelve_month_start[:3]+"'"+str(twelve_month_year_start)
    
    months = 1 
    next_month_end, next_year_end = get_month_and_year(val_end_date_omm,months)
    start_2=next_month_end[:3]+"'"+str(next_year_end)

    row_3=["Performance Months", ["Each observation snapshot month looks at next 12 months of default performance.",
                                 f"First observation snapshot, {omm_start_date}, looks at performance in {start_1} to {end_1}.",
                                 f"Last observation snapshot, {omm_end_date}, looks at performance in {start_2} to {imm_end_date}."]
          ]
    
    for i,x in enumerate(row_3):   
        if i==0:
            fill_table(table,3,i,False,x,0,0,0,0)  
        else:
            for j in range(len(row_3[1])):
                if(j!=0):
                    paragraph=table.cell(3,i).add_paragraph()
                    paragraph.style='List Bullet'
                    paragraph.style.font.size=Pt(8)
                    paragraph.style.bold=True
                    run=paragraph.add_run(f"{row_3[1][j]}")
                    paragraph.alignment = 0
                    run.font.size=Pt(9)
                    run.font.name='Arial'
                else:
                    paragraph=table.cell(3,i).paragraphs[0]
                    paragraph.style='List Bullet'
                    paragraph.style.font.size=Pt(8)
                    paragraph.style.bold=True
                    run=paragraph.add_run(f"{row_3[1][j]}")
                    paragraph.alignment = 0
                    run.font.size=Pt(9)
                    run.font.name='Arial'  
        
    table_width=[1.81, 4.68]                
    set_table_col_width(table,table_width) 
    
    add_auto_numbered_heading(level=2,text="Scope of Validation",bold=True,font_pt=12,font_color=(0,0,0),font_name='Arial')
    
    text_=f"This section explains the scope of {str(val_end_date_imm.year)} annual validation exercise."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    bullet_points=[
    # f"First time validation i.e., ‘New Model validation’ of the {full_form} B-Score was performed by Auronova consulting in 2020. ",

# f"Last Annual validation i.e., {str(val_end_date_imm.year -1 )} Annual Validation of the {full_form} B-Score was performed by FAB Model Validation team in {str(val_end_date_imm.year -1 )}." ,

f"This report captures the ‘Annual Model validation’ results of {str(val_end_date_imm.year)}. The annual validation was performed by FAB Model Validation team.",

"Quantitative validation was performed in accordance with the policy i.e., all B-score related KPIs viz. GINI, PSI, HCI, MAPE were covered as part of validation exercise. Additionally, variable level analysis was also performed.",

"Qualitative validation of data quality, exclusions were also performed. Other qualitative validation parameters, which are applied during ‘new model validation’ viz. review of model development document, model methodology etc. were out of scope of this validation.", 
  
"Status of the Action items raised during previous validation were also reviewed. Fresh action items were issued based on this year’s validation results.", 

"Model implementation testing i.e., sample checks on B-score generated manually vis-à-vis B-score generated in production system on the validation data were also performed."
        
]
    for comment in bullet_points:
        paragraph=doc.add_paragraph()
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        
        paragraph.paragraph_format.left_indent = Inches(0.5)
        paragraph.paragraph_format.right_indent = Inches(0.5)        
        paragraph.style='List Bullet'
        paragraph.style.font.size=Pt(8)
        paragraph.style.bold=True
        run=paragraph.add_run(f"{comment}")
        run.font.size=Pt(10)
        run.font.color.rgb=word_rgb(0,0,0)
        run.font.name='Arial'
        
        paragraph.paragraph_format.space_after = Pt(8)
        paragraph.paragraph_format.space_before = Pt(8) 
        
    text_=f"As part of this exercise, {str(val_end_date_imm.year -1 )} validation report was used as reference:"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')    
    
    bullet_points=["{To be manually updated}"]
    
    for comment in bullet_points:
        paragraph=doc.add_paragraph()
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        
        paragraph.paragraph_format.left_indent = Inches(0.5)
        paragraph.paragraph_format.right_indent = Inches(0.5)        
        paragraph.style='List Bullet'
        paragraph.style.font.size=Pt(8)
        paragraph.style.bold=True
        run=paragraph.add_run(f"{comment}")
        run.font.size=Pt(10)
        run.font.color.rgb=word_rgb(0,0,0)
        run.font.name='Arial'
        
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)
        
    add_auto_numbered_heading(level=2,text="Summary of Validation Results",bold=True,font_pt=12,font_color=(0,0,0),font_name='Arial') 
    text_="The summary of model validation is provided in the table below:"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    table=create_table(10,2)
    add_table_caption("Validation Summary")     
    for row in table.rows:
        row.height = Inches(0.2)   
        
    col_0=[
           "Model Name",
           "Validation type",
           "Model Tier",
           "Model Performance RAG status",
           "Model Risk",
           "Summary comment",
           "Validator name",
           "Reviewer name",
           "Approver",
           "Next review date"
         ]
        
    for i,x in enumerate(col_0):
        fill_table(table,i,0,True,x,255,255,255,0)
        fill_table_color(table,i,0,"0070C0") 
        
    if portfolio_code=="pl":
        model_tier="HIGH"
    if portfolio_code=="al":
        model_tier="MEDIUM"
    if portfolio_code=="nl":
        model_tier="HIGH"
        
    col_1=[
           f"Retail {full_form} B-Score/TTC-PD Model",
           "Annual Model Validation",
           model_tier,
           " ",
           " ",
           " ",
           " ",
           " ",
           "Model Risk Management Committee",
           " "
         ]
    for i,x in enumerate(col_1):
        fill_table(table,i,1,False,x,0,0,0,0)
        fill_table_color(table,i,1,"FFFFFF")
        
    table_width=[2.03,4.47]                
    set_table_col_width(table,table_width) 
   
    add_auto_numbered_heading(level=3,text="Qualitative Validation Results Summary",bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    text_="The following table presents summary of the qualitative model validation results."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')  
    
    table=create_table(3,3)
    add_table_caption("Qualitative validation results summary") 
    
    for row in table.rows:
        row.height = Inches(0.3)
    
    row_0=["Category","Observation","Conclusion"]
    for i,x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)
        fill_table_color(table,0,i,"0070C0")
        
    col_0=["Data adequacy and quality","Exclusions"]
    for i,x in enumerate(col_0):
        fill_table(table,i+1,0,True,x,0,0,0,0)
        fill_table_color(table,i+1,0,"FFFFFF")  
        
        
    col_1=[
           ["Exposure trends in model validation data reconciled closely with the portfolio business reports for all the quarters.",
           "No significant gap observed in the quarterly trends across DPD buckets in the validation data.",
            "Conclusion: Based on evidence, Model Validation Team considered the validation data to be Fit for Purpose."
           ],
        [
            "No major data anomaly was found in the number of observations in the exclusions.",
            "After removing exclusions, number of observations in each segment was adequate for performing model validation."   ,
            "Conclusion: Based on evidence, Model Validation Team considered the validation data to be Fit for Purpose."
        ]
    ]
    
    for i in range(len(col_1)):
        for j in range(len(col_1[i])):
            if(j==2):
                paragraph=table.cell(i+1,1).add_paragraph()
                run=paragraph.add_run(f"{col_1[i][j]}")
                paragraph.alignment = 0
                run.font.size=Pt(9)
                run.font.name='Arial' 
            else:
                
                if(j!=0):
                    paragraph=table.cell(i+1,1).add_paragraph()
                    paragraph.style='List Bullet'
                    paragraph.style.font.size=Pt(8)
                    paragraph.style.bold=True
                    run=paragraph.add_run(f"{col_1[i][j]}")
                    paragraph.alignment = 0
                    run.font.size=Pt(9)
                    run.font.name='Arial'
                else:
                    paragraph=table.cell(i+1,1).paragraphs[0]
                    paragraph.style='List Bullet'
                    paragraph.style.font.size=Pt(8)
                    paragraph.style.bold=True
                    run=paragraph.add_run(f"{col_1[i][j]}")
                    paragraph.alignment = 0
                    run.font.size=Pt(9)
                    run.font.name='Arial'                    
        
    col_2=["Compliant","Compliant"]
    for i,x in enumerate(col_2):
        fill_table(table,i+1,2,True,x,255,255,255,1)
        fill_table_color(table,i+1,2,"00B050")    
        
    table_width=[1.09,4.61,0.80]                
    set_table_col_width(table,table_width) 
    
    text_="Conclusion (Qualitative Validation)"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')     


    comment="No major concerns on the validation of the data quality and applied exclusions. Model Validation Team considers the validation data to be Fit-for-Purpose."
    paragraph=doc.add_paragraph()
    
    paragraph.paragraph_format.left_indent = Inches(0.5)
    paragraph.paragraph_format.right_indent = Inches(0.5)        
    paragraph.style='List Bullet'
    paragraph.style.font.size=Pt(8)
    paragraph.style.bold=True
    run=paragraph.add_run(f"{comment}")
    run.font.size=Pt(10)
    run.font.color.rgb=word_rgb(0,0,0)
    run.font.name='Arial'
    
    paragraph.paragraph_format.space_after = Pt(6)
    paragraph.paragraph_format.space_before = Pt(6)
    
    add_auto_numbered_heading(level=3,text="Quantitative Validation Results Summary",bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial')   
    
    text_=f"The quantitative validation summary for the {no_segments} {full_form} model segments is as follows."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')  
    ####################### added new
    if no_segments >=2:
        table=create_table(5,1+no_segments)
        add_table_caption("Quantitative validation - Data Summary")
    else:
        table=create_table(5,1+no_segments)
        add_table_caption("Quantitative validation - Data Summary")

    for row in table.rows:
        row.height = Inches(0.3)  
        
    # if portfolio_code=="pl":
    #     row_0=["Data Statistics","STL-Current","STL-Delinquent"]
    # if portfolio_code=="al":
    #     row_0=["Data Statistics","Current Banking","Current Non-Banking","Delinquent"]
    if portfolio_code=="nl":
        if no_segments ==1 :
            row_0=["Data Statistics","NL-Current "]
        else:
            row_0=["Data Statistics","NL-Current ","NL-Delinquent"]
        


        
        
    for i,x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)
        fill_table_color(table,0,i,"0070C0")
    
    col_0=[ "Validation Period",
            "Observation Snapshot months",
            "Observations in Validation data",
            "Observed Default Rate"
            ]
    for i,x in enumerate(col_0):
        fill_table(table,i+1,0,True,x,0,0,0,0)
        fill_table_color(table,i+1,0,"FFFFFF") 
        
    temp_df=exclusion_df[(exclusion_df['SCORE_DATE']<=val_end_date_output) & (exclusion_df['SCORE_DATE']>=val_start_date_output) & (exclusion_df['BSCORE_EXCLUSION_REASON'].str.contains("POST_EXCLUSION"))]
    temp_df["TARGET_12"].replace(99,0,inplace=True)
    temp_df_grp=temp_df.groupby("SCORE_MODEL_CODE").agg(count=pd.NamedAgg(column='ACCOUNT_ID',aggfunc='count'),
                                           bads=pd.NamedAgg(column='TARGET_12',aggfunc='sum')).reset_index(drop=True)
    temp_df_grp["bad_rate"]=temp_df_grp["bads"]/temp_df_grp["count"]
    del temp_df_grp["bads"]
    temp_df_grp["count"]=temp_df_grp["count"].apply(lambda x: f"{x:,}")
    temp_df_grp["bad_rate"]=temp_df_grp["bad_rate"].apply(lambda x: f"{x:2%}")
    temp_df_grp["bad_rate"]=temp_df_grp["bad_rate"].str[:-5]+"%"
         
    for j in range(no_segments):
        
        x_temp=temp_df_grp.iloc[j,:].to_list()
        col_1=[ f"{omm_start_date} to {omm_end_date}",
                f"Monthly snapshots from {omm_start_date} to {omm_end_date}",
                ]    
        col_1=col_1+x_temp
        for i,x in enumerate(col_1):
            fill_table(table,i+1,j+1,False,x,0,0,0,1)
            fill_table_color(table,i+1,j+1,"FFFFFF")        
    
    
    text_=f"The below table shows overall validation test result summary for the {no_segments} model segments,"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    table=create_table(2+no_segments,7)
    add_table_caption("Segment level - Quantitative Test Result") 
    
    for row in table.rows:
        row.height = Inches(0.3)
        
    table.cell(0,0).merge(table.cell(1,0))
    table.cell(0,1).merge(table.cell(0,3))
    
    table.cell(0,4).merge(table.cell(1,4))
    table.cell(0,5).merge(table.cell(1,5))
    table.cell(0,6).merge(table.cell(1,6))
    
    row_0=["Segment","GINI","PSI","HCI","MAPE"]
    
    for i,x in enumerate(row_0):
        if (i==0 or i==1):
            fill_table(table,0,i,True,x,255,255,255,1)
            fill_table_color(table,0,i,"0070C0") 
        else:
            fill_table(table,0,i+2,True,x,255,255,255,1)
            fill_table_color(table,0,i+2,"0070C0")             
    
    row_1=["Benchmark","Validation Data","GINI Drop"]
    
    for i,x in enumerate(row_1):
        fill_table(table,1,i+1,True,x,255,255,255,1)
        fill_table_color(table,1,i+1,"0070C0")
        
    # if portfolio_code=="pl":
    #     col_0=["Segment 1: STL- Current","Segment 2: STL- Delinquent"]
    #     row_1=[]
    #     row_2=[]
    #     rows_0=[row_1,row_2]

    if portfolio_code=="nl":
        if no_segments == 1:
            col_0=["Segment 1: Current"]
            row_1=[]
            rows_0=[row_1]
        elif no_segments == 2:
            col_0=["Segment 1: Current","Segment 2: Delinquent"]
            row_1=[]
            row_2=[]
            rows_0=[row_1,row_2]
        
        
    for i,x in enumerate(col_0):
        fill_table(table,i+2,0,True,x,0,0,0,0)
        fill_table_color(table,i+2,0,"FFFFFF") 
        
    scorecard_available=portfolio_wise_score_model_code(segment=segment)
    score_dict={}
    for j,x in enumerate(scorecard_available):
        rag_color_list=[]
        psi_hci=psi_hci_list[x]
        y_psi_hci="PSI"
        df_1_psi,color_list_psi,document_color_psi=df_psi_hci_chart(df=psi_hci,y=y_psi_hci) 
        val_value_psi=round(df_1_psi["PSI"].to_list()[-1],4)
        y_psi_hci="HCI"
        df_1_hci,color_list_hci,document_color_hci=df_psi_hci_chart(df=psi_hci,y=y_psi_hci)    
        val_value_hci=round(df_1_hci["HCI"].to_list()[-1],4)  
        
        gini=gini_list[x]
        gini.rename({"QTR":"QUARTER","GINI_STATISTIC":"GINI"},axis=1,inplace=True)
        df,color_list,per_change_gini,document_color_gini=df_gini_chart(df=gini,score_model_code=x,score_model_wise_bm_gini_dict=score_model_wise_bm_gini_dict,z=z)
        bm_value_gini=round(df["GINI"].to_list()[-1],4) 
        val_value_gini=round(df["GINI"].to_list()[-2],4) 
        per_change_gini=per_change_gini
         
        df_mape=mape_list[x]
        df_1_1,df_1_2=df_mape_chart(df=df_mape,y="Percent(%)",score_model_code=x,score_model_wise_imp_pd_dict=score_model_wise_imp_pd_dict,x=x1)
        # display(df_1_1)
        lt_odr=df_1_2[df_1_2['variable']=='LT_ODR']["Percent(%)"].to_list()[0]
        lt_edr=df_1_2[df_1_2['variable']=='LT_EDR']["Percent(%)"].to_list()[0]

       
        ini_date=sorted(df_1_2[x1].unique())[0]
        fin_date=sorted(df_1_2[x1].unique())[-1]
    
        mape=round(abs((lt_edr-lt_odr)/lt_odr),4)
        
        if (mape<=0.2):
            document_color_mape="#00B050"
        elif (mape<=0.3):
            document_color_mape="#FFC000"
        else:
            document_color_mape="#FF0000"
        
        rows_0[j].append("{:2%}".format(bm_value_gini)[:-5]+"%")
        rows_0[j].append("{:2%}".format(val_value_gini)[:-5]+"%")
        
        if (val_value_gini<.3):
            rows_0[j].append("Absolute Segment GINI <30%")
            document_color_gini="#FF0000"
        else:
            if per_change_gini>0:
                rows_0[j].append("Decreased by {:2%}".format(per_change_gini)[:-5]+"%")
            else:
                text="Increased by "
                text+="{:2%}".format(per_change_gini)[1:-5]+"%"
                rows_0[j].append(text)
            
        rows_0[j].append("{:2%}".format(val_value_psi)[:-5]+"%")
        rows_0[j].append("{:2%}".format(val_value_hci)[:-5]+"%")
        rows_0[j].append("{:2%}".format(mape)[:-5]+"%")
        
        for i,text in enumerate(rows_0[j]):
            fill_table(table,j+2,1+i,False,text,0,0,0,1)
            if i<2:
                fill_table_color(table,j+2,1+i,"FFFFFF") 
            elif i==2:
                fill_table_color(table,j+2,1+i,document_color_gini)
                rag_color_list.append(document_color_gini)
            elif i==3:
                fill_table_color(table,j+2,1+i,document_color_psi) 
                rag_color_list.append(document_color_psi)
            elif i==4:
                fill_table_color(table,j+2,1+i,document_color_hci)
                rag_color_list.append(document_color_hci)
            else:
                fill_table_color(table,j+2,1+i,document_color_mape)  
                rag_color_list.append(document_color_mape)
                
        score_dict[x]=rag_color_list
    
    table_width=[1.82,0.83,.78,1.39,0.57,0.57,0.57]

    # img_path = r'table_image.png'
    # save_table_as_image(table,img_path)    
               
    set_table_col_width(table,table_width) 
    
    text_="Conclusion (Quantitative Validation)"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')  
    
    bullet_points=["GINI (Model Discrimination)",
                  "PSI (Population Stability)",
                  "HCI (Score Concentration)",
                  "MAPE (Model Forecast Error)"]
    
    for comment in bullet_points:
        paragraph=doc.add_paragraph()
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        
        paragraph.paragraph_format.left_indent = Inches(0.5)
        paragraph.paragraph_format.right_indent = Inches(0.5)        
        paragraph.style='List Bullet'
        paragraph.style.font.size=Pt(8)
        paragraph.style.bold=True
        run=paragraph.add_run(f"{comment}")
        run.font.size=Pt(10)
        run.bold=True
        run.font.color.rgb=word_rgb(0,0,0)
        run.font.name='Arial'
        
        paragraph.paragraph_format.space_after = Pt(8)
        paragraph.paragraph_format.space_before = Pt(8) 
    
    add_auto_numbered_heading(level=3,text="Model RAG status",bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    text_="Model Segment Level RAG status"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')     
    
    text_1="Based on the criteria outlined in the Bank’s MRM risk quantification framework, the RAG status for each segmented model and portfolio has been computed as per the metrics provided below using validation results (Refer "
    text_2="Section 4.2 "
    text_3="for details)."
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=True,italic=False,font_pt=10,font_color=(5,99,193),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_3,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')

    table=create_table(2+no_segments,7)
    add_table_caption("Segment level - RAG Status") 
    
    for row in table.rows:
        row.height = Inches(0.3) 
        
    table.cell(0,1).merge(table.cell(0,5))   
    table.cell(0,6).merge(table.cell(1,6)) 
    
    row_0=["","Risk Scores","","","","","Final RAG"]
    for i,x in enumerate(row_0):
        if (i==1 or i==6):
            fill_table(table,0,i,True,x,255,255,255,1)
            fill_table_color(table,0,i,"0070C0") 
        else:
            fill_table_color(table,0,i,"0070C0")      
    
    row_1=["Model Segment","GINI drop Score","PSI Score","HCI Score","MAPE Score","Score Point = 70% GINI + 10% PSI +10% HCI + 10% MAPE"]
    
    for i,x in enumerate(row_1):
        #print(x)
        fill_table(table,1,i,True,x,255,255,255,1)
        fill_table_color(table,1,i,"0070C0") 
    

    # if portfolio_code=="pl":
    #     col_0=["Segment 1: STL- Current","Segment 2: STL- Delinquent"]
    #     row_1=[]
    #     row_2=[]
    #     rows_0=[row_1,row_2]
        
    if portfolio_code=="nl":
        # col_0=["Segment 1: Current","Segment 2: Delinquent"]
        # row_1=[]
        # row_2=[]
        # rows_0=[row_1,row_2]
        if no_segments == 1:
            col_0=["Segment 1: Current"]
            row_1=[]
            rows_0=[row_1]
        elif no_segments == 2:
            col_0=["Segment 1: Current","Segment 2: Delinquent"]
            row_1=[]
            row_2=[]
            rows_0=[row_1,row_2]
        
    # elif portfolio_code=="al":
    #     col_0=[]
    # else:
    #     col_0=[]
        
    for i,x in enumerate(col_0):
        fill_table(table,i+2,0,True,x,0,0,0,0)
        fill_table_color(table,i+2,0,"FFFFFF") 
    
    def rag_score(color):
        if color=="#00B050":
            score=1
        elif color=="#FFC000":
            score=2
        else:
            score=3
        return score
    
    def final_score(row_list):
        final_score=0
        for i in range(len(row_list)):
            if i==0:
                final_score=final_score+0.7*row_list[i]
            else:
                final_score=final_score+0.1*row_list[i]
                
        return final_score
    
    overall_model_rag_score={}
    
    for j,x in enumerate(scorecard_available):
        row_=score_dict[x]
        row_j=[]
        for i in range(len(row_)):
            row_j.append(rag_score(row_[i]))
        
        final_rag_score=final_score(row_j)
        final_rag_score_=round(final_rag_score,2)
        l_1=[]
        overall_model_rag_score[x]=l_1
        l_1.append(final_rag_score_)
        row_j.append(str(final_rag_score_))
        if final_rag_score<=1.3:
            final_rag_color="#00B050"
            final_rag_status="Green"
            l_1.append(final_rag_status)
            l_1.append(final_rag_color)
        elif final_rag_score<=2.35:
            final_rag_color="#FFC000"
            final_rag_status="Amber"
            l_1.append(final_rag_status)
            l_1.append(final_rag_color)
        else:
            final_rag_color="#FF0000"
            final_rag_status="Red"
            l_1.append(final_rag_status)
            l_1.append(final_rag_color)
        
        row_j.append(final_rag_status)
        color=final_rag_color
        for i,k in enumerate(row_j):
            if i!=len(row_j)-1:
                k=str(k)
                fill_table(table,j+2,1+i,False,k,0,0,0,1)
                fill_table_color(table,j+2,1+i,"FFFFFF")
            else:
                fill_table(table,j+2,1+i,False,k,0,0,0,1)
                fill_table_color(table,j+2,1+i,color)
      
    
    table_width=[2,0.75,0.59,0.58,0.59,1.56,0.57]                
    set_table_col_width(table,table_width)  
    
    text_="Conclusion: Model Segment Level RAG status"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')  
    
    text_="In accordance with MRM policy,"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')    
    
    text_="Overall Model RAG status"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    text_=f"The table below shows the overall Model Risk of the {full_form} B-score model that is defined as exposure weighted RAG status of the model segments:"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')     
    
    table=create_table(3+no_segments,5)
    add_table_caption("Overall Model - RAG Status") 
    
    for row in table.rows:
        row.height = Inches(0.3) 
        
    table.cell(0,0).merge(table.cell(0,4))   
    table.cell(2+no_segments,0).merge(table.cell(2+no_segments,2))   
    
    row_0=[f"Aggregate ({full_form} Portfolio)"]
    for i,x in enumerate(row_0):
        fill_table(table,0,0,True,x,255,255,255,1)
        fill_table_color(table,0,0,"0070C0")
    
    row_1=["Segment","Weights (Exposure %)","Model Risk Score","RAG status","Exposure Weighted Score Points"]
    for i,x in enumerate(row_1):
        fill_table(table,1,i,True,x,255,255,255,1)
        fill_table_color(table,1,i,"0070C0") 

    # if portfolio_code=="pl":
    #     col_0=["Segment 1: STL- Current","Segment 2: STL- Delinquent"]
    #     row_1=[]
    #     row_2=[]
    #     rows_0=[row_1,row_2]

    if portfolio_code=="nl":
        if no_segments==1:
            col_0=["Segment 1: Current"]
            row_1=[]
            rows_0=[row_1]
        elif no_segments==2:
            col_0=["Segment 1: Current","Segment 2: Delinquent"]
            row_1=[] 
            row_2=[]
            rows_0=[row_1,row_2]
        
    # elif portfolio_code=="al":
    #     col_0=[]
    # else:
    #     col_0=[]   
    
        
    for i,x in enumerate(col_0):
        fill_table(table,i+2,0,True,x,0,0,0,0)
        fill_table_color(table,i+2,0,"FFFFFF") 
        
    def rag_status_based_color(status):
        if status=="Red":
            color="#FF0000"
        elif status=="Amber":
            color="#FFC000"
        else:
            color="#00B050"
        return color
          
    scorecard_available=portfolio_wise_score_model_code(segment)

    overall_rag_temp_df["%exposure"]=overall_rag_temp_df["exposure"]/overall_rag_temp_df["exposure"].sum()
    overall_rag_exposure_sum=overall_rag_temp_df["exposure"].sum()
    overall_rag_temp_df_temp=overall_rag_temp_df["%exposure"].to_list()
    for i,x in enumerate(overall_rag_temp_df_temp):
        x=round(x,4)
        x="{:2%}".format(x)[:-5]+"%"
        fill_table(table,i+2,1,False,x,0,0,0,1)
        fill_table_color(table,i+2,1,"FFFFFF")         
    
    total_exposure_wtd_rag_score=0
    for count_seg,x in enumerate(scorecard_available):
        
        exposure_wtd_score=overall_model_rag_score[x][0]*overall_rag_temp_df_temp[count_seg]
        exposure_wtd_score=round(exposure_wtd_score,2)
        total_exposure_wtd_rag_score+=exposure_wtd_score
        overall_model_rag_score[x].insert(2,exposure_wtd_score)
        for i,j in enumerate(overall_model_rag_score[x][:-1]):
            j=str(j)
            fill_table(table,2+count_seg,2+i,False,j,0,0,0,1)
            color=overall_model_rag_score[x][-1]
            if i==1:
                fill_table_color(table,2+count_seg,2+i,color)
            else:
                fill_table_color(table,2+count_seg,2+i,"FFFFFF")
        
    row_last=["Weighted Score","",""]
    row_last.append(str(round(total_exposure_wtd_rag_score,2)))
    
    if total_exposure_wtd_rag_score<=1.3:
        final_rag_color="#00B050"
        final_rag_status="Green"
        
    elif total_exposure_wtd_rag_score<=2.35:
        final_rag_color="#FFC000"
        final_rag_status="Amber"

    else:
        final_rag_color="#FF0000"
        final_rag_status="Red"

    row_last.insert(3,final_rag_status)
    for i,x in enumerate(row_last):
        if i==3:
            fill_table(table,2+no_segments,i,True,x,0,0,0,1)
            fill_table_color(table,2+no_segments,i,final_rag_color) 
        else:
            fill_table(table,2+no_segments,i,True,x,0,0,0,1)
            fill_table_color(table,2+no_segments,i,"FFFFFF")             
        
    table_width=[2,1.1,1.2,1.2,1]                
    set_table_col_width(table,table_width) 
    
    text_="Conclusion: Overall Model RAG status"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    text_1="Based on this score, overall Model RAG status is "
    if (final_rag_status=="Green"):
        text_2="GREEN – Fit for Purpose "
        color=(0,176,80)
    elif(final_rag_status=="Amber"):
        text_2="AMBER – Not Fit for Purpose "  
        color=(255,192,0)
    else:
        text_2="RED – Not Fit for Purpose "
        color=(255,0,0)
    text_3="with action items."

    
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=True,underline=False,italic=False,font_pt=10,font_color=color,font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_3,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')

        
    add_auto_numbered_heading(level=3,text="Action Items",bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    text_="Previous Action Items"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    text_=f"Validation Team reviewed the status of the Action items issued during {str(val_end_date_imm.year-1)} Annual Validation. Refer to Table 13 for more details."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')   
    
    table=create_table(2,7)
    add_table_caption(f"Previous Action Items") 
    
    for row in table.rows:
        row.height = Inches(0.3) 
        
    row_0=["S.No.","Validation","Segment","Issue","Action Item","Severity","Status"]
    for i,x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)
        fill_table_color(table,0,i,"0070C0") 
        
    table_width=[0.4,0.75,0.9,1.31,1.85,.65,0.55]               
    set_table_col_width(table,table_width)
        
    text_="Conclusion: Previous Action items"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')    
    
    text_="New Action items"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    text_=f"New set of action items were raised during this year i.e., {str(val_end_date_imm.year)} annual validation. Refer to Table 14 for more details."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')   

    table=create_table(2,7)
    add_table_caption(f"New Action Items") 
    
    for row in table.rows:
        row.height = Inches(0.3) 
        
    row_0=["S.No.", "Segment",	"Issue",	"Action Item","Severity","Remediation Plan",	"Action item Owner"]
    for i,x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)
        fill_table_color(table,0,i,"0070C0") 
        
    table_width=[0.45,0.85,0.85,1.85,0.7,0.95,0.9]                
    set_table_col_width(table,table_width)
    
    text_="Conclusion: New Action items"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    add_auto_numbered_heading(level=2,text="Conclusion",bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    doc.add_page_break()
    
    return

#...................................................................................................................................................................................#

def document_part_2(exclusion_df,portfolio_code,val_start_date_imm,val_start_date_omm,val_end_date_imm,val_end_date_omm,path_folder,formatted_date=None):
    
    def calculate_and_plot_portfolio_summary(df):
        temp=df[df["MIS_RECON_EXCLUSION_REASON"].str.contains("POST_EXCLUSION")]
        temp_=temp.groupby("SCORE_DATE").agg(accounts=pd.NamedAgg(column="ACCOUNT_ID",aggfunc='count'),
                                      exposure=pd.NamedAgg(column="LOAN_BALANCE_00",aggfunc='sum')).reset_index(drop=True)
        temp_["exposure"]=round(temp_["exposure"]/1000000,2)
        temp_.insert(0,'SCORE_DATE',exclusion_df['SCORE_DATE'].unique())
        temp_["SCORE_DATE"]=pd.to_datetime(temp_["SCORE_DATE"])
        temp_["SCORE_DATE"]=temp_["SCORE_DATE"].apply(lambda x: x.month_name()[:3]+"-"+str(x.year)[2:])

        accounts_min_max={}
        accounts_min_max["minimum"]='{:,}'.format(temp_['accounts'].min())
        accounts_min_max["maximum"]='{:,}'.format(temp_['accounts'].max())
        exposure_min_max={}
        exposure_min_max["minimum"]='{:,}'.format(int(round(temp_['exposure'].min(),0)))
        exposure_min_max["maximum"]='{:,}'.format(int(round(temp_['exposure'].max(),0)))   
        
        sns.set(style="white", palette="dark",font_scale=0.9,rc={"figure.figsize":(20, 6)})
        fig,ax1=plt.subplots()
        bar_plot= sns.barplot(data=temp_,x="SCORE_DATE", y='accounts',ax=ax1,
                              color='#002060',width=0.4,gap=0.15,linewidth=0.75,label="No. of Accounts")
        
        ax2=ax1.twinx()
        
        line_plot = sns.lineplot(data=temp_,x="SCORE_DATE",
                                 y='exposure',ax=ax2,color="#C00000",linewidth=5,label="Exposure(AED Mn.)")
        
        # for index, row in df_1_2.iterrows():
        
        #     plt.annotate(f'{row[y]:.2f}%',(row[x],row[y]),textcoords="offset points",xytext=(2,4), ha='center',fontsize=8,color='#002060')
        # line_plot = sns.lineplot(data=df_1_1,x=x, y=y,hue='QTR',palette="flare",linestyle='dashed',linewidth=0.8)
        # logger.debug(f"Line graph plotted for {len(df_1_2[x])} score bands.")
        plt.title('Portfolio Summary', {'fontsize':15,'color':"#002060",'fontweight':1000})        
        ax1.set_xlabel('SCORE_DATE',{'fontsize':15,'color':'#002060','fontweight':1000},labelpad=10)        
        ax1.set_ylabel('No. of Accounts',{'fontsize':15,'color':'#002060','fontweight':1000},labelpad=10)        
        ax2.set_ylabel('Exposure(AED Mn.)',{'fontsize':15,'color':'#002060','fontweight':1000},labelpad=10)   
        ax1.legend(loc='upper right',fontsize=12,bbox_to_anchor=(0.50,0.95))        
        ax2.legend(loc='upper right',fontsize=12,bbox_to_anchor=(0.65,0.95))
        bar_plot.set_yticklabels(bar_plot.get_yticklabels(),fontsize=13)        
        line_plot.set_yticklabels(line_plot.get_yticklabels(),fontsize=13)
        
        #ax2.set_yticklabels(np.arange(0,1000,100),fontsize=9)
        bar_plot.set_xticklabels(bar_plot.get_xticklabels(),fontsize=13,rotation=90)        
        ax1.set_ylim(top=max(temp_['accounts'])+3000) 
        #ax2.set_ylim(0, round((max(temp_['exposure']) + 200) / 100 ) * 100)
        ax2.set_ylim(bottom=min(temp_['exposure']) - 2, top=max(temp_['exposure']) + 2)
        ax2.set_yticks(ax2.get_yticks())

      
        #ax2.set_ylim(top=min(temp_['exposure'])+2000)
        #ax2.set_ylim(bottom=min(temp_['exposure'])-200)
        for spine in line_plot.spines.values():        
            spine.set_edgecolor("#BFBFBF")        
            spine.set_linewidth(0.75)
        
        # plt.show()
        
        plt.savefig(f'{path_folder}/nl/{portfolio_code.upper()}_PORTFOLIO_SUMMARY.png',bbox_inches='tight')
                
        return accounts_min_max,exposure_min_max
    
    accounts_min_max,exposure_min_max=calculate_and_plot_portfolio_summary(df=exclusion_df)
    
    def calculate_and_plot_odr_trend(df):
        
        temp_df=df[(df["SCORE_DATE"]<=val_end_date_omm) & (exclusion_df['BSCORE_EXCLUSION_REASON'].str.contains("POST_EXCLUSION"))]
        temp_df["TARGET_12"].replace(99,0,inplace=True)
        if portfolio_code=="nl":
            temp_df=temp_df[~temp_df["SCORE_DATE"].isin(["2020-04-30","2020-05-31","2020-06-30","2020-07-31","2021-05-31"])]
        
        temp_df_grp=temp_df.groupby("SCORE_DATE").agg(count=pd.NamedAgg(column='ACCOUNT_ID',aggfunc='count'),
                                               bads=pd.NamedAgg(column='TARGET_12',aggfunc='sum')).reset_index(drop=True)
        temp_df_grp["bad_rate"]=round(temp_df_grp["bads"]*100/temp_df_grp["count"],2)
        temp_df_grp.insert(0,'SCORE_DATE',temp_df['SCORE_DATE'].unique())
        temp_df_grp["SCORE_DATE"]=pd.to_datetime(temp_df_grp["SCORE_DATE"])
        temp_df_grp["SCORE_DATE"]=temp_df_grp["SCORE_DATE"].apply(lambda x: x.month_name()[:3]+"-"+str(x.year)[2:])
        del temp_df_grp["bads"]
        del temp_df_grp["count"]
        
        temp_df_grp_avg=round(temp_df_grp['bad_rate'].mean(),2)
        # logger.info(f"Plotting of Rank Order graph begins for {score_model_code} scorecard.")
        
        sns.set(style="white", palette="dark",font_scale=0.9,rc={"figure.figsize":(20, 6)})        
        fig,ax1=plt.subplots()        
        line_plot = sns.lineplot(data=temp_df_grp,x="SCORE_DATE", y='bad_rate',ax=ax1,color="#C00000",linewidth=5,label="ODR(%)")
        # for index, row in temp_df_grp.iterrows():
            # line_plot.annotate(str(row['bad_rate'])+"%",(row['SCORE_DATE'],row['bad_rate']),textcoords="offset points",xytext=(6,20), ha='center',fontsize=10,color='#002060')
                 
        # logger.debug(f"Line graph plotted for {len(df_1_2[x])} score bands.")
        plt.title('ODR Trend', {'fontsize':15,'color':"#002060",'fontweight':1000})
        ax1.set_xlabel('SCORE_DATE',{'fontsize':15,'color':'#002060','fontweight':1000},labelpad=10)
        ax1.set_ylabel('ODR(%)',{'fontsize':15,'color':'#002060','fontweight':1000},labelpad=10)        
        ax1.legend(loc='upper right',fontsize=12,bbox_to_anchor=(1,0.95))        
        # bar_plot.set_yticklabels(bar_plot.get_yticklabels(),fontsize=9)
        line_plot.set_yticklabels(line_plot.get_yticklabels(),fontsize=13)
        # ax2.set_yticklabels(np.arange(0,25000,2000),fontsize=9)        
        line_plot.set_xticklabels(line_plot.get_xticklabels(),fontsize=13,rotation=90)
        ax1.set_ylim(bottom=min(temp_df_grp['bad_rate'])-0.35,top=max(temp_df_grp['bad_rate'])+0.35)
        ax1.set_yticks(ax1.get_yticks())
        # plt.ylim(0,max(temp_df_grp['bad_rate']+0.2))
        # plt.yticks(range(0,1.0,0.1))
        for spine in line_plot.spines.values():
            spine.set_edgecolor("#BFBFBF")
            spine.set_linewidth(0.75)
        # plt.show()
        plt.savefig(f'{path_folder}/nl/{portfolio_code.upper()}_ODR_Trend.png',bbox_inches='tight')        
        # logger.info(f"Plotting of Rank Order graph ends for {score_model_code} scoremodel.")
        # logger.debug("#######################################################################################")
        
        return temp_df_grp_avg
    
    temp_df_grp_avg=calculate_and_plot_odr_trend(df=exclusion_df)
            
    def paragraph_add(space_after,space_before,text,bold,font_pt,font_color,font_name):    
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        run=paragraph.add_run()
        run.text= text
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return
    
    def mixed_paragraph_add(space_after,space_before):
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        
        return paragraph
    
    def mixed_run_add(paragraph,text,bold,underline,italic,font_pt,font_color,font_name):
        run=paragraph.add_run(text)
        run.bold=bold
        run.underline=underline
        run.italic=italic
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return

    def b_u_paragraph_add(space_after,space_before,text,bold,underline,italic,font_pt,font_color,font_name):
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        run=paragraph.add_run()
        run.text= text
        run.bold=bold
        run.underline=underline
        run.italic=italic
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return
        
    def add_auto_numbered_heading(level, text,bold,font_pt,font_color,font_name):
        if level==1:
            heading_counters[1]=0
        if level==2:
            heading_counters[2]=0
        if level==3:
            heading_counters[3]=0
        heading_counters[level - 1] += 1 # Increment the counter for the current level
        number = '.'.join(str(heading_counters[i]) for i in range(level)) # Generate the numbering
        heading = doc.add_paragraph()
        run = heading.add_run()
        run.text=f"{number}  {text}"
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        heading.style = f'Heading {level}'
        heading.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        if level==1:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(18)
        else:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(6)            
        return 

    def create_table(n_rows,n_cols):
        table=doc.add_table(rows=n_rows, cols=n_cols)
        table.style='Table Grid' 
        table.alignment=WD_TABLE_ALIGNMENT.CENTER
        table.autofit=True 
        
        for i in range(len(table.rows)):
            for j in range(len(table.columns)):
        
                table.cell(i,j).paragraphs[0].alignment  = WD_ALIGN_VERTICAL.CENTER
                table.cell(i,j).vertical_alignment=WD_CELL_VERTICAL_ALIGNMENT.CENTER
                
        return table
    
    def fill_table_color(table,row,col,color="FFFFFF"):
        shading_elm_1 = parse_xml(r'<w:shd {} w:fill="{}"/>'.format(nsdecls('w'),color))
        table.rows[row].cells[col]._tc.get_or_add_tcPr().append(shading_elm_1)
        return

    
    def fill_table(table,row,col,bold,text,r,g,b,alignment):
        
        cell_1=table.cell(row,col).paragraphs[0]
        run=cell_1.add_run()
        run.text=text
        run.bold=bold
        run.font.size=Pt(9)
        run.font.name='Arial'
        run.font.color.rgb=word_rgb(r,g,b)
        cell_1.alignment=alignment
        return
    
    def add_table_caption(caption):
        paragraph = doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)
        run=paragraph.add_run()
        run.text=f"Table {len(doc.tables)}: {caption}"
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.style = 'Caption'
        run.font.name='Arial'
        run.bold=True
        run.font.size=Pt(8)
        run.font.color.rgb=word_rgb(0,0,0)  
        
        return
    
    def set_table_col_width(table,width_arr):
        for i,x in enumerate(width_arr):
            for col in table.columns[i].cells:
                col.width=Inches(x)
        return   
    
    
    def add_picture_caption(caption):
        paragraph = doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)
        run=paragraph.add_run()
        run.text=f"Figure {len(doc.inline_shapes)}: {caption}"
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.style = 'Caption'
        run.font.name='Arial'
        run.bold=True
        run.font.size=Pt(8)
        run.font.color.rgb=word_rgb(0,0,0)  
        
        return 
    
    val_start_date_imm=pd.to_datetime(val_start_date_imm)
    val_start_date_omm=pd.to_datetime(val_start_date_omm)    
    val_end_date_imm=pd.to_datetime(val_end_date_imm)
    val_end_date_omm=pd.to_datetime(val_end_date_omm)
         
    sections=doc.sections
    for section in sections:
        section.top_margin=Inches(1)
        section.bottom_margin=Inches(1)
        section.left_margin=Inches(1)
        section.right_margin=Inches(1)
        
    add_auto_numbered_heading(level=1,text="Portfolio Overview",bold=True,font_pt=14,font_color=(0,0,0),font_name='Arial')    
    
    full_form=portfolio_wise_full_form(portfolio_code=portfolio_code)
    if portfolio_code !="nl":
        text_=f"FAB offers {full_form} to both UAE Nationals as well as expats."
    else:
        text_=f"FAB offers {full_form} to UAE Nationals only. These loans have tenure as high as 20 years with options of balloon payment at maturity."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    n_months=len(exclusion_df['SCORE_DATE'].unique()) 
    imm_end_date=val_end_date_imm.month_name()[:3]+"'"+str(val_end_date_imm.year)
    text_=f"The graph below provides a summary of the {full_form.lower()} portfolio across {n_months} months starting {formatted_date} to {imm_end_date}. The portfolio provided below includes the existing defaulted accounts not written-off at each observation month excluding only the policy exclusions."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    
    img_stream_1 = f'{path_folder}/nl/{portfolio_code.upper()}_PORTFOLIO_SUMMARY.png'
    left = Inches(1.5)
    top = Inches(1.15)
    pic = doc.add_picture(img_stream_1,width=Inches(7), height=Inches(3.1))
    last_paragraph=doc.paragraphs[-1]
    last_paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.CENTER
    add_picture_caption("Portfolio Overview (exposure, no of accounts)")
    
    # text_=f"Exposure trend in the validation data is mostly stable with a decreasing trend (ranging between AED {exposure_min_max['minimum']} Mn. to AED {exposure_min_max['maximum']} Mn.). Trends of active accounts volume in the validation data is also mostly stable with a decreasing trend (ranging between {accounts_min_max['minimum']} to {accounts_min_max['maximum']})."
    prompt =  "Provide consice summary in one paragraph, Analyze the chart (containing bank customer data), offering technical insights and interpretations of the numerical data provided. Note that in this context, 'exposure' refers to the balance of loans with customers. Your understanding and insights will help illuminate the significance of the information presented."
    text_ = get_insights(img_stream_1,prompt)
    text_ = " ".join(text_.split())
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    img_stream_1 = f'{path_folder}/nl/{portfolio_code.upper()}_ODR_Trend.png'
    left = Inches(1.5)
    top = Inches(1.15)
    pic = doc.add_picture(img_stream_1,width=Inches(6.65), height=Inches(3.1))
    last_paragraph=doc.paragraphs[-1]
    last_paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.CENTER
    add_picture_caption("Portfolio Overview (ODR)")    
    
    omm_end_date=val_end_date_omm.month_name()[:3]+"'"+str(val_end_date_omm.year)[2:]
    if portfolio_code=="pl":
        text_=f"Average default rate for the entire portfolio from Jan'20 to {omm_end_date} is approximately {temp_df_grp_avg}%. The observed default rate has been volatile in the second half of 2020. Sharp increase in Q3 2020 was followed by a sharp decline in Q4 2020. From Q1 2021 onwards we see a decreasing trend in ODR, however this has stabilised in the recent quarters. This indicates that portfolio quality has improved over time. "
        paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 

    elif portfolio_code=="nl":
        # prompt =  f"Provide consice summary in one paragraph, Analyze the image, offering technical insights and interpretations of the numerical data provided. and also talk about Average default rate for the entire portfolio {temp_df_grp_avg}% please note ODR is observed default rate"
        prompt = f"Could you analyze a line graph (containing bank ODR data) displaying the Observed Default Rate (ODR) for a specified period or across different groups? The graph illustrates the percentage or frequency of observed defaults within a given context, such as a portfolio of loans or credit cards. Consider trends over time, comparative analysis between segments, seasonal patterns, outliers, implications for risk management, and future outlook. Please give inference in one paragraph, supported by data and analysis and Average default rate for the entire portfolio is {temp_df_grp_avg}%"
        text_1 = get_insights(img_stream_1,prompt)
        text_ = " ".join(text_.split())
        paragraph=mixed_paragraph_add(space_after=8,space_before=0)
        mixed_run_add(paragraph=paragraph,text=text_1,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
        # mixed_run_add(paragraph=paragraph,text=text_2,bold=True,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
        # mixed_run_add(paragraph=paragraph,text=text_3,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
        
    else:
        text=""
        paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 

    doc.add_page_break()
    return

#...................................................................................................................................................................................#

def document_part_3(exclusion_df,portfolio_code,val_start_date_imm,val_start_date_omm,val_end_date_imm,val_end_date_omm,segment,formatted_date=None,model_dev_benchmark_date=None,model_dev_snap_date=None,formatted_date_1=None):
    
    def paragraph_add(space_after,space_before,text,bold,font_pt,font_color,font_name):    
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        run=paragraph.add_run()
        run.text= text
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return

    def mixed_paragraph_add(space_after,space_before):
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        
        return paragraph
    
    def mixed_run_add(paragraph,text,bold,underline,italic,font_pt,font_color,font_name):
        run=paragraph.add_run(text)
        run.bold=bold
        run.underline=underline
        run.italic=italic
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return
    
    def b_u_paragraph_add(space_after,space_before,text,bold,underline,italic,font_pt,font_color,font_name):
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        run=paragraph.add_run()
        run.text= text
        run.bold=bold
        run.underline=underline
        run.italic=italic
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return
        
    def add_auto_numbered_heading(level, text,bold,font_pt,font_color,font_name):
        if level==1:
            heading_counters[1]=0
        if level==2:
            heading_counters[2]=0
        if level==3:
            heading_counters[3]=0
        heading_counters[level - 1] += 1 # Increment the counter for the current level
        number = '.'.join(str(heading_counters[i]) for i in range(level)) # Generate the numbering
        heading = doc.add_paragraph()
        run = heading.add_run()
        run.text=f"{number}  {text}"
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        heading.style = f'Heading {level}'
        heading.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        if level==1:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(18)
        else:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(6)            
        return 

    def create_table(n_rows,n_cols):
        table=doc.add_table(rows=n_rows, cols=n_cols)
        table.style='Table Grid' 
        table.alignment=WD_TABLE_ALIGNMENT.CENTER
        table.autofit=True 
        
        for i in range(len(table.rows)):
            for j in range(len(table.columns)):
        
                table.cell(i,j).paragraphs[0].alignment  = WD_ALIGN_VERTICAL.CENTER
                table.cell(i,j).vertical_alignment=WD_CELL_VERTICAL_ALIGNMENT.CENTER
                
        return table
    
    def fill_table_color(table,row,col,color="FFFFFF"):
        shading_elm_1 = parse_xml(r'<w:shd {} w:fill="{}"/>'.format(nsdecls('w'),color))
        table.rows[row].cells[col]._tc.get_or_add_tcPr().append(shading_elm_1)
        return

    
    def fill_table(table,row,col,bold,text,r,g,b,alignment):
        
        cell_1=table.cell(row,col).paragraphs[0]
        run=cell_1.add_run()
        run.text=text
        run.bold=bold
        run.font.size=Pt(9)
        run.font.name='Arial'
        run.font.color.rgb=word_rgb(r,g,b)
        cell_1.alignment=alignment
        return
    def add_table_caption(caption):
        paragraph = doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)
        run=paragraph.add_run()
        run.text=f"Table {len(doc.tables)}: {caption}"
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.style = 'Caption'
        run.font.name='Arial'
        run.bold=True
        run.font.size=Pt(8)
        run.font.color.rgb=word_rgb(0,0,0)  
        
        return
    
    def set_table_col_width(table,width_arr):
        for i,x in enumerate(width_arr):
            for col in table.columns[i].cells:
                col.width=Inches(x)
        return       
 
    
    val_start_date_imm=pd.to_datetime(val_start_date_imm)
    val_start_date_omm=pd.to_datetime(val_start_date_omm)    
    val_end_date_imm=pd.to_datetime(val_end_date_imm)
    val_end_date_omm=pd.to_datetime(val_end_date_omm)
    
    imm_end_date=val_end_date_imm.month_name()[:3]+"'"+str(val_end_date_imm.year)
    omm_start_date=val_start_date_omm.month_name()[:3]+"'"+str(val_start_date_omm.year)
    omm_end_date=val_end_date_omm.month_name()[:3]+"'"+str(val_end_date_omm.year)
         
    sections=doc.sections
    for section in sections:
        section.top_margin=Inches(1)
        section.bottom_margin=Inches(1)
        section.left_margin=Inches(1)
        section.right_margin=Inches(1)

     
    add_auto_numbered_heading(level=1,text="Data Review",bold=True,font_pt=14,font_color=(0,0,0),font_name='Arial')   
    
    # paragraph=doc.add_paragraph()
    # paragraph.paragraph_format.space_after = Pt(0)
    # paragraph.paragraph_format.space_before = Pt(8)
    text_= f"The validation team performed a detailed review of the data encompassing a reconciliation of the data with the MIS reports, assessment of data assumptions, exclusions, and performed data quality checks."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    
    add_auto_numbered_heading(level=2,text="Data Preparation",bold=True,font_pt=12,font_color=(0,0,0),font_name='Arial') 
    
    if portfolio_code!="nl":
        text= model_dev_snap_date
    else:
        
        text=model_dev_snap_date # asked -> input from model development team
        
    text_=f"The Validation Team was provided with benchmark data spanning {text} (NL-Current) by the Model Development Team for discriminatory power and model forecast accuracy testing. Additionally, population stability and concentration tests were conducted using the same {text} benchmark data. Additional Model output data ({formatted_date} to {imm_end_date}) was independently sourced by the Validation team with the help of the IT team. The details of data preparation are present in Table."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')    
        
    table=create_table(4,3)
    add_table_caption("Data Summary")
    for row in table.rows:
        row.height = Inches(0.3)
    
    row_0=["Data Statistics",	"Benchmark",	"Monitoring"]
    for i,x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)
        fill_table_color(table,0,i,"0070C0")
        
    col_0=["Total data period available",
            "Observation window (For GINI, MAPE)",
            "Observation window (For PSI, HCI)"
            ]
    for i,x in enumerate(col_0):
        fill_table(table,i+1,0,False,x,0,0,0,1)
        fill_table_color(table,i+1,0,"FFFFFF")
   
    # if portfolio_code!="nl":
    #     snap="Monthly snapshots from Jan 2018 to Dec 2018"
    # else:
    #     snap="NL-Current - Monthly snapshots from Jan 2018 to Dec 2018 NL-Delinquent - Monthly snapshots from Jan 2015 to Dec 2017"
    col_1=[
            model_dev_benchmark_date,
            model_dev_snap_date ,    # asked
            model_dev_snap_date
            ]
    
    for i,x in enumerate(col_1):
        fill_table(table,i+1,1,False,x,0,0,0,1)
        fill_table_color(table,i+1,1,"FFFFFF")  
        
    imm_start=val_start_date_imm.month_name()[:3]+"'"+str(val_start_date_imm.year) 
    imm_end=val_end_date_imm.month_name()[:3]+"'"+str(val_end_date_imm.year)
    omm_start=val_start_date_omm.month_name()[:3]+"'"+str(val_start_date_omm.year)
    omm_end=val_end_date_omm.month_name()[:3]+"'"+str(val_end_date_omm.year)
    
    col_2=[
            f"{formatted_date_1} to {imm_end}",
            f"Monthly snapshots from {omm_start} to {omm_end}",
            f"{imm_start} to {imm_end}"
            ]
        
    for i,x in enumerate(col_2):
        fill_table(table,i+1,2,False,x,0,0,0,1)
        fill_table_color(table,i+1,2,"FFFFFF")         
    
    # paragraph=doc.add_paragraph()
    # paragraph.paragraph_format.space_after = Pt(8)
    # paragraph.paragraph_format.space_before = Pt(8)
    
    text_="Data Preparation"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    x=val_start_date_omm.month_name()[:3]+"'"+str(val_start_date_omm.year)
    y=val_end_date_omm.month_name()[:3]+"'"+str(val_end_date_omm.year)
    text_=f"The following steps were taken to create out of time data sample from monthly snapshots {x} to {y}."
    paragraph_add(space_after=8,space_before=8,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    scorecard_available=portfolio_wise_score_model_code(segment=segment)
    no=len(scorecard_available)
    
    table=create_table(2+no,6)
    add_table_caption("Steps of Data preparation Summary")
    
    for row in table.rows:
        row.height = Inches(0.3)
        
    table.cell(1,0).merge(table.cell(no,0))
    table.cell(1,1).merge(table.cell(no,1))
    table.cell(1,3).merge(table.cell(no,3))
    table.cell(1,4).merge(table.cell(no,4))
    

    row_0=["Steps",	"File Type","File Name",	"Action",	"Data Time period",	"Number of Observations"]    
    for i,x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)
        fill_table_color(table,0,i,"0070C0")
        
    col_0=["1","2"]
    for i,x in enumerate(col_0):
        fill_table(table,i+no,0,False,x,0,0,0,1)
        fill_table_color(table,i+no,0,"FFFFFF")    
    
    col_1=["Benchmark Data",
           "Monitoring Data"
           ]
    for i,x in enumerate(col_1):
        fill_table(table,i+no,1,False,x,0,0,0,1)
        fill_table_color(table,i+no,1,"FFFFFF") 
        
    # if portfolio_code=="al":
    #     col_2=["Segment 1. NL-Current","Segment 2. AL_CURR_NSTL","Segment 3. AL_DELQ","L3 (IT Output) datasets","Post-exclusion datasets"]
    # elif portfolio_code=="pl":
    #     col_2=["Segment 1. NL-Current","Segment 2. PL-STL-Delinquent","L3 (IT Output) datasets","Post-exclusion datasets"]
    if portfolio_code=="nl":
        if no ==1:
            col_2=["Segment 1. NL-Current","L3 (IT Output) datasets"]
        else:
            col_2=["Segment 1. NL-Current","Segment 2. NL-Delinquent","L3 (IT Output) datasets"]

        
        
    for i,x in enumerate(col_2):
        fill_table(table,i+1,2,False,x,0,0,0,1)
        fill_table_color(table,i+1,2,"FFFFFF")  

    
    col_3=[["Standardizing the benchmark data." , "Renaming the model variables-based on development document",
                "Use Filter DROP_FLAG = 0"],
           " "
              ]
    for i,x in enumerate(col_3):
        if(i!=0):
            fill_table(table,i+no,3,False,x,0,0,0,0)
            fill_table_color(table,i+no,3,"FFFFFF")  
        else:
            for j in range(len(x)):
                if (j!=0):
                    paragraph=table.cell(i+no,3).add_paragraph()
                    paragraph.style='List Bullet'
                    paragraph.style.font.size=Pt(8)
                    paragraph.style.bold=True
                    run=paragraph.add_run(f"{col_3[i][j]}")
                    paragraph.alignment = 0
                    run.font.size=Pt(9)
                    run.font.name='Arial'
                else:
                    paragraph=table.cell(i+no,3).paragraphs[0]
                    run=paragraph.add_run(f"{col_3[i][j]}")
                    paragraph.alignment = 0
                    run.font.size=Pt(9)
                    run.font.name='Arial'  
    
    y=val_end_date_imm.month_name()[:3]+"'"+str(val_end_date_imm.year)
    
    col_4=[model_dev_snap_date,  # asked
            f"{formatted_date} to {y}",  # asked f"{formatted_date_1}- {imm_end}"
          
            ]
    for i,x in enumerate(col_4):
        fill_table(table,i+no,4,False,x,0,0,0,1)
        fill_table_color(table,i+no,4,"FFFFFF") 
      
    pre_exc_accts=exclusion_df.shape[0]
    if portfolio_code!="nl":
        post_exc_accts=exclusion_df[(exclusion_df['BSCORE_EXCLUSION_REASON'].str.contains("POST_EXCLUSION"))].shape[0]
    else:
        post_exc_accts=exclusion_df[(~exclusion_df['SCORE_DATE'].isin(["2020-04-30","2020-05-31","2020-06-30","2020-07-31","2021-05-31"])) & (exclusion_df['BSCORE_EXCLUSION_REASON'].str.contains("POST_EXCLUSION"))].shape[0]
        
    pre_exc=f'{pre_exc_accts:,}'
    post_exc=f'{post_exc_accts:,}'
    # if portfolio_code=="al":
    #     col_5=["144,079","8,070",pre_exc,post_exc]
    # elif portfolio_code=="pl":
    #     col_5=["144,079","18,380",pre_exc,post_exc]
    if portfolio_code=="nl":
        if no ==1:
            col_5=[" ",pre_exc,
                #    post_exc
                   ]   
        else:
            col_5=[" "," ",pre_exc,
                #    post_exc
                   ]   
             
    for i,x in enumerate(col_5):
        fill_table(table,i+1,5,False,x,0,0,0,1)
        fill_table_color(table,i+1,5,"FFFFFF")    

    table.allow_autofit=True
                
    table_width=[0.53,0.78,1.33,2.25,0.69,0.95]                
    set_table_col_width(table,table_width)
        
    # paragraph=doc.add_paragraph()
    # paragraph.paragraph_format.space_after = Pt(8)
    # paragraph.paragraph_format.space_before = Pt(8)        
        
    text_1=f"Representativeness of validation data was checked by reconciling them with Business MIS reports dating from end of {formatted_date} till the {y}. The details of the reconciliation can be checked in the "
    text_2="section 3.2 "
    text_3="of the current report."

    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=True,italic=False,font_pt=10,font_color=(5,99,193),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_3,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')

    text_1= "Post the reconciliation with the MIS reports, exclusions were applied to the data in line with the development document. The details of the exclusion applied, and the corresponding logic are in the "
    text_2="section 3.4 "
    text_3="of the current report."
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=True,italic=False,font_pt=10,font_color=(5,99,193),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_3,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    add_auto_numbered_heading(level=2,text="Data Reconciliation",bold=True,font_pt=12,font_color=(0,0,0),font_name='Arial')

    # paragraph=doc.add_paragraph(style='Heading 2')
    # paragraph.paragraph_format.space_after = Pt(6)
    # paragraph.paragraph_format.space_before = Pt(6)
    # text_="Data Reconciliation"
    # paragraph_add(text=text_,bold=True,font_pt=12,font_color=(0,0,0),font_name='Arial') 
    
    # paragraph=doc.add_paragraph()
    # paragraph.paragraph_format.space_after = Pt(8)
    # paragraph.paragraph_format.space_before = Pt(0)

    text_= "The business MIS reports form an independent check to evaluate the veracity of the validation data. Validation Team sourced the MIS reports from {to be updated} and observed a gap of {to be updated} between the monthly Exposure from Business MIS reports and the data used for Model Validation."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')   
    
    # paragraph=doc.add_paragraph()
    # paragraph.paragraph_format.space_after = Pt(8)
    # paragraph.paragraph_format.space_before = Pt(0)

    text_1= "This suggest that validation data closely reconciles with the Business reports. Detailed analysis is present in the Appendix "
    text_2="Section 8.3"
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=True,italic=False,font_pt=10,font_color=(5,99,193),font_name='Arial')
    
    add_auto_numbered_heading(level=2,text="Data Quality",bold=True,font_pt=12,font_color=(0,0,0),font_name='Arial')

    text_= "Before starting with the validation, basic data quality checks were applied from the perspective of scorecard validation data request. The data was checked for following issues:"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    bullet_points = ["Missing/ Garbage values",
                     "Coverage of all the key scorecard fields",
                     "Availability of fields required for data exclusions if any",
                     "Unique ID for mapping validation and performance data",
                     "Range of DPD values and its intuitiveness",
                     "Presence of duplicate rows in data"
                     ]
    
    for comment in bullet_points:
        paragraph=doc.add_paragraph()
        
        paragraph.paragraph_format.left_indent = Inches(0.5)
        paragraph.paragraph_format.right_indent = Inches(0.5)        
        paragraph.style='List Bullet'
        paragraph.style.font.size=Pt(8)
        paragraph.style.bold=True
        run=paragraph.add_run(f"{comment}")
        run.font.size=Pt(10)
        run.font.color.rgb=word_rgb(0,0,0)
        run.font.name='Arial'
        
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)    
    
    # paragraph=doc.add_paragraph()
    # paragraph.paragraph_format.space_after = Pt(8)
    # paragraph.paragraph_format.space_before = Pt(0)

    text_= "Overall, the data quality was good and in line with what was requested for scorecard validation. All the key data fields had data in acceptable shape and form."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')     
    
    # paragraph=doc.add_paragraph()
    # paragraph.paragraph_format.space_after = Pt(8)
    # paragraph.paragraph_format.space_before = Pt(0)

    text_1= "The details of the data quality checks have been captured in Appendix "
    text_2="Section 8.4"
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=True,italic=False,font_pt=10,font_color=(5,99,193),font_name='Arial')
    
    add_auto_numbered_heading(level=2,text="Waterfall of Exclusions",bold=True,font_pt=12,font_color=(0,0,0),font_name='Arial')

    # paragraph=doc.add_paragraph(style='Heading 2')
    # paragraph.paragraph_format.space_after = Pt(6)
    # paragraph.paragraph_format.space_before = Pt(6)
    # text_="Waterfall of Exclusions"
    # paragraph_add(text=text_,bold=True,font_pt=12,font_color=(0,0,0),font_name='Arial')
    
    # paragraph=doc.add_paragraph()
    # paragraph.paragraph_format.space_after = Pt(8)
    # paragraph.paragraph_format.space_before = Pt(0)

    text_= f"For validation, Model Validation team sourced the benchmark dataset for assessing model discriminatory power from the 2020 new model validation. Model Validation team further sourced additional L3 output data from {formatted_date} to {y} for the purpose of Out of time validation with the help of IT."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
   
    # paragraph=doc.add_paragraph()
    # paragraph.paragraph_format.space_after = Pt(8)
    # paragraph.paragraph_format.space_before = Pt(0)
    portfolio=portfolio_wise_full_form(portfolio_code=portfolio_code)
    text_= f"The below table summarizes the exclusion waterfall across observation window and performance window for {portfolio} portfolio:"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    
    a,b,c=make_waterfall_exclusion(df=exclusion_df,waterfall_type="imm",
                                   val_start_date_omm=val_start_date_omm,val_end_date_omm=val_end_date_omm)
    
    for i,seg in enumerate(c["SCORE_MODEL_CODE"].to_list()):
        c.iloc[i,0]=f"Segment {i+1}: "+c.iloc[i,0]

    val_start_date_omm=pd.to_datetime(val_start_date_omm)
    val_end_date_omm=pd.to_datetime(val_end_date_omm)
    
    scorecard_available=portfolio_wise_score_model_code(segment=segment)
    n_rows=3+b.shape[0]+len(scorecard_available)
    n_cols=a.shape[0]+1
    
    table=create_table(n_rows,n_cols)
    add_table_caption(f"Validation Observation Data Exclusion Waterfall {formatted_date} to {y}")
    for row in table.rows[:1]:
        row.height = Inches(0.3)
        
    table.cell(0,0).merge(table.cell(0,n_cols-1))
    text=f"Validation data waterfall for Stability and Concentration statistics {formatted_date} to {y}"
    fill_table(table,0,0,True,text,255,255,255,1)
    fill_table_color(table,0,0,color="#0070C0")
    
    row_1=sorted(a['YEAR_SCORE'].astype('str').unique())
    for i,x in enumerate(row_1):
        if(i!=(len(row_1)-1)):
            x=f"Jan'{x}-Dec'{x}"
        else:
            x=f"Jan'{x}-{y}"
        fill_table(table,1,i+1,True,x,0,0,0,1)
        fill_table_color(table,1,i+1,color="#DBDBDB") 
    
    fill_table_color(table,1,0,color="#DBDBDB")
    
    row_2=["Total Observations"] + a["Total Observations"].to_list()
    for i,x in enumerate(row_2):
        if(i==0):
            fill_table(table,2,i,False,x,0,0,0,0)
        else:
            fill_table(table,2,i,False,x,0,0,0,2)        
        fill_table_color(table,2,i,color="#FFFFFF")  
        
    # b['BSCORE_EXCLUSION_REASON']="(-) "+b['BSCORE_EXCLUSION_REASON'].str[:]
    
    for j in range(b.shape[0]):
        row_=b.iloc[j,:].to_list()
        if(j!=b.shape[0]-1):
            for i,x in enumerate(row_):
                if(i==0):
                    fill_table(table,j+3,i,False,x,0,0,0,0)
                else:
                    fill_table(table,j+3,i,False,x,0,0,0,2)        
                fill_table_color(table,j+3,i,color="#FFFFFF") 
        else:
            for i,x in enumerate(row_):
                if(i==0):
                    fill_table(table,j+3,i,True,x,0,0,0,0)
                else:
                    fill_table(table,j+3,i,True,x,0,0,0,2)        
                fill_table_color(table,j+3,i,color="#DBDBDB")        
            
    for j in range(c.shape[0]):
        row_=c.iloc[j,:].to_list()  
        for i,x in enumerate(row_):
            if(i==0):
                fill_table(table,j+3+b.shape[0],i,False,x,0,0,0,0)
            else:
                fill_table(table,j+3+b.shape[0],i,False,x,0,0,0,2)        
            fill_table_color(table,j+3+b.shape[0],i,color="#FFFFFF") 
    
        
    #set_table_col_width(table,[2.5,1,1,1,1])
    
    # paragraph=doc.add_paragraph()
    # paragraph.paragraph_format.space_after = Pt(0)
    # paragraph.paragraph_format.space_before = Pt(8)
    
    a,b,c=make_waterfall_exclusion(df=exclusion_df,waterfall_type="omm", val_start_date_omm=val_start_date_omm,val_end_date_omm=val_end_date_omm)
    
    exclusion_df_waterfall=exclusion_df[exclusion_df["BSCORE_EXCLUSION_REASON"].str.contains("POST_EXCLUSION")]
    exclusion_df_waterfall=exclusion_df_waterfall[(exclusion_df_waterfall["SCORE_DATE"]>=val_start_date_omm)&(exclusion_df_waterfall["SCORE_DATE"]<=val_end_date_omm)]
    exclusion_df_waterfall["TARGET_12"].replace(99,0,inplace=True)
    bad_rate_total=len(exclusion_df_waterfall[exclusion_df_waterfall["TARGET_12"]==1])/exclusion_df_waterfall.shape[0]
    bad_=f"{bad_rate_total:.2%}"
    
    b.iloc[-1,1]=b.iloc[-1,1]+ f" ({bad_})"
    b.iloc[-1,0]=b.iloc[-1,0]+ " (Bad_Rate)"
    
    for i,seg in enumerate(c["SCORE_MODEL_CODE"].to_list()):
        c.iloc[i,0]=f"Segment {i+1}: "+c.iloc[i,0]+" (Bad_Rate)"
        bad_rate_total=len(exclusion_df_waterfall[(exclusion_df_waterfall["SCORE_MODEL_CODE"]==seg) & (exclusion_df_waterfall["TARGET_12"]==1)])/exclusion_df_waterfall[(exclusion_df_waterfall["SCORE_MODEL_CODE"]==seg)].shape[0]
        bad_=f"{bad_rate_total:.2%}"
        c.iloc[i,1]=c.iloc[i,1]+ f" ({bad_})"
        
    val_waterfall_start=val_start_date_omm.month_name()[:3]+"'"+str(val_start_date_omm.year)
    val_waterfall_end=val_end_date_omm.month_name()[:3]+"'"+str(val_end_date_omm.year)
    n_rows=2+b.shape[0]+len(scorecard_available)
    n_cols=2
    
    
    table=create_table(n_rows,n_cols)
    add_table_caption(f"Validation Performance Data Exclusion Waterfall {val_waterfall_start} to {val_waterfall_end}")
    for row in table.rows[:1]:
        row.height = Inches(0.3)   
    
    table.cell(0,0).merge(table.cell(0,n_cols-1))
    text=f"Validation data waterfall for Discriminatory Power and Calibration statistics {val_waterfall_start} to {val_waterfall_end}"
    fill_table(table,0,0,True,text,255,255,255,1)
    fill_table_color(table,0,0,color="#0070C0")
    
    row_1=["Total Observations"] + [a]
    for i,x in enumerate(row_1):
        if(i==0):
            fill_table(table,1,i,False,x,0,0,0,0)
        else:
            fill_table(table,1,i,False,x,0,0,0,2)        
        fill_table_color(table,1,i,color="#FFFFFF")
    
    for j in range(b.shape[0]):
        row_=b.iloc[j,:].to_list()
        if(j!=b.shape[0]-1):
            for i,x in enumerate(row_):
                if(i==0):
                    fill_table(table,j+2,i,False,x,0,0,0,0)
                else:
                    fill_table(table,j+2,i,False,x,0,0,0,2)        
                fill_table_color(table,j+2,i,color="#FFFFFF") 
        else:
            for i,x in enumerate(row_):
                if(i==0):
                    fill_table(table,j+2,i,True,x,0,0,0,0)
                else:
                    fill_table(table,j+2,i,True,x,0,0,0,2)        
                fill_table_color(table,j+2,i,color="#DBDBDB")
                
    for j in range(c.shape[0]):
        row_=c.iloc[j,:].to_list()  
        for i,x in enumerate(row_):
            if(i==0):
                fill_table(table,j+2+b.shape[0],i,False,x,0,0,0,0)
            else:
                fill_table(table,j+2+b.shape[0],i,False,x,0,0,0,2)        
            fill_table_color(table,j+2+b.shape[0],i,color="#FFFFFF") 
    
    set_table_col_width(table,[4.75,1.74])
    
    # paragraph=doc.add_paragraph()
    # paragraph.paragraph_format.space_after = Pt(8)
    # paragraph.paragraph_format.space_before = Pt(0)
    text_1="The replication of the exclusion logic from the model development was completed by reconciling the numbers based on the codes. No logical inconsistency was found in their codes. The details of the exclusion analysis can be captured in Appendix "
    text_2="Section 8.4"
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=True,italic=False,font_pt=10,font_color=(5,99,193),font_name='Arial')
    
    add_auto_numbered_heading(level=2,text="Model Implementation results",bold=True,font_pt=12,font_color=(0,0,0),font_name='Arial')     
    
    text_="As part of the annual model validation exercise, the model validation team looked at reconciling L3 generated b-scores vis-à-vis manual scoring. Following were the steps involved."
    paragraph_add(space_after=8,space_before=8,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')    
    
    bullet_points=["Raw scores were calculated from variable weights for each segment based on the equation of model from development document.",
"The raw score was further recalibrated based on model equation from recalibration document.", 
"The manually generated scores were compared with system generated b-scores for the validation data."
]
    for comment in bullet_points:
        paragraph=doc.add_paragraph()
        
        paragraph.paragraph_format.left_indent = Inches(0.5)
        paragraph.paragraph_format.right_indent = Inches(0.5)        
        paragraph.style='List Number'
        paragraph.style.font.size=Pt(10)
        paragraph.style.font.name="Arial"
        paragraph.style.bold=True
        run=paragraph.add_run(f"{comment}")
        run.font.size=Pt(10)
        run.font.color.rgb=word_rgb(0,0,0)
        run.font.name='Arial'
        
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)   
    
    scorecard_available=portfolio_wise_score_model_code(segment=segment)
    no_segments=len(scorecard_available)
    table=create_table(1+no_segments,3)
    add_table_caption("Model Implementation Summary")  
    
    table.cell(1,0).merge(table.cell(no_segments,0))
    
    for row in table.rows:
        row.height = Inches(0.3)
        
    row_0=["Implementation Check",	"Scorecard Segment",	"KPIs"]
    for i,x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)
        fill_table_color(table,0,i,"0070C0")
        
    col_0=["Reconciliation of B-Score (Manually calculated) with B-Score (System calculated in ODS (L3))"]
    for i,x in enumerate(col_0):
        fill_table(table,i+1,0,False,x,0,0,0,1)
        fill_table_color(table,i+1,0,"FFFFFF") 
       
    if portfolio_code=="nl":
        if no_segments==1:
            col_1=["NL-Current"]
        else:
            col_1=["NL-Current","NL-Delinquent"]

        
    # elif portfolio_code=="pl":
    #     col_1=["PL-STL-Current","PL-STL-Delinquent"]
    # elif portfolio_code=="al":
    #     col_1=["",""]
    for i,x in enumerate(col_1):
        fill_table(table,i+1,1,False,x,0,0,0,1)
        fill_table_color(table,i+1,1,"FFFFFF")

        if no_segments==1:
            col_2=[f"100% match between {formatted_date} to {imm_end_date}."]
  
        else:
            col_2=[f"100% match between {formatted_date} to {imm_end_date}.", f"100% match between {formatted_date} to {imm_end_date}."]
  
        
    for i,x in enumerate(col_2):
        fill_table(table,i+1,2,False,x,0,0,0,1)
        fill_table_color(table,i+1,2,"FFFFFF")

    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    
    text_1="Conclusion :"
    text_2=" The replication produced a 100% match."
    
    mixed_run_add(paragraph=paragraph,text=text_1,bold=True,underline=False,italic=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    text_1="The details of the model replication results are in the Appendix "
    text_2="Section 8.5"
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=True,italic=False,font_pt=10,font_color=(5,99,193),font_name='Arial')
    
    # add_auto_numbered_heading(level=3,text="Model Implementation for MOB LT6",bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial')    
    
    # text_="Implemented Approach for accounts less than 6 MOB (LT6)"
    # paragraph_add(space_after=8,space_before=0,text=text_,bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    # text_="The currently implemented approach for accounts having less than 6 months on books (MOB LT6) is built around the following points:"
    # paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
#     bullet_points=[
#     "TTC_PD Current implies the average 6th month B Score PD for good accounts with MOB less than 6 in last 12 months.", 
# "TTC_PD Delinquent implies the average 6th month B Score PD for delinquent accounts with MOB less than 6 in last 12 months.",
# "Average of the above scores is used based on current or delinquent status in the latest month to score MOB less than 6 accounts."
# ]
#     for comment in bullet_points:
#         paragraph=doc.add_paragraph()
        
#         paragraph.paragraph_format.left_indent = Inches(0.5)
#         paragraph.paragraph_format.right_indent = Inches(0.5)        
#         paragraph.style='List Bullet'
#         paragraph.style.font.size=Pt(8)
#         paragraph.style.bold=True
#         run=paragraph.add_run(f"{comment}")
#         run.font.size=Pt(10)
#         run.font.color.rgb=word_rgb(0,0,0)
#         run.font.name='Arial'
        
#         paragraph.paragraph_format.space_after = Pt(6)
#         paragraph.paragraph_format.space_before = Pt(6)   
        
#     text_="Conclusion"
#     b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
        
#     text_="Validation replicated the approach for MOB LT6 accounts and reconciled the numbers with production. It was observed that:"
#     paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
        
#     bullet_points=[
#     "There is a clear separation between PD’s assigned to current and delinquent segment and the rank order is maintained.", 
# "Numbers reconciled closely with the currently implemented production PD figure for Current Segment.",
# "Although there is some deviation from the production PD figure for Delinquent Segment, this can be attributed to the difference in time windows since validation average is based on data from Aug’22 to Jul’23. An action item has been raised to update the production PD."
# ]
#     for comment in bullet_points:
#         paragraph=doc.add_paragraph()
        
#         paragraph.paragraph_format.left_indent = Inches(0.5)
#         paragraph.paragraph_format.right_indent = Inches(0.5)        
#         paragraph.style='List Bullet'
#         paragraph.style.font.size=Pt(8)
#         paragraph.style.bold=True
#         run=paragraph.add_run(f"{comment}")
#         run.font.size=Pt(10)
#         run.font.color.rgb=word_rgb(0,0,0)
#         run.font.name='Arial'
        
#         paragraph.paragraph_format.space_after = Pt(6)
#         paragraph.paragraph_format.space_before = Pt(6)   
    
    # text_="A comparison of long-term average of monthly PD estimate (Aug’22-Jul’23) with production PD is presented in Table 21."
    # paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    # table=create_table(1+no_segments,4)
    # add_table_caption("Comparison of production PD with average PD for MOB LT6")  
    
    # for row in table.rows:
    #     row.height = Inches(0.4)
        
    # row_0=["Segment","Production PD (as of Feb’23)","Average of monthly PD estimates (Aug’22-Jul’23)","Difference"]
    # for i,x in enumerate(row_0):
    #     fill_table(table,0,i,True,x,255,255,255,1)
    #     fill_table_color(table,0,i,"0070C0")
        
    # # if portfolio_code=="pl":
    # #     col_0=["PL-STL-Current","PL-STL-Delinquent"]
    # #     col_1=["1.63%","19.64%"]
    # #     col_2=["2.51%","9.33%"]
    # #     col_3=["88 bips","1031 bips"]
        
    # if portfolio_code=="nl":
    #     if no_segments ==1:
    #         col_0=["NL-Current"]
    #         col_1=["1.51%"]
    #         col_2=["1.60%"]
    #         col_3=["9 bips"]
    #     else:
    #         col_0=["NL-Current","NL-Delinquent"]
    #         col_1=["1.51%","22.26%"]
    #         col_2=["1.60%","19.13%"]
    #         col_3=["9 bips","313 bips"]

        
    
    # for i,x in enumerate(col_0):
    #     fill_table(table,i+1,0,True,x,0,0,0,1)
    #     fill_table_color(table,i+1,0,"FFFFFF")  
        
    # for i,x in enumerate(col_1):
    #     fill_table(table,i+1,1,False,x,0,0,0,1)
    #     fill_table_color(table,i+1,1,"FFFFFF")
        
    # for i,x in enumerate(col_2):
    #     fill_table(table,i+1,2,False,x,0,0,0,1)
    #     fill_table_color(table,i+1,2,"FFFFFF")
        
    # for i,x in enumerate(col_3):
    #     fill_table(table,i+1,3,False,x,0,0,0,1)
    #     fill_table_color(table,i+1,3,"FFFFFF")
        
    # text_1="Details are attached in Appendix "
    # text_2="Section 8.8"
    # paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    # mixed_run_add(paragraph=paragraph,text=text_1,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    # mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=True,italic=False,font_pt=10,font_color=(5,99,193),font_name='Arial')
    
    doc.add_page_break()
    return
#...................................................................................................................................................................................#
    
def document_part_4(exclusion_df,portfolio_code,val_start_date_imm,val_start_date_omm,val_end_date_imm,val_end_date_omm,segment,formatted_year=None):
    
    def paragraph_add(space_after,space_before,text,bold,font_pt,font_color,font_name):    
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        run=paragraph.add_run()
        run.text= text
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return
    
    def mixed_paragraph_add(space_after,space_before):
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        
        return paragraph
    
    def mixed_run_add(paragraph,text,bold,underline,italic,font_pt,font_color,font_name):
        run=paragraph.add_run(text)
        run.bold=bold
        run.underline=underline
        run.italic=italic
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return

    
    def b_u_paragraph_add(space_after,space_before,text,bold,underline,italic,font_pt,font_color,font_name):
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        run=paragraph.add_run()
        run.text= text
        run.bold=bold
        run.underline=underline
        run.italic=italic
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
    def add_auto_numbered_heading(level, text,bold,font_pt,font_color,font_name):
        if level==1:
            heading_counters[1]=0
        if level==2:
            heading_counters[2]=0
        if level==3:
            heading_counters[3]=0
        heading_counters[level - 1] += 1 # Increment the counter for the current level
        number = '.'.join(str(heading_counters[i]) for i in range(level)) # Generate the numbering
        heading = doc.add_paragraph()
        run = heading.add_run()
        run.text=f"{number}  {text}"
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        heading.style = f'Heading {level}'
        heading.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        if level==1:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(18)
        else:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(6)            
        return 

    def create_table(n_rows,n_cols):
        table=doc.add_table(rows=n_rows, cols=n_cols)
        table.style='Table Grid' 
        table.alignment=WD_TABLE_ALIGNMENT.CENTER
        table.autofit=True 
        
        for i in range(len(table.rows)):
            for j in range(len(table.columns)):
        
                table.cell(i,j).paragraphs[0].alignment  = WD_ALIGN_VERTICAL.CENTER
                table.cell(i,j).vertical_alignment=WD_CELL_VERTICAL_ALIGNMENT.CENTER
                
        return table
    
    def fill_table_color(table,row,col,color="FFFFFF"):
        shading_elm_1 = parse_xml(r'<w:shd {} w:fill="{}"/>'.format(nsdecls('w'),color))
        table.rows[row].cells[col]._tc.get_or_add_tcPr().append(shading_elm_1)
        return

    
    def fill_table(table,row,col,bold,text,r,g,b,alignment):
        
        cell_1=table.cell(row,col).paragraphs[0]
        run=cell_1.add_run()
        run.text=text
        run.bold=bold
        run.font.size=Pt(9)
        run.font.name='Arial'
        run.font.color.rgb=word_rgb(r,g,b)
        cell_1.alignment=alignment
        return
    
    def add_table_caption(caption):
        paragraph = doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)
        run=paragraph.add_run()
        run.text=f"Table {len(doc.tables)}: {caption}"
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.style = 'Caption'
        run.font.name='Arial'
        run.bold=True
        run.font.size=Pt(8)
        run.font.color.rgb=word_rgb(0,0,0)  
        
        return
    
    def set_table_col_width(table,width_arr):
        for i,x in enumerate(width_arr):
            for col in table.columns[i].cells:
                col.width=Inches(x)
        return       
 

    scorecard_available=portfolio_wise_score_model_code(segment=segment)
    no_segments=len(scorecard_available)

    val_start_date_imm=pd.to_datetime(val_start_date_imm)
    val_start_date_omm=pd.to_datetime(val_start_date_omm)    
    val_end_date_imm=pd.to_datetime(val_end_date_imm)
    val_end_date_omm=pd.to_datetime(val_end_date_omm)
         
    sections=doc.sections
    for section in sections:
        section.top_margin=Inches(1)
        section.bottom_margin=Inches(1)
        section.left_margin=Inches(1)
        section.right_margin=Inches(1)
    
    add_auto_numbered_heading(level=1,text="Validation Results",bold=True,font_pt=14,font_color=(0,0,0),font_name='Arial')    

    portfolio=portfolio_wise_full_form(portfolio_code=portfolio_code)
    text_= f"This section describes the detailed monitoring results. Validation of {portfolio} B-score is a yearly model validation exercise which consists of 2 parts."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    bullet_points = ["Qualitative validation",
                     "Quantitative validation",
                     ]
    
    for comment in bullet_points:
        paragraph=doc.add_paragraph()
        
        paragraph.paragraph_format.left_indent = Inches(0.5)
        paragraph.paragraph_format.right_indent = Inches(0.5)        
        paragraph.style='List Bullet'
        paragraph.style.font.size=Pt(9)
        paragraph.style.font.name='Arial'
        paragraph.style.bold=True
        run=paragraph.add_run(f"{comment}")
        run.font.size=Pt(10)
        run.font.color.rgb=word_rgb(0,0,0)
        run.font.name='Arial'
        
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)    
    
    add_auto_numbered_heading(level=2,text="Qualitative monitoring",bold=True,font_pt=12,font_color=(0,0,0),font_name='Arial')     

    # paragraph=doc.add_paragraph()
    # paragraph.paragraph_format.space_after = Pt(6)
    # paragraph.paragraph_format.space_before = Pt(6)
    text_=f"Considering {portfolio_code.upper()} B-score is an existing model and detailed qualitative validation of model methodology and model fitting were done during the new model validation in {formatted_year}, for this quarterly monitoring exercise only the relevant qualitative monitoring on data quality and modelling exclusions were performed."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')    
    
    # paragraph=doc.add_paragraph()
    # paragraph.paragraph_format.space_after = Pt(6)
    # paragraph.paragraph_format.space_before = Pt(6)
    text_="The table below describes the various qualitative monitoring criteria and their assessment description."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')   

    table=create_table(3,2)
    add_table_caption("Qualitative monitoring criteria")    
    # for row in table.rows[0].cells:
    #     row.height = Inches(0.3)

    for row in table.rows:
        row.height = Inches(0.3)
    
    row_0=["Qualitative Monitoring","Key checks"]
    
    for i,x in enumerate(row_0):   
        fill_table(table,0,i,True,x,255,255,255,1)  
        fill_table_color(table,row=0,col=i,color="0070C0")
        
    row_1=["Review of Monitoring Data Quality",
           ["Are there significant gaps in the data?",
           "Are there unexplained variations in the data?",
            "Is the modelling data reconciling closely with business reports?"
           ]
    ]
    
    for i in range(len(row_1)):
        if(i==0):
            fill_table(table,1,i,False,row_1[i],0,0,0,0)
        else:

            for j in range(len(row_1[1])):
                if(j!=0):
                    paragraph=table.cell(1,i).add_paragraph()
                    paragraph.style='List Bullet'
                    paragraph.style.font.size=Pt(8)
                    paragraph.style.bold=True
                    run=paragraph.add_run(f"{row_1[1][j]}")
                    paragraph.alignment = 0
                    run.font.size=Pt(9)
                    run.font.name='Arial'
                else:
                    paragraph=table.cell(1,i).paragraphs[0]
                    paragraph.style='List Bullet'
                    paragraph.style.font.size=Pt(8)
                    paragraph.style.bold=True
                    run=paragraph.add_run(f"{row_1[1][j]}")
                    paragraph.alignment = 0
                    run.font.size=Pt(9)
                    run.font.name='Arial'                    
                
                             
                
    row_2=["Review of Data and Modelling Exclusions", "Are there any explained dips / rise in the volume of Modelling exclusions?"]
    
    for i,x in enumerate(row_2):   
        if i==0:
            fill_table(table,2,i,False,x,0,0,0,0)  
        else:
            paragraph=table.cell(2,i).paragraphs[0]
            run=paragraph.add_run(f"{x}")
            paragraph.alignment = 0
            paragraph.style='List Bullet'
            paragraph.style.font.size=Pt(8)
            paragraph.style.bold=True
            run.font.size=Pt(9)
            run.font.name='Arial'     
      
    # paragraph=doc.add_paragraph()
    # paragraph.paragraph_format.space_after = Pt(8)
    # paragraph.paragraph_format.space_before = Pt(0)
    text_="The above 2 qualitative validation points are examined in detail as follows."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    add_auto_numbered_heading(level=3,text="Review of validation data quality ",bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    bullet_points=[
    "In Section 2: Portfolio Overview, overall exposure and account volume in the validation data was found to be stable.", 
"In Section 3.2: Data Reconciliation, the validation data reconciled well with portfolio business reports.",
"In  Section 3.3: Data Quality, no major data issue was observed during data quality checks.",
"In Section 3.5: Model Implementation results: B-score generated in L3 production system reconciled 100% with manually calculated B-score."
]
    
    for comment in bullet_points:
        paragraph=doc.add_paragraph()
        
        paragraph.paragraph_format.left_indent = Inches(0.5)
        paragraph.paragraph_format.right_indent = Inches(0.5)        
        paragraph.style='List Bullet'
        paragraph.style.font.size=Pt(8)
        paragraph.style.bold=True
        run=paragraph.add_run(f"{comment}")
        run.font.size=Pt(10)
        run.font.color.rgb=word_rgb(0,0,0)
        run.font.name='Arial'
        
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)
    
    add_auto_numbered_heading(level=3,text="Scope Definition and Exclusions",bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    bullet_points=[
    "In Section 3.4: Waterfall of exclusions, no major data anomaly was found in the number of observations in the exclusions. The number of observations in each segment for performing model validation were adequate after removing exclusions." ]
    
    for comment in bullet_points:
        paragraph=doc.add_paragraph()
        
        paragraph.paragraph_format.left_indent = Inches(0.5)
        paragraph.paragraph_format.right_indent = Inches(0.5)        
        paragraph.style='List Bullet'
        paragraph.style.font.size=Pt(8)
        paragraph.style.bold=True
        run=paragraph.add_run(f"{comment}")
        run.font.size=Pt(10)
        run.font.color.rgb=word_rgb(0,0,0)
        run.font.name='Arial'
        
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)
    
    paragraph=doc.add_paragraph()
    text_1="Conclusion :"
    text_2=" Based on the above results, model validation data was found to be Fit-for-Purpose."
    
    mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=True,underline=False,italic=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')

    add_auto_numbered_heading(level=2,text="Quantitative monitoring",bold=True,font_pt=12,font_color=(0,0,0),font_name='Arial') 

    text_=f"In accordance with MRM policy, following KPIs are used for evaluating the performance of the {portfolio_code.upper()} B-Score model of the overall model and individual segments."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 

    
    bullet_points = ["Population Stability Index i.e., PSI (Stability Analysis)",
                     "Highest Concentration Index i.e., HCI (Concentration Analysis)",
                     "% Drop in GINI (drop in discriminatory power)",
                     "Mean Absolute Percent Error i.e., MAPE (Accuracy analysis)"
                     ]
    
    for comment in bullet_points:
        paragraph=doc.add_paragraph()
        
        paragraph.paragraph_format.left_indent = Inches(0.5)
        paragraph.paragraph_format.right_indent = Inches(0.5)        
        paragraph.style='List Bullet'
        paragraph.style.font.size=Pt(8)
        paragraph.style.bold=True
        run=paragraph.add_run(f"{comment}")
        run.font.size=Pt(10)
        run.font.color.rgb=word_rgb(0,0,0)
        run.font.name='Arial'
        
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)
    
    text_2="The Red-Amber-Green cut-offs applied for each of these tests are as listed below:"
        
    paragraph=doc.add_paragraph()
    paragraph.paragraph_format.space_after = Pt(8)
    paragraph.paragraph_format.space_before = Pt(0)
    run = paragraph.add_run()
    run.text=text_2
    run.font.size=Pt(10)
    run.font.color.rgb=word_rgb(0,0,0)
    run.font.name='Arial'
    

    table=create_table(6,4)
    add_table_caption("RAG cut-offs for Validation Test")
    for i in range(len(table.rows)):
        for j in range(len(table.columns)):
            table.cell(i,j).height=Inches(.5)
            table.cell(i,j).width=Inches(1.6)
            
    col_width=Inches(0.2)
    row_height=Inches(1.5)
    # for col in table.columns:
    #     col.width=col_width
    #     col.height=row_height
        
    for row in table.rows:
        row.height=col_width
        
        
    table.cell(0,0).merge(table.cell(1,0))
    table.cell(0,1).merge(table.cell(0,3))
    
    # def fill_table_color(row,col,color="FFFFFF"):
    #     shading_elm_1 = parse_xml(r'<w:shd {} w:fill="{}"/>'.format(nsdecls('w'),color))
    #     table.rows[row].cells[col]._tc.get_or_add_tcPr().append(shading_elm_1)
        
    fill_table_color(table=table,row=0,col=0,color="0070C0")
    fill_table_color(table=table,row=0,col=1,color="0070C0")
    
    for i in range(3):
        fill_table_color(table,1,i+1,"0070C0")
    
    table.style='Table Grid' 
    
    name_list_1=["Test","Validation Status"]
    name_list_2=["Red","Amber","Green"]
    
    for i,x in enumerate(name_list_1):   
        fill_table(table,0,i,True,x,255,255,255,1)
        
    for i,x in enumerate(name_list_2):
        fill_table(table,1,i+1,True,x,255,255,255,1)
      
    col_0=[
        "% Drop in GINI",
        "Mean Absolute Percent Error (MAPE)",
        "Population Stability Index (PSI)",
        "Highest Concentration Index (HCI)"
            ]
    
    col_1=[
        "GINI drop >30% from benchmark",
        "MAPE > 30%",
        "PSI> 25%",
        "HCI> 30%"
            ]
    
    col_2=[
        "GINI drop >20% and <=30% from benchmark",
        "MAPE >20% and MAPE <=30%",
        "PSI> 10% and PSI<= 25%",
        "HCI> 20% and HCI<= 30%"
            ]
    
    col_3=[
        "GINI drop <=20% from benchmark",
        "MAPE<=20%",
        "PSI<= 10%",
        "HCI<= 20%"
            ]
    
    for i, x in enumerate(col_0):
        fill_table(table,i+2,0,False,x,0,0,0,1)
        fill_table_color(table,row=i+2,col=0)
        
    for i, x in enumerate(col_1):
        fill_table(table,i+2,1,False,x,0,0,0,1)
        fill_table_color(table,row=i+2,col=1,color="FF0000")
        
    for i, x in enumerate(col_2):
        fill_table(table,i+2,2,False,x,0,0,0,1)
        fill_table_color(table,row=i+2,col=2,color="FFC000")
        
    for i, x in enumerate(col_3):
        fill_table(table,i+2,3,False,x,0,0,0,1)
        fill_table_color(table,row=i+2,col=3,color="00B050")

    text_="In Addition, Key performance indicators are also checked at model variable level. These are as follows."
    
    bullet_points=[
    "Character Stability Index",
    "Rank Ordering",
    "Multicollinearity (VIF)",
    "Information Value (IV)"
]
    
    for comment in bullet_points:
        paragraph=doc.add_paragraph()
        
        paragraph.paragraph_format.left_indent = Inches(0.5)
        paragraph.paragraph_format.right_indent = Inches(0.5)        
        paragraph.style='List Bullet'
        paragraph.style.font.size=Pt(9)
        paragraph.style.font.name='Arial'
        paragraph.style.bold=True
        run=paragraph.add_run(f"{comment}")
        run.font.size=Pt(10)
        run.font.color.rgb=word_rgb(0,0,0)
        run.font.name='Arial'
        
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6) 
        
    text_1="The results of quantitative validation as per the above criteria is captured in detail in "
    text_2="section 4.2 "
    text_3="of the document."
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=True,italic=False,font_pt=10,font_color=(5,99,193),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_3,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    text_="To check the model performance according to the existing criteria set in Bank’s MRM framework following approach is used to find model’s final RAG status. Score point allocation based on RAG status of metrics is as below:"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 

    table=create_table(4,2)
    add_table_caption("Score Points")
    for row in table.rows:
        row.height = Inches(0.2)
    
    row_0=[
    "RAG Status","Score Points"]
    
    row_1=["Green",	"1"]
    row_2=["Amber",	"2"]    
    row_3=["Red",	"3"]
    
    for i, x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)
        fill_table_color(table,row=0,col=i,color="0070C0")# 
        
    for i, x in enumerate(row_1):
        fill_table(table,1,i,False,x,0,0,0,1)
        fill_table_color(table,row=1,col=i,color="00B050")
        
    for i, x in enumerate(row_2):
        fill_table(table,2,i,False,x,0,0,0,1)
        fill_table_color(table,row=2,col=i,color="FFC000")
        
    for i, x in enumerate(row_3):
        fill_table(table,3,i,False,x,0,0,0,1)
        fill_table_color(table,row=3,col=i,color="FF0000")
        
    set_table_col_width(table,[2.03,2.03])
    
    text_="Final model score weight is computed as:"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    text_1="Weighted Score = 70%GINI + 10% PSI + 10% HCI + 10% MAPE"
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.CENTER
    mixed_run_add(paragraph=paragraph,text=text_1,bold=True,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Cambria Math')
    
    text_="Based on this weight final model RAG status is computed as below:"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    table=create_table(4,2)
    add_table_caption("Criteria for final RAG Status") 
    for row in table.rows:
        row.height = Inches(0.2)
    
    row_0=[
    "RAG Status","Weighted Score Cut-off"]
    
    row_1=["Green",	"<= 1.3"]
    row_2=["Amber",	"> 1.3 and <= 2.35"]    
    row_3=["Red",	"> 2.35"]
    
    for i, x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)
        fill_table_color(table,row=0,col=i,color="0070C0")# 
        
    for i, x in enumerate(row_1):
        fill_table(table,1,i,False,x,0,0,0,1)
        fill_table_color(table,row=1,col=i,color="00B050")
        
    for i, x in enumerate(row_2):
        fill_table(table,2,i,False,x,0,0,0,1)
        fill_table_color(table,row=2,col=i,color="FFC000")
        
    for i, x in enumerate(row_3):
        fill_table(table,3,i,False,x,0,0,0,1)
        fill_table_color(table,row=3,col=i,color="FF0000")
    
    set_table_col_width(table,[2.03,2.03])
    
    text_1="The results of final model RAG status as per the bank’s MRM framework is captured in detail in "
    text_2="section 1.4.3."
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=True,italic=False,font_pt=10,font_color=(5,99,193),font_name='Arial')
    
    text_="Percentage drop in GINI"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')    
    
    text_="The GINI Coefficient, also called the ‘Accuracy Ratio’, is a measure for evaluating the discriminatory power of a scorecard. By ordering Score ranges from worst to best, it is possible to measure the extent to which the scorecard differentiates between Good and Bad accounts."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    text_="Population Stability Index"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')     
    
    text_="Population Stability Index (PSI) is used for monitoring the changes in distribution of a score between validation sample and development sample."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    text_="When a model deteriorates in performance, checking distributional changes can help with identifying possible causes. If at least one variable has changed significantly or if several variables have moved to some extent, it might be time to recalibrate or to rebuild the model. "
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    text_="The PSI is calculated using the formula below:"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    text_1="PSI= ∑ [(Validation %-Development %)*ln((Validation %)/(Development %))]"
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.CENTER
    mixed_run_add(paragraph=paragraph,text=text_1,bold=True,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Cambria Math') 
    
    text_="The Population Stability Index (PSI) / Character stability Index (CSI) can be interpreted as follows:"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 

    bullet_points=[
    "0% to 10% implies that the population distribution is satisfactory (Green)",
    "10% to 25% implies that population distribution is changing and needs closer monitoring (Amber)",
    "> 25% implies that population distribution has changed significantly (Red)"
]
      
    for comment in bullet_points:
        paragraph=doc.add_paragraph()
        
        paragraph.paragraph_format.left_indent = Inches(0.5)
        paragraph.paragraph_format.right_indent = Inches(0.5)        
        paragraph.style='List Bullet'
        paragraph.style.font.size=Pt(9)
        paragraph.style.font.name='Arial'
        paragraph.style.bold=True
        run=paragraph.add_run(f"{comment}")
        run.font.size=Pt(10)
        run.font.color.rgb=word_rgb(0,0,0)
        run.font.name='Arial'
        
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)     
        
    text_="Concentration Testing"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')  
    
    text_="As per the Bank’s MRM framework, Highest Concentration Index (HCI) measures the highest population concentration bucket of the portfolio. To compute HCI, same buckets are used as used during PSI computation using population from validation data."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    text_="Rank Order Analysis"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    text_="The rank ordering denotes the ability of the model to maintain the monotonic nature of the actual default rate with descending trend. Break of rank order in the top deciles imply that scorecard’s predicted and actual results are not in line."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    text_="Model Forecast Accuracy (MAPE) - Long Term ODR vs Forecasted TTC PD"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')  
    
    text_="This section explains the model forecast accuracy check of the long-term average default rate or the Through the cycle (TTC) PD. Forecasted TTC PD is used along with macro-economic forecast in the estimation of IFRS9 models."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 

    text_="Characteristic Stability Analysis & Risk Discrimination Power (Rank Ordering)"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')  
    
    text_1=f"CSI provides more granular view of the reasons behind changes in overall population distribution. It measures how much the population distribution has changed across the bins of various variables used in the model. Please refer to "
    text_2="section 8.6 "
    text_3=f"in the appendix for detailed CSI results of all the {no_segments} segments."
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=True,italic=False,font_pt=10,font_color=(5,99,193),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_3,bold=False,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    text_="Multicollinearity"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    text_="Multicollinearity is a phenomenon in which one independent variable in a multiple regression model can be linearly predicted by another independent variables with a substantial degree of accuracy. The parameter estimates of such multiple regression models can change erratically in response to small changes in model or data."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')  
    
    text_="Although having Multicollinearity does not reduce model’s predictive power, it however leads to inaccurate assessment of significance of individual parameters. One way to assess multicollinearity is to use Variance Inflation Factor."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')  
        
    text_="VIF regresses each independent parameter with the remaining independent parameters in a multiple regression model to assess the correlation. The following formula is used to compute VIF for each parameter in the model."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')  
    
    text_1="VIF=1⁄((1-R^2))"
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.CENTER
    mixed_run_add(paragraph=paragraph,text=text_1,bold=True,underline=False,italic=False,font_pt=10,font_color=(0,0,0),font_name='Cambria Math')
    
    return
#...................................................................................................................................................................................#


def pasting_imm_monitoring_charts_word(df,y_psi,y_hci,title,portfolio_code,score_model_code,path_folder,bm_year=2020):
    
    ''' 
        This function pastes charts in the Word.
    
        input params:
           
        @ title_name - This takes string variable denoting the Title of the slide.
        @ score_model_code - This takes string variable and accounts for the scorecard whose charts are to be pasted.i.e., "PB_UAE_BSCR04"
        @ ppt_file_path_name - This takes string variable for the path and file name of the ppt.Ex- "ppt/final_ppt_1.pptx"
        
        output params:
        
        It returns a saved presentation at given file path.
    
    '''
    # logger.info(f"Pasting of monitoring charts in ppt begins for {score_model_code} scorecard")

    # from docx.oxml import OxmlElement
    # from docx.shared import Inches,Pt, RGBColor
    # from docx.enum.text import WD_PARAGRAPH_ALIGNMENT    
    
    sections=doc.sections
    for section in sections:
        section.top_margin=Inches(1)
        section.bottom_margin=Inches(1)
        section.left_margin=Inches(1)
        section.right_margin=Inches(1)
        
    def paragraph_add(space_after,space_before,text,bold,font_pt,font_color,font_name):    
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        run=paragraph.add_run()
        run.text= text
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return

    def mixed_paragraph_add(space_after,space_before):
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        
        return paragraph
    
    def mixed_run_add(paragraph,text,bold,underline,italic,font_pt,font_color,font_name):
        run=paragraph.add_run(text)
        run.bold=bold
        run.underline=underline
        run.italic=italic
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return    
    
        
    def add_picture_caption(caption):
        paragraph = doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)
        run=paragraph.add_run()
        run.text=f"Figure {len(doc.inline_shapes)}: {caption}"
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.style = 'Caption'
        run.font.name='Arial'
        run.bold=True
        run.font.size=Pt(7)
        run.font.color.rgb=word_rgb(0,0,0)  
        
        return
        
    # heading_counters=heading_counter()
    def add_auto_numbered_heading(level, text,bold,font_pt,font_color,font_name):
        if level==1:
            heading_counters[1]=0
        if level==2:
            heading_counters[2]=0
        if level==3:
            heading_counters[3]=0
        heading_counters[level - 1] += 1 # Increment the counter for the current level
        number = '.'.join(str(heading_counters[i]) for i in range(level)) # Generate the numbering
        heading = doc.add_paragraph()
        run = heading.add_run()
        run.text=f"{number} {text}"
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        heading.style = f'Heading {level}'
        heading.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        if level==1:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(18)
        else:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(6)            
        return 

    y=title.split()[:2]+title.split()[3:]
    y=" ".join(y)
    add_auto_numbered_heading(level=3,text=y,bold=True,font_pt=12,font_color=(0,0,0),font_name='Arial')
    
   
    # paragraph=doc.add_paragraph()
    # paragraph.paragraph_format.left_indent = Inches(0.5)
    # run = paragraph.add_run()
    # run.text="1. Population Stability Index"
    # run.bold=True
    # run.font.size=Pt(10)
    # run.font.color.rgb=word_rgb(0,0,0)
    # run.font.name='Arial'    
    
    add_auto_numbered_heading(level=4,text="Population Stability Index",bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    full_form=portfolio_wise_full_form(portfolio_code=portfolio_code)
    text_=f"The graphs below show population stability index for out of time validation sample based on Jan'{str(bm_year)} to Dec'{str(bm_year)} period as benchmark for {full_form} portfolio:"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    img_stream_1 = f"{path_folder}/nl/PSI_{score_model_code}.png"
    
    left = Inches(1.5)
    top = Inches(1.15)
    pic = doc.add_picture(img_stream_1,width=Inches(5.65), height=Inches(3.1))
    last_paragraph=doc.paragraphs[-1]
    last_paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.CENTER
    add_picture_caption("PSI")
    
    df_1_psi,color_list_psi,document_color=df_psi_hci_chart(df=df,y=y_psi) 
    df_1_hci,color_list_hci,document_color=df_psi_hci_chart(df=df,y=y_hci)

    val_value_psi=df_1_psi["PSI"].to_list()[-1]

    val_value_hci=df_1_hci["HCI"].to_list()[-1]

    ini_qtr=df_1_psi["YYYYMM"].to_list()[0]

    fin_qtr=df_1_psi["YYYYMM"].to_list()[-1]

    

    if (val_value_psi<=0.1):

        status_psi="LOW"
        status_psi_color=(0,176,80)
        text_psi="The PSI is within the MRM policy’s prescribed threshold of 10%."

    elif (val_value_psi<=.25):

        status_psi="MEDIUM"
        status_psi_color=(255,192,0)       
        text_psi="The PSI is within the MRM policy’s prescribed threshold of 10%-25%."

    else:

        status_psi="HIGH"
        status_psi_color=(255,0,0)
        text_psi="The PSI has breached the MRM policy’s prescribed threshold of 25%."
        

    if (val_value_hci<=0.2):

        status_hci="LOW"
        status_hci_color=(0,176,80)
        text_hci="The HCI is within the MRM policy’s prescribed threshold of 20%."

    elif (val_value_hci<=0.3):

        status_hci="MEDIUM"
        status_hci_color=(255,192,0)
        text_hci="The HCI is within the MRM policy’s prescribed threshold of 20%-30%."

    else:

        status_hci="HIGH"
        status_hci_color=(255,0,0)
        text_hci="The HCI has breached the MRM policy’s prescribed threshold of 30%."
      
    # first_point=f"PSI is calculated using 4 quarters rolling window appraoch. Hence, sample of {fin_qtr} include samples from ({ini_qtr} to {fin_qtr})."
    # second_point=f"Accordingly, for the entire input monitoring time period i.e. ({ini_qtr} to {fin_qtr}), PSI value is {round(val_value_psi*100,2)}%."
    # bullet_points = [first_point,
    #                 second_point
    #             ]

    # for comment in bullet_points:
    #     paragraph=doc.add_paragraph()
    #     paragraph.paragraph_format.left_indent = Inches(0.5)
    #     paragraph.paragraph_format.right_indent = Inches(0.5)        
    #     paragraph.style='List Bullet'
    #     paragraph.style.font.size=Pt(8)
    #     paragraph.style.bold=True
    #     run=paragraph.add_run(f"{comment}")
    #     run.font.size=Pt(10)
    #     run.font.color.rgb=word_rgb(0,0,0)
    #     run.font.name='Arial'
        
    #     paragraph.paragraph_format.space_after = Pt(6)
    #     paragraph.paragraph_format.space_before = Pt(6)
    prompt =  f"Could you interpret an graph depicting PSI (Population Stability Index) numbers? The PSI values represent the population stability within different segments or groups. Please analyze the image to identify trends or variations in PSI across ({ini_qtr} to {fin_qtr}). Consider the implications of PSI values for population stability and provide insights based on the analysis. Please Note if PSI is less than 10% means variable distribution is stable vice versa and  and dont give general information"
    text_ = get_insights(img_stream_1,prompt)
    text_ = " ".join(text_.split())
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    text_1="Conclusion : "
    text_2=text_psi+" Based on the evidence, PSI for the model has been assigned a RAG status of"
    text_3=f" {status_psi}."
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=True,underline=False,italic=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=False,italic=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_3,bold=True,underline=False,italic=True,font_pt=10,font_color=status_psi_color,font_name='Arial')
     
    add_auto_numbered_heading(level=4,text="Concentration Testing",bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    text_=f"The graphs below show HCI for out of time validation sample based on Jan’{str(bm_year)} to Dec’{str(bm_year)} period as benchmark for {full_form} portfolio:"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    img_stream_2 = F"{path_folder}/nl/HCI_{score_model_code}.png"
    left = Inches(1.15)
    top = Inches(4.15)
    pic = doc.add_picture(img_stream_2, width=Inches(5.65), height=Inches(3.1))
    
    last_paragraph=doc.paragraphs[-1]
    last_paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.CENTER
    add_picture_caption("HCI")

    # first_point=f"HCI is calculated using 4 quarters rolling window appraoch. Hence, sample of {fin_qtr} include samples from ({ini_qtr} to {fin_qtr})."
    # second_point=f"Accordingly, for the entire input monitoring time period i.e. ({ini_qtr} to {fin_qtr}), HCI value is {round(val_value_hci*100,2)}%."

 
    # bullet_points = [first_point,
    #                 second_point
    #                 ]

    # for bullet in bullet_points:
    #     paragraph=doc.add_paragraph()
        
    #     paragraph.paragraph_format.left_indent = Inches(0.5)
    #     paragraph.paragraph_format.right_indent = Inches(0.5)
    #     paragraph.style='List Bullet'
    #     paragraph.style.font.size=Pt(8)
    #     paragraph.style.bold=True
    #     run=paragraph.add_run(f"{bullet}")
    #     run.font.size=Pt(10)
    #     run.font.color.rgb=word_rgb(0,0,0)
    #     run.font.name='Arial'
    prompt =  f"""Could you analyze an graph depicting HCI (High Concentration Index) numbers? The HCI values range from -1 to +1 and indicate the level of concentration or inequality within different segments or categories. Please interpret the graph to identify trends or variations in HCI across ({ini_qtr} to {fin_qtr}). Consider the implications of HCI values for concentration and inequality. Highest concentration index (HCI) measures the highest population concentration bucket of the portfolio, and provide insights based on the analysis PLEASE consider RAG cut of 20% if below 20% its Low else High and dont give general information"""
    text_ = get_insights(img_stream_2,prompt)
    text_ = " ".join(text_.split())
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    text_1="Conclusion : "
    text_2=text_hci+" Based on the evidence, HCI for the model has been assigned a RAG status of"
    text_3=f" {status_hci}."

    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=True,underline=False,italic=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=False,italic=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_3,bold=True,underline=False,italic=True,font_pt=10,font_color=status_hci_color,font_name='Arial')
    # logger.debug("HCI graph pasted")
    # logger.debug("###############################################")
    return

#...................................................................................................................................................................................#


def pasting_omm_monitoring_charts_1_word(df_gini,df_rankorder,title,portfolio_code,score_model_code,path_folder,score_model_wise_bm_gini_dict,z=None,model_dev_snap_date=None):
    
    ''' 
        This function pastes charts in the Word.
    
        input params:
           
        @ title_name - This takes string variable denoting the Title of the slide.
        @ score_model_code - This takes string variable and accounts for the scorecard whose charts are to be pasted.i.e., "PB_UAE_BSCR04"
        @ ppt_file_path_name - This takes string variable for the path and file name of the ppt.Ex- "ppt/final_ppt_1.pptx"
        
        output params:
        
        It returns a saved presentation at given file path.
    
    '''

    
    sections=doc.sections
    for section in sections:
        section.top_margin=Inches(1)
        section.bottom_margin=Inches(1)
        section.left_margin=Inches(1)
        section.right_margin=Inches(1)
        
    def paragraph_add(space_after,space_before,text,bold,font_pt,font_color,font_name):    
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        run=paragraph.add_run()
        run.text= text
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return
       
    def mixed_paragraph_add(space_after,space_before):
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        
        return paragraph
    
    def mixed_run_add(paragraph,text,bold,underline,italic,font_pt,font_color,font_name):
        run=paragraph.add_run(text)
        run.bold=bold
        run.underline=underline
        run.italic=italic
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return
    
    def add_picture_caption(caption):
        paragraph = doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)
        run=paragraph.add_run()
        run.text=f"Figure {len(doc.inline_shapes)}: {caption}"
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.style = 'Caption'
        run.font.name='Arial'
        run.bold=True
        run.font.size=Pt(8)
        run.font.color.rgb=word_rgb(0,0,0)  
        
        return
    
    # heading_counters=heading_counter()
    def add_auto_numbered_heading(level, text,bold,font_pt,font_color,font_name):
        if level==1:
            heading_counters[1]=0
        if level==2:
            heading_counters[2]=0
        if level==3:
            heading_counters[3]=0
        heading_counters[level - 1] += 1 # Increment the counter for the current level
        number = '.'.join(str(heading_counters[i]) for i in range(level)) # Generate the numbering
        heading = doc.add_paragraph()
        run = heading.add_run()
        run.text=f"{number} {text}"
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        heading.style = f'Heading {level}'
        heading.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        if level==1:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(18)
        else:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(6)            
        return 
   
    add_auto_numbered_heading(level=4,text="Percentage drop in coefficient of GINI",bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    full_form=portfolio_wise_full_form(portfolio_code=portfolio_code)    
    text_=f"The graphs show discriminatory power test summary for this segment of {full_form} portfolio:"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
 
    # Add first graph
    img_stream_1 = f"{path_folder}/nl/GINI_{score_model_code}.png"
    left = Inches(1.5)
    top = Inches(1.15)
    pic = doc.add_picture(img_stream_1,width=Inches(5.65), height=Inches(3.1))
    last_paragraph=doc.paragraphs[-1]
    last_paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.CENTER
    add_picture_caption("GINI")

    gini=df_gini
    gini.rename({"QTR":"QUARTER","GINI_STATISTIC":"GINI"},axis=1,inplace=True)
    df,color_list,per_change_gini,document_color=df_gini_chart(df=gini,score_model_code=score_model_code,score_model_wise_bm_gini_dict=score_model_wise_bm_gini_dict,z=z)
    
    if (per_change_gini<0):
        result_= "an increase"
        per_change_gini=0-per_change_gini
        text_gini="The GINI has significantly increased from benchmark period"
        status_gini_color=(0,176,80)
    else:
        result_="a decrease"

    if (per_change_gini<=0.2):
        status="LOW"
        text_gini="The Drop in GINI is within the MRM policy’s prescribed threshold of 20%."
        status_gini_color=(0,176,80)
    elif (per_change_gini<=0.3):
        status="MEDIUM"
        text_gini="The Drop in GINI is within the MRM policy’s prescribed threshold of 20%-30%."
        status_gini_color=(255,192,0)
    else:
        status="HIGH"  
        text_gini="The Drop in GINI has breached the MRM policy’s prescribed threshold of 30%."
        status_gini_color=(255,0,0)
    
    ini_qtr=df[z].to_list()[0]
    fin_qtr=df[z].to_list()[-3]
    
    if ((portfolio_code=="nl") & (score_model_code=="PB_UAE_BSCR12")):
        time_period=model_dev_snap_date  
    else:
        time_period=model_dev_snap_date
        
    # first_point=f"The Benchmark {time_period} GINI was {round(df['GINI'].to_list()[-1]*100,2)}%."
    # second_point=f"The GINI for the full Validation sample ({ini_qtr} to {fin_qtr}) is {round(df['GINI'].to_list()[-2]*100,2)}%. Hence, there is {result_} of {round(per_change_gini*100,2)}% in GINI from the benchmark."
    # # third_point=f"Based on the evidence, GINI for the model has been assigned a RAG status of {status}."

    # bullet_points = [first_point,
    #                 second_point
    #                 ]

    # for comment in bullet_points:
    #     paragraph=doc.add_paragraph()
        
    #     paragraph.paragraph_format.left_indent = Inches(0.5)
    #     paragraph.paragraph_format.right_indent = Inches(0.5)        
    #     paragraph.style='List Bullet'
    #     paragraph.style.font.size=Pt(8)
    #     paragraph.style.bold=True
    #     run=paragraph.add_run(f"{comment}")
    #     run.font.size=Pt(10)
    #     run.font.color.rgb=word_rgb(0,0,0)
    #     run.font.name='Arial'
        
    #     paragraph.paragraph_format.space_after = Pt(6)
    #     paragraph.paragraph_format.space_before = Pt(6)

    prompt =  f"""Could you analyze an graph depicting GINI index and also consider following points 1) The Benchmark {time_period} GINI was {round(df['GINI'].to_list()[-1]*100,2)}%.
     and The GINI for the full Validation sample ({ini_qtr} to {fin_qtr}) is {round(df['GINI'].to_list()[-2]*100,2)}%. Hence, there is {result_} of {round(per_change_gini*100,2)}% in GINI from the benchmark"""
    text_ = get_insights(img_stream_1,prompt)
    text_ = " ".join(text_.split())
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 

    text_1="Conclusion : "
    text_2=text_gini+" Based on the evidence, GINI for the model has been assigned a RAG status of"
    text_3=f" {status}."

    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=True,underline=False,italic=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=False,italic=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_3,bold=True,underline=False,italic=True,font_pt=10,font_color=status_gini_color,font_name='Arial')
    
    add_auto_numbered_heading(level=4,text="Rank Order Analysis",bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    full_form=portfolio_wise_full_form(portfolio_code=portfolio_code)    
    text_=f"The graphs show Rank Order Analysis for this segment of {full_form} portfolio:"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    img_stream_2 = F"{path_folder}/nl/RANK_ORDER_{score_model_code}.png"
    left = Inches(1)
    top = Inches(4.15)
    pic = doc.add_picture(img_stream_2, width=Inches(6.65), height=Inches(3.1))
    
    last_paragraph=doc.paragraphs[-1]
    last_paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.CENTER
    add_picture_caption("Rank Order")
    
    val_df=df_rank_order(df=df_rankorder)[1]
    val_df["%change"]=val_df['BAD_RATE'].pct_change()
    no_breaks=len(val_df[val_df['%change']>0])
    
    if (no_breaks<=0):
        first_point=f"Strong rank order demonstrated by the scorecard throughout high, medium and low risk bands in all the quarters under monitoring ({ini_qtr} to {fin_qtr})."
        text_rankorder="Expected Rank Order demonstrated by the scorecard. Hence, no concern."
        
    elif(no_breaks<=2):
        first_point=f"Mostly stable and satisfactory risk Rank ordering demonstrated by the scorecard in all the quarters under monitoring ({ini_qtr} to {fin_qtr})."
        text_rankorder="Some minor trend breaks were observed but this does not raise a concern."
    
    else:
        first_point=f"Poor rank order demonstrated by the scorecard. Multiple trend breaks have been observed in all the quarters under monitoring ({ini_qtr} to {fin_qtr})."
        text_rankorder="Rank Ordering is weak and not satisfactory."
    
    # bullet_points = [first_point,
    # ]    

    # for bullet in bullet_points:
    #     paragraph=doc.add_paragraph()
        
    #     paragraph.paragraph_format.left_indent = Inches(0.5)
    #     paragraph.paragraph_format.right_indent = Inches(0.5)
    #     paragraph.style='List Bullet'
    #     paragraph.style.font.size=Pt(8)
    #     paragraph.style.bold=True
    #     run=paragraph.add_run(f"{bullet}")
    #     run.font.size=Pt(10)
    #     run.font.color.rgb=word_rgb(0,0,0)
    #     run.font.name='Arial'
    prompt =  f"""Could you analyze an graph depicting RANK order and also consider following points {first_point}"""
    text_ = get_insights(img_stream_2,prompt)
    text_ = " ".join(text_.split())
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 

    # logger.debug("HCI graph pasted")
    # logger.debug("###############################################")
    text_1="Conclusion : "
    text_2=text_rankorder
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=True,underline=False,italic=True,font_pt=10,font_color=(0,0,0),font_name='Arial')    
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=False,italic=True,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    return

#...................................................................................................................................................................................#


def pasting_omm_monitoring_charts_2_word(df,y,title,portfolio_code,score_model_code,path_folder,score_model_wise_imp_pd_dict,x1=None):
    
    ''' 
        This function pastes charts in the Word.
    
        input params:
           
        @ title_name - This takes string variable denoting the Title of the slide.
        @ score_model_code - This takes string variable and accounts for the scorecard whose charts are to be pasted.i.e., "PB_UAE_BSCR04"
        @ ppt_file_path_name - This takes string variable for the path and file name of the ppt.Ex- "ppt/final_ppt_1.pptx"
        
        output params:
        
        It returns a saved presentation at given file path.
    
    '''
    sections=doc.sections
    for section in sections:
        section.top_margin=Inches(1)
        section.bottom_margin=Inches(1)
        section.left_margin=Inches(1)
        section.right_margin=Inches(1)
        
    def paragraph_add(space_after,space_before,text,bold,font_pt,font_color,font_name):    
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        run=paragraph.add_run()
        run.text= text
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return
    
    def mixed_paragraph_add(space_after,space_before):
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        
        return paragraph
    
    def mixed_run_add(paragraph,text,bold,underline,italic,font_pt,font_color,font_name):
        run=paragraph.add_run(text)
        run.bold=bold
        run.underline=underline
        run.italic=italic
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return       

    def add_picture_caption(caption):
        paragraph = doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)
        run=paragraph.add_run()
        run.text=f"Figure {len(doc.inline_shapes)}: {caption}"
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.style = 'Caption'
        run.font.name='Arial'
        run.bold=True
        run.font.size=Pt(8)
        run.font.color.rgb=word_rgb(0,0,0)  
        
        return
    
    # heading_counters=heading_counter()
    def add_auto_numbered_heading(level, text,bold,font_pt,font_color,font_name):
        if level==1:
            heading_counters[1]=0
        if level==2:
            heading_counters[2]=0
        if level==3:
            heading_counters[3]=0
        heading_counters[level - 1] += 1 # Increment the counter for the current level
        number = '.'.join(str(heading_counters[i]) for i in range(level)) # Generate the numbering
        heading = doc.add_paragraph()
        run = heading.add_run()
        run.text=f"{number} {text}"
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        heading.style = f'Heading {level}'
        heading.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        if level==1:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(18)
        else:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(6)            
        return 

    add_auto_numbered_heading(level=4,text="Model Forecast Accuracy (MAPE) - Long Term ODR vs Forecasted TTC PD",bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    full_form=portfolio_wise_full_form(portfolio_code=portfolio_code)    
    text_=f"The graphs below show Long Term Observed Default Rate(LT_ODR), Long Term Expected Default Rate(LT_EDR) for {full_form} portfolio:"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    # Add first graph
    img_stream_1 = f"{path_folder}/nl/MAPE_{score_model_code}.png"
    left = Inches(1.5)
    top = Inches(1.15)
    pic = doc.add_picture(img_stream_1,width=Inches(6.65), height=Inches(3.1))
    last_paragraph=doc.paragraphs[-1]
    last_paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.CENTER
    add_picture_caption("MAPE")
    
    df_1_1,df_1_2=df_mape_chart(df=df,y=y,score_model_code=score_model_code,score_model_wise_imp_pd_dict=score_model_wise_imp_pd_dict,x=x1)

    lt_odr=df_1_2[df_1_2['variable']=='LT_ODR'][y].to_list()[0]
    lt_edr=df_1_2[df_1_2['variable']=='LT_EDR'][y].to_list()[0]

    if x1 == "QTR":
        ini_date=sorted(df_1_2[x1].unique())[0]
    
        fin_date=sorted(df_1_2[x1].unique())[-1]
    
    else:
        ini_date=sorted(df_1_2[x1].unique())[0]
        date_object = datetime.strptime(ini_date, '%Y-%m-%d')
        ini_date = date_object.strftime("%b'%Y")
        fin_date=sorted(df_1_2[x1].unique())[-1]
        date_object = datetime.strptime(fin_date, '%Y-%m-%d')
        fin_date = date_object.strftime("%b'%Y")
    
    mape=abs((lt_edr-lt_odr)/lt_odr)
    
    if (mape<=0.2):
        status="LOW"
        text_mape="Model accuracy remains strong for this segment."
        status_mape_color=(0,176,80)

    elif (mape<=0.3):
        status="MEDIUM"
        text_mape="Model accuracy remains moderate for this segment."
        status_mape_color=(255,192,0)
        
    else:
        status="HIGH"
        text_mape="Model accuracy remains poor for this segment."
        status_mape_color=(255,0,0)
    
    first_point=f"Long term observed default rate is calculated from {ini_date} onwards till {fin_date}, which looks at next 12 months of default performance."
    second_point=f"This segment has an Observed Long term default Rate of {round(lt_odr,2)}% while the Predicted Long term default rate for the same period is {round(lt_edr,2)}%."
    third_point=f"The MAPE against the TTC PD is {round(mape*100,2)}%."
    # fourth_point=f"Based on the evidence, MAPE for the model has been assigned a RAG status of {status}."
    if portfolio_code=="nl":
        fourth_point="The observed gap in months from April-2020 to July-2020 and May-2021 is due to incorrect DPD assignment leading to incorrect scorecard tagging. Due to this reason, data for these 5 months have been excluded."
    else:
        fourth_point=""
    bullet_points = [first_point,
    second_point,
    third_point,
    # fourth_point
    ]

    # for comment in bullet_points:
    #     paragraph=doc.add_paragraph()
        
    #     paragraph.paragraph_format.left_indent = Inches(0.5)
    #     paragraph.paragraph_format.right_indent = Inches(0.5)        
    #     paragraph.style='List Bullet'
    #     paragraph.style.font.size=Pt(8)
    #     paragraph.style.bold=True
    #     run=paragraph.add_run(f"{comment}")
    #     run.font.size=Pt(10)
    #     run.font.color.rgb=word_rgb(0,0,0)
    #     run.font.name='Arial'
        
    #     paragraph.paragraph_format.space_after = Pt(6)
    #     paragraph.paragraph_format.space_before = Pt(6)

    prompt =  f"""Could you analyze an graph depicting RANK order and also consider following points 1) {first_point}  2) {second_point} 3){third_point}"""
    text_ = get_insights(img_stream_1,prompt)
    text_ = " ".join(text_.split())
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 

    text_1="Conclusion : "
    text_2=text_mape+" Based on the evidence, MAPE for the model has been assigned a RAG status of"
    text_3=f" {status}."
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=True,underline=False,italic=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=False,italic=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_3,bold=True,underline=False,italic=True,font_pt=10,font_color=status_mape_color,font_name='Arial')
    
    return

#...................................................................................................................................................................................#

def pasting_variable_analysis_charts(df_csi,df_iv,df_vif,current_qtr,y,title,portfolio_code,score_model_code,var_name_list,var_full_form_list,val_start_date_imm,val_start_date_omm,val_end_date_imm,val_end_date_omm,path_folder,segment,model_dev_snap_date=None,model_dev_benchmark_date=None):
    
    sections=doc.sections
    for section in sections:
        section.top_margin=Inches(1)
        section.bottom_margin=Inches(1)
        section.left_margin=Inches(1)
        section.right_margin=Inches(1)
        
    def paragraph_add(space_after,space_before,text,bold,font_pt,font_color,font_name):    
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        run=paragraph.add_run()
        run.text= text
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return
    
    def mixed_paragraph_add(space_after,space_before):
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        
        return paragraph
    
    def mixed_run_add(paragraph,text,bold,underline,italic,font_pt,font_color,font_name):
        run=paragraph.add_run(text)
        run.bold=bold
        run.underline=underline
        run.italic=italic
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return 
    
    def b_u_paragraph_add(space_after,space_before,text,bold,underline,italic,font_pt,font_color,font_name):
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        run=paragraph.add_run()
        run.text= text
        run.bold=bold
        run.underline=underline
        run.italic=italic
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
    
    def add_picture_caption(caption):
        paragraph = doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)
        run=paragraph.add_run()
        run.text=f"Figure {len(doc.inline_shapes)}: {caption}"
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.style = 'Caption'
        run.font.name='Arial'
        run.bold=True
        run.font.size=Pt(8)
        run.font.color.rgb=word_rgb(0,0,0)  
        
        return
    
    def add_auto_numbered_heading(level, text,bold,font_pt,font_color,font_name):
        if level==1:
            heading_counters[1]=0
        if level==2:
            heading_counters[2]=0
        if level==3:
            heading_counters[3]=0
        heading_counters[level - 1] += 1 # Increment the counter for the current level
        number = '.'.join(str(heading_counters[i]) for i in range(level)) # Generate the numbering
        heading = doc.add_paragraph()
        run = heading.add_run()
        run.text=f"{number} {text}"
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        heading.style = f'Heading {level}'
        heading.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        if level==1:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(18)
        else:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(6)            
        return 

    def create_table(n_rows,n_cols):
        table=doc.add_table(rows=n_rows, cols=n_cols)
        table.style='Table Grid' 
        table.alignment=WD_TABLE_ALIGNMENT.CENTER
        table.autofit=True 
        
        for i in range(len(table.rows)):
            for j in range(len(table.columns)):
        
                table.cell(i,j).paragraphs[0].alignment  = WD_ALIGN_VERTICAL.CENTER
                table.cell(i,j).vertical_alignment=WD_CELL_VERTICAL_ALIGNMENT.CENTER
                
        return table
    
    def fill_table_color(table,row,col,color="FFFFFF"):
        shading_elm_1 = parse_xml(r'<w:shd {} w:fill="{}"/>'.format(nsdecls('w'),color))
        table.rows[row].cells[col]._tc.get_or_add_tcPr().append(shading_elm_1)
        return

    
    def fill_table(table,row,col,bold,text,r,g,b,alignment):
        
        cell_1=table.cell(row,col).paragraphs[0]
        run=cell_1.add_run()
        run.text=text
        run.bold=bold
        run.font.size=Pt(9)
        run.font.name='Arial'
        run.font.color.rgb=word_rgb(r,g,b)
        cell_1.alignment=alignment
        return
    
    def add_table_caption(caption):
        paragraph = doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)
        run=paragraph.add_run()
        run.text=f"Table {len(doc.tables)}: {caption}"
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.style = 'Caption'
        run.font.name='Arial'
        run.bold=True
        run.font.size=Pt(8)
        run.font.color.rgb=word_rgb(0,0,0)  
        
        return
    
    def set_table_col_width(table,width_arr):
        for i,x in enumerate(width_arr):
            for col in table.columns[i].cells:
                col.width=Inches(x)
        return      
    

    val_start_date_imm=pd.to_datetime(val_start_date_imm)
    val_start_date_omm=pd.to_datetime(val_start_date_omm)    
    val_end_date_imm=pd.to_datetime(val_end_date_imm)
    val_end_date_omm=pd.to_datetime(val_end_date_omm)  
    
    csi_1=df_csi[df_csi["YYYYMM"]=="Validation_Sample"].reset_index(drop=True)
    csi_1['VARIABLE_NAME']=csi_1['VARIABLE_NAME'].str[:-4]
    req_cols=["SCORE_CARD","VARIABLE_NAME","VALUE","BUCKET_SCORE","CONCENTRATION_BENCHMARK","CSI"]
    csi_1=csi_1[req_cols]
    csi_1["VALUE"]=csi_1["VALUE"].astype('str')
    csi_1['FIRST_VAL']=csi_1['VALUE'].apply(lambda x: x.split(',')[0][1:])
    csi_1['FIRST_VAL']=csi_1['FIRST_VAL'].replace('-','0')
    csi_1['FIRST_VAL']=csi_1['FIRST_VAL'].replace('.0','0')
    csi_1['FIRST_VAL'] = csi_1['FIRST_VAL'].replace('', '0') ### added extra
    print(csi_1['FIRST_VAL'].value_counts())
    csi_1['FIRST_VAL']=csi_1['FIRST_VAL'].astype('float')
    csi_1['FIRST_VAL']=csi_1['FIRST_VAL'].astype('int64') # added astype('int')
    csi_1.sort_values(by=['VARIABLE_NAME','FIRST_VAL','BUCKET_SCORE'],inplace=True)
    csi_1_RANK= csi_1.groupby(["VARIABLE_NAME","BUCKET_SCORE"]).agg(                   
                        dummy_score = pd.NamedAgg(column='BUCKET_SCORE', aggfunc='count'),
                        CONCENTRATION_BENCHMARK = pd.NamedAgg(column='CONCENTRATION_BENCHMARK', aggfunc='sum'),
                        VALUE = pd.NamedAgg(column='VALUE', aggfunc='sum')).reset_index(drop=False)
    
    csi_1_RANK['FIRST_VAL']=csi_1_RANK['VALUE'].apply(lambda x: x.split(',')[0][1:])
    #print(csi_1_RANK['FIRST_VAL'].unique())
    csi_1_RANK['FIRST_VAL']=csi_1_RANK['FIRST_VAL'].replace('-','0')
    csi_1_RANK['FIRST_VAL']=csi_1_RANK['FIRST_VAL'].str.replace('.03.0', '0.30', regex=False)
    csi_1_RANK['FIRST_VAL']=csi_1_RANK['FIRST_VAL'].replace('.0','0')
    csi_1_RANK['FIRST_VAL']=csi_1_RANK['FIRST_VAL'].replace('','0')
    csi_1_RANK['FIRST_VAL']=csi_1_RANK['FIRST_VAL'].astype('float')
    csi_1_RANK['FIRST_VAL']=csi_1_RANK['FIRST_VAL'].astype('int64') # added.astype('int')
    csi_1_RANK["VALUE"]=np.where(csi_1_RANK["VALUE"].str.len()<=50,csi_1_RANK["VALUE"],"[Neg "+csi_1_RANK["VALUE"].str[7:11]+"Neg "+csi_1_RANK["VALUE"].str[17:20]+" , "+"[Neg "+csi_1_RANK["VALUE"].str[114:118]+"Neg "+csi_1_RANK["VALUE"].str[124:127])
    csi_1_RANK.sort_values(by=['VARIABLE_NAME','FIRST_VAL','BUCKET_SCORE'],inplace=True)
                             
    def merge_buckets(bucket):
        x=bucket
        x=str(x)
        y=x.split(")[")
        if (len(y)==1):
            return x
            
        else:
            a=y[0].split(" ")[0]
            b=y[(len(y)-1)].split(" ")[-1]  
            return a+b
    csi_1_RANK["VALUE"]=np.where(csi_1_RANK["dummy_score"]==1,csi_1_RANK["VALUE"],csi_1_RANK["VALUE"].apply(merge_buckets))
    
    # csi_1_RANK=csi_1_RANK.sort_values(by=["VARIABLE_NAME","VALUE"])
    # print(f"shape is {csi_1_RANK.shape[0]} and {csi_1_RANK.shape[1]} for {score_model_code}")
    
    iv=df_iv
    iv["BAD_RATE"]=iv['Bad_count']/iv['Total_count']
    req_cols=["SCORE_CARD","VARIABLE_NAME","BUCKET_SCORE","PERC_POP","PERC_GOOD","PERC_BAD","BAD_RATE","IV"]
    iv_1=iv[req_cols]
    iv_1['VARIABLE_NAME']=iv_1['VARIABLE_NAME'].str[:-3]

    final_df=pd.merge(csi_1_RANK,iv_1,on=["VARIABLE_NAME","BUCKET_SCORE"])
    final_df=final_df[["VARIABLE_NAME","VALUE","BUCKET_SCORE","CONCENTRATION_BENCHMARK","PERC_POP","PERC_GOOD","PERC_BAD","BAD_RATE","IV"]]
    final_df["VARIABLE_NAME"]=final_df["VARIABLE_NAME"].str[:]+"_RNG"
    
    list_l=["CONCENTRATION_BENCHMARK","PERC_POP","PERC_GOOD","PERC_BAD","BAD_RATE","IV"]
    for x in list_l:
        final_df[x]=final_df[x].apply(lambda x: f"{x:.2%}")
    
    full_form=portfolio_wise_full_form(portfolio_code=portfolio_code)
    scorecard_available=portfolio_wise_score_model_code(segment)
    no_segments=len(scorecard_available)
    y=title.split()[:2]+title.split()[3:]
    y=" ".join(y)
    y1=" ".join(title.split()[3:])
    

    add_auto_numbered_heading(level=4,text="Variable Level Analysis",bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial')  
    
    text_="This section describes the model variable level analysis."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    text_="Characteristic Stability Analysis & Risk Discrimination Power (Rank Ordering)"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
        
    text_=f"The tables below present the CSI levels of variables of the {y}. There are {len(var_name_list)} model variables in the {y1}."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')        
    for i, var_name in enumerate(var_name_list):
        
        var_full_name=var_full_form_list[i]
        csi_1=csi_1[['VARIABLE_NAME','CSI']].drop_duplicates()
        variable=var_name[:-4]
        csi_var_val=csi_1[csi_1['VARIABLE_NAME']==variable]['CSI']
        if float(csi_var_val)<0.1:
            text_var_csi=f"Variable has very low CSI indicating it hasn’t witnessed any significant population shift over time from the benchmark period {model_dev_snap_date}."
        elif float(csi_var_val)<0.25:
            text_var_csi=f"Variable has moderate CSI indicating it hasn’t witnessed much significant population shift over time from the benchmark period {model_dev_snap_date}."
        else:
            text_var_csi=f"Variable has high CSI indicating it has witnessed a significant population shift over time from the benchmark period {model_dev_snap_date}."
        
        if i!=0:
            text_=" "
            paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 

        text_=f"Var {i+1}: {var_full_name}"
        b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
            
        
        img_stream_1 = f"{path_folder}/nl/CSI_{score_model_code}_{var_name}.png"
        left = Inches(1.5)
        top = Inches(1.15)
        pic = doc.add_picture(img_stream_1,width=Inches(6.65), height=Inches(3.1))
        last_paragraph=doc.paragraphs[-1]
        last_paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.CENTER
        x="/"
        add_picture_caption(f"CSI – {y1} ({str((i+1))}{x}{str(len(var_name_list))})") 
        
        prompt =  f"""Could you analyze an graph depicting CSI (Characteristic Stability index) and also consider following points 1) {text_var_csi}"""
        text_ = get_insights(img_stream_1,prompt)
        text_ = " ".join(text_.split())
        paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 

        # paragraph=doc.add_paragraph()
        
        # paragraph.paragraph_format.left_indent = Inches(0.5)
        # paragraph.paragraph_format.right_indent = Inches(0.5)        
        # paragraph.style='List Bullet'
        # paragraph.style.font.size=Pt(8)
        # paragraph.style.bold=True
        # run=paragraph.add_run(f"{text_var_csi}")
        # run.font.size=Pt(10)
        # run.font.color.rgb=word_rgb(0,0,0)
        # run.font.name='Arial'
        
        # paragraph.paragraph_format.space_after = Pt(6)
        # paragraph.paragraph_format.space_before = Pt(6)        

        df_rank_order_var=final_df[final_df["VARIABLE_NAME"]==var_name]
        
        df_rank_order_var_conclusion_text=df_rank_order_var.copy(deep=True)
        df_rank_order_var_conclusion_text["BAD_RATE"]=df_rank_order_var_conclusion_text["BAD_RATE"].str[:-1].astype("float")
        df_rank_order_var_conclusion_text["BUCKET_SCORE"]=df_rank_order_var_conclusion_text["BUCKET_SCORE"].astype("int")
        df_rank_order_var_conclusion_text["pct_change_score"]=df_rank_order_var_conclusion_text["BUCKET_SCORE"].pct_change()
        df_rank_order_var_conclusion_text["pct_change_bad_rate"]=df_rank_order_var_conclusion_text["BAD_RATE"].pct_change()
        df_rank_order_var_conclusion_text["product_pct_change_score_bad_rate"]=df_rank_order_var_conclusion_text["pct_change_score"]*df_rank_order_var_conclusion_text["pct_change_bad_rate"]

        no_breaks=df_rank_order_var_conclusion_text[df_rank_order_var_conclusion_text['product_pct_change_score_bad_rate']>0].shape[0]
    

        if no_breaks==0:
            text_rank_order_conclusion="Variable follows expected rank order based on bin scores."
            
        elif ((no_breaks>0)&(no_breaks<=df_rank_order_var_conclusion_text.shape[0]//2)):
            text_rank_order_conclusion="Variable follows expected rank order based on bin scores, but minor trend break has been observed in rank order."
            
        else:
            text_rank_order_conclusion="Variable doesn't follow expected rank order based on bin scores. Some concerns."

        # text_rank_order_conclusion_list.append(text_rank_order_conclusion)
        n_rows=df_rank_order_var.shape[0]
        table=create_table(2+n_rows,7)
        add_table_caption(f"Rank Order Summary – {y1} ({str((i+1))}{x}{str(len(var_name_list))})")
        
        for row in table.rows:
            row.height = Inches(0.3)
            
        table.cell(0,0).merge(table.cell(1,0))
        table.cell(0,1).merge(table.cell(1,1))
        table.cell(0,2).merge(table.cell(1,2))
        table.cell(0,3).merge(table.cell(0,5))
        table.cell(0,6).merge(table.cell(1,6))
        
        omm_start_date=val_start_date_omm.month_name()[:3]+"'"+str(val_start_date_omm.year)[2:]
        omm_end_date=val_end_date_omm.month_name()[:3]+"'"+str(val_end_date_omm.year)[2:]
        row_0= ["BIN","SCORE",f"Benchmark {model_dev_snap_date} Population", f"Performance ({omm_start_date}-{omm_end_date})","Validation Bad Rate"]
        for k,x in enumerate(row_0):  
            
            if(k!=len(row_0)-1):
                fill_table(table,0,k,True,x,255,255,255,1)  
                fill_table_color(table,row=0,col=k,color="0070C0")
            else:
                fill_table(table,0,k+2,True,x,255,255,255,1)  
                fill_table_color(table,row=0,col=k+2,color="0070C0") 
                
        row_1=["%Population","%Good", "%Bad" ]
        for k,x in enumerate(row_1): 
            fill_table(table,1,k+3,True,x,255,255,255,1)  
            fill_table_color(table,row=1,col=k+3,color="0070C0")  
        
        dis_power=0
        dict_pop_bads={}
        
        for k in range(n_rows): 
            values=df_rank_order_var.iloc[k,1:-1].to_list()
            
            list_pop_bads=[]
            dict_pop_bads[float(values[1])]=list_pop_bads
            list_pop_bads.append(str(round(float(values[3][:4]),0))[:-2])
            list_pop_bads.append(str(round(float(values[5][:4]),0))[:-2])
            iv=df_rank_order_var.iloc[k,-1]
            
            dis_power+=abs(float(values[4][:-1])-float(values[5][:-1]))
            if dis_power<20:
                text_disc="satisfactory"
            elif dis_power<30:
                text_disc="moderate"
            elif dis_power<50:
                text_disc="good"
            else:
                text_disc="very strong"
            
            for j,x in enumerate(values): 
                # print(x, str(x))
                fill_table(table,2+k,j,False,str(x),0,0,0,1)  
                fill_table_color(table,row=2+k,col=j,color="FFFFFF")      
                
        min_val_dict_pop_bads_score=min(dict_pop_bads.keys())
        pop_no=dict_pop_bads[min_val_dict_pop_bads_score][0]
        bads_no=dict_pop_bads[min_val_dict_pop_bads_score][1]
        
        table_width=[0.9,0.9,0.9,2.88,0.9]                
        set_table_col_width(table,table_width)
        
        # val_df["%change"]=val_df['BAD_RATE'].pct_change()
        # no_breaks=len(val_df[val_df['%change']>0])     
    
        text_1="Conclusion : "
        paragraph=mixed_paragraph_add(space_after=8,space_before=0)
        mixed_run_add(paragraph=paragraph,text=text_1,bold=True,underline=False,italic=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
        
        bullet_points=[f"Variable has {text_disc} discriminating power.",
                      f"The Variable has an Information Value of {iv} in the validation period.",
                      text_rank_order_conclusion,
                      f"The highest-risk (lowest score) bin contains ~{pop_no}% population and captures ~{bads_no}% of Bads."]
        
        for comments in bullet_points:
            paragraph=doc.add_paragraph()
            
            paragraph.paragraph_format.left_indent = Inches(0.5)
            paragraph.paragraph_format.right_indent = Inches(0.5)        
            paragraph.style='List Bullet'
            paragraph.style.font.size=Pt(8)
            paragraph.style.bold=True
            run=paragraph.add_run(f"{comments}")
            run.font.size=Pt(10)
            run.font.color.rgb=word_rgb(0,0,0)
            run.font.name='Arial'
            
            paragraph.paragraph_format.space_after = Pt(6)
            paragraph.paragraph_format.space_before = Pt(6)        
        
        if (i==(len(var_name_list)-1)):

            paragraph=doc.add_paragraph()
            text_1="Segment Conclusion : "
            paragraph=mixed_paragraph_add(space_after=8,space_before=0)    
            mixed_run_add(paragraph=paragraph,text=text_1,bold=True,underline=True,italic=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
           
            bullet_points=[
            "Overall Variable performance is good.",
            "Variables’ discrimination power is also good."
                    ]
            
            for comment in bullet_points:
                paragraph=doc.add_paragraph()
                
                paragraph.paragraph_format.left_indent = Inches(0.5)
                paragraph.paragraph_format.right_indent = Inches(0.5)        
                paragraph.style='List Bullet'
                paragraph.style.font.size=Pt(8)
                paragraph.style.bold=True
                run=paragraph.add_run(f"{comment}")
                run.font.size=Pt(10)
                run.font.color.rgb=word_rgb(0,0,0)
                run.font.name='Arial'
                
                paragraph.paragraph_format.space_after = Pt(6)
                paragraph.paragraph_format.space_before = Pt(6)
                
        doc.add_page_break()
        
    text_="Multicollinearity"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial')  
    
    text_="A general rule of thumb, any parameter with VIF greater 5 indicates high multicollinearity and should be dropped from the model. The table below presents the VIF for variables during development and in out of time validation sample."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')  
    
    df_1=df_vif.iloc[1:,:]
    df_1["Variable"]=df_1["Variable"].str[:-3]+"_RNG"
    df_1=df_1.sort_values(by="Variable")
    df_1=df_1.reset_index(drop=True)
    df_1["VIF"]=round(df_1["VIF"],3)
    max_vif=df_1["VIF"].max()
    df_1["VIF"]=df_1["VIF"].astype('str')
    
    n_rows=df_1.shape[0]
    n_cols=df_1.shape[1]
    y_seg=y.split()
    y_seg=y_seg[2:]
    y_seg=" ".join(y_seg)
    y_seg=portfolio_code.upper()+" "+y_seg  
    table=create_table(2+n_rows,2)
    add_table_caption(f"VIF Summary – {y_seg}")
    
    for row in table.rows:
        row.height = Inches(0.3)
        
    table.cell(0,0).merge(table.cell(0,1))

    row_0=[y_seg]
    for i,x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)  
        fill_table_color(table,row=0,col=i,color="0070C0")
    
    row_1=["Variable Name","OOT Sample"]
    for i,x in enumerate(row_1):
        fill_table(table,1,i,True,x,255,255,255,1)  
        fill_table_color(table,row=1,col=i,color="0070C0") 
        
    for j in range(n_cols):
        if (j!=0):
            temp_list=df_1.iloc[:,j].to_list()   
            for i,x in enumerate(temp_list):
                fill_table(table,2+i,j,False,x,0,0,0,1)  
                fill_table_color(table,row=2+i,col=j,color="FFFFFF") 
        else:
            temp_list=df_1.iloc[:,j].to_list()
            size=len(temp_list)
            print(">>>>>>>>>>>>size",size)
            print(">>>>>>>>>>>>>>>>>>var_full_form_list",var_full_form_list)
            temp_list=[]
            for i in range(size):
                var_full_name=var_full_form_list[i]
                temp_list.append(var_full_name)
            for i,x in enumerate(temp_list):
                fill_table(table,2+i,j,False,x,0,0,0,0)  
                fill_table_color(table,row=2+i,col=j,color="FFFFFF") 
            
    table_width=[5.45,1.05]                
    set_table_col_width(table,table_width)
    
    if max_vif<=5:
        text_vif="As can be observed from the table above that VIF for all parameters in out of time validation sample is below 5. No significant concern of multi-collinearity."
    else:
        text_vif="As can be observed from the table above that VIF for some parameters in out of time validation sample are above 5, there are some concerns of multi-collinearity."
        
    text_1="Conclusion : "
    text_2=text_vif
    paragraph=mixed_paragraph_add(space_after=8,space_before=0)
    mixed_run_add(paragraph=paragraph,text=text_1,bold=True,underline=False,italic=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    mixed_run_add(paragraph=paragraph,text=text_2,bold=False,underline=False,italic=True,font_pt=10,font_color=(0,0,0),font_name='Arial')  
        
    return
      
    
    
#...................................................................................................................................................................................#


def document_part_5_6_7_8(exclusion_df,portfolio_code,val_start_date_imm,val_start_date_omm,val_end_date_imm,val_end_date_omm,path_folder,bank_name=None):
    
    def paragraph_add(space_after,space_before,text,bold,font_pt,font_color,font_name):    
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        run=paragraph.add_run()
        run.text= text
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return
    
    def mixed_paragraph_add(space_after,space_before):
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        
        return paragraph
    
    def mixed_run_add(paragraph,text,bold,underline,italic,font_pt,font_color,font_name):
        run=paragraph.add_run(text)
        run.bold=bold
        run.underline=underline
        run.italic=italic
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
        return
    
    def b_u_paragraph_add(space_after,space_before,text,bold,underline,italic,font_pt,font_color,font_name):
        paragraph=doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(space_after)
        paragraph.paragraph_format.space_before = Pt(space_before)
        paragraph.alignment=WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        run=paragraph.add_run()
        run.text= text
        run.bold=bold
        run.underline=underline
        run.italic=italic
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        
    def add_auto_numbered_heading(level, text,bold,font_pt,font_color,font_name):
        if level==1:
            heading_counters[1]=0
        if level==2:
            heading_counters[2]=0
        if level==3:
            heading_counters[3]=0
        heading_counters[level - 1] += 1 # Increment the counter for the current level
        number = '.'.join(str(heading_counters[i]) for i in range(level)) # Generate the numbering
        heading = doc.add_paragraph()
        run = heading.add_run()
        run.text=f"{number}  {text}"
        run.bold=bold
        run.font.size=Pt(font_pt)
        run.font.color.rgb=word_rgb(font_color[0],font_color[1],font_color[2])
        run.font.name=font_name 
        heading.style = f'Heading {level}'
        heading.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        if level==1:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(18)
        else:
            heading.paragraph_format.space_after = Pt(6)
            heading.paragraph_format.space_before = Pt(6)            
        return 

    def create_table(n_rows,n_cols):
        table=doc.add_table(rows=n_rows, cols=n_cols)
        table.style='Table Grid' 
        table.alignment=WD_TABLE_ALIGNMENT.CENTER
        table.autofit=True 
        
        for i in range(len(table.rows)):
            for j in range(len(table.columns)):
        
                table.cell(i,j).paragraphs[0].alignment  = WD_ALIGN_VERTICAL.CENTER
                table.cell(i,j).vertical_alignment=WD_CELL_VERTICAL_ALIGNMENT.CENTER
                
        return table
    
    def fill_table_color(table,row,col,color="FFFFFF"):
        shading_elm_1 = parse_xml(r'<w:shd {} w:fill="{}"/>'.format(nsdecls('w'),color))
        table.rows[row].cells[col]._tc.get_or_add_tcPr().append(shading_elm_1)
        return

    
    def fill_table(table,row,col,bold,text,r,g,b,alignment):
        
        cell_1=table.cell(row,col).paragraphs[0]
        run=cell_1.add_run()
        run.text=text
        run.bold=bold
        run.font.size=Pt(9)
        run.font.name='Arial'
        run.font.color.rgb=word_rgb(r,g,b)
        cell_1.alignment=alignment
        return
    
    def add_table_caption(caption):
        paragraph = doc.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)
        run=paragraph.add_run()
        run.text=f"Table {len(doc.tables)}: {caption}"
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.style = 'Caption'
        run.font.name='Arial'
        run.bold=True
        run.font.size=Pt(8)
        run.font.color.rgb=word_rgb(0,0,0)  
        
        return
    
    def set_table_col_width(table,width_arr):
        for i,x in enumerate(width_arr):
            for col in table.columns[i].cells:
                col.width=Inches(x)
        return       
 
    
    val_start_date_imm=pd.to_datetime(val_start_date_imm)
    val_start_date_omm=pd.to_datetime(val_start_date_omm)    
    val_end_date_imm=pd.to_datetime(val_end_date_imm)
    val_end_date_omm=pd.to_datetime(val_end_date_omm)
         
    sections=doc.sections
    for section in sections:
        section.top_margin=Inches(1)
        section.bottom_margin=Inches(1)
        section.left_margin=Inches(1)
        section.right_margin=Inches(1)
        
    add_auto_numbered_heading(level=1,text="Model Implementation Testing",bold=True,font_pt=14,font_color=(0,0,0),font_name='Arial')    

    text_= "Validation of Model implementation was performed by reconciling the system generated score in L3 with independently calculated B-score."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')

    text_= "As explained in Section 3.5: Model Implementation results, B-score generated in L3 production system reconciled 100% with manually calculated B-score.  "
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    add_auto_numbered_heading(level=1,text="Key Model Findings",bold=True,font_pt=14,font_color=(0,0,0),font_name='Arial')    
    
    full_form=portfolio_wise_full_form(portfolio_code=portfolio_code)
    text_=f"Following are the key model validation findings for {full_form} B-score annual validation. "
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
#     if portfolio_code!="nl":
#         text_no_1="Validation Data Quality: There were no major concerns on the validation data quality. Data quality and volume was adequate for performing a robust quantitative validation."    
#     else:
#         text_no_1="Validation Data Quality: There was some concern on the validation data quality i.e., Incorrect DPD assignment leading to incorrect scorecard assignment was observed in Apr’20 to Jul’20 and in May’21. However, these months were excluded from the validation exercise for both Current and Delinquent segments. Overall data quality and volume was adequate for performing a robust quantitative validation."
        
#          "B-score generated from the production system (L3) matched 100% with B-score calculated independently by the Validation Team using each model variable score weight from the model development document.  Hence, no concern on model implementation. "       
    
    bullet_points=[
   "Validation Data Quality:",
"Model Implementation:",
        "Model Discrimination:",
        "Population Shift:",
        "Score Concentration:",
        "Model forecast accuracy of Long Term ODR:",
        "Overall:"
]
    
    for comment in bullet_points:
        paragraph=doc.add_paragraph()
        
        paragraph.paragraph_format.left_indent = Inches(0.5)
        paragraph.paragraph_format.right_indent = Inches(0.5)        
        paragraph.style='List Bullet'
        paragraph.style.font.size=Pt(8)
        paragraph.style.bold=True
        run=paragraph.add_run(f"{comment}")
        run.font.size=Pt(10)
        run.bold=True
        run.font.color.rgb=word_rgb(0,0,0)
        run.font.name='Arial'
        
        paragraph.paragraph_format.space_after = Pt(6)
        paragraph.paragraph_format.space_before = Pt(6)
        
    add_auto_numbered_heading(level=1,text="Action Items and Recommendations",bold=True,font_pt=14,font_color=(0,0,0),font_name='Arial')    
        
    text_= "This section describes the review of previous validation action items as well new action items issued for the current validation."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')

    text_="Previous Action items"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')  
    
    table=create_table(2,7)
    add_table_caption(f"Previous Action Items") 
    
    for row in table.rows:
        row.height = Inches(0.3) 
        
    row_0=["S.No.","Validation","Segment","Issue","Action Item","Severity","Status"]
    for i,x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)
        fill_table_color(table,0,i,"0070C0") 
        
    table_width=[0.4,0.75,0.9,1.31,1.85,.65,0.55]               
    set_table_col_width(table,table_width)
        
    text_="Conclusion: Previous Action items"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial')    
    
    text_="New Action items"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 
    
    text_=f"New set of action items were raised during this year i.e., {str(val_end_date_imm.year)} annual validation. Refer to below table for more details."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')   

    table=create_table(2,7)
    add_table_caption(f"New Action Items") 
    
    for row in table.rows:
        row.height = Inches(0.3) 
        
    row_0=["S.No.", "Segment",	"Issue",	"Action Item","Severity","Remediation Plan",	"Action item Owner"]
    for i,x in enumerate(row_0):
        fill_table(table,0,i,True,x,255,255,255,1)
        fill_table_color(table,0,i,"0070C0") 
        
    table_width=[0.45,0.85,0.85,1.85,0.7,0.95,0.9]                
    set_table_col_width(table,table_width)
    
    text_="Conclusion: New Action items"
    b_u_paragraph_add(space_after=8,space_before=0,text=text_,bold=True,underline=True,italic=False,font_pt=10,font_color=(0,0,0),font_name='Arial') 

    add_auto_numbered_heading(level=1,text="Appendix",bold=True,font_pt=14,font_color=(0,0,0),font_name='Arial')    
        
    text_= "End of the Report"
    paragraph_add(space_after=8,space_before=0,text=text_,bold=True,font_pt=10,font_color=(0,0,0),font_name='Arial')
    
    text_ ="This document has been prepared specifically as part of Annual Validation exercise by the FAB Model Validation team based on the defined scope of the engagement."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')

    text_ ="All the supporting analyses for conclusions drawn in this document/data have been provided to the bank, in the form of excel templates, presentations etc., to understand the rationale. This document should be read in conjunction with all these supporting documents."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')

    text_ ="Any observations with respect to this document should be addressed to the ‘Enterprise Risk Solutions’ of the bank."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')

    text_ ="The contents of this document are confidential and shall not be reproduced without the explicit consent of First Abu Dhabi Bank."
    paragraph_add(space_after=8,space_before=0,text=text_,bold=False,font_pt=10,font_color=(0,0,0),font_name='Arial')

    new_directory = f"{path_folder}\\{portfolio_code}"
    new_path = os.path.join(new_directory, f"{bank_name}.docx")
    
    doc.save(new_path)
    

    return
# table
# graph
        